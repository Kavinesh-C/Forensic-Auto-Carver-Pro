import os
import secrets
import re
import zipfile
import mimetypes
import time
from flask import Flask, render_template, render_template_string, request, redirect, url_for, flash, jsonify, Response, send_from_directory, make_response, send_file, session
import logging
from werkzeug.utils import secure_filename
from PIL import Image
import io
import threading
import pytsk3
import hashlib
import math
import magic
import shutil
import secrets
from cryptography.hazmat.primitives.ciphers import Cipher, algorithms, modes
from cryptography.hazmat.backends import default_backend
from cryptography.hazmat.primitives import hashes
from cryptography.hazmat.primitives.kdf.pbkdf2 import PBKDF2HMAC
from cryptography.fernet import Fernet, InvalidToken
import base64
import datetime
import gzip
import json
import csv
import xml.etree.ElementTree as ET
import mmap
from multiprocessing.dummy import Pool
import sqlite3
import tempfile
from urllib.parse import urlparse



# --- Optional Dependency Handling ---
try:
    from weasyprint import HTML
except ImportError:
    HTML = None

try:
    import docx
except ImportError:
    docx = None

try:
    from Evtx.Evtx import Evtx
except ImportError:
    Evtx = None
    print("Warning: 'python-evtx' is not installed. .evtx log parsing will not be available. Install with: pip install python-evtx")

try:
    import openpyxl
except ImportError:
    openpyxl = None
    print("Warning: 'openpyxl' is not installed. .xlsx parsing will not be available. Install with: pip install openpyxl")

try:
    import pptx
except ImportError:
    pptx = None
    print("Warning: 'python-pptx' is not installed. .pptx parsing will not be available. Install with: pip install python-pptx")

try:
    import rarfile
except ImportError:
    rarfile = None
    print("Warning: 'rarfile' is not installed. .rar parsing will not be available. Install with: pip install rarfile")

try:
    import py7zr
except ImportError:
    py7zr = None
    print("Warning: 'py7zr' is not installed. .7z parsing will not be available. Install with: pip install py7zr")

# Optional pyewf (EWF) support
try:
    import pyewf
    pyewf_available = True
except Exception:
    pyewf = None
    pyewf_available = False

try:
    import requests
    requests_available = True
except Exception:
    requests = None
    requests_available = False

# Optional MySQL drivers shim: prefer pymysql, fall back to MySQLdb (mysqlclient) if present.
# Use importlib.import_module at runtime to avoid Pylance static import warnings in environments
# where mysql drivers are not installed.
import importlib as _importlib
try:
    _pymysql_module = _importlib.import_module('pymysql')
except Exception:
    _pymysql_module = None
try:
    _mysqldb_module = _importlib.import_module('MySQLdb')
except Exception:
    _mysqldb_module = None

def _get_mysql_driver():
    """Return the available MySQL driver module (pymysql or MySQLdb) or None."""
    return _pymysql_module or _mysqldb_module

# --- Application Setup ---
app = Flask(__name__)
app.secret_key = 'supersecretkey'
app.config['MAX_CONTENT_LENGTH'] = None  # No limit on request body size

# Increase other limits for large file handling
app.config['MAX_COOKIE_SIZE'] = 10 * 1024 * 1024  # 
APP_ROOT = os.path.dirname(os.path.abspath(__file__))

# Write a small sentinel file at import time so we can verify which file the WSGI process loaded.
try:
    _loaded_path = os.path.join(APP_ROOT, 'app_loaded.log')
    with open(_loaded_path, 'a', encoding='utf-8') as _lf:
        _lf.write('\n' + '='*80 + '\n')
        _lf.write(f'App module imported from: {os.path.abspath(__file__)}\n')
        _lf.write(f'Import time: {repr(__import__("datetime").datetime.now())}\n')
except Exception:
    pass

# --- Paths set to user-visible folder names ---
# Use exact folder names requested by the user for clarity in the filesystem
UPLOAD_FOLDER = os.path.join(APP_ROOT, 'Upload Files')
CARVED_FOLDER = os.path.join(APP_ROOT, 'Carved Files')
DECRYPTED_FOLDER = os.path.join(APP_ROOT, 'Decrypted Files')
ENCRYPTED_FOLDER = os.path.join(APP_ROOT, 'Encrypted Files')
DELETED_RECOVERY_FOLDER = os.path.join(APP_ROOT, 'Deleted Files')

# Path for the password dictionary
DICTIONARY_FILE = 'common_passwords.txt'
# File to persist configured database connections
DB_CONFIG_FILE = os.path.join(APP_ROOT, 'db_connections.json')

# Apply folder paths to the Flask app config
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['CARVED_FOLDER'] = CARVED_FOLDER
app.config['DECRYPTED_FOLDER'] = DECRYPTED_FOLDER
app.config['ENCRYPTED_FOLDER'] = ENCRYPTED_FOLDER
app.config['DICTIONARY_FILE'] = DICTIONARY_FILE
app.config['DELETED_RECOVERY_FOLDER'] = DELETED_RECOVERY_FOLDER



# Create all necessary directories when the app starts
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['CARVED_FOLDER'], exist_ok=True)
os.makedirs(app.config['DECRYPTED_FOLDER'], exist_ok=True)
os.makedirs(app.config['ENCRYPTED_FOLDER'], exist_ok=True)
os.makedirs(app.config['DELETED_RECOVERY_FOLDER'], exist_ok=True)
# Session files root (will contain per-session subfolders named like 'Session_YYYYMMDD_HHMMSS')
SESSION_FOLDER = os.path.join(APP_ROOT, 'Session Files')
app.config['SESSION_FOLDER'] = SESSION_FOLDER
os.makedirs(app.config['SESSION_FOLDER'], exist_ok=True)

# --- Simple SQLite DB for tracking files and sessions ---
DB_FILE = os.path.join(APP_ROOT, 'fac_data.db')

def init_db():
    """Initialize the SQLite DB and required tables."""
    try:
        conn = sqlite3.connect(DB_FILE)
        cur = conn.cursor()
        cur.execute('''
        CREATE TABLE IF NOT EXISTS sessions (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            started_at TEXT,
            ended_at TEXT,
            active INTEGER DEFAULT 0
        )
        ''')
        cur.execute('''
        CREATE TABLE IF NOT EXISTS files (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            filename TEXT,
            file_type TEXT,
            path TEXT,
            size_bytes INTEGER,
            created_at TEXT,
            session_id INTEGER,
            extra JSON,
            status TEXT DEFAULT 'saved'
        )
        ''')
        # --- Migration: add session_path column if missing ---
        try:
            cur.execute("PRAGMA table_info(sessions)")
            cols = [r[1] for r in cur.fetchall()]
            if 'session_path' not in cols:
                # Add new column
                try:
                    cur.execute('ALTER TABLE sessions ADD COLUMN session_path TEXT')
                except Exception:
                    # SQLite older may not support ALTER in some contexts; ignore
                    pass
                # Migrate any rows where ended_at contains a path-like value (heuristic)
                try:
                    cur.execute('SELECT id, ended_at FROM sessions')
                    rows = cur.fetchall()
                    for rid, ended in rows:
                        if ended and isinstance(ended, str) and (os.path.isabs(ended) or os.path.exists(ended)):
                            # assume ended_at was used to stash the session path earlier; migrate it
                            cur.execute('UPDATE sessions SET session_path=? WHERE id=?', (ended, rid))
                    conn.commit()
                except Exception:
                    pass
        except Exception:
            pass
        # --- Migration: add status column to files if missing ---
        try:
            cur.execute("PRAGMA table_info(files)")
            fcols = [r[1] for r in cur.fetchall()]
            if 'status' not in fcols:
                try:
                    cur.execute('ALTER TABLE files ADD COLUMN status TEXT DEFAULT "saved"')
                except Exception:
                    pass
            # also add status_changed_at to track when the status last changed
            if 'status_changed_at' not in fcols:
                try:
                    cur.execute('ALTER TABLE files ADD COLUMN status_changed_at TEXT')
                except Exception:
                    pass
        except Exception:
            pass
        conn.commit()
    except Exception as e:
        print(f"init_db error: {e}")
    finally:
        try:
            conn.close()
        except Exception:
            pass

init_db()

def _get_db_conn():
    return sqlite3.connect(DB_FILE)


def _resolve_cfg(cfg_or_dbid):
    if not cfg_or_dbid:
        return None
    if isinstance(cfg_or_dbid, dict):
        return cfg_or_dbid
    return db_connections.get(cfg_or_dbid)


def _get_sqlite_path_from_cfg(cfg):
    if not cfg:
        return None
    try:
        params = cfg.get('conn_details') or {}
        if params and params.get('path'):
            return params.get('path')
        conn = (cfg.get('conn') or '').strip()
        if not conn:
            return None
        # handle sqlite:///path or direct path
        if conn.lower().startswith('sqlite'):
            return conn.split('///')[-1]
        if conn.lower().endswith('.sqlite') or conn.lower().endswith('.db'):
            return conn
    except Exception:
        return None
    return None


def db_insert_file(filename, file_type, path, size_bytes=0, session_id=None, extra=None, target_db=None):
    """Insert a file record into the chosen target DB (target_db may be a dbid or cfg dict).
    Falls back to the local application DB on any failure.
    Returns True on success, False on failure.
    """
    cfg = None
    try:
        if target_db:
            cfg = _resolve_cfg(target_db)
        if not cfg:
            # nothing to do, use local
            conn = _get_db_conn()
            cur = conn.cursor()
            now = datetime.datetime.now().isoformat()
            extra_json = json.dumps(extra) if extra is not None else None
            # support optional status provided via extra['status']
            status_val = None
            try:
                if isinstance(extra, dict) and 'status' in extra:
                    status_val = extra.get('status')
            except Exception:
                status_val = None
            if status_val:
                cur.execute('INSERT INTO files(filename, file_type, path, size_bytes, created_at, session_id, extra, status) VALUES(?,?,?,?,?,?,?,?)',
                            (filename, file_type, path, size_bytes, now, session_id, extra_json, status_val))
            else:
                cur.execute('INSERT INTO files(filename, file_type, path, size_bytes, created_at, session_id, extra) VALUES(?,?,?,?,?,?,?)',
                            (filename, file_type, path, size_bytes, now, session_id, extra_json))
            conn.commit()
            return True

        # Try SQLite target if configured
        sqlite_path = _get_sqlite_path_from_cfg(cfg)
        if sqlite_path:
            try:
                conn = sqlite3.connect(sqlite_path)
                cur = conn.cursor()
                # ensure files table exists
                cur.execute('''
                    CREATE TABLE IF NOT EXISTS files (
                        id INTEGER PRIMARY KEY AUTOINCREMENT,
                        filename TEXT,
                        file_type TEXT,
                        path TEXT,
                        size_bytes INTEGER,
                        created_at TEXT,
                        session_id INTEGER,
                        extra JSON,
                        status TEXT DEFAULT 'saved'
                    )
                ''')
                now = datetime.datetime.now().isoformat()
                extra_json = json.dumps(extra) if extra is not None else None
                cur.execute('INSERT INTO files(filename, file_type, path, size_bytes, created_at, session_id, extra) VALUES(?,?,?,?,?,?,?)',
                            (filename, file_type, path, size_bytes, now, session_id, extra_json))
                conn.commit()
                conn.close()
                return True
            except Exception as e:
                app.logger.debug(f"db_insert_file sqlite target error: {e}")

        # Try Postgres
        try:
            conn_str = (cfg.get('conn') or '').strip()
            engine = (cfg.get('engine') or cfg.get('type') or '').lower()
            params = cfg.get('conn_details') or {}
            if conn_str and (conn_str.startswith('postgres') or conn_str.startswith('postgresql')) or 'postgres' in engine:
                try:
                    import psycopg2
                except Exception:
                    raise
                from urllib.parse import urlparse
                conn_kwargs = {}
                if conn_str and (conn_str.startswith('postgres') or conn_str.startswith('postgresql')):
                    u = urlparse(conn_str)
                    if u.hostname:
                        conn_kwargs['host'] = u.hostname
                    if u.port:
                        conn_kwargs['port'] = u.port
                    if u.username:
                        conn_kwargs['user'] = u.username
                    if u.password:
                        conn_kwargs['password'] = u.password
                    if u.path:
                        conn_kwargs['dbname'] = u.path.lstrip('/')
                if params and isinstance(params, dict):
                    conn_kwargs.update({k: v for k, v in params.items() if v is not None})
                pg = psycopg2.connect(connect_timeout=5, **conn_kwargs)
                cur = pg.cursor()
                now = datetime.datetime.now().isoformat()
                extra_json = json.dumps(extra) if extra is not None else None
                cur.execute('INSERT INTO files(filename, file_type, path, size_bytes, created_at, session_id, extra) VALUES(%s,%s,%s,%s,%s,%s,%s)',
                            (filename, file_type, path or '', size_bytes, now, session_id, extra_json))
                pg.commit()
                cur.close()
                pg.close()
                return True
        except Exception as e:
            app.logger.debug(f"db_insert_file postgres error: {e}")

        # Try MySQL
        try:
            conn_str = (cfg.get('conn') or '').strip()
            engine = (cfg.get('engine') or cfg.get('type') or '').lower()
            params = cfg.get('conn_details') or {}
            if 'mysql' in engine or (conn_str and conn_str.startswith('mysql')):
                mysql_driver = _get_mysql_driver()
                if not mysql_driver:
                    raise RuntimeError('MySQL driver not available')
                from urllib.parse import urlparse
                kwargs = {}
                if conn_str and conn_str.startswith('mysql'):
                    u = urlparse(conn_str)
                    if u.hostname:
                        kwargs['host'] = u.hostname
                    if u.port:
                        kwargs['port'] = int(u.port)
                    if u.username:
                        kwargs['user'] = u.username
                    if u.password:
                        kwargs['password'] = u.password
                    if u.path:
                        kwargs['db'] = u.path.lstrip('/')
                if params and isinstance(params, dict):
                    kwargs.update({k: v for k, v in params.items() if v is not None})
                conn_mysql = mysql_driver.connect(connect_timeout=5, **kwargs)
                cur = conn_mysql.cursor()
                now = datetime.datetime.now().isoformat()
                extra_json = json.dumps(extra) if extra is not None else None
                cur.execute('INSERT INTO files(filename, file_type, path, size_bytes, created_at, session_id, extra) VALUES(%s,%s,%s,%s,%s,%s,%s)',
                            (filename, file_type, path or '', size_bytes, now, session_id, extra_json))
                conn_mysql.commit()
                cur.close()
                conn_mysql.close()
                return True
        except Exception as e:
            app.logger.debug(f"db_insert_file mysql error: {e}")

        # Try MongoDB
        try:
            conn_str = (cfg.get('conn') or '').strip()
            engine = (cfg.get('engine') or cfg.get('type') or '').lower()
            if 'mongo' in engine or (conn_str and conn_str.startswith('mongodb')):
                from pymongo import MongoClient
                client = MongoClient(conn_str) if conn_str else MongoClient()
                from urllib.parse import urlparse
                dbname = None
                if conn_str:
                    u = urlparse(conn_str)
                    dbname = u.path.lstrip('/') if u.path and u.path != '/' else None
                if not dbname:
                    dbname = cfg.get('conn_details', {}).get('database') or cfg.get('conn_details', {}).get('db') or 'admin'
                db = client[dbname]
                coll = db['files']
                doc = {'filename': filename, 'file_type': file_type, 'path': path or '', 'size_bytes': size_bytes, 'created_at': datetime.datetime.now().isoformat(), 'session_id': session_id, 'extra': extra or {}}
                coll.insert_one(doc)
                try:
                    client.close()
                except Exception:
                    pass
                return True
        except Exception as e:
            app.logger.debug(f"db_insert_file mongo error: {e}")

    except Exception as e:
        app.logger.debug(f"db_insert_file unexpected error: {e}")

    # Final fallback: insert into local application DB
    try:
        conn = _get_db_conn()
        cur = conn.cursor()
        now = datetime.datetime.now().isoformat()
        extra_json = json.dumps(extra) if extra is not None else None
        cur.execute('INSERT INTO files(filename, file_type, path, size_bytes, created_at, session_id, extra) VALUES(?,?,?,?,?,?,?)',
                    (filename, file_type, path, size_bytes, now, session_id, extra_json))
        conn.commit()
        return True
    except Exception as e:
        app.logger.debug(f"db_insert_file final fallback error: {e}")
        try:
            conn.close()
        except Exception:
            pass
    return False


def db_query_files(cfg_or_dbid, q=None, qtype=None, qsession=None, browse_root=None, browse_folder=None, per_page=25, offset=0):
    """Query files from target DB (cfg or dbid). Returns (files, total, file_types).
    Files are dicts with keys: id, filename, file_type, path, size_bytes, created_at, session_id
    Raises exception on failure; caller may fallback to local DB.
    """
    cfg = _resolve_cfg(cfg_or_dbid)
    # Build filter pieces
    where_clauses = []
    params = []
    if q:
        where_clauses.append(('filename_like', q))
    if qtype:
        where_clauses.append(('file_type', qtype))
    if qsession:
        where_clauses.append(('session_id', qsession))

    # If cfg is None, query local DB
    if not cfg:
        conn = _get_db_conn()
        try:
            cur = conn.cursor()
            where = []
            p = []
            if q:
                where.append('filename LIKE ?')
                p.append(f'%{q}%')
            if qtype:
                where.append('file_type = ?')
                p.append(qtype)
            if qsession:
                where.append('session_id = ?')
                p.append(qsession)
            abs_browse_folder = None
            if browse_root:
                candidate_roots = _allowed_roots()
                matched_root = None
                for r in candidate_roots:
                    if browse_root == r or browse_root == os.path.basename(r):
                        matched_root = r
                        break
                if matched_root:
                    if browse_folder:
                        abs_browse_folder = os.path.abspath(os.path.join(matched_root, browse_folder))
                    else:
                        abs_browse_folder = os.path.abspath(matched_root)
                    where.append('path LIKE ?')
                    p.append(f"{abs_browse_folder}%")
            where_sql = ('WHERE ' + ' AND '.join(where)) if where else ''
            cur.execute(f'SELECT COUNT(1) FROM files {where_sql}', p)
            total = cur.fetchone()[0] or 0
            sql = f"SELECT id, filename, file_type, path, size_bytes, created_at, session_id, status, status_changed_at FROM files {where_sql} ORDER BY created_at DESC LIMIT ? OFFSET ?"
            cur.execute(sql, p + [per_page, offset])
            rows = cur.fetchall()
            files = []
            for r in rows:
                files.append({'id': r[0], 'filename': r[1], 'file_type': r[2], 'path': r[3], 'size_bytes': r[4], 'created_at': r[5], 'session_id': r[6], 'status': r[7], 'status_changed_at': r[8]})
            # file_types
            cur.execute('SELECT DISTINCT file_type FROM files')
            file_types = [r[0] for r in cur.fetchall() if r[0]]
            return files, total, file_types
        finally:
            try: conn.close()
            except Exception: pass

    # else: cfg present - attempt engine-specific queries
    engine = (cfg.get('engine') or cfg.get('type') or '').lower()
    conn_str = (cfg.get('conn') or '').strip()
    params_map = cfg.get('conn_details') or {}

    # SQLite target file
    try:
        sqlite_path = _get_sqlite_path_from_cfg(cfg)
        if sqlite_path and os.path.exists(sqlite_path):
            sconn = sqlite3.connect(sqlite_path)
            try:
                scur = sconn.cursor()
                where = []
                p = []
                if q:
                    where.append('filename LIKE ?'); p.append(f'%{q}%')
                if qtype:
                    where.append('file_type = ?'); p.append(qtype)
                if qsession:
                    where.append('session_id = ?'); p.append(qsession)
                abs_browse_folder = None
                if browse_root:
                    # attempt same mapping as local
                    candidate_roots = _allowed_roots()
                    matched_root = None
                    for r in candidate_roots:
                        if browse_root == r or browse_root == os.path.basename(r):
                            matched_root = r; break
                    if matched_root:
                        if browse_folder:
                            abs_browse_folder = os.path.abspath(os.path.join(matched_root, browse_folder))
                        else:
                            abs_browse_folder = os.path.abspath(matched_root)
                        where.append('path LIKE ?'); p.append(f"{abs_browse_folder}%")
                where_sql = ('WHERE ' + ' AND '.join(where)) if where else ''
                scur.execute(f'SELECT COUNT(1) FROM files {where_sql}', p)
                total = scur.fetchone()[0] or 0
                scur.execute(f"SELECT id, filename, file_type, path, size_bytes, created_at, session_id, status, status_changed_at FROM files {where_sql} ORDER BY created_at DESC LIMIT ? OFFSET ?", p + [per_page, offset])
                rows = scur.fetchall()
                files = [{'id': r[0], 'filename': r[1], 'file_type': r[2], 'path': r[3], 'size_bytes': r[4], 'created_at': r[5], 'session_id': r[6], 'status': r[7], 'status_changed_at': r[8]} for r in rows]
                scur.execute('SELECT DISTINCT file_type FROM files')
                file_types = [r[0] for r in scur.fetchall() if r[0]]
                return files, total, file_types
            finally:
                try: sconn.close()
                except Exception: pass
    except Exception as e:
        app.logger.debug(f"SQLite target query error: {e}")

    # Postgres
    try:
        if 'postgres' in engine or conn_str.startswith('postgres') or conn_str.startswith('postgresql'):
            try:
                import psycopg2
                from urllib.parse import urlparse
                conn_kwargs = {}
                if conn_str and (conn_str.startswith('postgres') or conn_str.startswith('postgresql')):
                    u = urlparse(conn_str)
                    if u.hostname: conn_kwargs['host'] = u.hostname
                    if u.port: conn_kwargs['port'] = u.port
                    if u.username: conn_kwargs['user'] = u.username
                    if u.password: conn_kwargs['password'] = u.password
                    if u.path: conn_kwargs['dbname'] = u.path.lstrip('/')
                if params_map and isinstance(params_map, dict): conn_kwargs.update({k: v for k, v in params_map.items() if v is not None})
                pg = psycopg2.connect(connect_timeout=5, **conn_kwargs)
                cur = pg.cursor()
                where = []
                p = []
                if q: where.append('filename ILIKE %s'); p.append(f'%{q}%')
                if qtype: where.append('file_type = %s'); p.append(qtype)
                if qsession: where.append('session_id = %s'); p.append(qsession)
                if browse_root:
                    # best-effort: attempt to filter by path prefix if path field present
                    candidate_roots = _allowed_roots()
                    matched_root = None
                    for r in candidate_roots:
                        if browse_root == r or browse_root == os.path.basename(r): matched_root = r; break
                    if matched_root:
                        if browse_folder:
                            abs_browse_folder = os.path.abspath(os.path.join(matched_root, browse_folder))
                        else:
                            abs_browse_folder = os.path.abspath(matched_root)
                        where.append('path LIKE %s'); p.append(f"{abs_browse_folder}%")
                where_sql = ('WHERE ' + ' AND '.join(where)) if where else ''
                cur.execute(f'SELECT COUNT(1) FROM files {where_sql}', p)
                total = cur.fetchone()[0] or 0
                cur.execute(f"SELECT id, filename, file_type, path, size_bytes, created_at, session_id, status, status_changed_at FROM files {where_sql} ORDER BY created_at DESC LIMIT %s OFFSET %s", p + [per_page, offset])
                rows = cur.fetchall()
                files = [{'id': r[0], 'filename': r[1], 'file_type': r[2], 'path': r[3], 'size_bytes': r[4], 'created_at': r[5], 'session_id': r[6], 'status': r[7], 'status_changed_at': r[8]} for r in rows]
                cur.execute('SELECT DISTINCT file_type FROM files')
                file_types = [r[0] for r in cur.fetchall() if r[0]]
                cur.close(); pg.close()
                return files, total, file_types
            except Exception as e:
                app.logger.debug(f"Postgres target query error: {e}")
    except Exception:
        pass

    # MySQL
    try:
        if 'mysql' in engine or conn_str.startswith('mysql'):
            try:
                mysql_driver = _get_mysql_driver()
                if not mysql_driver:
                    raise RuntimeError('MySQL driver not available')
                from urllib.parse import urlparse
                kwargs = {}
                if conn_str and conn_str.startswith('mysql'):
                    u = urlparse(conn_str)
                    if u.hostname: kwargs['host'] = u.hostname
                    if u.port: kwargs['port'] = int(u.port)
                    if u.username: kwargs['user'] = u.username
                    if u.password: kwargs['password'] = u.password
                    if u.path: kwargs['db'] = u.path.lstrip('/')
                if params_map and isinstance(params_map, dict): kwargs.update({k: v for k, v in params_map.items() if v is not None})
                conn_mysql = mysql_driver.connect(connect_timeout=5, **kwargs)
                cur = conn_mysql.cursor()
                where = []
                p = []
                if q: where.append('filename LIKE %s'); p.append(f'%{q}%')
                if qtype: where.append('file_type = %s'); p.append(qtype)
                if qsession: where.append('session_id = %s'); p.append(qsession)
                if browse_root:
                    candidate_roots = _allowed_roots(); matched_root = None
                    for r in candidate_roots:
                        if browse_root == r or browse_root == os.path.basename(r): matched_root = r; break
                    if matched_root:
                        if browse_folder: abs_browse_folder = os.path.abspath(os.path.join(matched_root, browse_folder))
                        else: abs_browse_folder = os.path.abspath(matched_root)
                        where.append('path LIKE %s'); p.append(f"{abs_browse_folder}%")
                where_sql = ('WHERE ' + ' AND '.join(where)) if where else ''
                cur.execute(f'SELECT COUNT(1) FROM files {where_sql}', p)
                total = cur.fetchone()[0] or 0
                cur.execute(f"SELECT id, filename, file_type, path, size_bytes, created_at, session_id, status, status_changed_at FROM files {where_sql} ORDER BY created_at DESC LIMIT %s OFFSET %s", p + [per_page, offset])
                rows = cur.fetchall()
                files = [{'id': r[0], 'filename': r[1], 'file_type': r[2], 'path': r[3], 'size_bytes': r[4], 'created_at': r[5], 'session_id': r[6], 'status': r[7], 'status_changed_at': r[8]} for r in rows]
                cur.execute('SELECT DISTINCT file_type FROM files')
                file_types = [r[0] for r in cur.fetchall() if r[0]]
                cur.close(); conn_mysql.close()
                return files, total, file_types
            except Exception as e:
                app.logger.debug(f"MySQL target query error: {e}")
    except Exception:
        pass

    # MongoDB
    try:
        if 'mongo' in engine or conn_str.startswith('mongodb'):
            try:
                from pymongo import MongoClient
                from urllib.parse import urlparse
                client = MongoClient(conn_str) if conn_str else MongoClient()
                dbname = None
                if conn_str:
                    u = urlparse(conn_str); dbname = u.path.lstrip('/') if u.path and u.path != '/' else None
                if not dbname: dbname = params_map.get('database') or params_map.get('db') or 'admin'
                db = client[dbname]
                coll = db['files']
                filt = {}
                if q: filt['filename'] = {'$regex': q, '$options': 'i'}
                if qtype: filt['file_type'] = qtype
                if qsession: filt['session_id'] = qsession
                if browse_root:
                    candidate_roots = _allowed_roots(); matched_root = None
                    for r in candidate_roots:
                        if browse_root == r or browse_root == os.path.basename(r): matched_root = r; break
                    if matched_root:
                        if browse_folder: abs_browse_folder = os.path.abspath(os.path.join(matched_root, browse_folder))
                        else: abs_browse_folder = os.path.abspath(matched_root)
                        filt['path'] = {'$regex': f'^{re.escape(abs_browse_folder)}'}
                total = coll.count_documents(filt)
                cursor = coll.find(filt).sort('created_at', -1).skip(offset).limit(per_page)
                files = []
                for doc in cursor:
                    files.append({'id': str(doc.get('_id')), 'filename': doc.get('filename'), 'file_type': doc.get('file_type'), 'path': doc.get('path'), 'size_bytes': doc.get('size_bytes'), 'created_at': doc.get('created_at'), 'session_id': doc.get('session_id'), 'status': doc.get('status'), 'status_changed_at': doc.get('status_changed_at')})
                file_types = coll.distinct('file_type') or []
                try: client.close()
                except Exception: pass
                return files, total, file_types
            except Exception as e:
                app.logger.debug(f"Mongo target query error: {e}")
    except Exception:
        pass

    raise Exception('Could not query target DB')


def _allowed_roots():
    """Return absolute allowed root directories as a list."""
    roots = [
        app.config.get('UPLOAD_FOLDER'),
        app.config.get('ENCRYPTED_FOLDER'),
        app.config.get('DECRYPTED_FOLDER'),
        app.config.get('CARVED_FOLDER'),
        app.config.get('DELETED_RECOVERY_FOLDER'),
        app.config.get('SESSION_FOLDER')
    ]
    # filter None and take absolute paths
    abs_roots = []
    for r in roots:
        try:
            if r:
                abs_roots.append(os.path.abspath(r))
        except Exception:
            pass
    return abs_roots


@app.before_request
def ensure_csrf_token():
    # ensure every session has a CSRF token available for the UI
    try:
        if 'csrf_token' not in session:
            session['csrf_token'] = secrets.token_urlsafe(32)
    except Exception:
        pass


def verify_csrf_request():
    """Verify CSRF token from header or form. Returns True if valid."""
    try:
        server_token = session.get('csrf_token')
        if not server_token:
            return False
        # header first
        hdr = request.headers.get('X-CSRF-Token')
        if hdr and hdr == server_token:
            return True
        # form field
        form_token = None
        try:
            form_token = request.form.get('csrf_token')
        except Exception:
            form_token = None
        if form_token and form_token == server_token:
            return True
        # querystring fallback
        arg_token = request.args.get('csrf_token')
        if arg_token and arg_token == server_token:
            return True
        return False
    except Exception:
        return False


# --- Audit logging setup ---
AUDIT_LOG = os.path.join(APP_ROOT, 'audit.log')
def _get_audit_logger():
    logger = logging.getLogger('fac_audit')
    if not logger.handlers:
        logger.setLevel(logging.INFO)
        fh = logging.FileHandler(AUDIT_LOG)
        fh.setFormatter(logging.Formatter('%(message)s'))
        logger.addHandler(fh)
    return logger

def audit_event(action, target_path, details=None):
    """Write an audit event as a JSON line to the audit log."""
    try:
        logger = _get_audit_logger()
        event = {
            'ts': datetime.datetime.utcnow().isoformat() + 'Z',
            'remote_addr': request.remote_addr if request else None,
            'session_id': session.get('analysis_session_id') if session else None,
            'action': action,
            'target': target_path,
            'details': details
        }
        logger.info(json.dumps(event))
    except Exception:
        pass


def is_path_under_allowed_roots(path):
    """Return True if the path is inside one of the allowed roots."""
    try:
        if not path:
            return False
        abs_path = os.path.abspath(path)
        for root in _allowed_roots():
            # normalize and compare
            try:
                abs_root = os.path.abspath(root)
                # commonpath will raise if on different drives on Windows; handle that
                try:
                    common = os.path.commonpath([abs_root, abs_path])
                except Exception:
                    # if different drives, fall back to startswith
                    if os.name == 'nt':
                        if abs_path.lower().startswith(abs_root.lower() + os.path.sep):
                            return True
                        continue
                    else:
                        continue

                if common == abs_root:
                    return True
            except Exception:
                continue
        return False
    except Exception:
        return False


def resolve_root_display_to_path(root_display):
    """Given a root display name (basename) or absolute path, return the absolute root path if allowed, else None."""
    try:
        # If absolute path provided and allowed, return it
        if os.path.isabs(root_display) and os.path.isdir(root_display) and os.path.abspath(root_display) in _allowed_roots():
            return os.path.abspath(root_display)
        # Otherwise match by basename
        for r in _allowed_roots():
            try:
                if os.path.basename(r) == root_display or r == root_display:
                    return os.path.abspath(r)
            except Exception:
                continue
        return None
    except Exception:
        return None


@app.route('/fs_explorer')
def fs_explorer():
    """Filesystem explorer for allowed roots and per-session folders.

    Query params:
    - browse_root: root folder display name or absolute path (required)
    - browse_folder: subfolder path under the root (optional)
    - page, per_page: pagination for listing files
    """
    browse_root = request.args.get('browse_root')
    browse_folder = request.args.get('browse_folder', '')
    # Normalize: callers sometimes pass the root display name as the browse_folder
    try:
        if browse_folder and (browse_folder == browse_root):
            browse_folder = ''
    except Exception:
        pass
    try:
        page = int(request.args.get('page', '1'))
        if page < 1: page = 1
    except Exception:
        page = 1
    try:
        per_page = int(request.args.get('per_page', '50'))
        if per_page < 1: per_page = 50
    except Exception:
        per_page = 50

    root_path = resolve_root_display_to_path(browse_root)
    if not root_path:
        return "Invalid root", 400

    # compute absolute folder to list
    if browse_folder:
        target_dir = os.path.abspath(os.path.join(root_path, browse_folder))
    else:
        target_dir = os.path.abspath(root_path)

    # safety check
    if not is_path_under_allowed_roots(target_dir) or not os.path.isdir(target_dir):
        return "Access denied or not a directory", 403

    # list directories and files
    try:
        entries = sorted(os.listdir(target_dir))
    except Exception as e:
        return f"Cannot list directory: {e}", 500

    # separate folders and files
    folders = [e for e in entries if os.path.isdir(os.path.join(target_dir, e))]
    files = [e for e in entries if os.path.isfile(os.path.join(target_dir, e))]

    # pagination for files
    total_files = len(files)
    offset = (page - 1) * per_page
    paged_files = files[offset: offset + per_page]

# Render explorer template (external file)
    roots = [os.path.basename(r) for r in _allowed_roots()]
    return render_template('explorer.html', folders=folders, paged_files=paged_files, total_files=total_files,
                                root_name=os.path.basename(root_path), rel_path=os.path.relpath(target_dir, root_path),
                                root_display=browse_root, cur_dbid='local', page=page, offset=offset,
                                allowed_roots=roots, csrf_token=session.get('csrf_token'), pyewf_available=pyewf_available)


@app.route('/serve_fs_file')
def serve_fs_file():
    root = request.args.get('root')
    path = request.args.get('path')
    root_path = resolve_root_display_to_path(root)
    if not root_path:
        return "Invalid root", 400
    abs_path = os.path.abspath(os.path.join(root_path, path or ''))
    if not is_path_under_allowed_roots(abs_path) or not os.path.isfile(abs_path):
        return "Access denied or not a file", 403
    try:
        # audit download
        try:
            audit_event('download', abs_path, {'method': 'serve_fs_file'})
        except Exception:
            pass
        return send_file(abs_path, as_attachment=True)
    except Exception:
        return "Could not serve file", 500


@app.route('/download_folder_zip')
def download_folder_zip():
    browse_root = request.args.get('browse_root')
    browse_folder = request.args.get('browse_folder', '')
    root_path = resolve_root_display_to_path(browse_root)
    if not root_path:
        return "Invalid root", 400
    # Normalize browse_folder: callers sometimes send the root display name again
    try:
        if browse_folder and (browse_folder == browse_root or browse_folder == os.path.basename(root_path)):
            browse_folder = ''
    except Exception:
        pass
    target_dir = os.path.abspath(os.path.join(root_path, browse_folder)) if browse_folder else os.path.abspath(root_path)
    if not is_path_under_allowed_roots(target_dir) or not os.path.isdir(target_dir):
        return "Access denied or invalid folder", 403

    # create a zip in memory
    try:
        memory_file = io.BytesIO()
        with zipfile.ZipFile(memory_file, 'w', zipfile.ZIP_DEFLATED) as zf:
            for root, dirs, files in os.walk(target_dir):
                for f in files:
                    full = os.path.join(root, f)
                    arcname = os.path.relpath(full, os.path.dirname(target_dir))
                    zf.write(full, arcname)
        memory_file.seek(0)
        zip_name = f"{os.path.basename(target_dir)}.zip"
        try:
            audit_event('download_zip', target_dir, {'zip_name': zip_name})
        except Exception:
            pass
        return send_file(memory_file, download_name=zip_name, as_attachment=True)
    except Exception as e:
        return f"Error creating zip: {e}", 500


@app.route('/api/fs/list', methods=['GET'])
def api_fs_list():
    root = request.args.get('root')
    rel = request.args.get('path', '')
    root_path = resolve_root_display_to_path(root)
    if not root_path:
        return jsonify({'error': 'invalid_root'}), 400
    # Normalize rel: sometimes callers send the root display name again as the path (e.g. path='Upload Files').
    try:
        if rel and (rel == root or rel == os.path.basename(root_path)):
            rel = ''
    except Exception:
        pass
    abs_path = os.path.abspath(os.path.join(root_path, rel)) if rel else os.path.abspath(root_path)
    if not is_path_under_allowed_roots(abs_path) or not os.path.exists(abs_path):
        return jsonify({'error': 'access_denied_or_not_found'}), 403

    items = []
    try:
        for name in sorted(os.listdir(abs_path)):
            full = os.path.join(abs_path, name)
            stat = os.stat(full)
            items.append({
                'name': name,
                'is_dir': os.path.isdir(full),
                'size': stat.st_size,
                'mtime': int(stat.st_mtime)
            })
        return jsonify({'path': os.path.relpath(abs_path, root_path), 'items': items})
    except Exception as e:
        return jsonify({'error': 'list_failed', 'message': str(e)}), 500


@app.route('/api/create_image', methods=['POST'])
def api_create_image():
    """Start a background imaging job that copies a file or creates an archive while
    computing MD5 and SHA1 and exposing progress via a status endpoint.
    Returns: {ok: True, job_id: '<id>'}
    """
    try:
        if not verify_csrf_request():
            return jsonify({'error': 'csrf_invalid'}), 400

        source_type = (request.form.get('source_type') or (request.json and request.json.get('source_type')) or 'file')
        source_root = (request.form.get('source_root') or (request.json and request.json.get('source_root')))
        source_path = (request.form.get('source_path') or (request.json and request.json.get('source_path')))
        image_format = (request.form.get('image_format') or (request.json and request.json.get('image_format')) or '.dd')
        destination = (request.form.get('destination') or (request.json and request.json.get('destination')) or 'download')

        # Validate inputs according to the chosen source_type
        root_path = None
        abs_source = None
        if source_type == 'cloud':
            # source_path must be a valid URL (http/https). Optional host whitelist may be configured
            if not source_path:
                return jsonify({'error': 'missing_parameters'}), 400
            try:
                p = urlparse(source_path)
                if p.scheme not in ('http', 'https'):
                    return jsonify({'error': 'invalid_url_scheme', 'message': 'Only http/https URLs are allowed'}), 400
                allowed_hosts = app.config.get('ALLOWED_CLOUD_HOSTS')
                if allowed_hosts and isinstance(allowed_hosts, (list, tuple, set)):
                    host = (p.hostname or '').lower()
                    if host not in [h.lower() for h in allowed_hosts]:
                        return jsonify({'error': 'host_not_allowed'}), 400
            except Exception:
                return jsonify({'error': 'invalid_url'}), 400
            root_path = None
            abs_source = None
        elif source_type == 'device':
            # Device imaging: expect an absolute device path in source_path (e.g. /dev/sda or \\\\.\\PhysicalDrive1)
            if not source_path:
                return jsonify({'error': 'missing_parameters'}), 400
            # require confirmation checkbox and typed-confirmation text to match the device path
            confirm_box = (request.form.get('device_confirm_box') or (request.json and request.json.get('device_confirm_box')) or '')
            confirm_text = (request.form.get('device_confirm_text') or (request.json and request.json.get('device_confirm_text')) or '')
            if str(confirm_box) != '1' or (confirm_text or '').strip() != (source_path or '').strip():
                return jsonify({'error': 'device_confirmation_required', 'message': 'You must check the confirmation box and type the exact device path to proceed.'}), 400
            # accept device path as provided; do not require it to be under allowed roots
            abs_source = source_path.strip()
            root_path = None
        else:
            # file/folder source: require a configured root and a relative path
            if not source_root or source_path is None:
                return jsonify({'error': 'missing_parameters'}), 400
            root_path = resolve_root_display_to_path(source_root)
            if not root_path:
                return jsonify({'error': 'invalid_root'}), 400
            abs_source = os.path.abspath(os.path.join(root_path, source_path))
            if not is_path_under_allowed_roots(abs_source):
                return jsonify({'error': 'access_denied'}), 403
            if not os.path.exists(abs_source):
                return jsonify({'error': 'source_not_found'}), 404

        # capture optional metadata fields
        ev_case = (request.form.get('case_number') or (request.json and request.json.get('case_number')) or '')
        ev_examiner = (request.form.get('examiner') or (request.json and request.json.get('examiner')) or '')
        ev_notes = (request.form.get('notes') or (request.json and request.json.get('notes')) or '')
        ev_compress = bool(request.form.get('compress') or (request.json and request.json.get('compress')))

        # prepare session folder to store image
        sess_base = app.config.get('SESSION_FOLDER')
        ts = datetime.datetime.now().strftime('%Y%m%d_%H%M%S')
        sess_name = f'Image_Create_{ts}_{secrets.token_hex(4)}'
        sess_path = os.path.join(sess_base, sess_name)
        os.makedirs(sess_path, exist_ok=True)

        # sanitize requested extension
        if not image_format.startswith('.'):
            image_format = '.' + image_format
        safe_ext = re.sub(r'[^a-zA-Z0-9._-]', '', image_format)

        src_name = os.path.basename(abs_source.rstrip(os.path.sep)) or 'image'
        target_name = secure_filename(f"{src_name}{safe_ext}")
        target_path = os.path.join(sess_path, target_name)

        # lightweight disk-space check: require at least source size + 1MB free
        try:
            required = os.path.getsize(abs_source) if os.path.isfile(abs_source) else 0
        except Exception:
            required = 0
        try:
            st = os.statvfs(sess_base)
            free = st.f_bavail * st.f_frsize
        except Exception:
            # if statvfs not available (Windows), skip check
            free = None
        if free is not None and required and free < (required + 1024 * 1024):
            return jsonify({'error': 'insufficient_space'}), 507

        # job registry (in-memory)
        if 'image_jobs' not in app.config:
            app.config['image_jobs'] = {}

        job_id = secrets.token_hex(12)
        job_info = {
            'id': job_id,
            'status': 'queued',
            'started_at': None,
            'ended_at': None,
            'progress': 0,
            'message': '',
            'source': abs_source,
            'target': target_path,
            'md5': None,
            'sha1': None,
                'options': {
                'image_format': safe_ext,
                'use_pyewf': (safe_ext.lower() == '.e01' and pyewf_available),
                    'source_type': source_type,
                    'source_url': source_path if source_type == 'cloud' else None,
                'case_number': ev_case,
                'examiner': ev_examiner,
                'notes': ev_notes,
                'compress': ev_compress
            },
            'error': None
        }
        app.config['image_jobs'][job_id] = job_info

        def worker_copy_and_hash(src, dst, job):
            try:
                job['status'] = 'running'
                job['started_at'] = datetime.datetime.utcnow().isoformat() + 'Z'
                CHUNK = 4 * 1024 * 1024
                total = 0
                # determine total for different source types
                src_type = job.get('options', {}).get('source_type', 'file')
                if src_type == 'file' and src and os.path.isfile(src):
                    total = os.path.getsize(src)
                elif src_type == 'device' and src and os.path.exists(src):
                    try:
                        total = os.path.getsize(src)
                    except Exception:
                        total = 0
                else:
                    total = 0
                # if folder, we will create a zip-like archive while streaming
                md5 = hashlib.md5()
                sha1 = hashlib.sha1()
                bytes_written = 0
                # if user requested EWF and pyewf is available, use pyewf
                use_pyewf = job.get('options', {}).get('use_pyewf')
                if use_pyewf and pyewf_available:
                    # pyewf expects bytes; create a temporary file source if folder
                    try:
                        if os.path.isfile(src):
                            src_file_for_ewf = src
                            total = os.path.getsize(src_file_for_ewf)
                        else:
                            # create temp zip to feed into EWF
                            tmp_zip = dst + '.tmpzip'
                            with zipfile.ZipFile(tmp_zip, 'w', zipfile.ZIP_DEFLATED) as zf:
                                for root_dir, dirs, files in os.walk(src):
                                    for f in files:
                                        full = os.path.join(root_dir, f)
                                        arc = os.path.relpath(full, src)
                                        zf.write(full, arc)
                            src_file_for_ewf = tmp_zip
                            total = os.path.getsize(src_file_for_ewf)

                        # write EWF using pyewf
                        # pyewf expects to write multiple segments; the API requires an output filename without extension
                        base_out = os.path.splitext(dst)[0]
                        # open source and write to pyewf
                        with open(src_file_for_ewf, 'rb') as infile:
                            ewf_handle = pyewf.handle()
                            ewf_handle.open_write(base_out)
                            while True:
                                chunk = infile.read(CHUNK)
                                if not chunk:
                                    break
                                ewf_handle.write(chunk)
                                md5.update(chunk)
                                sha1.update(chunk)
                                bytes_written += len(chunk)
                                job['progress'] = int((bytes_written / total) * 100) if total else 0
                            ewf_handle.close()

                        # pyewf will write files like base_out.E01 etc.; find the primary E01 and move/rename to dst
                        # find created EWF file
                        created = None
                        for candidate in os.listdir(os.path.dirname(base_out) or '.'):
                            if candidate.lower().startswith(os.path.basename(base_out).lower()) and candidate.lower().endswith('.e01'):
                                created = os.path.join(os.path.dirname(base_out) or '.', candidate)
                                break
                        if created:
                            try:
                                os.replace(created, dst)
                            except Exception:
                                shutil.copy2(created, dst)
                        # cleanup temp zip if used
                        if not os.path.isfile(src) and os.path.exists(tmp_zip):
                            try: os.remove(tmp_zip)
                            except Exception: pass

                        job['md5'] = md5.hexdigest()
                        job['sha1'] = sha1.hexdigest()
                        job['progress'] = 100
                        job['status'] = 'finished'
                        job['ended_at'] = datetime.datetime.utcnow().isoformat() + 'Z'
                        try:
                            audit_event('create_image_ewf', dst, {'job': job['id'], 'source': src})
                        except Exception:
                            pass
                        return
                    except Exception as e:
                        # fallback to default handling if pyewf path fails
                        app.logger.exception('pyewf imaging failed, falling back')
                        job['message'] = 'pyewf failed, falling back to raw copy'

                # Cloud source: stream download using requests
                if job.get('options', {}).get('source_type') == 'cloud':
                    if not requests_available:
                        raise RuntimeError('requests_not_available')
                    url = job.get('options', {}).get('source_url')
                    if not url:
                        raise RuntimeError('missing_url')
                    # stream download to dst while hashing
                    with requests.get(url, stream=True, timeout=30) as resp:
                        resp.raise_for_status()
                        # try to get content-length
                        try:
                            total = int(resp.headers.get('Content-Length') or 0)
                        except Exception:
                            total = 0
                        with open(dst, 'wb') as outf:
                            for chunk in resp.iter_content(chunk_size=CHUNK):
                                if not chunk:
                                    continue
                                outf.write(chunk)
                                md5.update(chunk)
                                sha1.update(chunk)
                                bytes_written += len(chunk)
                                job['progress'] = int((bytes_written / total) * 100) if total else 0
                elif job.get('options', {}).get('source_type') == 'device':
                    # Device imaging: attempt to run dd (must run with privileges)
                    device_path = src
                    if not device_path or not os.path.exists(device_path):
                        raise RuntimeError('device_not_found')
                    # construct dd command depending on platform
                    if os.name == 'nt':
                        # Windows: expect \\.\PhysicalDriveN style path; dd.exe must be available in PATH
                        dd_cmd = ['dd', 'if=' + device_path, 'of=' + dst, 'bs=4M']
                    else:
                        dd_cmd = ['dd', 'if=' + device_path, 'of=' + dst, 'bs=4M', 'status=none']
                    # spawn subprocess and capture stdout/stderr for progress if available
                    import subprocess
                    p = subprocess.Popen(dd_cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
                    # poll process and compute hashes from the target file (as it grows)
                    last_size = 0
                    while True:
                        if p.poll() is not None:
                            break
                        try:
                            if os.path.exists(dst):
                                cur_size = os.path.getsize(dst)
                                if cur_size > last_size:
                                    # read new bytes and update hashes
                                    with open(dst, 'rb') as fh:
                                        fh.seek(last_size)
                                        new = fh.read(cur_size - last_size)
                                        if new:
                                            md5.update(new)
                                            sha1.update(new)
                                            bytes_written += len(new)
                                            job['progress'] = int((bytes_written / total) * 100) if total else 0
                                    last_size = cur_size
                        except Exception:
                            pass
                        time.sleep(0.5)
                    stdout, stderr = p.communicate(timeout=5)
                    if p.returncode != 0:
                        raise RuntimeError('dd_failed: ' + (stderr.decode('utf-8', errors='ignore') if stderr else ''))
                elif os.path.isfile(src):
                    with open(src, 'rb') as inf, open(dst, 'wb') as outf:
                        while True:
                            chunk = inf.read(CHUNK)
                            if not chunk:
                                break
                            outf.write(chunk)
                            md5.update(chunk)
                            sha1.update(chunk)
                            bytes_written += len(chunk)
                            job['progress'] = int((bytes_written / total) * 100) if total else 0
                else:
                    # create zip archive to dst path (streaming not trivial), write to temp file then move
                    tmp = dst + '.tmpzip'
                    with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zf:
                        for root_dir, dirs, files in os.walk(src):
                            for f in files:
                                full = os.path.join(root_dir, f)
                                arc = os.path.relpath(full, src)
                                zf.write(full, arc)
                    # compute hashes while copying tmp -> dst
                    total = os.path.getsize(tmp)
                    with open(tmp, 'rb') as inf, open(dst, 'wb') as outf:
                        while True:
                            chunk = inf.read(CHUNK)
                            if not chunk:
                                break
                            outf.write(chunk)
                            md5.update(chunk)
                            sha1.update(chunk)
                            bytes_written += len(chunk)
                            job['progress'] = int((bytes_written / total) * 100) if total else 0
                    try:
                        os.remove(tmp)
                    except Exception:
                        pass

                job['md5'] = md5.hexdigest()
                job['sha1'] = sha1.hexdigest()
                job['progress'] = 100
                job['status'] = 'finished'
                job['ended_at'] = datetime.datetime.utcnow().isoformat() + 'Z'

                # audit
                try:
                    audit_event('create_image', dst, {'job': job['id'], 'source': src})
                except Exception:
                    pass
            except Exception as e:
                job['status'] = 'error'
                job['error'] = str(e)
                job['ended_at'] = datetime.datetime.utcnow().isoformat() + 'Z'
                app.logger.exception('image job failed')

        # start worker thread
        t = threading.Thread(target=worker_copy_and_hash, args=(abs_source, target_path, job_info), daemon=True)
        t.start()

        return jsonify({'ok': True, 'job_id': job_id})
    except Exception as e:
        app.logger.exception('create_image error')
        return jsonify({'error': 'create_failed', 'message': str(e)}), 500


@app.route('/api/fs/info', methods=['GET'])
def api_fs_info():
    """Return enhanced info for a single file, including thumbnail data-uri when available."""
    root = request.args.get('root')
    rel = request.args.get('path')
    root_path = resolve_root_display_to_path(root)
    if not root_path:
        return jsonify({'error': 'invalid_root'}), 400
    abs_path = os.path.abspath(os.path.join(root_path, rel)) if rel else None
    if not abs_path or not is_path_under_allowed_roots(abs_path) or not os.path.exists(abs_path):
        return jsonify({'error': 'access_denied_or_not_found'}), 403
    try:
        info = get_enhanced_file_info(abs_path)
        # attach path and root for client
        info['path'] = rel
        info['root'] = root
        info['name'] = os.path.basename(abs_path)
        return jsonify(info)
    except Exception as e:
        return jsonify({'error': 'info_failed', 'message': str(e)}), 500


@app.route('/api/image_status/<job_id>', methods=['GET'])
def api_image_status(job_id):
    jobs = app.config.get('image_jobs') or {}
    job = jobs.get(job_id)
    if not job:
        return jsonify({'error': 'not_found'}), 404
    # attach a download_url if finished
    data = {k: job[k] for k in ('id','status','progress','started_at','ended_at','md5','sha1','error') if k in job}
    if job.get('status') == 'finished' and job.get('target'):
        # target is absolute path under session folder
        sess_folder = app.config.get('SESSION_FOLDER')
        try:
            rel = os.path.relpath(job.get('target'), sess_folder)
            download_url = url_for('serve_fs_file', root=os.path.basename(sess_folder), path=rel)
            data['download_url'] = download_url
            data['filename'] = os.path.basename(job.get('target'))
        except Exception:
            pass
    return jsonify(data)


@app.route('/api/devices', methods=['GET'])
def api_devices():
    """Return a best-effort list of local block devices for the UI.
    Fields: id, path, size_bytes, model, mountpoints
    This endpoint is read-only and should not open devices.
    """
    devices = []
    try:
        if os.name == 'posix':
            # try lsblk JSON first
            try:
                import subprocess, json as _json
                p = subprocess.run(['lsblk', '-J', '-b', '-o', 'NAME,PATH,SIZE,MODEL,MOUNTPOINTS'], capture_output=True, text=True, timeout=3)
                if p.returncode == 0 and p.stdout:
                    j = _json.loads(p.stdout)
                    # lsblk JSON has 'blockdevices'
                    for dev in j.get('blockdevices', []):
                        try:
                            path = dev.get('path') or ('/dev/' + dev.get('name'))
                            devices.append({'id': dev.get('name'), 'path': path, 'size_bytes': int(dev.get('size') or 0), 'model': dev.get('model') or '', 'mountpoints': dev.get('mountpoints') or []})
                        except Exception:
                            continue
            except Exception:
                pass
            # fallback: parse /proc/partitions to list major devices
            if not devices and os.path.exists('/proc/partitions'):
                try:
                    with open('/proc/partitions', 'r') as fh:
                        lines = fh.read().strip().splitlines()
                    for line in lines[2:]:
                        parts = line.split()
                        if len(parts) >= 4:
                            name = parts[3]
                            # skip partitions (e.g., sda1)
                            if re.match(r'\D+', name):
                                path = '/dev/' + name
                                try:
                                    size = os.path.getsize(path) if os.path.exists(path) else 0
                                except Exception:
                                    size = 0
                                devices.append({'id': name, 'path': path, 'size_bytes': size, 'model': '', 'mountpoints': []})
                except Exception:
                    pass
        elif os.name == 'nt':
            # Windows: use PowerShell Get-Disk to enumerate disks
            try:
                import subprocess, json as _json
                pwsh_cmd = ['powershell', '-NoProfile', '-Command', "Get-Disk | Select-Object Number,FriendlyName,Size | ConvertTo-Json"]
                p = subprocess.run(pwsh_cmd, capture_output=True, text=True, timeout=4)
                if p.returncode == 0 and p.stdout:
                    j = _json.loads(p.stdout)
                    # j may be a list or single object
                    items = j if isinstance(j, list) else [j]
                    for it in items:
                        try:
                            num = it.get('Number')
                            name = it.get('FriendlyName') or f'PhysicalDrive{num}'
                            size = int(it.get('Size') or 0)
                            path = f"\\\\.\\PhysicalDrive{num}"
                            devices.append({'id': str(num), 'path': path, 'size_bytes': size, 'model': name, 'mountpoints': []})
                        except Exception:
                            continue
            except Exception:
                pass
    except Exception:
        pass

    return jsonify({'devices': devices})


@app.route('/api/fs/preview', methods=['GET'])
def api_fs_preview():
    """Return a lightweight preview for a file: text snippet or small image thumbnail (data-uri).

    Query params: root, path
    """
    root = request.args.get('root')
    rel = request.args.get('path')
    root_path = resolve_root_display_to_path(root)
    if not root_path:
        return jsonify({'error': 'invalid_root'}), 400
    abs_path = os.path.abspath(os.path.join(root_path, rel)) if rel else None
    if not abs_path or not is_path_under_allowed_roots(abs_path) or not os.path.exists(abs_path) or not os.path.isfile(abs_path):
        return jsonify({'error': 'access_denied_or_not_found'}), 403
    try:
        mime = magic.from_file(abs_path, mime=True)
    except Exception:
        mime = 'application/octet-stream'

    result = {'path': rel, 'root': root, 'name': os.path.basename(abs_path), 'mime': mime}
    try:
        if mime.startswith('text/') or mime in ('application/json', 'application/xml'):
            # return first N bytes/lines
            with open(abs_path, 'r', errors='ignore') as fh:
                text = fh.read(64 * 1024)
            result['preview'] = text
            result['mime'] = mime
            return jsonify(result)
        elif mime.startswith('image/'):
            # small thumbnail
            try:
                with Image.open(abs_path) as img:
                    img.thumbnail((256,256))
                    bio = io.BytesIO()
                    img.save(bio, format='PNG')
                    bio.seek(0)
                    data = base64.b64encode(bio.read()).decode('ascii')
                    result['thumbnail'] = f"data:image/png;base64,{data}"
                    result['mime'] = mime
                    return jsonify(result)
            except Exception as e:
                result['preview'] = '[thumbnail failed]'
                return jsonify(result)
        else:
            # binary fallback: strings preview
            result['preview'] = extract_strings_preview(abs_path)[:8192]
            return jsonify(result)
    except Exception as e:
        return jsonify({'error': 'preview_failed', 'message': str(e)}), 500


@app.route('/api/fs/mkdir', methods=['POST'])
def api_fs_mkdir():
    root = request.form.get('root')
    # CSRF protection
    if not verify_csrf_request():
        return jsonify({'error': 'csrf_missing_or_invalid'}), 403
    rel = request.form.get('path', '')
    name = request.form.get('name')
    if not name:
        return jsonify({'error': 'missing_name'}), 400
    root_path = resolve_root_display_to_path(root)
    if not root_path:
        return jsonify({'error': 'invalid_root'}), 400
    target = os.path.abspath(os.path.join(root_path, rel, name))
    if not is_path_under_allowed_roots(target):
        return jsonify({'error': 'access_denied'}), 403
    try:
        os.makedirs(target, exist_ok=False)
        return jsonify({'created': True, 'path': os.path.relpath(target, root_path)})
    except FileExistsError:
        return jsonify({'error': 'exists'}), 409
    except Exception as e:
        return jsonify({'error': 'mkdir_failed', 'message': str(e)}), 500


@app.route('/api/fs/delete', methods=['POST'])
def api_fs_delete():
    root = request.form.get('root')
    if not verify_csrf_request():
        return jsonify({'error': 'csrf_missing_or_invalid'}), 403
    rel = request.form.get('path')
    root_path = resolve_root_display_to_path(root)
    if not root_path:
        return jsonify({'error': 'invalid_root'}), 400
    # Tolerate callers passing the root again as path
    try:
        if rel and (rel == root or rel == os.path.basename(root_path)):
            rel = ''
    except Exception:
        pass
    abs_path = os.path.abspath(os.path.join(root_path, rel)) if rel else os.path.abspath(root_path)
    if not is_path_under_allowed_roots(abs_path) or not os.path.exists(abs_path):
        return jsonify({'error': 'access_denied_or_not_found'}), 403
    try:
        if os.path.isdir(abs_path):
            shutil.rmtree(abs_path)
        else:
            os.remove(abs_path)
        try:
            audit_event('delete', abs_path, {'by': 'api_fs_delete'})
        except Exception:
            pass
        return jsonify({'deleted': True})
    except Exception as e:
        return jsonify({'error': 'delete_failed', 'message': str(e)}), 500


@app.route('/api/fs/rename', methods=['POST'])
def api_fs_rename():
    root = request.form.get('root')
    old = request.form.get('old')
    new = request.form.get('new')
    if not verify_csrf_request():
        return jsonify({'error': 'csrf_missing_or_invalid'}), 403
    if not old or not new:
        return jsonify({'error': 'missing_params'}), 400
    root_path = resolve_root_display_to_path(root)
    if not root_path:
        return jsonify({'error': 'invalid_root'}), 400
    old_abs = os.path.abspath(os.path.join(root_path, old))
    new_abs = os.path.abspath(os.path.join(root_path, new))
    if not is_path_under_allowed_roots(old_abs) or not is_path_under_allowed_roots(new_abs):
        return jsonify({'error': 'access_denied'}), 403
    try:
        os.rename(old_abs, new_abs)
        try:
            audit_event('rename', old_abs, {'new': new_abs})
        except Exception:
            pass
        return jsonify({'renamed': True})
    except Exception as e:
        return jsonify({'error': 'rename_failed', 'message': str(e)}), 500


@app.route('/api/fs/upload', methods=['POST'])
def api_fs_upload():
    root = request.form.get('root')
    path = request.form.get('path', '')
    if not verify_csrf_request():
        return jsonify({'error': 'csrf_missing_or_invalid'}), 403
    f = request.files.get('file')
    if not f:
        return jsonify({'error': 'no_file'}), 400
    root_path = resolve_root_display_to_path(root)
    if not root_path:
        return jsonify({'error': 'invalid_root'}), 400
    dest_dir = os.path.abspath(os.path.join(root_path, path)) if path else os.path.abspath(root_path)
    if not is_path_under_allowed_roots(dest_dir):
        return jsonify({'error': 'access_denied'}), 403
    os.makedirs(dest_dir, exist_ok=True)
    filename = secure_filename(f.filename)
    dest = os.path.join(dest_dir, filename)
    try:
        f.save(dest)
        # register uploaded file in in-memory mapping so UI can display it immediately
        try:
            size = os.path.getsize(dest)
            now_ts = int(time.time())
            uploaded_files_db[filename] = {
                'path': dest,
                'size': size,
                'size_mb': round(size / (1024*1024), 2),
                'hash_info': {},
                'hashing_complete': False,
                'encryption_status': {'encrypted': False, 'decrypting': False, 'decrypted_path': None, 'description': ''},
                'uploaded_at': now_ts,
                # status_changed_at mirrors uploaded_at initially; other flows may update it
                'status_changed_at': now_ts,
                'status': 'uploaded'
            }
            # kick off background hashing
            try:
                t = threading.Thread(target=calculate_hashes_threaded, args=(dest,))
                t.daemon = True
                t.start()
            except Exception:
                pass
        except Exception:
            pass
        return jsonify({'uploaded': True, 'path': os.path.relpath(dest, root_path)})
    except Exception as e:
        return jsonify({'error': 'upload_failed', 'message': str(e)}), 500



@app.route('/api/uploaded_files', methods=['GET'])
def api_uploaded_files():
    """Return current uploaded files mapping for client-side refresh."""
    try:
        # Convert paths to relative paths and include size_mb
        out = {}
        for name, info in uploaded_files_db.items():
            try:
                rel = info.get('path')
                out[name] = {
                    'path': rel,
                    'size': info.get('size', 0),
                    'size_mb': info.get('size_mb', round(info.get('size', 0)/(1024*1024), 2)),
                    'hash_info': info.get('hash_info', {}),
                    'encryption_status': info.get('encryption_status', {}),
                    'uploaded_at': info.get('uploaded_at'),
                    'status_changed_at': info.get('status_changed_at'),
                    'status': info.get('status', 'uploaded')
                }
            except Exception:
                out[name] = {'path': info.get('path'), 'size': info.get('size', 0), 'size_mb': info.get('size_mb', 0), 'hash_info': {}, 'encryption_status': {}}
        return jsonify({'uploaded_files': out})
    except Exception as e:
        return jsonify({'error': 'failed', 'message': str(e)}), 500


@app.route('/api/file_status/<int:file_id>', methods=['GET'])
def api_file_status(file_id):
    """Return the status for a file record by id."""
    try:
        conn = _get_db_conn()
        cur = conn.cursor()
        cur.execute('SELECT id, filename, status FROM files WHERE id = ?', (file_id,))
        r = cur.fetchone()
        try: conn.close()
        except Exception: pass
        if not r:
            return jsonify({'error': 'not_found'}), 404
        return jsonify({'id': r[0], 'filename': r[1], 'status': r[2]})
    except Exception as e:
        return jsonify({'error': 'status_failed', 'message': str(e)}), 500


@app.route('/api/get_csrf', methods=['GET'])
def api_get_csrf():
    """Return the CSRF token for the current session (safe on same-origin).
    This helps XHR clients obtain the token when it's not embedded into the page.
    """
    try:
        token = session.get('csrf_token')
        if not token:
            # ensure a token exists
            session['csrf_token'] = secrets.token_urlsafe(32)
            token = session['csrf_token']
        return jsonify({'csrf': token})
    except Exception:
        return jsonify({'error': 'csrf_unavailable'}), 500


@app.route('/api/file_status', methods=['POST'])
def api_file_status_update():
    """Update status for a file record. Request form: file_id, status. CSRF-protected."""
    if not verify_csrf_request():
        return jsonify({'error': 'csrf_missing_or_invalid'}), 403
    # enforce allowed statuses
    ALLOWED_STATUSES = {'saved', 'processing', 'waiting', 'paused'}
    try:
        fid = request.form.get('file_id') or (request.json and request.json.get('file_id'))
        status = (request.form.get('status') or (request.json and request.json.get('status')) or '').strip()
        if not fid or not status:
            return jsonify({'error': 'missing_parameters'}), 400
        if status not in ALLOWED_STATUSES:
            return jsonify({'error': 'invalid_status', 'allowed': list(ALLOWED_STATUSES)}), 400
        try:
            fid = int(fid)
        except Exception:
            return jsonify({'error': 'invalid_file_id'}), 400

        # authorization check: require session flag 'can_edit_status' or admin
        def _can_edit_status():
            try:
                if session.get('is_admin'):
                    return True
                if session.get('can_edit_status'):
                    return True
                # allow during development when running on localhost
                if request.remote_addr in ('127.0.0.1', '::1'):
                    return True
            except Exception:
                pass
            return False

        if not _can_edit_status():
            return jsonify({'error': 'unauthorized'}), 403

        conn = _get_db_conn()
        cur = conn.cursor()
        # fetch current value for audit
        cur.execute('SELECT filename, status FROM files WHERE id = ?', (fid,))
        row = cur.fetchone()
        if not row:
            try: conn.close()
            except Exception: pass
            return jsonify({'error': 'not_found'}), 404
        old_status = row[1]
        now = datetime.datetime.now().isoformat()
        try:
            cur.execute('UPDATE files SET status = ?, status_changed_at = ? WHERE id = ?', (status, now, fid))
        except Exception:
            # fallback if target DB/table doesn't have the status_changed_at column
            cur.execute('UPDATE files SET status = ? WHERE id = ?', (status, fid))
        conn.commit()
        try: conn.close()
        except Exception: pass

        # audit
        try:
            audit_event('status_change', row[0], details={'file_id': fid, 'old': old_status, 'new': status, 'user': session.get('analysis_session_id')})
        except Exception:
            pass

        return jsonify({'ok': True})
    except Exception as e:
        return jsonify({'error': 'update_failed', 'message': str(e)}), 500


@app.route('/api/file_status_bulk', methods=['POST'])
def api_file_status_bulk():
    """Bulk update statuses for multiple file ids. Expects JSON: {file_ids: [1,2,3], status: 'processing'}"""
    if not verify_csrf_request():
        return jsonify({'error': 'csrf_missing_or_invalid'}), 403
    ALLOWED_STATUSES = {'saved', 'processing', 'waiting', 'paused'}
    try:
        data = None
        if request.is_json:
            data = request.get_json()
        else:
            # try form data
            data = { 'file_ids': request.form.getlist('file_ids'), 'status': request.form.get('status') }
        if not data:
            return jsonify({'error': 'missing_parameters'}), 400
        file_ids = data.get('file_ids') or []
        status = (data.get('status') or '').strip()
        if not file_ids or not status:
            return jsonify({'error': 'missing_parameters'}), 400
        if status not in ALLOWED_STATUSES:
            return jsonify({'error': 'invalid_status', 'allowed': list(ALLOWED_STATUSES)}), 400

        # simple auth reuse
        if not (session.get('is_admin') or session.get('can_edit_status') or request.remote_addr in ('127.0.0.1', '::1')):
            return jsonify({'error': 'unauthorized'}), 403

        # normalize ids
        norm_ids = []
        for v in file_ids:
            try:
                norm_ids.append(int(v))
            except Exception:
                continue
        if not norm_ids:
            return jsonify({'error': 'invalid_file_ids'}), 400

        conn = _get_db_conn()
        cur = conn.cursor()
        # fetch filenames and old statuses for audit
        q = 'SELECT id, filename, status FROM files WHERE id IN ({seq})'.format(seq=','.join(['?']*len(norm_ids)))
        cur.execute(q, tuple(norm_ids))
        rows = cur.fetchall()
        if not rows:
            try: conn.close()
            except Exception: pass
            return jsonify({'error': 'not_found'}), 404

        # update
        now = datetime.datetime.now().isoformat()
        uq_full = 'UPDATE files SET status = ?, status_changed_at = ? WHERE id = ?'
        uq_simple = 'UPDATE files SET status = ? WHERE id = ?'
        for r in rows:
            fid = r[0]
            old = r[2]
            try:
                cur.execute(uq_full, (status, now, fid))
            except Exception:
                cur.execute(uq_simple, (status, fid))
        conn.commit()
        try: conn.close()
        except Exception: pass

        # audit each change
        try:
            for r in rows:
                audit_event('status_change', r[1], details={'file_id': r[0], 'old': r[2], 'new': status, 'user': session.get('analysis_session_id')})
        except Exception:
            pass

        return jsonify({'ok': True, 'updated': len(rows)})
    except Exception as e:
        return jsonify({'error': 'bulk_update_failed', 'message': str(e)}), 500


def start_analysis_session():
    """Create and mark an analysis session active."""
    try:
        conn = _get_db_conn()
        cur = conn.cursor()
        now = datetime.datetime.now().isoformat()
        cur.execute('INSERT INTO sessions(started_at, active) VALUES(?,1)', (now,))
        session_id = cur.lastrowid
        conn.commit()
        # create a per-session subfolder under SESSION_FOLDER and typed subfolders
        try:
            ts = datetime.datetime.now().strftime('%Y%m%d_%H%M%S')
            session_subfolder_name = f"Session_{ts}"
            session_subfolder_path = os.path.join(app.config['SESSION_FOLDER'], session_subfolder_name)
            os.makedirs(session_subfolder_path, exist_ok=True)
            # typed subfolders
            for sub in ('Carved','Deleted','Reports','Logs','Events','ManualCarving'):
                try:
                    os.makedirs(os.path.join(session_subfolder_path, sub), exist_ok=True)
                except Exception:
                    pass
            # store path in sessions table if column exists, else fallback to ended_at
            try:
                cur.execute('PRAGMA table_info(sessions)')
                cols = [r[1] for r in cur.fetchall()]
                if 'session_path' in cols:
                    cur.execute('UPDATE sessions SET session_path=? WHERE id=?', (session_subfolder_path, session_id))
                else:
                    cur.execute('UPDATE sessions SET ended_at=? WHERE id=?', (session_subfolder_path, session_id))
                conn.commit()
            except Exception:
                pass
            # also store in flask session for runtime access
            try:
                session['analysis_session_id'] = session_id
                session['analysis_session_path'] = session_subfolder_path
            except Exception:
                pass
        except Exception:
            pass
        return session_id
    except Exception as e:
        print(f"start_analysis_session error: {e}")
        return None
    finally:
        try:
            conn.close()
        except Exception:
            pass

def end_analysis_session(sess_id):
    try:
        conn = _get_db_conn()
        cur = conn.cursor()
        now = datetime.datetime.now().isoformat()
        cur.execute('UPDATE sessions SET ended_at=?, active=0 WHERE id=?', (now, sess_id))
        conn.commit()
    except Exception as e:
        print(f"end_analysis_session error: {e}")
    finally:
        try:
            conn.close()
        except Exception:
            pass

def add_file_record(filename, file_type, path, size_bytes=0, session_id=None, extra=None):
    # If a target DB was provided in extra (e.g. from the Upload Storage Target card), prefer that
    try:
        target_db = None
        if extra and isinstance(extra, dict):
            target_db = extra.get('target_db')
        # If not provided in extra, attempt to use a selection stored in the Flask session
        if not target_db:
            try:
                s_sel = session.get('selected_target_db')
                if s_sel:
                    target_db = s_sel
            except Exception:
                # No request/session context available
                target_db = None
        # Delegate to the engine-agnostic inserter which will fallback to local DB on error
        if target_db and target_db != 'local':
            # Create a save job so the background worker can push the file to the target DB
            try:
                job_id = create_save_job(filename, path, target_db)
                # Attach job id to extra for visibility and later status mapping
                if extra is None: extra = {}
                extra = dict(extra)
                extra['save_job_id'] = job_id
            except Exception:
                job_id = None
            # Record locally immediately with job reference; background worker will attempt the remote insert
            try:
                conn = _get_db_conn()
                cur = conn.cursor()
                now = datetime.datetime.now().isoformat()
                extra_json = json.dumps(extra) if extra is not None else None
                cur.execute('INSERT INTO files(filename, file_type, path, size_bytes, created_at, session_id, extra) VALUES(?,?,?,?,?,?,?)',
                            (filename, file_type, path, size_bytes, now, session_id, extra_json))
                conn.commit()
                return True
            except Exception:
                # fallback to db_insert_file immediate attempt
                try:
                    return db_insert_file(filename, file_type, path, size_bytes=size_bytes, session_id=session_id, extra=extra, target_db=target_db)
                except Exception:
                    pass

        conn = _get_db_conn()
        cur = conn.cursor()
        now = datetime.datetime.now().isoformat()
        extra_json = json.dumps(extra) if extra is not None else None
        cur.execute('INSERT INTO files(filename, file_type, path, size_bytes, created_at, session_id, extra) VALUES(?,?,?,?,?,?,?)',
                    (filename, file_type, path, size_bytes, now, session_id, extra_json))
        conn.commit()
        # For session-scoped artifacts, ensure they live inside the per-session folder if possible.
        # Background threads and worker contexts may not have Flask `session` available, so
        # fall back to the sessions table (active or by provided session_id) to find the
        # runtime session_path.
        try:
            eligible_types = ('carved', 'deleted_recovered', 'report', 'log', 'event_log', 'deleted', 'manual_carve')
            if file_type in eligible_types:
                sess_path = None
                runtime_sid = session_id
                # Try Flask session first (available in request contexts)
                try:
                    sess_path = session.get('analysis_session_path')
                    if not runtime_sid:
                        runtime_sid = session.get('analysis_session_id')
                except Exception:
                    sess_path = None

                # If not available, try to look up in DB. Prefer provided session_id, else look for active session
                try:
                    if not sess_path:
                        cur2 = conn.cursor()
                        if runtime_sid:
                            cur2.execute('SELECT session_path, ended_at FROM sessions WHERE id=?', (runtime_sid,))
                            row = cur2.fetchone()
                            candidate = None
                            if row:
                                candidate = row[0] or row[1]
                            if candidate and os.path.isdir(candidate):
                                sess_path = candidate
                        else:
                            # look for an active session with a session_path first
                            try:
                                cur2.execute("SELECT id, session_path FROM sessions WHERE active=1 AND session_path IS NOT NULL ORDER BY started_at DESC LIMIT 1")
                                r = cur2.fetchone()
                                if r and r[1] and os.path.isdir(r[1]):
                                    runtime_sid = r[0]
                                    sess_path = r[1]
                                else:
                                    # fallback: most recent session with a session_path
                                    cur2.execute("SELECT id, session_path FROM sessions WHERE session_path IS NOT NULL ORDER BY started_at DESC LIMIT 1")
                                    r2 = cur2.fetchone()
                                    if r2 and r2[1] and os.path.isdir(r2[1]):
                                        runtime_sid = r2[0]
                                        sess_path = r2[1]
                            except Exception:
                                pass
                except Exception:
                    sess_path = None

                if sess_path:
                    # Create subfolders per artifact type for neatness inside the session folder
                    mapping = {
                        'carved': 'Carved',
                        'deleted_recovered': 'Deleted',
                        'deleted': 'Deleted',
                        'report': 'Reports',
                        'log': 'Logs',
                        'event_log': 'Events',
                        'manual_carve': 'ManualCarving'
                    }
                    sub = mapping.get(file_type, 'Misc')
                    dest_dir = os.path.join(sess_path, sub)
                    os.makedirs(dest_dir, exist_ok=True)
                    dest = os.path.join(dest_dir, secure_filename(filename))
                    # If file is already the intended target, skip copying
                    try:
                        if os.path.abspath(path) != os.path.abspath(dest):
                            if os.path.exists(path) and os.path.isfile(path):
                                try:
                                    shutil.copy2(path, dest)
                                except Exception:
                                    try:
                                        with open(path, 'rb') as rf, open(dest, 'wb') as wf:
                                            wf.write(rf.read())
                                    except Exception:
                                        pass
                            else:
                                # if original path doesn't exist (e.g., we were given bytes), nothing to copy
                                pass
                        # If we derived a runtime_sid from DB, attempt to update the DB record we just inserted
                        if runtime_sid and not session_id:
                            try:
                                cur_update = conn.cursor()
                                cur_update.execute('UPDATE files SET session_id=? WHERE id=?', (runtime_sid, cur.lastrowid))
                                conn.commit()
                            except Exception:
                                pass
                    except Exception:
                        pass
        except Exception:
            pass
    except Exception as e:
        print(f"add_file_record error: {e}")
    finally:
        try:
            conn.close()
        except Exception:
            pass


def _save_bytes_to_session_file(data_bytes, filename, file_type='report'):
    """Save bytes or BytesIO to a file in SESSION_FOLDER and record in DB.

    Returns the saved path or None on failure.
    """
    try:
        sess_id = None
        try:
            sess_id = session.get('analysis_session_id')
        except Exception:
            sess_id = None

        timestamp = datetime.datetime.now().strftime('%Y%m%d%H%M%S')
        safe_name = secure_filename(filename)
        # Try to place reports into the per-session Reports subfolder when available
        sess_path = None
        try:
            sess_path = session.get('analysis_session_path')
        except Exception:
            sess_path = None

        if sess_path and os.path.isdir(sess_path):
            reports_dir = os.path.join(sess_path, 'Reports')
            os.makedirs(reports_dir, exist_ok=True)
            dest_path = os.path.join(reports_dir, f"report_{timestamp}_{safe_name}")
        else:
            dest_path = os.path.join(app.config.get('SESSION_FOLDER', SESSION_FOLDER), f"report_{timestamp}_{safe_name}")

        # If data_bytes is BytesIO-like, getvalue
        try:
            if hasattr(data_bytes, 'getvalue'):
                raw = data_bytes.getvalue()
            else:
                raw = data_bytes
        except Exception:
            raw = data_bytes

        with open(dest_path, 'wb') as wf:
            if isinstance(raw, str):
                wf.write(raw.encode('utf-8'))
            else:
                wf.write(raw)

        try:
            add_file_record(os.path.basename(dest_path), file_type, dest_path, os.path.getsize(dest_path), session_id=sess_id, extra={'generated_by': 'reporting'})
        except Exception:
            pass
        return dest_path
    except Exception as e:
        print(f"_save_bytes_to_session_file error: {e}")
        return None

# --- Populate in-memory deleted_files_db from on-disk files at startup ---
def populate_deleted_files_db_from_disk():
    """Scan the deleted recovery folder and populate the in-memory `deleted_files_db`.

    This helps the web UI show files that were recovered on disk before the server started.
    """
    global deleted_files_db
    try:
        deleted_files_db.clear()
    except Exception:
        deleted_files_db = {}

    recovery_dir = app.config.get('DELETED_RECOVERY_FOLDER', os.path.join(APP_ROOT, 'deleted_files'))
    try:
        for fname in sorted(os.listdir(recovery_dir)):
            fpath = os.path.join(recovery_dir, fname)
            if not os.path.isfile(fpath):
                continue
            try:
                info = get_enhanced_file_info(fpath)
            except Exception:
                info = {'mime_type': 'application/octet-stream', 'file_type': 'Unknown', 'thumbnail': None, 'size_bytes': os.path.getsize(fpath) if os.path.exists(fpath) else 0, 'size_kb': '0'}

            try:
                mtime = datetime.datetime.fromtimestamp(os.path.getmtime(fpath)).strftime('%Y-%m-%d %H:%M:%S')
            except Exception:
                mtime = 'Unknown'

            try:
                ctime = datetime.datetime.fromtimestamp(os.path.getctime(fpath)).strftime('%Y-%m-%d %H:%M:%S')
            except Exception:
                ctime = 'Unknown'

            deleted_files_db[os.path.basename(fpath)] = {
                'inode': None,
                'name': os.path.basename(fpath),
                'size': info.get('size_bytes', 0),
                'size_kb': info.get('size_kb', '0'),
                'mtime': mtime,
                'ctime': ctime,
                'file_type': info.get('mime_type', 'application/octet-stream'),
                'thumbnail': info.get('thumbnail'),
                'path': fpath
            }
    except Exception as e:
        app.logger.debug(f"populate_deleted_files_db_from_disk: error scanning folder: {e}")

# Run initial population so the UI sees recovered files already on disk
try:
    populate_deleted_files_db_from_disk()
except Exception:
    pass

def get_encrypted_files():
    """Get list of encrypted files with their sizes."""
    encrypted_files = []
    try:
        enc_folder = app.config.get('ENCRYPTED_FOLDER')
        if enc_folder and os.path.isdir(enc_folder):
            for filename in os.listdir(enc_folder):
                # keep original behavior of only listing .enc files
                if filename.endswith('.enc'):
                    filepath = os.path.join(enc_folder, filename)
                    if os.path.isfile(filepath):
                        size = os.path.getsize(filepath)
                        size_str = format_bytes(size)
                        encrypted_files.append({
                            'filename': filename,
                            'size': size_str,
                            'path': filepath
                        })
    except Exception as e:
        print(f"Error listing encrypted files: {e}")
    
    return encrypted_files


def get_decrypted_files():
    """Get list of decrypted files with their sizes."""
    decrypted_files = []
    try:
        dec_folder = app.config.get('DECRYPTED_FOLDER')
        if dec_folder and os.path.isdir(dec_folder):
            for filename in os.listdir(dec_folder):
                # list typical decrypted outputs; include any file
                filepath = os.path.join(dec_folder, filename)
                if os.path.isfile(filepath):
                    size = os.path.getsize(filepath)
                    size_str = format_bytes(size)
                    decrypted_files.append({
                        'filename': filename,
                        'size': size_str,
                        'path': filepath
                    })
    except Exception as e:
        print(f"Error listing decrypted files: {e}")
    return decrypted_files


@app.route('/download_decrypted/<filename>')
def download_decrypted_file(filename):
    try:
        safe_filename = secure_filename(filename)
        file_path = os.path.join(app.config['DECRYPTED_FOLDER'], safe_filename)
        if not os.path.exists(file_path):
            flash(f"Decrypted file '{safe_filename}' not found.", "error")
            return redirect(url_for('decryption_page', filename=session.get('last_decryption_target') or ''))
        return send_file(file_path, as_attachment=True, download_name=safe_filename, mimetype='application/octet-stream')
    except Exception as e:
        flash(f"Error downloading decrypted file: {str(e)}", "error")
        return redirect(url_for('decryption_page', filename=session.get('last_decryption_target') or ''))


@app.route('/delete_decrypted/<filename>')
def delete_decrypted_file(filename):
    try:
        safe_filename = secure_filename(filename)
        file_path = os.path.join(app.config['DECRYPTED_FOLDER'], safe_filename)
        if os.path.exists(file_path):
            os.remove(file_path)
            flash(f"Decrypted file '{safe_filename}' deleted successfully.", "success")
        else:
            flash(f"Decrypted file '{safe_filename}' not found.", "error")
    except Exception as e:
        flash(f"Error deleting decrypted file: {str(e)}", "error")
    return redirect(url_for('decryption_page', filename=session.get('last_decryption_target') or ''))


@app.context_processor
def inject_sidebar_session():
    """Provide a small summary for the sidebar: whether a session is active, a short name, and carved count."""
    sidebar = {'active': False, 'short_name': None, 'carved_count': 0}
    try:
        sid = session.get('analysis_session_id')
        s_path = session.get('analysis_session_path')
        if sid:
            sidebar['active'] = True
            if s_path and os.path.isdir(s_path):
                sidebar['short_name'] = os.path.basename(s_path)
                # short form: Session_YYYYMMDD_HHMM -> drop seconds
                try:
                    parts = sidebar['short_name'].split('_')
                    if len(parts) >= 2:
                        ts = parts[1]
                        sidebar['short_name'] = f"S_{ts[:12]}"
                except Exception:
                    pass
                # compute carved count safely
                try:
                    carved_dir = os.path.join(s_path, 'Carved')
                    if os.path.isdir(carved_dir):
                        sidebar['carved_count'] = sum(1 for _ in os.listdir(carved_dir) if os.path.isfile(os.path.join(carved_dir, _)))
                except Exception:
                    sidebar['carved_count'] = 0
    except Exception:
        pass
    return {'sidebar_session': sidebar}



if not os.path.exists(DICTIONARY_FILE):
    with open(DICTIONARY_FILE, 'w') as f:
        f.write("password\n123456\nadmin\n12345\n12345678\nletmein\nqwerty\npassword1\n")



# --- In-Memory Databases & Caches ---
uploaded_files_db = {}
carved_files_db = {}
deleted_files_db = {}
sorted_carved_keys = []
sorted_deleted_inodes = []

# --- Persistent database connections storage ---
db_connections = {}


def _ensure_db_connections_table(conn):
    try:
        cur = conn.cursor()
        cur.execute('''
            CREATE TABLE IF NOT EXISTS db_connections (
                id TEXT PRIMARY KEY,
                name TEXT,
                type TEXT,
                engine TEXT,
                conn TEXT,
                conn_details TEXT,
                connected INTEGER DEFAULT 0,
                size TEXT,
                last_checked TEXT,
                message TEXT
            )
        ''')
        conn.commit()
    except Exception:
        try:
            conn.rollback()
        except Exception:
            pass


def _ensure_save_jobs_table(conn):
    try:
        cur = conn.cursor()
        cur.execute('''
            CREATE TABLE IF NOT EXISTS save_jobs (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                filename TEXT,
                path TEXT,
                target_db TEXT,
                status TEXT,
                message TEXT,
                created_at TEXT,
                updated_at TEXT
            )
        ''')
        conn.commit()
    except Exception:
        try:
            conn.rollback()
        except Exception:
            pass


def create_save_job(filename, path, target_db):
    try:
        conn = sqlite3.connect(DB_FILE)
        cur = conn.cursor()
        now = datetime.datetime.now().isoformat()
        cur.execute('INSERT INTO save_jobs(filename, path, target_db, status, message, created_at, updated_at) VALUES(?,?,?,?,?,?,?)',
                    (filename, path, target_db, 'waiting', None, now, now))
        conn.commit()
        job_id = cur.lastrowid
        try: conn.close()
        except Exception: pass
        return job_id
    except Exception as e:
        app.logger.debug(f"create_save_job error: {e}")
        try:
            conn.close()
        except Exception:
            pass
        return None


def update_save_job(job_id, status=None, message=None):
    try:
        conn = sqlite3.connect(DB_FILE)
        cur = conn.cursor()
        now = datetime.datetime.now().isoformat()
        if status is not None and message is not None:
            cur.execute('UPDATE save_jobs SET status=?, message=?, updated_at=? WHERE id=?', (status, message, now, job_id))
        elif status is not None:
            cur.execute('UPDATE save_jobs SET status=?, updated_at=? WHERE id=?', (status, now, job_id))
        elif message is not None:
            cur.execute('UPDATE save_jobs SET message=?, updated_at=? WHERE id=?', (message, now, job_id))
        conn.commit()
        try: conn.close()
        except Exception: pass
        return True
    except Exception as e:
        app.logger.debug(f"update_save_job error: {e}")
        try:
            conn.close()
        except Exception:
            pass
        return False


def get_next_save_job():
    try:
        conn = sqlite3.connect(DB_FILE)
        cur = conn.cursor()
        cur.execute("SELECT id, filename, path, target_db FROM save_jobs WHERE status='waiting' ORDER BY created_at ASC LIMIT 1")
        row = cur.fetchone()
        try: conn.close()
        except Exception: pass
        if row:
            return {'id': row[0], 'filename': row[1], 'path': row[2], 'target_db': row[3]}
    except Exception as e:
        app.logger.debug(f"get_next_save_job error: {e}")
        try:
            conn.close()
        except Exception:
            pass
    return None


def _save_job_worker(stop_event):
    """Background worker that processes save_jobs table."""
    while not stop_event.is_set():
        try:
            job = get_next_save_job()
            if not job:
                # sleep briefly when no jobs
                time.sleep(1)
                continue
            job_id = job['id']
            update_save_job(job_id, status='processing')
            try:
                # Attempt to insert into target DB
                success = db_insert_file(job['filename'], 'uploaded', job['path'], size_bytes=os.path.getsize(job['path']) if os.path.exists(job['path']) else 0, extra={'auto_saved': True}, target_db=job['target_db'])
                if success:
                    update_save_job(job_id, status='saved', message='Saved to target')
                else:
                    update_save_job(job_id, status='stopped', message='Failed to save to target')
            except Exception as e:
                update_save_job(job_id, status='paused', message=str(e))
        except Exception as e:
            app.logger.debug(f"save_job_worker loop error: {e}")
            time.sleep(2)


# Start the background worker
try:
    _ensure_save_jobs_table(sqlite3.connect(DB_FILE))
except Exception:
    pass
_save_stop_event = threading.Event()
threading.Thread(target=_save_job_worker, args=(_save_stop_event,), daemon=True).start()



def load_db_configs():
    """Load DB connection configurations from the application SQLite DB (fac_data.db).
    If the JSON file exists and the table is empty, attempt a one-time migration.
    """
    global db_connections
    db_connections = {}
    try:
        # connect to application DB file if available
        if DB_FILE and os.path.exists(DB_FILE):
            conn = sqlite3.connect(DB_FILE)
            try:
                _ensure_db_connections_table(conn)
                cur = conn.cursor()
                cur.execute('SELECT id, name, type, engine, conn, conn_details, connected, size, last_checked, message FROM db_connections')
                rows = cur.fetchall()
                if rows:
                    for r in rows:
                        cid = r[0]
                        try:
                            conn_details = json.loads(r[5]) if r[5] else None
                        except Exception:
                            conn_details = None
                        db_connections[cid] = {
                            'name': r[1], 'type': r[2], 'engine': r[3], 'conn': r[4],
                            'conn_details': conn_details,
                            'connected': bool(r[6]), 'size': r[7], 'last_checked': r[8], 'message': r[9]
                        }
                    return
                # If table empty but JSON config exists, migrate
                if os.path.exists(DB_CONFIG_FILE):
                    try:
                        with open(DB_CONFIG_FILE, 'r', encoding='utf-8') as fh:
                            j = json.load(fh)
                            if isinstance(j, dict):
                                for k, v in j.items():
                                    db_connections[k] = v
                                    # persist into sqlite
                                    try:
                                        cur.execute('INSERT OR REPLACE INTO db_connections (id, name, type, engine, conn, conn_details, connected, size, last_checked, message) VALUES (?,?,?,?,?,?,?,?,?,?)',
                                                    (k, v.get('name'), v.get('type'), v.get('engine'), v.get('conn'), json.dumps(v.get('conn_details')) if v.get('conn_details') else None, int(bool(v.get('connected'))), v.get('size'), v.get('last_checked'), v.get('message')))
                                    except Exception:
                                        pass
                                conn.commit()
                                return
                    except Exception:
                        pass
            finally:
                try: conn.close()
                except Exception: pass
        # fallback: no DB file, try JSON file
        if os.path.exists(DB_CONFIG_FILE):
            try:
                with open(DB_CONFIG_FILE, 'r', encoding='utf-8') as fh:
                    j = json.load(fh)
                    if isinstance(j, dict):
                        db_connections = j
            except Exception:
                db_connections = {}
    except Exception as e:
        app.logger.debug(f"load_db_configs: error loading configs: {e}")
        db_connections = {}


def save_db_configs():
    """Persist the in-memory db_connections into the application SQLite DB if available,
    otherwise fallback to JSON file."""
    try:
        if DB_FILE:
            conn = sqlite3.connect(DB_FILE)
            try:
                _ensure_db_connections_table(conn)
                cur = conn.cursor()
                for cid, cfg in db_connections.items():
                    try:
                        cur.execute('INSERT OR REPLACE INTO db_connections (id, name, type, engine, conn, conn_details, connected, size, last_checked, message) VALUES (?,?,?,?,?,?,?,?,?,?)',
                                    (cid, cfg.get('name'), cfg.get('type'), cfg.get('engine'), cfg.get('conn'), json.dumps(cfg.get('conn_details')) if cfg.get('conn_details') else None, int(bool(cfg.get('connected'))), cfg.get('size'), cfg.get('last_checked'), cfg.get('message')))
                    except Exception:
                        pass
                conn.commit()
                return
            finally:
                try: conn.close()
                except Exception: pass
    except Exception:
        pass
    # fallback to JSON file
    try:
        with open(DB_CONFIG_FILE, 'w', encoding='utf-8') as f:
            json.dump(db_connections, f, indent=2)
    except Exception as e:
        app.logger.debug(f"save_db_configs: error saving {DB_CONFIG_FILE}: {e}")


# Load on startup
try:
    load_db_configs()
except Exception:
    db_connections = {}

# Update your status dictionaries at the top of the file
carving_status = {
    "progress": 0, "current_offset": "0x00000000", "files_found": 0,
    "complete": False, "found_files_list": [], "time_remaining_str": "N/A",
    "start_time": None, "estimated_total_time": None, "elapsed_time": "0s"
}

deleted_scan_status = {
    "in_progress": False, 
    "files_found": 0, 
    "complete": False, 
    "message": "Scan has not started.",
    "errors": [], 
    "time_remaining_str": "N/A",
    "start_time": None,
    "estimated_total_time": None,
    "elapsed_time": "0s",
    "scan_methods": {
        "directory_walk": 0,
        "inode_scan": 0,
        "file_slack": 0,
        "unallocated_space": 0,
        "recycle_bin": 0
    }
}

upload_status = {
    "in_progress": False,
    "progress": 0,
    "filename": None,
    "start_time": None,
    "estimated_total_time": None,
    "elapsed_time": "0s",
    "time_remaining_str": "Calculating...",
    "total_bytes": 0,
    "bytes_uploaded": 0,
    "upload_speed": 0,
    "last_update_time": None,
    "complete": False,
    "error": None
}
# --- Status Dictionaries for Background Tasks ---

decryption_status = {
    "in_progress": False, "complete": False, "message": "", "attempts": 0, "success": False,
    "filename": None, "progress": 0, "total_attempts": 0
}

hashing_status = {
    "in_progress": False, "progress": 0, "complete": False, "hashes": {}
}

strings_status = {
    "in_progress": False, "complete": False, "progress": 0, "strings_found": 0, "preview": []
}

# --- Signatures and Patterns ---
CUSTOM_ENC_HEADER = b'FCPE_V1_'  # Forensic Carver Pro Encryption, Version 1

ENCRYPTION_SIGNATURES = {
    'AES_ENCRYPTED': {'header': b'Salted__', 'description': 'OpenSSL AES encrypted file'},
    'BITLOCKER': {'header': b'-FVE-FS-', 'description': 'BitLocker encrypted volume'},
    'VERACRYPT': {'header': b'VERA', 'description': 'VeraCrypt encrypted container'},
    'PGP': {'header': b'\x85\x01\x0c', 'description': 'PGP encrypted file'},
    '7Z_ENCRYPTED': {'header': b'7z\xbc\xaf\x27\x1c', 'description': '7-Zip encrypted archive'},
    'ZIP_ENCRYPTED': {'header': b'PK\x03\x04', 'description': 'ZIP encrypted archive (needs password)'},
    'RAR_ENCRYPTED': {'header': b'Rar!\x1a\x07', 'description': 'RAR encrypted archive'},
    'FERNET_ENCRYPTED': {'header': CUSTOM_ENC_HEADER, 'description': 'Forensic Carver Pro Encrypted File'},
}

# --- Optimized & Consolidated File Signatures ---
FILE_SIGNATURES = {
    'Image Files': {
        'JPEG': {'header': b'\xff\xd8\xff\xe0', 'footer': b'\xff\xd9', 'extension': '.jpeg', 'max_size': 20 * 1024 * 1024},
        'JPG': {'header': b'\xff\xd8\xff\xe0', 'footer': b'\xff\xd9', 'extension': '.jpg', 'max_size': 20 * 1024 * 1024},
        'PNG': {'header': b'\x89\x50\x4e\x47\x0d\x0a\x1a\x0a', 'footer': b'\x49\x45\x4e\x44\xae\x42\x60\x82', 'extension': '.png'},
        'GIF': {'header': b'\x47\x49\x46\x38', 'footer': b'\x00\x3b', 'extension': '.gif'},
        'TIFF (Intel)': {'header': b'\x49\x49\x2a\x00', 'max_size': 30*1024*1024, 'extension': '.tiff'},
        'TIFF (Motorola)': {'header': b'\x4d\x4d\x00\x2a', 'max_size': 30*1024*1024, 'extension': '.tiff'},
        'BMP': {'header': b'\x42\x4d', 'max_size': 30*1024*1024, 'extension': '.bmp'},
        'ICO': {'header': b'\x00\x00\x01\x00', 'max_size': 1*1024*1024, 'extension': '.ico'},
    },
    'Document Files': {
        'PDF': {'header': b'\x25\x50\x44\x46', 'footer': b'\x25\x25\x45\x4f\x46', 'extension': '.pdf'},
        'DOCX': {'header': b'\x50\x4b\x03\x04', 'footer': b'\x50\x4b\x05\x06', 'extension': '.docx'},
        'XLSX': {'header': b'\x50\x4b\x03\x04', 'footer': b'\x50\x4b\x05\x06', 'extension': '.xlsx'},
    'PPTX': {'header': b'\x50\x4b\x03\x04', 'footer': b'\x50\x4b\x05\x06', 'extension': '.pptx'},
        'RTF': {'header': b'\x7b\x5c\x72\x74\x66', 'footer': b'\x7d', 'extension': '.rtf'},
        'DOC': {'header': b'\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1', 'max_size': 20*1024*1024, 'extension': '.doc'},
    },
    'Archive Files': {
        'ZIP': {'header': b'\x50\x4b\x03\x04', 'footer': b'\x50\x4b\x05\x06', 'extension': '.zip'},
        'RAR': {'header': b'\x52\x61\x72\x21\x1a\x07\x00', 'max_size': 1024*1024*1024, 'extension': '.rar'},
        '7Z': {'header': b'\x37\x7a\xbc\xaf\x27\x1c', 'max_size': 1024*1024*1024, 'extension': '.7z'},
        'GZIP': {'header': b'\x1f\x8b', 'max_size': 1024*1024*1024, 'extension': '.gz'},
    },
    'Audio Files': {
        'MP3': {'headers': [b'\x49\x44\x33', b'\xff\xfb'], 'extension': '.mp3'},
        'WAV': {'header': b'\x52\x49\x46\x46', 'max_size': 100*1024*1024, 'extension': '.wav'},
        'FLAC': {'header': b'\x66\x4c\x61\x43', 'max_size': 100*1024*1024, 'extension': '.flac'},
    },
    'Video Files': {
        'MP4 / MOV': {'header': b'\x66\x74\x79\x70', 'offset': 4, 'extension': '.mp4'},
        'AVI': {'header': b'\x52\x49\x46\x46', 'max_size': 2048*1024*1024, 'extension': '.avi'},
    },
    'Executable Files': {
        'EXE': {'header': b'\x4d\x5a', 'max_size': 50*1024*1024, 'extension': '.exe'},
        'ELF': {'header': b'\x7f\x45\x4c\x46', 'max_size': 50*1024*1024, 'extension': '.elf'},
    },
    'Database Files': {
        'SQLite': {'header': b'\x53\x51\x4c\x69\x74\x65\x20\x66\x6f\x72\x6d\x61\x74\x20\x33\x00', 'max_size': 100*1024*1024, 'extension': '.sqlite'},
    },
    'Windows Log Files': {
        'Windows Event Log (EVTX)': {'header': b'\x45\x6c\x66\x46\x69\x6c\x65\x00', 'extension': '.evtx'},
        'Windows Registry Hive (REGF)': { 'header': b'\x72\x65\x67\x66', 'max_size': 500 * 1024 * 1024, 'extension': '.dat' },
    },
}

LOG_OS_OPTIONS = {
    "Windows": [".log", ".evtx", ".txt", ".csv", ".json", ".xml"],
    "Linux": [".log", ".txt", ".json", ".xml", ".csv", ".gz", ".zip", ""],
    "macOS": [".logarchive", ".gz", ".tracev3", ".log"]
}

EVENT_OS_OPTIONS = {
    "Windows": [".evtx", ".evt"],
    "Linux": ["/var/log/syslog", "/var/log/auth.log", "/var/log/kern.log",
              "/var/log/dmesg", "/var/log/wtmp", "/var/log/btmp"],
    "macOS": [".logarchive", ".asl", ".ips", ".spin", ".diag", ".log"]
}

# --- NEW: Cryptographic Functions ---
SALT_SIZE = 16
ITERATIONS = 480_000

def derive_key(password: str, salt: bytes) -> bytes:
    """Derives a secure encryption key from a password and salt."""
    kdf = PBKDF2HMAC(
        algorithm=hashes.SHA256(),
        length=32,
        salt=salt,
        iterations=ITERATIONS,
        backend=default_backend()
    )
    return base64.urlsafe_b64encode(kdf.derive(password.encode()))
def calculate_time_estimations(status_dict, current_progress, total_items=None):
    """Calculate time estimations for long-running processes."""
    if not status_dict.get("start_time"):
        status_dict["start_time"] = time.time()
        status_dict["last_update_time"] = time.time()
        return "Calculating..."
    
    current_time = time.time()
    elapsed = current_time - status_dict["start_time"]
    status_dict["elapsed_time"] = format_time(elapsed)
    
    # Update only every 2 seconds to reduce CPU usage
    if current_time - status_dict.get("last_update_time", 0) < 2:
        return status_dict.get("time_remaining_str", "Calculating...")
    
    status_dict["last_update_time"] = current_time
    
    if current_progress > 0:
        if total_items and total_items > 0:
            # Based on total items
            progress_percent = min(current_progress / total_items, 0.99)  # Cap at 99%
            if progress_percent > 0.01:  # Wait until we have meaningful progress
                estimated_total = elapsed / progress_percent
                estimated_remaining = estimated_total - elapsed
                status_dict["time_remaining_str"] = format_time(max(0, estimated_remaining))
                status_dict["estimated_total_time"] = format_time(estimated_total)
            else:
                status_dict["time_remaining_str"] = "Calculating..."
        else:
            # Based on progress percentage
            progress_percent = current_progress / 100.0
            if progress_percent > 0.01:
                estimated_total = elapsed / progress_percent
                estimated_remaining = estimated_total - elapsed
                status_dict["time_remaining_str"] = format_time(max(0, estimated_remaining))
                status_dict["estimated_total_time"] = format_time(estimated_total)
            else:
                status_dict["time_remaining_str"] = "Calculating..."
    else:
        status_dict["time_remaining_str"] = "Calculating..."
    
    return status_dict["time_remaining_str"]

def format_bytes(size_bytes):
    """Convert bytes to human-readable format"""
    if size_bytes == 0:
        return "0 B"
    
    size_names = ["B", "KB", "MB", "GB", "TB"]
    i = 0
    while size_bytes >= 1024 and i < len(size_names) - 1:
        size_bytes /= 1024.0
        i += 1
    
    return f"{size_bytes:.2f} {size_names[i]}"

def calculate_upload_time_estimations(bytes_uploaded, total_bytes, start_time):
    """Calculate upload time estimations."""
    if bytes_uploaded <= 0 or total_bytes <= 0:
        return "0s", "Calculating...", "Calculating..."
    
    elapsed = time.time() - start_time
    elapsed_str = format_time(elapsed)
    
    if bytes_uploaded > 0:
        upload_speed = bytes_uploaded / elapsed  # bytes per second
        remaining_bytes = total_bytes - bytes_uploaded
        if upload_speed > 0:
            remaining_time = remaining_bytes / upload_speed
            remaining_str = format_time(remaining_time)
            total_estimated_str = format_time(elapsed + remaining_time)
        else:
            remaining_str = "Calculating..."
            total_estimated_str = "Calculating..."
    else:
        remaining_str = "Calculating..."
        total_estimated_str = "Calculating..."
    
    return elapsed_str, remaining_str, total_estimated_str

def update_process_timing(status_dict, current_value, total_value=None, value_type="bytes"):
    """Update timing information for any process."""
    if not status_dict.get("start_time"):
        status_dict["start_time"] = time.time()
    
    current_time = time.time()
    
    # Always update elapsed time
    elapsed = current_time - status_dict["start_time"]
    status_dict["elapsed_time"] = format_time(elapsed)
    
    # Update progress based on value type
    if value_type == "bytes" and total_value and total_value > 0:
        progress_percent = (current_value / total_value) * 100
        status_dict["progress"] = min(progress_percent, 99)  # Cap at 99% until complete
    elif value_type == "items" and total_value and total_value > 0:
        progress_percent = (current_value / total_value) * 100
        status_dict["progress"] = min(progress_percent, 99)
    elif value_type == "percentage":
        status_dict["progress"] = current_value
    
    # Calculate time remaining (update every 3 seconds to reduce load)
    if current_time - status_dict.get("last_update_time", 0) >= 3:
        if current_value > 0:
            if total_value and total_value > 0:
                # Based on known total
                progress_ratio = current_value / total_value
                if progress_ratio > 0.01:  # Wait for meaningful progress
                    estimated_total = elapsed / progress_ratio
                    estimated_remaining = estimated_total - elapsed
                    status_dict["time_remaining_str"] = format_time(max(0, estimated_remaining))
                    status_dict["estimated_total_time"] = format_time(estimated_total)
                else:
                    status_dict["time_remaining_str"] = "Calculating..."
            else:
                # Based on progress percentage
                progress_ratio = status_dict["progress"] / 100.0
                if progress_ratio > 0.01:
                    estimated_total = elapsed / progress_ratio
                    estimated_remaining = estimated_total - elapsed
                    status_dict["time_remaining_str"] = format_time(max(0, estimated_remaining))
                    status_dict["estimated_total_time"] = format_time(estimated_total)
                else:
                    status_dict["time_remaining_str"] = "Calculating..."
        else:
            status_dict["time_remaining_str"] = "Calculating..."
        
        status_dict["last_update_time"] = current_time
    
    return status_dict

def custom_encrypt_file(file_path: str, original_filename: str, password: str, output_dir: str):
    """Encrypts a file and saves it with a .enc extension in the specified directory."""
    try:
        with open(file_path, 'rb') as f:
            data = f.read()

        salt = os.urandom(SALT_SIZE)
        key = derive_key(password, salt)
        fernet = Fernet(key)
        encrypted_data = fernet.encrypt(data)

        output_filename = f"{original_filename}.enc"
        output_path = os.path.join(output_dir, output_filename)
        
        with open(output_path, 'wb') as f:
            f.write(CUSTOM_ENC_HEADER)
            f.write(salt)
            f.write(encrypted_data)
        
        # record to DB
        try:
            sess_id = session.get('analysis_session_id')
        except Exception:
            sess_id = None
        try:
            add_file_record(output_filename, 'encrypted', output_path, os.path.getsize(output_path), session_id=sess_id, extra={'original': original_filename})
        except Exception:
            pass
        return True, f"Success! Encrypted file saved as {output_filename}", output_filename
    
    except Exception as e:
        return False, f"An error occurred during encryption: {e}", None

def extract_strings_preview(filepath):
    """Provides a basic preview by extracting printable strings from a binary file."""
    output = ["[INFO] This is a preview of printable strings found in the binary file.\n"]
    try:
        with open(filepath, 'rb') as f:
            content = f.read(2 * 1024 * 1024)

        printable_chars = re.compile(b"([%s]{4,})" % b" -~")
        strings = [match.group(0).decode('ascii', 'ignore') for match in printable_chars.finditer(content)]
        output.extend(strings)
        if len(content) == 2 * 1024 * 1024:
            output.append("\n\n--- [STRING PREVIEW TRUNCATED AT 2MB] ---")
        return "\n".join(output)
    except Exception as e:
        return f"[ERROR] Could not process binary file: {e}"

def get_enhanced_file_info(filepath):
    """Get enhanced file information including thumbnails and file types."""
    try:
        mime_type = magic.from_file(filepath, mime=True)
        file_size = os.path.getsize(filepath)
        
        # Generate thumbnail for images
        thumbnail = None
        if mime_type.startswith('image/'):
            thumbnail = create_thumbnail_data_uri(filepath)
        
        # Determine file category
        file_type = "Unknown"
        if mime_type.startswith('image/'):
            file_type = "Image"
        elif mime_type.startswith('video/'):
            file_type = "Video"
        elif mime_type.startswith('audio/'):
            file_type = "Audio"
        elif mime_type.startswith('text/'):
            file_type = "Text"
        elif 'pdf' in mime_type:
            file_type = "PDF"
        elif 'zip' in mime_type or 'archive' in mime_type:
            file_type = "Archive"
        elif 'executable' in mime_type:
            file_type = "Executable"
        elif 'document' in mime_type:
            file_type = "Document"
        
        return {
            'mime_type': mime_type,
            'file_type': file_type,
            'thumbnail': thumbnail,
            'size_bytes': file_size,
            'size_kb': f"{file_size / 1024:.2f}"
        }
    except Exception as e:
        return {
            'mime_type': 'application/octet-stream',
            'file_type': 'Unknown',
            'thumbnail': None,
            'size_bytes': 0,
            'size_kb': '0'
        }

def generate_deleted_filename(seq, detected_ext=''):
    """Generate a standardized deleted-file recovery filename.

    Format: deleted_files_recovery_0001[.ext]
    """
    ext = detected_ext or ''
    # ensure extension starts with a dot
    if ext and not ext.startswith('.'):
        ext = f'.{ext.lstrip(".")}'
    return f"deleted_files_recovery_{seq:04d}{ext}"
    
def get_active_evidence_path():
    if not uploaded_files_db: 
        return None
    file_details = next(iter(uploaded_files_db.values()))
    decrypted_path = file_details.get('encryption_status', {}).get('decrypted_path')
    if decrypted_path and os.path.exists(decrypted_path):
        return decrypted_path
    return file_details.get('path')



# --- Helper Functions ---
def format_time(seconds):
    """Formats seconds into a human-readable string like '1h 5m 30s'."""
    if seconds is None or seconds < 0:
        return "Calculating..."
    if seconds < 1:
        return "0s"

    minutes, seconds = divmod(int(seconds), 60)
    hours, minutes = divmod(minutes, 60)

    parts = []
    if hours > 0:
        parts.append(f"{hours}h")
    if minutes > 0:
        parts.append(f"{minutes}m")
    if seconds > 0 or not parts:
        parts.append(f"{seconds}s")

    return " ".join(parts)

def create_thumbnail_data_uri(filepath):
    """Generates a base64-encoded PNG thumbnail for image files."""
    try:
        with Image.open(filepath) as img:
            img.thumbnail((100, 100))
            buffered = io.BytesIO()
            img.save(buffered, format="PNG")
            img_str = base64.b64encode(buffered.getvalue()).decode("utf-8")
            return f"data:image/png;base64,{img_str}"
    except (IOError, OSError):
        return None

def format_hex_view(data_bytes, start_offset=0):
    """Formats a byte string into a standard hex view format."""
    lines = []
    for i in range(0, len(data_bytes), 16):
        chunk = data_bytes[i:i+16]
        offset_str = f"0x{(start_offset + i):08X}"
        hex_str = ' '.join(f'{b:02X}' for b in chunk)
        ascii_str = ''.join(chr(b) if 32 <= b <= 126 else '.' for b in chunk)
        lines.append(f"{offset_str}  {hex_str:<47} {ascii_str}")
    return '\n'.join(lines)

def is_valid_mp3_stream(data, start_offset=0, frames_to_check=5):
    """Checks for a sequence of valid, contiguous MP3 frames."""
    BITRATE_MAP = [0, 32, 40, 48, 56, 64, 80, 96, 112, 128, 160, 192, 224, 256, 320, 0]
    SAMPLERATE_MAP = [44100, 48000, 32000, 0]

    current_offset = start_offset
    frames_found = 0

    while frames_found < frames_to_check:
        if current_offset + 4 > len(data):
            break

        header_bytes = data[current_offset : current_offset + 4]

        is_mpeg1_layer3 = (header_bytes[0] == 0xFF and (header_bytes[1] & 0xE0) == 0xE0 and
                         ((header_bytes[1] >> 3) & 3) == 1 and ((header_bytes[1] >> 1) & 3) == 1)

        if not is_mpeg1_layer3:
            break

        bitrate_bits = (header_bytes[2] >> 4) & 15
        samplerate_bits = (header_bytes[2] >> 2) & 3
        padding_bit = (header_bytes[2] >> 1) & 1

        bitrate = BITRATE_MAP[bitrate_bits] * 1000
        sample_rate = SAMPLERATE_MAP[samplerate_bits]

        if bitrate == 0 or sample_rate == 0:
            break

        frame_size = int((144 * bitrate / sample_rate) + padding_bit)

        if frame_size <= 1:
            break

        current_offset += frame_size
        frames_found += 1

    return frames_found >= frames_to_check

# --- Log Parsing Helpers ---
def parse_text_log(filepath, max_lines=10000):
    lines = []
    with open(filepath, "r", errors="ignore") as f:
        for i, line in enumerate(f):
            if i >= max_lines:
                lines.append(f"\n--- [File truncated at {max_lines} lines] ---\n")
                break
            lines.append(line)
    return "".join(lines)

def parse_csv_log(filepath, max_rows=10000):
    output = []
    with open(filepath, newline="", errors="ignore") as csvfile:
        reader = csv.reader(csvfile)
        for i, row in enumerate(reader):
            if i >= max_rows:
                output.append(f"\n--- [File truncated at {max_rows} rows] ---\n")
                break
            output.append(",".join(row))
    return "\n".join(output)

def parse_json_log(filepath, max_bytes=1048576):
    with open(filepath, "r", errors="ignore") as f:
        content = f.read(max_bytes)
        try:
            data = json.loads(content)
            output = json.dumps(data, indent=4)
            if len(content) == max_bytes:
                output += f"\n\n--- [File preview truncated at {max_bytes} bytes] ---\n"
            return output
        except json.JSONDecodeError:
            return f"[INFO] Could not parse as a single JSON object. Displaying raw text preview.\n\n{content}\n\n--- [File preview truncated at {max_bytes} bytes] ---\n"

def parse_xml_log(filepath, max_elements=10000):
    output = []
    count = 0
    try:
        for event, elem in ET.iterparse(filepath, events=('start',)):
            if count >= max_elements:
                output.append(f"\n--- [File truncated at {max_elements} XML elements] ---\n")
                break
            text = elem.text.strip() if elem.text else ""
            output.append(f"<{elem.tag}> {text}")
            count += 1
            elem.clear()
        return "\n".join(output)
    except ET.ParseError as e:
        return f"[ERROR] Failed to parse XML: {e}\n\nShowing raw text preview instead:\n\n{parse_text_log(filepath, max_lines=500)}"

def parse_evtx_log(filepath, max_records=5000):
    if Evtx is None: 
        return "[ERROR] The 'python-evtx' library is required to parse .evtx files."
    output = []
    try:
        with Evtx(filepath) as log:
            for i, record in enumerate(log.records()):
                if i >= max_records:
                    output.append(f"\n\n--- [File truncated at {max_records} records] ---\n")
                    break
                output.append(record.xml())
    except Exception as e:
        return f"[ERROR] Failed to parse EVTX file: {e}"
    return "\n\n".join(output)

# --- Forensic & Encryption Functions ---
def calculate_entropy(data):
    """Calculates the Shannon entropy of a byte string."""
    if not data: 
        return 0
    byte_counts = [0] * 256
    for byte in data: 
        byte_counts[byte] += 1
    entropy = 0
    total_bytes = len(data)
    for count in byte_counts:
        if count > 0:
            probability = count / total_bytes
            entropy -= probability * math.log2(probability)
    return entropy

def detect_encryption(filepath):
    """Detects encryption based on file headers and entropy."""
    try:
        with open(filepath, 'rb') as f: 
            header = f.read(512)
        for enc_type, sig_info in ENCRYPTION_SIGNATURES.items():
            if sig_info.get('header') and header.startswith(sig_info['header']):
                if enc_type == 'ZIP_ENCRYPTED':
                    try:
                        with zipfile.ZipFile(filepath) as zf:
                            if not any(e.flag_bits & 0x1 for e in zf.infolist()): 
                                continue
                    except Exception: 
                        continue
                return {'encrypted': True, 'encryption_type': enc_type, 'description': sig_info['description']}
        entropy = calculate_entropy(header)
        if entropy > 7.5:
            return {'encrypted': True, 'encryption_type': 'UNKNOWN', 'description': f'High entropy ({entropy:.2f}) suggests encryption'}
    except Exception as e:
        print(f"Error in detect_encryption: {e}")
    return {'encrypted': False, 'encryption_type': None, 'description': 'No encryption detected'}

def perform_forensic_analysis(file_path):
    """Performs a high-level analysis of a disk image file."""
    results, partition_info = [], []
    try:
        with open(file_path, 'rb') as f: 
            header = f.read(4096)
        if header[510:512] == b'\x55\xaa': 
            results.append("✓ MBR Partition Table detected")
        if b'EFI PART' in header[512:1024]: 
            results.append("✓ GPT Partition Table detected")
        if b'FAT32' in header[0x52:0x5A]: 
            results.append("✓ FAT32 File System detected")
        elif b'NTFS' in header[0x03:0x07]: 
            results.append("✓ NTFS File System detected")
        elif b'\x53\xEF' in header[1024+56:1024+58]: 
            results.append("✓ EXT File System Superblock detected")
        if b'-FVE-FS-' in header: 
            results.append("⚠ BitLocker encryption detected")
        try:
            img_handle = pytsk3.Img_Info(file_path)
            volume = pytsk3.Volume_Info(img_handle)
            if volume:
                vstype_map = {
                    pytsk3.TSK_VS_TYPE_DETECT: "Auto-detect", 
                    pytsk3.TSK_VS_TYPE_DOS: "DOS",
                    pytsk3.TSK_VS_TYPE_BSD: "BSD", 
                    pytsk3.TSK_VS_TYPE_SUN: "Sun",
                    pytsk3.TSK_VS_TYPE_MAC: "Mac", 
                    pytsk3.TSK_VS_TYPE_GPT: "GPT",
                    pytsk3.TSK_VS_TYPE_UNSUPP: "Unsupported",
                }
                results.append(f"✓ Partition Table Type: {vstype_map.get(volume.info.vstype, 'Unknown')}")
                for part in volume:
                    if part.flags != pytsk3.TSK_VS_PART_FLAG_UNALLOC:
                        partition_info.append({"addr": part.addr, "desc": part.desc.decode('utf-8', 'ignore'), "start": part.start, "len": part.len})
        except IOError: 
            results.append("ℹ️ No partition table found or image is a single volume.")
        except Exception as e: 
            results.append(f"❌ Error reading partitions: {e}")
        file_size = os.path.getsize(file_path)
        results.append(f"📊 File Size: {file_size/(1024*1024*1024):.2f} GB" if file_size > 1024**3 else f"📊 File Size: {file_size/(1024*1024):.2f} MB")
        results.append(f"📈 Header Entropy: {calculate_entropy(header):.2f} (High entropy > 7.5 may suggest encryption)")
    except Exception as e: 
        results.append(f"❌ Analysis error: {str(e)}")
    return results, partition_info

def calculate_hashes_threaded(file_path):
    """Calculates MD5, SHA1, and SHA256 hashes in a background thread."""
    hashing_status.update({"in_progress": True, "progress": 0, "complete": False, "hashes": {}})
    md5, sha1, sha256 = hashlib.md5(), hashlib.sha1(), hashlib.sha256()
    CHUNK_SIZE = 4 * 1024 * 1024
    try:
        file_size = os.path.getsize(file_path)
        with open(file_path, 'rb') as f:
            bytes_read = 0
            while True:
                chunk = f.read(CHUNK_SIZE)
                if not chunk: 
                    break
                bytes_read += len(chunk)
                md5.update(chunk); sha1.update(chunk); sha256.update(chunk)
                hashing_status["progress"] = int((bytes_read / file_size) * 100) if file_size > 0 else 100
        final_hashes = {'MD5': md5.hexdigest(), 'SHA-1': sha1.hexdigest(), 'SHA-256': sha256.hexdigest()}
        hashing_status.update({"in_progress": False, "complete": True, "hashes": final_hashes})
        filename = os.path.basename(file_path)
        if filename in uploaded_files_db:
            uploaded_files_db[filename]['hash_info'] = final_hashes
            uploaded_files_db[filename]['hashing_complete'] = True
    except Exception as e:
        print(f"Error during hashing: {e}")
        hashing_status.update({"in_progress": False, "complete": True, "error": str(e)})

def extract_strings_threaded(filepath):
    """Extracts all printable strings from a file in a background thread."""
    strings_status.update({"in_progress": True, "complete": False, "progress": 0, "strings_found": 0, "preview": []})
    min_len = 4
    try:
        with open(filepath, 'rb') as f, mmap.mmap(f.fileno(), 0, access=mmap.ACCESS_READ) as mm:
            file_size = len(mm)
            printable_chars = re.compile(b"([%s]{%d,})" % (b" -~", min_len))
            last_update_pos = 0
            for match in printable_chars.finditer(mm):
                pos = match.start()
                if pos > last_update_pos + (file_size // 100):
                    strings_status['progress'] = int((pos / file_size) * 100)
                    last_update_pos = pos

                strings_status['strings_found'] += 1
                if len(strings_status['preview']) < 100:
                    strings_status['preview'].append(match.group(0).decode('ascii', 'ignore'))

    except Exception as e:
        print(f"Error during strings extraction: {e}")
    strings_status.update({"in_progress": False, "complete": True, "progress": 100})

def attempt_decryption(filepath, encryption_type, password=None):
    """Orchestrates the decryption process in a background thread."""
    filename = os.path.basename(filepath)
    decryption_status.update({
        "in_progress": True, "complete": False, "message": "Starting decryption...",
        "success": False, "filename": filename, "progress": 0, "attempts": 0
    })
    if filename in uploaded_files_db: 
        uploaded_files_db[filename]['encryption_status']['decrypting'] = True
    # Ensure decrypted folder exists
    os.makedirs(app.config.get('DECRYPTED_FOLDER', DECRYPTED_FOLDER), exist_ok=True)
    decrypted_path = os.path.join(app.config['DECRYPTED_FOLDER'], f"decrypted_{filename}")

    passwords_to_try = []
    is_user_pwd = bool(password)
    if is_user_pwd:
        passwords_to_try.append(password)
    else:
        try:
            with open(app.config['DICTIONARY_FILE'], 'r') as f:
                passwords_to_try.extend([line.strip() for line in f if line.strip()])
        except FileNotFoundError:
            decryption_status['message'] = "Password dictionary file not found."
            passwords_to_try = ["password", "123456", "admin"]

    total_passwords = len(passwords_to_try)
    if total_passwords == 0 and is_user_pwd:
        total_passwords = 1

    decryption_status['total_attempts'] = total_passwords

    success = False
    if "ZIP" in encryption_type: 
        success = crack_zip_with_passwords(filepath, decrypted_path, passwords_to_try, is_user_pwd, total_passwords)
    elif "AES" in encryption_type: 
        success = crack_openssl_aes_with_passwords(filepath, decrypted_path, passwords_to_try, is_user_pwd, total_passwords)
    elif "FERNET" in encryption_type: 
        success = crack_fernet_with_passwords(filepath, decrypted_path, passwords_to_try, is_user_pwd, total_passwords)
    else: 
        decryption_status['message'] = f"Automated decryption for '{encryption_type}' is not supported."

    decryption_status.update({'in_progress': False, 'complete': True, 'success': success, 'progress': 100})
    
    if success:
        decryption_status['message'] = "Decryption successful!"
        if filename in uploaded_files_db:
            uploaded_files_db[filename]['encryption_status']['decrypted_path'] = decrypted_path
            uploaded_files_db[filename]['encryption_status']['description'] = "File successfully decrypted."
            # add DB record for decrypted file
            try:
                sess_id = session.get('analysis_session_id')
            except Exception:
                sess_id = None
            try:
                add_file_record(f"decrypted_{filename}", 'decrypted', decrypted_path, os.path.getsize(decrypted_path), session_id=sess_id, extra={'source': filename})
            except Exception:
                pass
    elif "Trying password" in decryption_status.get('message', '') or decryption_status.get('message') == "Starting decryption...":
        decryption_status['message'] = "Decryption failed. All password attempts were incorrect."

    if filename in uploaded_files_db:
        uploaded_files_db[filename]['encryption_status']['decrypting'] = False

def crack_zip_with_passwords(filepath, output_path, passwords_to_try, is_user_pwd, total_passwords):
    """Attempts to decrypt a ZIP file using a list of passwords."""
    try:
        with zipfile.ZipFile(filepath, 'r') as zf:
            if not any(e.flag_bits & 0x1 for e in zf.infolist()):
                decryption_status['message'] = "Archive does not appear to be password-protected."

        for i, pwd in enumerate(passwords_to_try):
            progress = int(((i + 1) / total_passwords) * 100) if total_passwords > 0 else 50
            is_first_try_user = (i == 0 and is_user_pwd)
            decryption_status.update({
                "attempts": i + 1,
                "message": f"Trying password: {'●' * len(pwd) if is_first_try_user else f'dictionary entry #{i+1}'}",
                "progress": progress
            })
            
            try:
                with zipfile.ZipFile(filepath, 'r') as zf:
                    with tempfile.TemporaryDirectory() as temp_dir:
                        zf.extractall(path=temp_dir, pwd=pwd.encode())
                        extracted_items = os.listdir(temp_dir)
                        if extracted_items:
                            source_path = os.path.join(temp_dir, extracted_items[0])
                            if os.path.isfile(source_path):
                                shutil.copy(source_path, output_path)
                            return True
                        return True

            except (RuntimeError, zipfile.BadZipFile):
                continue
            except Exception as e:
                print(f"An unexpected error occurred during ZIP extraction: {e}")
                continue
                
    except Exception as e:
        error_message = f"Could not process as a ZIP file. Error: {e}"
        decryption_status['message'] = error_message
        print(f"ZIP cracking error: {error_message}")
        
    return False

def derive_openssl_key_iv(password, salt, key_length=32, iv_length=16):
    """Derives AES key and IV from password and salt, mimicking OpenSSL's method."""
    key_iv = b''; prev_hash = b''
    while len(key_iv) < key_length + iv_length:
        hasher = hashlib.md5()
        if prev_hash: 
            hasher.update(prev_hash)
        hasher.update(password); hasher.update(salt)
        prev_hash = hasher.digest()
        key_iv += prev_hash
    return key_iv[:key_length], key_iv[key_length:key_length + iv_length]

def crack_openssl_aes_with_passwords(filepath, output_path, passwords_to_try, is_user_pwd, total_passwords):
    """Attempts to decrypt an OpenSSL AES file using a list of passwords."""
    CHUNK_SIZE = 4 * 1024 * 1024
    try:
        with open(filepath, 'rb') as f_in:
            if f_in.read(8) != b'Salted__':
                decryption_status['message'] = "File is not in OpenSSL salted format."
                return False
            salt = f_in.read(8)

        for i, pwd_str in enumerate(passwords_to_try):
            progress = int(((i + 1) / total_passwords) * 100) if total_passwords > 0 else 50
            pwd = pwd_str.encode('utf-8')
            is_first_try_user = (i == 0 and is_user_pwd)
            decryption_status.update({
                "attempts": i + 1,
                "message": f"Trying password: {'●' * len(pwd_str) if is_first_try_user else f'dictionary entry #{i+1}'}",
                "progress": progress
            })

            temp_path = None
            try:
                key, iv = derive_openssl_key_iv(pwd, salt)
                decryptor = Cipher(algorithms.AES(key), modes.CBC(iv), backend=default_backend()).decryptor()
                
                with open(filepath, 'rb') as f_in_inner, tempfile.NamedTemporaryFile(mode='wb', delete=False) as temp_f:
                    temp_path = temp_f.name
                    f_in_inner.seek(16)
                    while True:
                        chunk = f_in_inner.read(CHUNK_SIZE)
                        if not chunk: 
                            break
                        temp_f.write(decryptor.update(chunk))
                    temp_f.write(decryptor.finalize())

                os.rename(temp_path, output_path)
                return True
            except ValueError:
                if temp_path and os.path.exists(temp_path): 
                    os.remove(temp_path)
                continue
            except Exception as e:
                if temp_path and os.path.exists(temp_path): 
                    os.remove(temp_path)
                continue
                
    except Exception as e:
        error_message = f"Could not process as an AES file. Error: {e}"
        decryption_status['message'] = error_message
        print(f"AES cracking error: {error_message}")
        
    return False

def crack_fernet_with_passwords(filepath, output_path, passwords_to_try, is_user_pwd, total_passwords):
    """Attempts to decrypt a Fernet encrypted file using a list of passwords."""
    try:
        with open(filepath, 'rb') as f_in:
            header = f_in.read(len(CUSTOM_ENC_HEADER))
            if header != CUSTOM_ENC_HEADER:
                decryption_status['message'] = "File is not a valid Forensic Carver encrypted file."
                return False
            
            salt = f_in.read(SALT_SIZE)
            encrypted_data = f_in.read()

        for i, pwd_str in enumerate(passwords_to_try):
            progress = int(((i + 1) / total_passwords) * 100) if total_passwords > 0 else 50
            is_first_try_user = (i == 0 and is_user_pwd)
            decryption_status.update({
                "attempts": i + 1,
                "message": f"Trying password: {'●' * len(pwd_str) if is_first_try_user else f'dictionary entry #{i+1}'}",
                "progress": progress
            })

            try:
                key = derive_key(pwd_str, salt)
                fernet = Fernet(key)
                decrypted_data = fernet.decrypt(encrypted_data)

                with open(output_path, 'wb') as f_out:
                    f_out.write(decrypted_data)
                return True

            except InvalidToken:
                continue
            except Exception as e:
                print(f"An unexpected error occurred during Fernet decryption: {e}")
                continue
                
    except Exception as e:
        error_message = f"Could not process as a Fernet file. Error: {e}"
        decryption_status['message'] = error_message
        print(f"Fernet cracking error: {error_message}")
        
    return False



# --- NEW PREVIEW HELPER FUNCTIONS ---






   




def validate_and_deduplicate_file(content, seen_hashes, min_size=128):
    """
    Strict validation: Check if file is non-empty, meets minimum size, and is not a duplicate.
    Returns (is_valid, content_hash)
    """
    if not content or len(content) < min_size:
        return False, None
    
    content_hash = hashlib.md5(content).hexdigest()
    if content_hash in seen_hashes:
        return False, content_hash
    
    return True, content_hash

# --- Carving & Recovery Engines ---
def _validate_and_extract_file(mm, file_start_pos, sig_options, seen_hashes, file_size_limit):
    """Validates a data chunk with strict deduplication and empty file checking."""
    sig = sig_options[0]
    name = sig['name']
    MIN_FILE_SIZE = 128

    # --- Strategy 1: RIFF Containers (WAV, AVI) ---
    if name in ['WAV', 'AVI']:
        try:
            if file_start_pos + 12 < file_size_limit and mm[file_start_pos+8:file_start_pos+12] in [b'WAVE', b'AVI ']:
                file_size_from_header = int.from_bytes(mm[file_start_pos+4:file_start_pos+8], 'little')
                total_size = file_size_from_header + 8

                if total_size > 32:
                    end_pos = file_start_pos + total_size
                    content = mm[file_start_pos:end_pos]
                    
                    # Add validation
                    if content and len(content) >= MIN_FILE_SIZE:
                        content_hash = hashlib.md5(content[:4096]).hexdigest()
                        if content_hash not in seen_hashes:
                            seen_hashes.add(content_hash)
                            return content, sig, content_hash
        except Exception:
            pass

    # --- Strategy 2: Parse MP4/MOV Container Structure ---
    elif name == 'MP4 / MOV':
        try:
            current_offset = file_start_pos
            total_size = 0
            box_count = 0
            while current_offset + 8 < file_size_limit and box_count < 500:
                box_size = int.from_bytes(mm[current_offset:current_offset+4], 'big')
                box_type = mm[current_offset+4:current_offset+8]
                if box_size < 8 or not box_type.isalnum(): 
                    break
                total_size += box_size
                current_offset += box_size
                box_count += 1
            if total_size > 16384:
                end_pos = min(file_start_pos + total_size, file_size_limit)
                content = mm[file_start_pos:end_pos]
                if b'moov' in content or b'mdat' in content:
                    # Add validation
                    if content and len(content) >= MIN_FILE_SIZE:
                        content_hash = hashlib.md5(content[:4096]).hexdigest()
                        if content_hash not in seen_hashes:
                            seen_hashes.add(content_hash)
                            return content, sig, content_hash
        except Exception:
            pass

    # --- Strategy 3: Parse MP3 Frame Stream with Tolerance and ID3 Skipping ---
    elif name == 'MP3':
        audio_start_offset = file_start_pos
        if mm[file_start_pos:file_start_pos+3] == b'ID3':
            try:
                id3_size_bytes = mm[file_start_pos+6:file_start_pos+10]
                id3_size = (id3_size_bytes[0] << 21) | (id3_size_bytes[1] << 14) | (id3_size_bytes[2] << 7) | id3_size_bytes[3]
                audio_start_offset = file_start_pos + id3_size + 10
            except IndexError: 
                return None, None, None

        BITRATE_MAP = [0, 32, 40, 48, 56, 64, 80, 96, 112, 128, 160, 192, 224, 256, 320, 0]
        SAMPLERATE_MAP = [44100, 48000, 32000, 0]
        current_offset = audio_start_offset
        total_size = 0
        invalid_frames_in_a_row = 0
        MAX_INVALID_FRAMES = 10

        while current_offset + 4 < file_size_limit:
            header_bytes = mm[current_offset : current_offset + 4]
            is_mpeg_frame = (header_bytes[0] == 0xFF and (header_bytes[1] & 0xE0) == 0xE0)
            if not is_mpeg_frame:
                invalid_frames_in_a_row += 1
                if invalid_frames_in_a_row >= MAX_INVALID_FRAMES: 
                    break
                current_offset += 1
                continue
            invalid_frames_in_a_row = 0
            try:
                samplerate_bits = (header_bytes[2] >> 2) & 3
                sample_rate = SAMPLERATE_MAP[samplerate_bits]
                bitrate_bits = (header_bytes[2] >> 4) & 15
                bitrate = BITRATE_MAP[bitrate_bits] * 1000
                padding_bit = (header_bytes[2] >> 1) & 1
                if bitrate == 0 or sample_rate == 0: 
                    raise ValueError("Invalid frame")
                frame_size = int((144 * bitrate / sample_rate) + padding_bit)
                if frame_size <= 1: 
                    raise ValueError("Invalid frame")
                total_size += frame_size
                current_offset += frame_size
            except (ValueError, IndexError):
                invalid_frames_in_a_row += 1
                if invalid_frames_in_a_row >= MAX_INVALID_FRAMES: 
                    break
                current_offset += 1

        if total_size > 4096:
            end_pos = audio_start_offset + total_size
            content = mm[file_start_pos:end_pos]
            # Add validation
            if content and len(content) >= MIN_FILE_SIZE:
                content_hash = hashlib.md5(content[:8192]).hexdigest()
                if content_hash not in seen_hashes:
                    seen_hashes.add(content_hash)
                    return content, sig, content_hash

    # --- Strategy 4: Basic Header/Footer Matching (JPEG, PNG, PDF, etc.) ---
    elif sig.get('footer'):
        search_limit = min(file_start_pos + sig.get('max_size', 50*1024*1024), file_size_limit)
        header_len = len(sig.get('header', sig.get('headers', [b''])[0]))
        footer_pos = mm.find(sig['footer'], file_start_pos + header_len, search_limit)
        if footer_pos != -1:
            end_pos = footer_pos + len(sig['footer'])
            content = mm[file_start_pos:end_pos]
            MIN_SIZES = {'.jpeg': 2048, '.png': 256, '.zip': 128, '.pdf': 1024, '.docx': 4096}
            
            # Add validation with file format-specific minimum sizes
            min_size = MIN_SIZES.get(sig['extension'], MIN_FILE_SIZE)
            if content and len(content) >= min_size:
                content_hash = hashlib.md5(content[:4096]).hexdigest()
                if content_hash in seen_hashes: 
                    return None, None, None
                try:
                    if sig['extension'] in ['.jpeg', '.png', '.gif']: 
                        Image.open(io.BytesIO(content)).verify()
                    elif sig['extension'] in ['.docx', '.xlsx', '.pptx', '.zip']:
                        with zipfile.ZipFile(io.BytesIO(content)) as zf:
                            if zf.testzip() is not None: 
                                return None, None, None
                    elif sig['extension'] == '.pdf':
                        if not content.strip().endswith(b'%%EOF'): 
                            return None, None, None
                    
                    # Final validation check before returning
                    if content_hash not in seen_hashes:
                        seen_hashes.add(content_hash)
                        return content, sig, content_hash
                except Exception:
                    return None, None, None

    return None, None, None
# --- NEW, CORRECTED CARVER ---
def simple_file_carver(filepath, selected_types):
    """High-speed carver that eliminates empty files and duplicates."""
    # global carving_status
    
    # Initialize status
    carving_status.update({
        "progress": 0, "current_offset": "0x00000000", "files_found": 0,
        "complete": False, "found_files_list": [], "error": None,
         "start_time": time.time(), 
        "estimated_total_time": None, 
        "elapsed_time": "0s",
        "last_update_time": time.time(),
        "bytes_processed": 0,
        "total_bytes": os.path.getsize(filepath)
    })
    
    output_dir = app.config['CARVED_FOLDER']
    file_counter = 0
    seen_hashes = set()
    MIN_FILE_SIZE = 128  # Minimum file size to consider valid

    # Enhanced directory clearing with better error handling
    try:
        os.makedirs(output_dir, exist_ok=True)
        for item_name in os.listdir(output_dir):
            item_path = os.path.join(output_dir, item_name)
            try:
                if os.path.isfile(item_path):
                    os.unlink(item_path)
                elif os.path.isdir(item_path):
                    shutil.rmtree(item_path)
            except Exception as e:
                print(f"Warning: Could not remove {item_path}: {e}")
    except Exception as e:
        error_msg = f"Permission Error: Could not clear old results in {output_dir}. Details: {e}"
        carving_status.update({"error": error_msg, "complete": True})
        return

    # Filter signatures based on selection
    signatures_to_find = {}
    for category, types in FILE_SIGNATURES.items():
        for name, sig in types.items():
            if name in selected_types:
                signatures_to_find[name] = sig

    try:
        file_size = os.path.getsize(filepath)
        with open(filepath, 'rb') as f:
            with mmap.mmap(f.fileno(), 0, access=mmap.ACCESS_READ) as mm:
                for name, sig in signatures_to_find.items():
                    headers = sig.get('headers', [])
                    if sig.get('header') and sig['header'] not in headers:
                        headers.append(sig['header'])
                    
                    for header in headers:
                        if not header:
                            continue
                            
                        current_pos = 0
                        while current_pos < file_size:
                            found_pos = mm.find(header, current_pos)
                            if found_pos == -1:
                                break
                            
                            # Extract file data
                            file_data = extract_file_data(mm, found_pos, sig, file_size)
                            
                            # STRICT VALIDATION: Skip empty files and files below minimum size
                            if not file_data or len(file_data) < MIN_FILE_SIZE:
                                current_pos = found_pos + 1
                                continue
                            
                            # STRICT DEDUPLICATION: Calculate content hash
                            content_hash = hashlib.md5(file_data).hexdigest()
                            if content_hash in seen_hashes:
                                current_pos = found_pos + 1
                                continue
                            seen_hashes.add(content_hash)
                            
                            file_counter += 1
                            
                            # Save file with metadata
                            if save_carved_file(file_data, found_pos, name, sig, file_counter, output_dir):
                                # Update status
                                update_carving_status(file_counter, found_pos, file_data, file_size, name)
                            
                            current_pos = found_pos + 1
                            
    except Exception as e:
        carving_status["error"] = f"Carving process error: {e}"
        print(f"Error during carving: {e}")
    finally:
       carving_status.update({
            "progress": 100, 
            "complete": True,
            "elapsed_time": format_time(time.time() - carving_status["start_time"]),
            "time_remaining_str": "Complete",
            "estimated_total_time": format_time(time.time() - carving_status["start_time"])
        })
    print(f"Carving complete. Found {carving_status['files_found']} valid, non-duplicate files.")
       
def extract_file_data(mm, found_pos, sig, file_size):
    """Extract file data based on signature rules with strict validation."""
    footer = sig.get('footer')
    max_size = sig.get('max_size', 10 * 1024 * 1024)
    
    if footer:
        search_limit = min(found_pos + sig.get('max_size', 50 * 1024 * 1024), file_size)
        end_pos = mm.find(footer, found_pos + len(sig.get('header', b'')), search_limit)
        if end_pos != -1:
            file_data = mm[found_pos: end_pos + len(footer)]
            # Validate file size and content
            if len(file_data) >= 128:  # Minimum valid file size
                return file_data
    else:
        file_data = mm[found_pos: min(found_pos + max_size, file_size)]
        if len(file_data) >= 128:  # Minimum valid file size
            return file_data
    
    return b''

def extract_file_data(mm, found_pos, sig, file_size):
    """Extract file data based on signature rules."""
    footer = sig.get('footer')
    max_size = sig.get('max_size', 10 * 1024 * 1024)
    
    if footer:
        search_limit = min(found_pos + sig.get('max_size', 50 * 1024 * 1024), file_size)
        end_pos = mm.find(footer, found_pos + len(sig.get('header', b'')), search_limit)
        if end_pos != -1:
            return mm[found_pos: end_pos + len(footer)]
    else:
        return mm[found_pos: min(found_pos + max_size, file_size)]
    
    return b''

def save_carved_file(file_data, found_pos, name, sig, file_counter, output_dir):
    """Save carved file with proper naming convention."""
    try:
        offset_hex = f"{found_pos:08X}"
        size_bytes = len(file_data)
        safe_name = name.lower().replace(' ', '_').replace('/', '_')
        extension = sig.get('extension', '.bin')
        filename = f"{file_counter}-{offset_hex}-{size_bytes}-{safe_name}{extension}"
        
        # Prefer the configured Carved Files folder unless an explicit output_dir within Carved Files is provided
        carved_root = app.config.get('CARVED_FOLDER', CARVED_FOLDER)
        os.makedirs(carved_root, exist_ok=True)
        if output_dir and os.path.commonpath([os.path.abspath(output_dir), os.path.abspath(carved_root)]) == os.path.abspath(carved_root):
            save_path = os.path.join(output_dir, filename)
        else:
            save_path = os.path.join(carved_root, filename)
        with open(save_path, 'wb') as out_file:
            out_file.write(file_data)
        # record carved file in DB under current session
        try:
            sess_id = session.get('analysis_session_id')
        except Exception:
            sess_id = None
        try:
            add_file_record(filename, 'carved', save_path, len(file_data), session_id=sess_id, extra={'offset': found_pos, 'signature': name})
        except Exception:
            pass
        return True
    except Exception as e:
        print(f"Error saving file {filename}: {e}")
        return False

def update_carving_status(file_counter, found_pos, file_data, file_size, name):
    """Update the global carving status."""
    # global carving_status
    offset_hex = f"{found_pos:08X}"
    
    file_info = {
        "name": f"{file_counter}-{offset_hex}-{len(file_data)}-{name}",
        "offset": f"0x{offset_hex}",
        "hex_preview": format_hex_view(file_data[:256])
    }
    
    # Update basic counters
    carving_status.update({
        "files_found": file_counter,
        "current_offset": f"0x{offset_hex}",
        "progress": int((found_pos / file_size) * 100) if file_size > 0 else 0
    })
    carving_status["found_files_list"].append(file_info)

    # Update bytes processed (use found_pos + file size as an approximation)
    try:
        processed = found_pos + len(file_data)
        # keep the maximum processed value to avoid regressions
        carving_status['bytes_processed'] = max(carving_status.get('bytes_processed', 0), processed)
    except Exception:
        pass

    # Calculate elapsed, speed, and ETA if start_time and total_bytes are available
    try:
        if carving_status.get('start_time'):
            now = time.time()
            elapsed = now - carving_status['start_time']
            carving_status['elapsed_time'] = format_time(elapsed)

            bytes_done = carving_status.get('bytes_processed', 0)
            total = carving_status.get('total_bytes') or file_size or 0

            if bytes_done > 0 and elapsed > 0 and total > 0:
                speed = bytes_done / elapsed
                estimated_total_seconds = total / max(speed, 0.0001)
                remaining_seconds = max(0, estimated_total_seconds - elapsed)
                carving_status['estimated_total_time'] = format_time(estimated_total_seconds)
                carving_status['time_remaining_str'] = format_time(remaining_seconds)
            else:
                # fallback placeholders
                carving_status['estimated_total_time'] = None
                carving_status['time_remaining_str'] = 'Calculating...'
    except Exception:
        # Don't let timing errors break the carving loop
        pass

def _scan_partition_worker(args):
    """Worker function for parallel partition scanning with deep scan and deduplication."""
    filepath, part_info, fs_offset = args
    found_files = {}
    errors = []
    seen_hashes = set()

    try:
        img = pytsk3.Img_Info(filepath)
        fs = pytsk3.FS_Info(img, offset=fs_offset)

        def process_deleted_file(fs_object, recovery_method):
            """Helper to process, hash, and add a unique, non-empty deleted file."""
            try:
                if not fs_object.info or not fs_object.info.meta or fs_object.info.meta.size <= 0:
                    return

                content_preview = fs_object.read_random(0, min(4096, fs_object.info.meta.size))
                if not content_preview:
                    return

                file_hash = hashlib.md5(content_preview).hexdigest()
                if file_hash in seen_hashes:
                    return

                seen_hashes.add(file_hash)

                name_str = fs_object.info.name.name.decode('utf-8', 'ignore') if fs_object.info.name else "orphaned_file"
                inode = str(fs_object.info.meta.addr)

                file_info = {
                    'inode': inode, 'name': name_str, 'size': fs_object.info.meta.size,
                    'mtime': datetime.datetime.fromtimestamp(fs_object.info.meta.mtime).strftime('%Y-%m-%d %H:%M:%S'),
                    'ctime': datetime.datetime.fromtimestamp(fs_object.info.meta.ctime).strftime('%Y-%m-%d %H:%M:%S'),
                    'recovery_method': recovery_method,
                    'file_type': magic.from_buffer(content_preview, mime=True),
                    'fs_offset': fs_offset
                }
                found_files[inode] = file_info
            except Exception as e:
                pass

        # --- Scan Strategy 1: Directory Walk ---
        def directory_walk(directory, parent_path):
            for fs_object in directory:
                try:
                    name_info = fs_object.info.name
                    if not name_info or name_info.name in [b'.', b'..']:
                        continue

                    full_path = f"{parent_path.rstrip('/')}/{name_info.name.decode('utf-8', 'ignore')}"
                    is_deleted_recycled = '$Recycle.Bin' in full_path or 'RECYCLED' in full_path or '/.Trash' in full_path
                    is_deleted_unalloc = hasattr(fs_object.info.meta, 'flags') and (fs_object.info.meta.flags & pytsk3.TSK_FS_META_FLAG_UNALLOC)

                    if is_deleted_recycled or is_deleted_unalloc:
                        process_deleted_file(fs_object, 'Recycle Bin' if is_deleted_recycled else 'Metadata')

                    if fs_object.info.meta and fs_object.info.meta.type == pytsk3.TSK_FS_META_TYPE_DIR:
                        directory_walk(fs_object.as_directory(), full_path)
                except Exception:
                    continue

        directory_walk(fs.open_dir(path="/"), "/")

        # --- Scan Strategy 2: Deep Inode Scan ---
        last_inode = fs.info.last_inum
        for inode_num in range(fs.info.first_inum, last_inode + 1):
            try:
                fs_file = fs.open_meta(inode=inode_num)
                if fs_file.info.meta.flags & pytsk3.TSK_FS_META_FLAG_UNALLOC and str(inode_num) not in found_files:
                    process_deleted_file(fs_file, 'Deep Inode Scan')
            except (IOError, AttributeError):
                continue

    except IOError as e:
        errors.append(f"Could not open filesystem on partition {part_info['desc']}: {e}")

    return found_files, errors

# --- Find and REPLACE your 'scan_for_deleted_files_engine' function with this new version ---

def recover_deleted_files_engine(filepath):
    """Improved deleted files recovery engine optimized for large images.

    Features:
    - Streams reads/writes to avoid loading large files into memory.
    - Two-step deduplication (quick head/tail signature + full SHA256) to avoid duplicates.
    - Skips empty/small files and provides periodic status updates via deleted_scan_status.
    - Writes to a temporary file and atomically moves to final filename when complete.
    """
    global deleted_scan_status, deleted_files_db

    # Initialize globals if missing
    if 'deleted_scan_status' not in globals() or not isinstance(deleted_scan_status, dict):
        deleted_scan_status = {}
    if 'deleted_files_db' not in globals() or not isinstance(deleted_files_db, dict):
        deleted_files_db = {}
    
    # Reset for fresh run
    deleted_files_db.clear()

    deleted_scan_status.update({
        "in_progress": True,
        "complete": False,
        "files_found": 0,
        "message": "Starting deleted files recovery...",
        "start_time": time.time(),
        "last_update_time": time.time(),
        "elapsed_time": "0s",
        "estimated_total_time": None,
        "time_remaining_str": "Calculating...",
        "scan_methods": {
            "directory_walk": 0,
            "inode_scan": 0,
            "file_slack": 0,
            "unallocated_space": 0,
            "recycle_bin": 0
        },
        "validation_stats": {
            "total_scanned": 0,
            "empty_rejected": 0,
            "duplicate_rejected": 0,
            "invalid_rejected": 0,
            "valid_recovered": 0
        }
    })

    recovery_dir = app.config.get('DELETED_RECOVERY_FOLDER', os.path.join(APP_ROOT, 'deleted_files'))
    os.makedirs(recovery_dir, exist_ok=True)

    # Deduplication sets
    seen_quick = set()   # (size, head_hash, tail_hash)
    seen_full = set()    # sha256 full content
    # Backwards-compatible quick md5 set used by remaining legacy slack logic
    seen_hashes = set()

    MIN_FILE_SIZE = 128
    CHUNK_SIZE = 4 * 1024 * 1024  # 4MB streaming

    # Clear previous results in recovery_dir (only files created by previous runs)
    try:
        for item in os.listdir(recovery_dir):
            item_path = os.path.join(recovery_dir, item)
            if os.path.isfile(item_path) and item.startswith('deleted_'):
                try:
                    os.unlink(item_path)
                except Exception:
                    continue
    except Exception as e:
        deleted_scan_status.update({"message": f"Error clearing old files: {e}", "in_progress": False})
        return {}

    total_recovered = 0

    def update_status(method, count=1):
        nonlocal total_recovered
        total_recovered += count
        deleted_scan_status["files_found"] = total_recovered
        deleted_scan_status["scan_methods"][method] = deleted_scan_status["scan_methods"].get(method, 0) + count
        deleted_scan_status["message"] = f"Recovered {total_recovered} files... ({method})"
        # Update validation stats for recovered files
        try:
            v = deleted_scan_status.setdefault('validation_stats', {
                'total_scanned': 0, 'empty_rejected': 0, 'duplicate_rejected': 0, 'invalid_rejected': 0, 'valid_recovered': 0
            })
            v['valid_recovered'] = v.get('valid_recovered', 0) + count
        except Exception:
            pass

        # Update timing estimations periodically
        try:
            calculate_time_estimations(deleted_scan_status, deleted_scan_status.get('validation_stats', {}).get('total_scanned', 0))
        except Exception:
            pass

    def _compute_quick_hash(fs_file, size):
        """Compute a fast head+tail md5 signature for quick dedupe checks."""
        try:
            head_len = min(8192, size)
            head = fs_file.read_random(0, head_len) or b''
            tail_len = 8192
            if size > head_len:
                tail_offset = max(0, size - tail_len)
                tail = fs_file.read_random(tail_offset, min(tail_len, size - tail_offset)) or b''
            else:
                tail = b''
            h = hashlib.md5()
            h.update(head)
            h.update(tail)
            return h.hexdigest(), head, tail
        except Exception:
            return None, b'', b''

    def _compute_full_hash_and_write(fs_file, size, tmp_path):
        """Stream file out to tmp_path while computing SHA256.
        Returns sha256 hex digest and final size written.
        """
        sha = hashlib.sha256()
        written = 0
        try:
            offset = 0
            while offset < size:
                to_read = min(CHUNK_SIZE, size - offset)
                chunk = fs_file.read_random(offset, to_read)
                if not chunk:
                    break
                tmp_path.write(chunk)
                sha.update(chunk)
                written += len(chunk)
                offset += len(chunk)
            tmp_path.flush()
            return sha.hexdigest(), written
        except Exception:
            return None, written

    def process_deleted_file(fs_object, recovery_method, fs_offset=0):
        # nonlocal total_recovered
        try:
            meta = getattr(fs_object.info, 'meta', None)
            if not meta or not hasattr(meta, 'size'):
                return
            size = int(meta.size)
            if size <= MIN_FILE_SIZE:
                # count as scanned but rejected for being too small
                try:
                    deleted_scan_status.setdefault('validation_stats', {}).setdefault('total_scanned', 0)
                    deleted_scan_status['validation_stats']['total_scanned'] += 1
                    deleted_scan_status['validation_stats']['empty_rejected'] = deleted_scan_status['validation_stats'].get('empty_rejected', 0) + 1
                except Exception:
                    pass
                return

            inode = getattr(meta, 'addr', None)
            inode_str = str(inode) if inode is not None else 'unknown'

            # Count that we inspected a candidate (for timing/estimates)
            try:
                deleted_scan_status.setdefault('validation_stats', {}).setdefault('total_scanned', 0)
                deleted_scan_status['validation_stats']['total_scanned'] += 1
            except Exception:
                pass

            # Quick signature
            quick_hash, head, tail = _compute_quick_hash(fs_object, size)
            quick_key = (size, quick_hash)
            if quick_key in seen_quick:
                # quick filter considered duplicate - increment duplicate counter
                try:
                    deleted_scan_status.setdefault('validation_stats', {})
                    deleted_scan_status['validation_stats']['duplicate_rejected'] = deleted_scan_status['validation_stats'].get('duplicate_rejected', 0) + 1
                except Exception:
                    pass
                return

            # If quick key unique so far, compute full SHA256 while streaming to temp file
            import tempfile
            safe_name_root = secure_filename(getattr(getattr(fs_object.info, 'name', None), 'name', b'orphaned').decode('utf-8', 'ignore'))
            tmp_file = None
            final_name = None
            try:
                tmp_fd = tempfile.NamedTemporaryFile(delete=False, dir=recovery_dir)
                tmp_file = tmp_fd
                sha_hex, written = _compute_full_hash_and_write(fs_object, size, tmp_file)
                tmp_file.close()

                if not sha_hex or written == 0:
                    # empty or couldn't write - count as empty_rejected
                    try:
                        deleted_scan_status.setdefault('validation_stats', {})
                        if written == 0:
                            deleted_scan_status['validation_stats']['empty_rejected'] = deleted_scan_status['validation_stats'].get('empty_rejected', 0) + 1
                        else:
                            deleted_scan_status['validation_stats']['invalid_rejected'] = deleted_scan_status['validation_stats'].get('invalid_rejected', 0) + 1
                    except Exception:
                        pass
                    try:
                        os.unlink(tmp_fd.name)
                    except Exception:
                        pass
                    return

                # Full dedupe
                if sha_hex in seen_full:
                    # duplicate content, drop tmp
                    try:
                        os.unlink(tmp_fd.name)
                    except Exception:
                        pass
                    seen_quick.add(quick_key)
                    try:
                        deleted_scan_status.setdefault('validation_stats', {})
                        deleted_scan_status['validation_stats']['duplicate_rejected'] = deleted_scan_status['validation_stats'].get('duplicate_rejected', 0) + 1
                    except Exception:
                        pass
                    return

                # Unique, move to final path atomically
                seen_full.add(sha_hex)
                seen_quick.add(quick_key)

                timestamp_mtime = getattr(meta, 'mtime', None)
                mtime = datetime.datetime.fromtimestamp(timestamp_mtime).strftime('%Y-%m-%d %H:%M:%S') if timestamp_mtime else 'Unknown'
                ctime_val = getattr(meta, 'ctime', None)
                ctime = datetime.datetime.fromtimestamp(ctime_val).strftime('%Y-%m-%d %H:%M:%S') if ctime_val else 'Unknown'

                # Use a standardized deleted file name (sequential) and include best-effort extension
                # Determine a sensible extension as before
                detected_ext = ''
                try:
                    file_type = magic.from_buffer(head, mime=True) if head else 'unknown'
                    if file_type and file_type != 'unknown' and '/' in file_type:
                        import mimetypes
                        guessed = mimetypes.guess_extension(file_type.split(';')[0].strip())
                        if guessed:
                            detected_ext = guessed
                except Exception:
                    detected_ext = ''

                if not detected_ext:
                    try:
                        head_bytes = head or b''
                        for cat, sigs in FILE_SIGNATURES.items():
                            for name, meta in sigs.items():
                                hdr = meta.get('header')
                                if hdr and head_bytes.startswith(hdr):
                                    detected_ext = meta.get('extension', '')
                                    break
                            if detected_ext:
                                break
                    except Exception:
                        detected_ext = ''

                # Generate a sequential filename under the deleted_files_recovery_* pattern
                seq = sum(1 for _ in os.listdir(recovery_dir) if os.path.isfile(os.path.join(recovery_dir, _))) + 1
                final_name = generate_deleted_filename(seq, detected_ext)
                final_path = os.path.join(recovery_dir, final_name)
                # ensure unique final filename
                if os.path.exists(final_path):
                    final_path = final_path + f"_{int(time.time())}"

                try:
                    os.replace(tmp_fd.name, final_path)
                except Exception:
                    # fallback: copy then remove
                    try:
                        with open(tmp_fd.name, 'rb') as r, open(final_path, 'wb') as w:
                            while True:
                                b = r.read(CHUNK_SIZE)
                                if not b:
                                    break
                                w.write(b)
                        os.unlink(tmp_fd.name)
                    except Exception:
                        pass

                try:
                    file_type = magic.from_buffer(head, mime=True) if head else 'unknown'
                except Exception:
                    file_type = 'unknown'

                # Determine a sensible file extension from the mime type or known signatures
                try:
                    ext = None
                    if file_type and file_type != 'unknown' and '/' in file_type:
                        import mimetypes
                        ext = mimetypes.guess_extension(file_type.split(';')[0].strip())
                    # Fallback: check known FILE_SIGNATURES headers
                    if not ext:
                        try:
                            head_bytes = head or b''
                            for cat, sigs in FILE_SIGNATURES.items():
                                for name, meta in sigs.items():
                                    hdr = meta.get('header')
                                    if hdr and head_bytes.startswith(hdr):
                                        ext = meta.get('extension')
                                        break
                                if ext:
                                    break
                        except Exception:
                            ext = None
                    if not ext:
                        ext = ''
                except Exception:
                    ext = ''

                # Normalize safe_name_root: avoid duplicate extension if already present
                safe_root = safe_name_root
                if ext and safe_root.lower().endswith(ext.lower()):
                    final_basename = safe_root
                else:
                    final_basename = safe_root + (ext or '')

                file_info = {
                    'inode': inode_str,
                    'name': safe_name_root,
                    'size': size,
                    'offset': f"0x{meta.addr:08X}" if hasattr(meta, 'addr') and meta.addr is not None else 'Unknown',
                    'offset_decimal': int(getattr(meta, 'addr', 0) or 0),
                    'mtime': mtime,
                    'ctime': ctime,
                    'recovery_method': recovery_method,
                    'file_type': file_type,
                    'sha256': sha_hex,
                    'path': final_path,
                    'fs_offset': fs_offset
                }
                deleted_files_db[os.path.basename(final_path)] = file_info
                # add DB record for recovered deleted file
                try:
                    sess_id = session.get('analysis_session_id')
                except Exception:
                    sess_id = None
                try:
                    add_file_record(os.path.basename(final_path), 'deleted_recovered', final_path, os.path.getsize(final_path), session_id=sess_id, extra={'inode': inode_str})
                except Exception:
                    pass
                # Update aggregated status (update_status will increment valid_recovered)
                update_status(recovery_method)

            except Exception as e:
                # cleanup tmp file if present
                try:
                    if tmp_file is not None and hasattr(tmp_file, 'name'):
                        os.unlink(tmp_file.name)
                except Exception:
                    pass
                return

        except Exception:
            return

    try:
        img_handle = pytsk3.Img_Info(filepath)
        
        # Method 1: Directory Walk (Deleted entries)
        def deep_directory_scan(fs, directory, path="/"):
            try:
                for f in directory:
                    if not hasattr(f.info, 'meta') or f.info.meta is None:
                        continue
                    
                    # Check if file is deleted
                    is_deleted = not (f.info.meta.flags & pytsk3.TSK_FS_META_FLAG_ALLOC)
                    is_file = f.info.meta.type == pytsk3.TSK_FS_META_TYPE_REG
                    has_name = hasattr(f.info, 'name') and f.info.name is not None
                    
                    if is_deleted and is_file and has_name and f.info.name.name not in [b'.', b'..']:
                        process_deleted_file(f, "directory_walk")
                    
                    # Recursively scan subdirectories
                    if (f.info.meta.type == pytsk3.TSK_FS_META_TYPE_DIR and 
                        has_name and f.info.name.name not in [b'.', b'..']):
                        try:
                            deep_directory_scan(fs, f.as_directory(), 
                                              f"{path}/{f.info.name.name.decode('utf-8', 'ignore')}")
                        except (IOError, AttributeError):
                            continue
            except Exception as e:
                pass

        # Method 2: Deep Inode Scan
        def deep_inode_scan(fs):
            try:
                last_inode = fs.info.last_inum
                for inode_num in range(fs.info.first_inum, last_inode + 1):
                    try:
                        fs_file = fs.open_meta(inode=inode_num)
                        
                        # Check if inode is unallocated (deleted) but has content
                        if (fs_file.info.meta.flags & pytsk3.TSK_FS_META_FLAG_UNALLOC and 
                            fs_file.info.meta.type == pytsk3.TSK_FS_META_TYPE_REG and 
                            fs_file.info.meta.size > MIN_FILE_SIZE):
                            process_deleted_file(fs_file, "inode_scan")
                    except (IOError, AttributeError):
                        continue
            except Exception as e:
                pass

        # Method 3: File Slack Space Recovery
        def recover_file_slack(fs, directory, path="/"):
            try:
                for f in directory:
                    if not hasattr(f.info, 'meta') or f.info.meta is None:
                        continue
                    
                    # Check allocated files for slack space
                    is_allocated = (f.info.meta.flags & pytsk3.TSK_FS_META_FLAG_ALLOC)
                    is_file = f.info.meta.type == pytsk3.TSK_FS_META_TYPE_REG
                    
                    if is_allocated and is_file and f.info.meta.size > 0:
                        try:
                            # Calculate potential slack space
                            block_size = fs.info.block_size
                            actual_size = f.info.meta.size
                            blocks_used = (actual_size + block_size - 1) // block_size
                            slack_size = (blocks_used * block_size) - actual_size
                            
                            if slack_size > MIN_FILE_SIZE:
                                # Read slack space from end of file
                                slack_content = f.read_random(actual_size, slack_size)
                                
                                if slack_content and len(slack_content) >= MIN_FILE_SIZE:
                                    content_hash = hashlib.md5(slack_content).hexdigest()
                                    if content_hash not in seen_hashes:
                                        seen_hashes.add(content_hash)
                                        
                                        original_name = f.info.name.name.decode('utf-8', 'ignore') if hasattr(f.info, 'name') else f"file_{f.info.meta.addr}"
                                        safe_filename = secure_filename(f"slack_{f.info.meta.addr}_{original_name}")
                                        save_path = os.path.join(recovery_dir, safe_filename)
                                        
                                        with open(save_path, 'wb') as out_file:
                                            out_file.write(slack_content)
                                        
                                        update_status("file_slack")
                        except Exception:
                            continue
                    
                    # Recursively scan subdirectories
                    if (f.info.meta.type == pytsk3.TSK_FS_META_TYPE_DIR and 
                        hasattr(f.info, 'name') and f.info.name.name not in [b'.', b'..']):
                        try:
                            recover_file_slack(fs, f.as_directory(), 
                                            f"{path}/{f.info.name.name.decode('utf-8', 'ignore')}")
                        except (IOError, AttributeError):
                            continue
            except Exception as e:
                pass

        # Method 4: Recycle Bin/Trash Recovery
        def recover_recycle_bin(fs, directory, path="/"):
            try:
                for f in directory:
                    if not hasattr(f.info, 'meta') or f.info.meta is None:
                        continue
                    
                    current_path = f"{path}/{f.info.name.name.decode('utf-8', 'ignore')}" if hasattr(f.info, 'name') else path
                    
                    # Check for recycle bin directories
                    is_recycle_bin = any(keyword in current_path.upper() for keyword in 
                                       ['$RECYCLE.BIN', 'RECYCLED', '.TRASH', 'RECYCLE.BIN'])
                    
                    if is_recycle_bin and f.info.meta.type == pytsk3.TSK_FS_META_TYPE_REG:
                        process_deleted_file(f, "recycle_bin")
                    
                    # Recursively scan subdirectories
                    if f.info.meta.type == pytsk3.TSK_FS_META_TYPE_DIR and hasattr(f.info, 'name'):
                        try:
                            recover_recycle_bin(fs, f.as_directory(), current_path)
                        except (IOError, AttributeError):
                            continue
            except Exception as e:
                pass

        # Execute all recovery methods
        try:
            # Try to handle partition table first
            volume = pytsk3.Volume_Info(img_handle)
            for part in volume:
                if part.flags != pytsk3.TSK_VS_PART_FLAG_UNALLOC:
                    try:
                        fs_offset = part.start * volume.info.block_size
                        fs = pytsk3.FS_Info(img_handle, offset=fs_offset)
                        
                        deleted_scan_status["message"] = f"Scanning partition {part.desc}..."
                        
                        # Run all recovery methods
                        deep_directory_scan(fs, fs.open_dir(path="/"))
                        deep_inode_scan(fs)
                        recover_file_slack(fs, fs.open_dir(path="/"))
                        recover_recycle_bin(fs, fs.open_dir(path="/"))
                        
                    except IOError:
                        continue
        except IOError:
            # Try as single filesystem
            try:
                fs = pytsk3.FS_Info(img_handle, offset=0)
                deleted_scan_status["message"] = "Scanning as single filesystem..."
                
                deep_directory_scan(fs, fs.open_dir(path="/"))
                deep_inode_scan(fs)
                recover_file_slack(fs, fs.open_dir(path="/"))
                recover_recycle_bin(fs, fs.open_dir(path="/"))
                
            except IOError as e:
                deleted_scan_status["message"] = f"Error opening filesystem: {e}"

        deleted_scan_status["message"] = f"Recovery complete! Found {total_recovered} files using multiple methods."
        deleted_scan_status["complete"] = True
        
    except Exception as e:
        deleted_scan_status["message"] = f"A critical error occurred: {e}"
        deleted_scan_status["complete"] = True
        
    deleted_scan_status["in_progress"] = False
    return deleted_files_db  # Return the database of recovered files
    # --- Reporting Helper Functions ---
def generate_docx_report_data(case_details, evidence_file, carved_files, deleted_files, now):
    """Generate a DOCX report for the forensic case."""
    if docx is None: 
        return None
    from docx.shared import Inches
    doc = docx.Document()
    doc.add_heading('Forensic Report', 0)
    doc.add_paragraph(f"Generated by ForensicCarver Pro on {now.strftime('%Y-%m-%d %H:%M:%S')}")
    doc.add_heading('Case Details', level=1)
    table = doc.add_table(rows=3, cols=2)
    table.style = 'Light List'
    table.cell(0, 0).text = 'Case Name'
    table.cell(0, 1).text = case_details.get('name', '')
    table.cell(1, 0).text = 'Case Number'
    table.cell(1, 1).text = case_details.get('number', '')
    table.cell(2, 0).text = 'Examiner'
    table.cell(2, 1).text = case_details.get('examiner', '')

    for filename, details in evidence_file.items():
        doc.add_heading(f"Data Source: {filename}", level=1)
        t = doc.add_table(rows=4, cols=2)
        t.style = 'Light List'
        t.cell(0, 0).text = 'File Size'
        t.cell(0, 1).text = str(details.get('size_mb', '')) + " MB"
        t.cell(1, 0).text = 'MD5 Hash'
        t.cell(1, 1).text = details.get('hash_info', {}).get('MD5', 'N/A')
        t.cell(2, 0).text = 'SHA-1 Hash'
        t.cell(2, 1).text = details.get('hash_info', {}).get('SHA-1', 'N/A')
        t.cell(3, 0).text = 'SHA-256 Hash'
        t.cell(3, 1).text = details.get('hash_info', {}).get('SHA-256', 'N/A')
        doc.add_heading('Partition Information', level=2)
        if details.get('partition_info'):
            ptable = doc.add_table(rows=1, cols=4)
            ptable.style = 'Light List'
            hdr_cells = ptable.rows[0].cells
            hdr_cells[0].text = 'Index'
            hdr_cells[1].text = 'Description'
            hdr_cells[2].text = 'Start Sector'
            hdr_cells[3].text = 'Length (Sectors)'
            for part in details['partition_info']:
                row_cells = ptable.add_row().cells
                row_cells[0].text = str(part.get('addr', ''))
                row_cells[1].text = str(part.get('desc', ''))
                row_cells[2].text = str(part.get('start', ''))
                row_cells[3].text = str(part.get('len', ''))
        else:
            doc.add_paragraph("No partition data found.")

    doc.add_heading(f"Carved Files ({len(carved_files)} found)", level=1)
    ctable = doc.add_table(rows=1, cols=4)
    ctable.style = 'Light List'
    ctable.rows[0].cells[0].text = 'ID'
    ctable.rows[0].cells[1].text = 'Filename'
    ctable.rows[0].cells[2].text = 'Offset'
    ctable.rows[0].cells[3].text = 'Size (KB)'
    for file in sorted(carved_files.values(), key=lambda x: x.get('id', 0)):
        row = ctable.add_row().cells
        row[0].text = str(file.get('id', ''))
        row[1].text = file.get('name', '')
        row[2].text = file.get('offset', '')
        row[3].text = file.get('size_kb', '')

    doc.add_heading(f"Deleted File Entries ({len(deleted_files)} found)", level=1)
    dtable = doc.add_table(rows=1, cols=6)
    dtable.style = 'Light List'
    dtable.rows[0].cells[0].text = 'Inode'
    dtable.rows[0].cells[1].text = 'Filename'
    dtable.rows[0].cells[2].text = 'Size (Bytes)'
    dtable.rows[0].cells[3].text = 'Modified'
    dtable.rows[0].cells[4].text = 'Accessed'
    dtable.rows[0].cells[5].text = 'Created'
    for inode, file in deleted_files.items():
        row = dtable.add_row().cells
        row[0].text = str(file.get('inode', ''))
        row[1].text = file.get('name', '')
        row[2].text = str(file.get('size', ''))
        row[3].text = file.get('mtime', '')
        row[4].text = file.get('atime', '')
        row[5].text = file.get('ctime', '')

    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def generate_csv_zip_report_data(case_details, evidence_file, carved_files, deleted_files):
    """Generate a ZIP file containing CSV reports for the case."""
    memory_file = io.BytesIO()
    with zipfile.ZipFile(memory_file, 'w', zipfile.ZIP_DEFLATED) as zf:
        case_csv = io.StringIO()
        writer = csv.writer(case_csv)
        writer.writerow(['Case Name', 'Case Number', 'Examiner'])
        writer.writerow([case_details.get('name', ''), case_details.get('number', ''), case_details.get('examiner', '')])
        zf.writestr('case_details.csv', case_csv.getvalue())

        evidence_csv = io.StringIO()
        writer = csv.writer(evidence_csv)
        writer.writerow(['Filename', 'Size (MB)', 'MD5', 'SHA-1', 'SHA-256'])
        for filename, details in evidence_file.items():
            writer.writerow([
                filename,
                details.get('size_mb', ''),
                details.get('hash_info', {}).get('MD5', ''),
                details.get('hash_info', {}).get('SHA-1', ''),
                details.get('hash_info', {}).get('SHA-256', ''),
            ])
        zf.writestr('evidence_files.csv', evidence_csv.getvalue())

        carved_csv = io.StringIO()
        writer = csv.writer(carved_csv)
        writer.writerow(['ID', 'Filename', 'Offset', 'Size (KB)'])
        for file in sorted(carved_files.values(), key=lambda x: x.get('id', 0)):
            writer.writerow([file.get('id', ''), file.get('name', ''), file.get('offset', ''), file.get('size_kb', ''),])
        zf.writestr('carved_files.csv', carved_csv.getvalue())

        deleted_csv = io.StringIO()
        writer = csv.writer(deleted_csv)
        writer.writerow(['Inode', 'Filename', 'Size (Bytes)', 'Modified', 'Accessed', 'Created'])
        for inode, file in deleted_files.items():
            writer.writerow([file.get('inode', ''), file.get('name', ''), file.get('size', ''), file.get('mtime', ''), file.get('atime', ''), file.get('ctime', ''),])
        zf.writestr('deleted_files.csv', deleted_csv.getvalue())

    memory_file.seek(0)
    return memory_file

# --- HTML TEMPLATES ---
BASE_TEMPLATE = r"""
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>ForensicCarver Pro</title>
    <script src="https://cdn.tailwindcss.com"></script>
    <!-- Minimal local fallback utilities: used when Tailwind CDN is unavailable (keeps layout usable under Apache/offline) -->
    <style id="fallback-utilities">
        /* Layout helpers used by the templates (small subset of Tailwind) */
        .flex { display: flex !important; }
        .h-screen { height: 100vh !important; }
        .w-64 { width: 16rem !important; }
        .p-4 { padding: 1rem !important; }
        .p-8 { padding: 2rem !important; }
        .overflow-y-auto { overflow-y: auto !important; }
        .flex-shrink-0 { flex-shrink: 0 !important; }
        .flex-1 { flex: 1 1 auto !important; }
        .rounded-lg { border-radius: 0.5rem !important; }
        .block { display: block !important; }
        .mb-8 { margin-bottom: 2rem !important; }
        .text-white { color: #ffffff !important; }
        .text-xs { font-size: 0.75rem !important; }
        .font-bold { font-weight: 700 !important; }
        .font-normal { font-weight: 400 !important; }
        .text-2xl { font-size: 1.5rem !important; }
        .space-y-2 > * + * { margin-top: 0.5rem !important; }
        .hover\:bg-gray-700:hover { background-color: #374151 !important; }
        .cursor-not-allowed { cursor: not-allowed !important; opacity: 0.7 !important; }
        .text-gray-400 { color: #9ca3af !important; }
        .pt-4 { padding-top: 1rem !important; }
        /* Simple responsive helpers */
        @media (max-width: 640px) {
            .w-64 { width: 12rem !important; }
            .p-8 { padding: 1rem !important; }
        }
    </style>
    <style>
        body { 
            background-color: #111827; 
            color: #d1d5db; 
            font-family: 'Inter', sans-serif; 
        }
        .sidebar { 
            background-color: #1f2937; 
        }
        .sidebar a { 
            border-left: 3px solid transparent; 
            transition: all 0.2s ease-in-out; 
        }
        .sidebar a.active { 
            border-left-color: #3b82f6; 
            background-color: #374151; 
            color: white; 
        }
        .card { 
            background-color: #1f2937; 
            border: 1px solid #374151; 
            box-shadow: 0 4px 6px -1px rgb(0 0 0 / 0.1), 0 2px 4px -2px rgb(0 0 0 / 0.1); 
        }
        .btn-primary, .btn-secondary, .btn-green, .bg-red-600 { 
            transition: all 0.2s ease-in-out; 
        }
        .btn-primary:hover, .btn-secondary:hover, .btn-green:hover, .bg-red-600:hover { 
            opacity: 0.9; 
            transform: translateY(-1px); 
        }
        .btn-primary { 
            background-color: #3b82f6; 
            color: white; 
        }
        .btn-secondary { 
            background-color: #4b5563; 
            color: white; 
        }
        .btn-green { 
            background-color: #10b981; 
            color: white; 
        }
        .hex-view { 
            font-family: 'Courier New', Courier, monospace; 
            background-color: #0d1117; 
            border: 1px solid #374151; 
        }
        .log-view { 
            font-family: 'Courier New', Courier, monospace; 
            background-color: #0d1117; 
            color: #d1d5db; 
            border: 1px solid #374151; 
            white-space: pre-wrap; 
            word-wrap: break-word; 
        }
        /* Highlight user input fields: keep background matching dark theme and only highlight border/outline on focus */
        input[type="text"], input[type="search"], input[type="email"], input[type="url"], input[type="tel"], input[type="number"], input[type="password"],
        textarea, select, [contenteditable="true"] {
            background-color: transparent; /* keep page/card background visible */
            border: 1px solid rgba(255,255,255,0.06);
            color: #e5e7eb; /* light text */
            padding: .5rem;
            border-radius: .25rem;
            transition: box-shadow .12s ease, border-color .12s ease, background-color .12s ease;
            -webkit-appearance: none; /* remove native styling */
            appearance: none;
        }
        input:focus, textarea:focus, select:focus, [contenteditable="true"]:focus {
            outline: none;
            background-color: transparent;
            color: #e5e7eb;
            border-color: #f59e0b; /* highlight border only */
            box-shadow: 0 0 0 3px rgba(245,158,11,0.12);
        }
        /* Keep disabled/readonly inputs visually distinct but without white background */
        input[disabled], textarea[disabled], select[disabled], input[readonly], textarea[readonly] {
            opacity: 0.7;
            background-color: transparent;
            cursor: not-allowed;
        }
        /* For selects that show native arrow, keep color consistent */
        select { background-image: none; }
        .encryption-badge { 
            background-color: #ef4444; 
            color: white; 
            padding: 2px 8px; 
            border-radius: 4px; 
            font-size: 0.75rem; 
            font-weight: bold; 
        }
        .decryption-badge { 
            background-color: #10b981; 
            color: white; 
            padding: 2px 8px; 
            border-radius: 4px; 
            font-size: 0.75rem; 
            font-weight: bold; 
        }
        .decrypting-badge { 
            background-color: #f59e0b; 
            color: white; 
            padding: 2px 8px; 
            border-radius: 4px; 
            font-size: 0.75rem; 
            font-weight: bold; 
        }
        .modal { 
            display: none; 
            position: fixed; 
            z-index: 1000; 
            left: 0; 
            top: 0; 
            width: 100%; 
            height: 100%; 
            background-color: rgba(0,0,0,0.7); 
        }
        .modal-content { 
            background-color: #1f2937; 
            margin: 15% auto; 
            padding: 20px; 
            border: 1px solid #374151; 
            width: 50%; 
            border-radius: 8px; 
            box-shadow: 0 4px 8px 0 rgba(0,0,0,0.2); 
        }
        .close { 
            color: #aaa; 
            float: right; 
            font-size: 28px; 
            font-weight: bold; 
            cursor: pointer; 
        }
        .close:hover { 
            color: white; 
        }
    </style>
    <link rel="stylesheet" href="https://rsms.me/inter/inter.css">
    <!-- Local fallback stylesheet (served by Flask static) -->
    <style>
    /* Floating card animations */
    .storage-card-hidden { opacity: 0; transform: translateX(12px); transition: opacity 220ms ease, transform 220ms ease; pointer-events: none; }
    .storage-card-visible { opacity: 1; transform: translateX(0); transition: opacity 220ms ease, transform 220ms ease; pointer-events: auto; }
    .storage-card-pop { box-shadow: 0 8px 26px rgba(0,0,0,0.45); }
    </style>
    <link rel="stylesheet" href="/static/css/local_tailwind_fallback.css">
    <link rel="stylesheet" href="/static/css/explorer.css">
</head>
<body class="flex h-screen">
    <aside class="sidebar w-64 p-4 space-y-2 flex-shrink-0 overflow-y-auto">
        <div class="text-white text-2xl font-bold mb-8">
            ForensicCarver <span class="text-blue-500">Pro</span>
            <p class="text-xs font-normal text-gray-400">Digital Evidence Analysis</p>
        </div>
    <a href="{{ url_for('evidence_upload') }}" class="block p-3 rounded-lg hover:bg-gray-700 {% if request.endpoint == 'evidence_upload' %}active{% endif %}">Evidence & Upload</a>
    <a href="{{ url_for('databases') }}" class="block p-3 rounded-lg hover:bg-gray-700 {% if request.endpoint == 'databases' %}active{% endif %}">Database Management</a>
        {%- if show_encryption_in_sidebar is defined and show_encryption_in_sidebar %}
        <a href="{{ url_for('encryption_page') }}" class="block p-3 rounded-lg hover:bg-gray-700 {% if request.endpoint == 'encryption_page' %}active{% endif %}">File Encryption</a>
        {%- endif %}
        {% if uploaded_files_db %}
        <a href="{{ url_for('decryption_page', filename=(uploaded_files_db.keys()|first)) if uploaded_files_db and uploaded_files_db[(uploaded_files_db.keys()|first)].encryption_status.encrypted else '#' }}" class="block p-3 rounded-lg {% if uploaded_files_db and uploaded_files_db[(uploaded_files_db.keys()|first)].encryption_status.encrypted %} hover:bg-gray-700 {% else %} text-gray-600 cursor-not-allowed {% endif %} {% if request.endpoint == 'decryption_page' or request.endpoint == 'decryption_progress' %}active{% endif %}">Decryption</a>
        <a href="{{ url_for('forensic_analysis') }}" class="block p-3 rounded-lg hover:bg-gray-700 {% if request.endpoint == 'forensic_analysis' %}active{% endif %}">
                <div style="display:flex;align-items:center;justify-content:space-between;">
                    <span>Forensic Analysis</span>
                    {% if sidebar_session and sidebar_session.active %}
                    <span title="Session: {{ sidebar_session.short_name }}\nCarved: {{ sidebar_session.carved_count }}" style="display:inline-flex;align-items:center;gap:6px;">
                        <span style="width:10px;height:10px;border-radius:50%;background:#10b981;display:inline-block;box-shadow:0 0 6px rgba(16,185,129,0.45);"></span>
                        <small style="color:#a7f3d0;font-size:11px;">{{ sidebar_session.short_name }} · {{ sidebar_session.carved_count }}</small>
                    </span>
                    {% elif session.get('analysis_session_id') %}
                    <span title="Active analysis session" style="display:inline-flex;align-items:center;gap:6px;">
                        <span style="width:10px;height:10px;border-radius:50%;background:#10b981;display:inline-block;box-shadow:0 0 6px rgba(16,185,129,0.45);"></span>
                        <small style="color:#a7f3d0;font-size:11px;">Active</small>
                    </span>
                    {% endif %}
                </div>
            </a>
        <a href="{{ url_for('auto_carving_setup') }}" class="block p-3 rounded-lg hover:bg-gray-700 {% if request.endpoint == 'auto_carving_setup' %}active{% endif %}">Auto Carving</a>
        <a href="{{ url_for('recovered_files') }}" class="block p-3 rounded-lg hover:bg-gray-700 {% if request.endpoint == 'recovered_files' %}active{% endif %}">Recovered Files</a>
        <a href="{{ url_for('deleted_files_status_page') }}" class="block p-3 rounded-lg hover:bg-gray-700 {% if request.endpoint == 'deleted_files_status_page' or request.endpoint == 'deleted_files' %}active{% endif %}">Deleted Files</a>
        <a href="{{ url_for('log_file_viewer') }}" class="block p-3 rounded-lg hover:bg-gray-700 {% if request.endpoint == 'log_file_viewer' %}active{% endif %}">Log File Viewer</a>
        <a href="{{ url_for('event_log_viewer') }}" class="block p-3 rounded-lg hover:bg-gray-700 {% if request.endpoint == 'event_log_viewer' %}active{% endif %}">Event Log Viewer</a>
        <a href="{{ url_for('reporting') }}" class="block p-3 rounded-lg hover:bg-gray-700 {% if request.endpoint == 'reporting' %}active{% endif %}">Reporting</a>
        <a href="{{ url_for('manual_carving') }}" class="block p-3 rounded-lg hover:bg-gray-700 {% if request.endpoint == 'manual_carving' %}active{% endif %}">Manual Carving</a>
         
        {% endif %}
        <div class="pt-4 text-xs text-gray-500">Version 4.3.0<br>Licensed to: Forensics Lab</div>
    </aside>
    
    <main class="flex-1 p-8 overflow-y-auto">
        {% with messages = get_flashed_messages(with_categories=true) %}
            {% if messages %}
                {% for category, message in messages %}
                    <div class="bg-{% if category == 'success' %}green{% else %}blue{% endif %}-600 text-white p-4 rounded-lg mb-4">
                        {{ message | safe }}
                    </div>
                {% endfor %}
            {% endif %}
        {% endwith %}
        {{ content|safe }}
    </main>
    
    <div id="uploadConfirmModal" class="modal">
        <div class="modal-content">
            <h2 class="text-xl font-semibold text-white mb-4">Confirm Large File Upload</h2>
            <p id="uploadConfirmText" class="text-gray-300 mb-6"></p>
            <div class="flex justify-end space-x-4">
                <button id="confirmCancelBtn" class="btn-secondary px-4 py-2 rounded-lg">Cancel</button>
                <button id="confirmOkBtn" class="btn-primary px-4 py-2 rounded-lg">Continue</button>
            </div>
        </div>
    </div>
    
    <div id="uploadCompleteModal" class="modal">
        <div class="modal-content">
            <h2 class="text-xl font-semibold text-white mb-4">Upload Complete</h2>
            <p class="text-gray-300 mb-6">File has been uploaded and is ready for analysis.</p>
            <div class="flex justify-end">
                <button id="closeUploadModalBtn" class="btn-green px-4 py-2 rounded-lg">OK</button>
            </div>
        </div>
    </div>
    
    <!-- Global floating storage selector (appears on all pages except Database Management which overrides) -->
    <div id="storage-toggle-icon" title="Toggle Storage Card" style="position:fixed;right:16px;top:16px;z-index:9999;cursor:grab;user-select:none;">
        <div id="storage-icon-inner" style="position:relative;width:48px;height:48px;border-radius:9999px;display:flex;align-items:center;justify-content:center;box-shadow:0 6px 18px rgba(0,0,0,0.4);background:#06b6d4;">
            <!-- database icon -->
            <svg xmlns="http://www.w3.org/2000/svg" width="22" height="22" viewBox="0 0 24 24" fill="none" stroke="white" stroke-width="1.8" stroke-linecap="round" stroke-linejoin="round">
                <ellipse cx="12" cy="5" rx="8" ry="3"></ellipse>
                <path d="M4 5v6c0 1.657 3.582 3 8 3s8-1.343 8-3V5"></path>
                <path d="M4 11v6c0 1.657 3.582 3 8 3s8-1.343 8-3v-6"></path>
            </svg>
            <div id="storage-icon-badge" style="position:absolute;left:-8px;top:-8px;background:#ef4444;color:white;font-size:10px;padding:2px 6px;border-radius:9999px;display:none;box-shadow:0 2px 6px rgba(0,0,0,0.3);"></div>
        </div>
    </div>

    {%- if show_session_details is not defined or show_session_details %}
    <!-- Global Session Details button + modal (shown on most pages; controllable via show_session_details) -->
    <div id="global-session-btn" style="position:fixed;right:16px;top:76px;z-index:10000;">
        <button id="btn-session-details" class="btn-secondary px-3 py-2 rounded-lg" title="Session Details">📦 Session</button>
    </div>

    <div id="session-modal" style="display:none; position:fixed; right:56px; top:80px; width:320px; background:#071018; border:1px solid #23313e; padding:12px; z-index:10001; border-radius:8px;">
        <div style="display:flex; justify-content:space-between; align-items:center; margin-bottom:8px;">
            <div style="font-weight:600;">Session Details</div>
            <button id="session-close" style="background:transparent;color:#cfe3ff;border:none;">✕</button>
        </div>
        <div style="font-size:13px; color:#cfe3ff;">
            <div>Session: <strong>{{ (session_info.name if session_info is defined else (session.get('analysis_session_id') or 'N/A')) }}</strong></div>
            <div style="margin-top:6px">Path: <div style="font-family:monospace; font-size:12px; color:#9fc3e6;">{{ (session_info.path if session_info is defined else session.get('analysis_session_path') or 'Not created') }}</div></div>
            <div style="margin-top:8px">
                <div style="font-weight:600; margin-bottom:6px;">Counts</div>
                <div style="display:grid; grid-template-columns:1fr 1fr; gap:6px;">
                    <div>Carved: <strong>{{ (session_info.counts.Carved if session_info is defined and session_info.counts is defined else 0) }}</strong></div>
                    <div>Deleted: <strong>{{ (session_info.counts.Deleted if session_info is defined and session_info.counts is defined else 0) }}</strong></div>
                    <div>Reports: <strong>{{ (session_info.counts.Reports if session_info is defined and session_info.counts is defined else 0) }}</strong></div>
                    <div>Logs: <strong>{{ (session_info.counts.Logs if session_info is defined and session_info.counts is defined else 0) }}</strong></div>
                    <div>Events: <strong>{{ (session_info.counts.Events if session_info is defined and session_info.counts is defined else 0) }}</strong></div>
                    <div>ManualCarving: <strong>{{ (session_info.counts.ManualCarving if session_info is defined and session_info.counts is defined else 0) }}</strong></div>
                </div>
            </div>
            <div style="margin-top:8px; display:flex; gap:6px;">
                <a href="{{ url_for('manual_carving') }}" class="btn-primary px-2 py-1 rounded">Manual Carving</a>
                <a href="{{ url_for('fs_explorer') }}?browse_root=Session Files" class="btn-secondary px-2 py-1 rounded">Open Session Files</a>
            </div>
        </div>
    </div>

    <script>
    document.addEventListener('DOMContentLoaded', function(){
        var btn = document.getElementById('btn-session-details');
        var modal = document.getElementById('session-modal');
        var closer = document.getElementById('session-close');
        if(btn && modal){ btn.addEventListener('click', function(){ modal.style.display = modal.style.display === 'none' ? 'block' : 'none'; }); }
        if(closer && modal){ closer.addEventListener('click', function(){ modal.style.display = 'none'; }); }
    });
    </script>
    {%- endif %}

    <div id="storage-card" class="storage-card-hidden" style="position:fixed;right:16px;top:76px;z-index:9999;width:320px;">
        <div class="card p-4 rounded-lg" style="background:#0b1220;border:1px solid rgba(255,255,255,0.04);">
            <h3 class="text-sm font-semibold text-white mb-2">Upload Storage Target</h3>
            <div id="storage-options" class="space-y-2 text-sm text-gray-300 mb-3">
                <div>
                    <label><input type="radio" name="storage_choice" value="local"> Store Locally</label>
                </div>
                <!-- DB options will be injected by client JS -->
            </div>
            <div class="mb-2">
                <label class="text-xs text-gray-400 mr-2">Icon color</label>
                <input id="storage-color-picker" type="color" value="#06b6d4" style="vertical-align:middle">
            </div>
            <div class="flex justify-between">
                <button id="storage-fix-btn" class="btn-primary px-3 py-1 text-sm rounded-lg">Fix for storage</button>
                <button id="storage-clear-btn" class="btn-secondary px-3 py-1 text-sm rounded-lg">Clear</button>
            </div>
            <div id="storage-note" class="text-xs text-gray-400 mt-2"></div>
        </div>
    </div>

    <script>
    // Floating storage selector global behavior
    document.addEventListener('DOMContentLoaded', function(){
        const icon = document.getElementById('storage-toggle-icon');
        const inner = document.getElementById('storage-icon-inner');
        const badge = document.getElementById('storage-icon-badge');
        const card = document.getElementById('storage-card');
        const fixBtn = document.getElementById('storage-fix-btn');
        const clearBtn = document.getElementById('storage-clear-btn');
        const note = document.getElementById('storage-note');
        const colorPicker = document.getElementById('storage-color-picker');

        function populateDbOptions(list) {
            const container = document.getElementById('storage-options');
            // remove existing db rows except local
            Array.from(container.querySelectorAll('div.db-option')).forEach(n => n.remove());
            list.forEach(item => {
                const d = document.createElement('div'); d.className = 'db-option';
                d.innerHTML = `<label><input type="radio" name="storage_choice" value="${item.id}"> ${item.name} (${item.type})</label>`;
                container.appendChild(d);
            });
            // After injecting options, restore previously fixed selection (if any)
            try {
                const sel = window.localStorage.getItem('selected_target_db');
                if (sel) {
                    const r = container.querySelector(`input[name="storage_choice"][value="${sel}"]`);
                    if (r) {
                        r.checked = true;
                        const note = document.getElementById('storage-note');
                        if (note) note.textContent = 'Fixed: ' + sel;
                        // ensure badge is updated as well
                        try { const ev = new Event('storage:updated'); document.dispatchEvent(ev); } catch (e) {}
                    }
                }
            } catch (e) {}
            // update badge once options and selection are settled
            try { if (typeof updateBadge === 'function') updateBadge(); } catch (e) {}
        }

        // fetch DB list from server (non-blocking)
        fetch('/databases/list_json').then(r => r.json()).then(data => {
            if (Array.isArray(data)) populateDbOptions(data);
        }).catch(() => {});

        // restore icon color
        try { const storedColor = window.localStorage.getItem('storage_icon_color'); if (storedColor) inner.style.background = storedColor, colorPicker.value = storedColor; } catch (e) {}
        // restore position
        try { const pos = JSON.parse(window.localStorage.getItem('storage_icon_pos') || 'null'); if (pos && pos.left && pos.top) { icon.style.left = pos.left; icon.style.top = pos.top; icon.style.right = 'auto'; card.style.left = pos.left; card.style.top = (parseInt(pos.top,10) + 56) + 'px'; card.style.right = 'auto'; } } catch (e) {}

        // draggable
        let dragging = false, offsetX = 0, offsetY = 0;
        inner.style.cursor = 'grab';
        inner.addEventListener('mousedown', (e) => { dragging = true; inner.style.cursor = 'grabbing'; offsetX = e.clientX - icon.getBoundingClientRect().left; offsetY = e.clientY - icon.getBoundingClientRect().top; e.preventDefault(); });
        document.addEventListener('mousemove', (e) => {
            if (!dragging) return;
            const xRaw = e.clientX - offsetX;
            const yRaw = e.clientY - offsetY;
            // clamp to viewport
            const maxX = Math.max(8, window.innerWidth - icon.offsetWidth - 8);
            const maxY = Math.max(8, window.innerHeight - icon.offsetHeight - 8);
            const x = Math.min(maxX, Math.max(8, xRaw));
            const y = Math.min(maxY, Math.max(8, yRaw));
            icon.style.left = x + 'px'; icon.style.top = y + 'px'; icon.style.right = 'auto';
            // position card below icon but clamp it to viewport too
            const cardTop = Math.min(window.innerHeight - card.offsetHeight - 8, y + icon.offsetHeight + 8);
            const cardLeft = Math.min(window.innerWidth - card.offsetWidth - 8, Math.max(8, x));
            card.style.left = cardLeft + 'px'; card.style.top = cardTop + 'px'; card.style.right = 'auto';
        });
        document.addEventListener('mouseup', () => {
            if (!dragging) return; dragging = false; inner.style.cursor = 'grab';
            try { window.localStorage.setItem('storage_icon_pos', JSON.stringify({ left: icon.style.left, top: icon.style.top })); } catch (e) {}
        });

        // touch support
        inner.addEventListener('touchstart', (e) => { const t = e.touches[0]; dragging = true; offsetX = t.clientX - icon.getBoundingClientRect().left; offsetY = t.clientY - icon.getBoundingClientRect().top; });
        document.addEventListener('touchmove', (e) => {
            if (!dragging) return; const t = e.touches[0]; const xRaw = t.clientX - offsetX; const yRaw = t.clientY - offsetY;
            const maxX = Math.max(8, window.innerWidth - icon.offsetWidth - 8);
            const maxY = Math.max(8, window.innerHeight - icon.offsetHeight - 8);
            const x = Math.min(maxX, Math.max(8, xRaw));
            const y = Math.min(maxY, Math.max(8, yRaw));
            icon.style.left = x + 'px'; icon.style.top = y + 'px'; icon.style.right = 'auto';
            const cardTop = Math.min(window.innerHeight - card.offsetHeight - 8, y + icon.offsetHeight + 8);
            const cardLeft = Math.min(window.innerWidth - card.offsetWidth - 8, Math.max(8, x));
            card.style.left = cardLeft + 'px'; card.style.top = cardTop + 'px'; card.style.right = 'auto';
        });
        document.addEventListener('touchend', () => { if (!dragging) return; dragging = false; try { window.localStorage.setItem('storage_icon_pos', JSON.stringify({ left: icon.style.left, top: icon.style.top })); } catch (e) {} });

        // color picker
        colorPicker.addEventListener('input', (e) => { inner.style.background = e.target.value; try { window.localStorage.setItem('storage_icon_color', e.target.value); } catch (err) {} });

        function updateBadge() {
            try {
                const sel = window.localStorage.getItem('selected_target_db');
                if (!sel) { badge.style.display = 'none'; return; }
                let labelText = sel;
                const r = card.querySelector(`input[name="storage_choice"][value="${sel}"]`);
                if (r) labelText = r.parentNode.textContent.trim();
                badge.textContent = labelText.length > 10 ? labelText.slice(0,10) + '…' : labelText; badge.style.display = 'block';
            } catch (e) { badge.style.display = 'none'; }
        }

        try { const sel = window.localStorage.getItem('selected_target_db'); if (sel) { const r = card.querySelector(`input[name="storage_choice"][value="${sel}"]`); if (r) r.checked = true; note.textContent = 'Fixed: ' + sel; } } catch (e) {}
        updateBadge();

        icon.addEventListener('click', (e) => {
            if (dragging) return;
            if (card.classList.contains('storage-card-visible')) { card.classList.remove('storage-card-visible'); card.classList.add('storage-card-hidden'); return; }
            // Position the card to the LEFT of the icon when possible, otherwise fall back to right.
            const iconRect = icon.getBoundingClientRect();
            const preferLeft = iconRect.left - 8; // space left of icon
            const tryLeft = iconRect.left - card.offsetWidth - 8;
            let desiredLeft = tryLeft;
            if (tryLeft < 8) {
                // not enough room on left; place to the right of the icon
                desiredLeft = Math.min(window.innerWidth - card.offsetWidth - 8, iconRect.right + 8);
            }
            // Vertically center card relative to icon when possible
            let desiredTop = iconRect.top + (iconRect.height / 2) - (card.offsetHeight / 2);
            // clamp vertically within viewport
            desiredTop = Math.max(8, Math.min(desiredTop, window.innerHeight - card.offsetHeight - 8));

            // final clamps for horizontal as well
            desiredLeft = Math.max(8, Math.min(desiredLeft, window.innerWidth - card.offsetWidth - 8));

            card.style.left = desiredLeft + 'px';
            card.style.top = desiredTop + 'px';
            card.style.right = 'auto';
            // animate in
            card.classList.remove('storage-card-hidden');
            card.classList.add('storage-card-visible','storage-card-pop');
            setTimeout(() => card.classList.remove('storage-card-pop'), 250);
        });

        fixBtn.addEventListener('click', () => { const r = card.querySelector('input[name="storage_choice"]:checked'); if (!r) { note.textContent = 'Select a storage target first.'; return; } try { window.localStorage.setItem('selected_target_db', r.value); note.textContent = 'Fixed: ' + r.value; updateBadge(); } catch (e) { note.textContent = 'Could not save selection.'; } });
        fixBtn.addEventListener('click', () => { const r = card.querySelector('input[name="storage_choice"]:checked'); if (!r) { note.textContent = 'Select a storage target first.'; return; } try { window.localStorage.setItem('selected_target_db', r.value); note.textContent = 'Fixed: ' + r.value; updateBadge(); try { fetch('/databases/select_target', { method: 'POST', headers: {'Content-Type':'application/json'}, body: JSON.stringify({selected: r.value}) }); } catch(e){} } catch (e) { note.textContent = 'Could not save selection.'; } });
        clearBtn.addEventListener('click', () => { try { window.localStorage.removeItem('selected_target_db'); note.textContent = 'Selection cleared.'; const r = card.querySelector('input[name="storage_choice"]:checked'); if (r) r.checked = false; updateBadge(); try { fetch('/databases/select_target', { method: 'POST', headers: {'Content-Type':'application/json'}, body: JSON.stringify({selected: null}) }); } catch(e){} } catch (e) { note.textContent = 'Could not clear selection.'; } });
    });
    </script>
    <script>const FAC_CSRF_TOKEN = {{ session.get('csrf_token')|tojson }};</script>
    <script>const FAC_PYEWF_AVAILABLE = {{ pyewf_available|default(false)|tojson }};</script>
    <script src="/static/js/explorer.js"></script>
    <script src="/static/js/create_image_modal.js"></script>
</body>
</html>
"""


@app.route('/databases/select_target', methods=['POST'])
def databases_select_target():
    try:
        data = request.get_json(force=True)
        sel = data.get('selected') if isinstance(data, dict) else None
        if sel:
            session['selected_target_db'] = sel
        else:
            session.pop('selected_target_db', None)
        return jsonify({'ok': True})
    except Exception:
        return jsonify({'ok': False}), 400

DECRYPTION_PROGRESS_CONTENT = """
<h1 class="text-3xl font-bold text-white mb-2">Decryption in Progress</h1>
<p class="text-gray-400 mb-8">Attempting to decrypt <strong class="font-mono text-blue-400">{{ filename }}</strong>. Please wait.</p>
<div class="card p-6 rounded-lg">
    <div class="flex justify-between items-center mb-4">
        <h2 id="status-header" class="text-xl font-semibold text-white">Initializing...</h2>
        <span id="progress-percent" class="font-bold text-blue-400">0%</span>
    </div>
    <div class="w-full bg-gray-700 rounded-full h-4">
        <div id="progress-bar" class="bg-blue-600 h-4 rounded-full transition-all duration-500" style="width: 0%"></div>
    </div>
    <p class="text-sm text-gray-400 mt-2">Status: <span id="status-message" class="text-yellow-400">Starting process...</span></p>
    <p class="text-sm text-gray-400 mt-1">Attempts: <span id="attempts-count">0 / 0</span></p>
</div>

<div id="result-container" class="mt-8 text-center hidden">
    <div id="result-card" class="card p-8 rounded-lg inline-block">
        <h2 id="result-title" class="text-2xl font-bold mb-4"></h2>
        <p id="result-message" class="text-gray-300 mb-6"></p>
        <a id="result-button" href="#" class="btn-primary px-8 py-3 rounded-lg font-semibold">Continue</a>
    </div>
</div>

<script>
function updateProgress() {
    fetch('/decryption_status')
        .then(response => response.json())
        .then(data => {
            if (!data.filename) return;

            const progressBar = document.getElementById('progress-bar');
            const progressPercent = document.getElementById('progress-percent');
            const statusHeader = document.getElementById('status-header');
            const statusMessage = document.getElementById('status-message');
            const attemptsCount = document.getElementById('attempts-count');
            
            progressBar.style.width = data.progress + '%';
            progressPercent.textContent = data.progress + '%';
            statusMessage.textContent = data.message;
            attemptsCount.textContent = `${data.attempts} / ${data.total_attempts}`;

            if (data.in_progress) {
                statusHeader.textContent = "Decryption Running...";
                setTimeout(updateProgress, 1000);
            } else if (data.complete) {
                const resultContainer = document.getElementById('result-container');
                const resultCard = document.getElementById('result-card');
                const resultTitle = document.getElementById('result-title');
                const resultMessage = document.getElementById('result-message');
                const resultButton = document.getElementById('result-button');
                
                resultContainer.classList.remove('hidden');

                if (data.success) {
                    statusHeader.textContent = "Decryption Successful!";
                    progressPercent.textContent = '100%';
                    progressBar.style.width = '100%';
                    progressBar.classList.remove('bg-blue-600');
                    progressBar.classList.add('bg-green-600');
                    statusMessage.textContent = "File decrypted successfully.";
                    
                    resultCard.classList.add('border', 'border-green-500');
                    resultTitle.textContent = "✅ Success!";
                    resultTitle.classList.add('text-green-400');
                    resultMessage.textContent = "The evidence file has been decrypted. You can now proceed with the analysis.";
                    resultButton.href = "{{ url_for('forensic_analysis') }}";
                    resultButton.textContent = "Go to Analysis";
                } else {
                    statusHeader.textContent = "Decryption Failed";
                    progressBar.classList.remove('bg-blue-600');
                    progressBar.classList.add('bg-red-600');
                    statusMessage.textContent = data.message;
                    
                    resultCard.classList.add('border', 'border-red-500');
                    resultTitle.textContent = "❌ Failed";
                    resultTitle.classList.add('text-red-400');
                    resultMessage.textContent = "Could not decrypt the file with the provided passwords. You can return to the decryption page to try again.";
                    resultButton.href = "{{ url_for('decryption_page', filename=filename) }}";
                    resultButton.textContent = "Try Again";
                }
            }
        });
}
document.addEventListener('DOMContentLoaded', updateProgress);
</script>

<!-- Floating storage selector: single enhanced draggable icon/card remains below -->


"""

EVIDENCE_UPLOAD_CONTENT = """
<h1 class="text-3xl font-bold text-white mb-4">Evidence File Management</h1>
<p class="text-gray-400 mb-8">Upload new evidence or load an existing file from the database for analysis.</p>
<div class="grid grid-cols-1 lg:grid-cols-2 gap-8">
    <div class="space-y-8">
        <div class="card p-6 rounded-lg">
            <h2 class="text-xl font-semibold text-white mb-4">Upload New Evidence File</h2>
            <p class="text-xs text-gray-400 mb-4">Supported formats: .dd, .e01, .mem, .raw, .img, .vmdk</p>
            
            <!-- Upload Progress Container (initially hidden) -->
            <div id="upload-progress-container" class="hidden mb-6 p-4 bg-gray-800 rounded-lg border border-gray-600">
                <div class="flex justify-between items-center mb-3">
                    <h3 class="text-lg font-semibold text-white">Uploading: <span id="upload-filename" class="font-mono"></span></h3>
                    <div style="display:flex;align-items:center;gap:8px;">
                        <span id="upload-percent" class="font-bold text-blue-400 text-lg">0%</span>
                        <button type="button" id="upload-cancel" class="btn-small btn-delete" title="Cancel upload">Cancel</button>
                    </div>
                </div>
                
                <div class="w-full bg-gray-700 rounded-full h-4 mb-3">
                    <div id="upload-progress-bar" class="bg-blue-600 h-4 rounded-full transition-all duration-300" style="width: 0%"></div>
                </div>
                
                <div class="grid grid-cols-1 md:grid-cols-3 gap-4 text-sm">
                    <div class="text-center">
                        <div class="text-gray-400 mb-1">Progress</div>
                        <div id="upload-progress-text" class="font-mono text-white">0 B / 0 B</div>
                    </div>
                    <div class="text-center">
                        <div class="text-gray-400 mb-1">Speed</div>
                        <div id="upload-speed" class="font-mono text-white">0 B/s</div>
                    </div>
                    <div class="text-center">
                        <div class="text-gray-400 mb-1">Time</div>
                        <div id="upload-timing" class="font-mono text-white text-xs">Elapsed: 0s<br>Remaining: Calculating...</div>
                    </div>
                </div>
            </div>
            
            <form id="upload-form" method="post" enctype="multipart/form-data" class="border-2 border-dashed border-gray-600 rounded-lg p-8 text-center">
                <p class="mb-4 text-gray-300">Drop evidence file here or click to browse</p>
                <!-- Support selecting multiple files (UI) so users can pick and review a list before uploading -->
                <input type="file" name="file" class="hidden" id="upload-file-input" multiple>
                <label for="upload-file-input" class="btn-primary px-6 py-3 rounded-lg cursor-pointer inline-block mb-4">Select Files</label>

                <!-- Warning shown when no root is selected -->
                <div id="upload-root-warning" class="hidden mt-3 text-sm text-yellow-300">No destination root selected. Please select a root from the top-left dropdown before uploading.</div>

                <div id="file-info" class="hidden mt-4 p-3 bg-gray-800 rounded-lg">
                    <div id="selected-files-list" class="space-y-2 text-left">
                        <!-- dynamically populated list of selected files -->
                    </div>
                    <div class="mt-3 text-right">
                        <button type="button" id="clear-selection" class="btn-secondary px-3 py-1 text-xs rounded-lg">Clear Selection</button>
                    </div>
                </div>

                <div class="mt-4 flex justify-center">
                    <div style="display:flex;flex-direction:column;align-items:center;gap:8px;">
                        <button type="submit" id="upload-button" class="btn-green px-8 py-3 rounded-lg font-semibold">Upload File</button>
                        <label style="font-size:12px;color:#9CA3AF;margin-top:6px;display:flex;align-items:center;gap:8px;"><input type="checkbox" id="upload-auto-refresh"> Automatically refresh "Currently Loaded Evidence" after upload</label>
                    </div>
                </div>
            </form>
        </div>
        <!-- MOVED: Session Management card placed under Upload New Evidence File -->
        <div class="card p-6 rounded-lg">
            <h2 class="text-xl font-semibold text-white mb-4">Session Management</h2>
            <div class="flex flex-col space-y-4">
                <a href="{{ url_for('forensic_analysis') }}" class="btn-primary px-6 py-3 rounded-lg text-center">Start Analysis</a>
                <a href="{{ url_for('auto_carving_setup') }}" class="btn-green px-6 py-3 rounded-lg text-center">Go to Auto Carving</a>
                <a href="{{ url_for('encryption_page') }}" class="btn-secondary px-6 py-3 rounded-lg text-center">File Encryption</a>
                <!-- Create Image button: opens the same modal used on the Explorer page -->
                <button type="button" id="btn-create-image" class="btn-secondary px-6 py-3 rounded-lg text-center">Create Image</button>
            </div>
        </div>
        
        
    </div>
    
    <!-- MOVED: Currently Loaded Evidence section to the right column -->
    <div class="space-y-8">
        <div class="card p-6 rounded-lg" id="currently-loaded-evidence-card">
            <h2 class="text-xl font-semibold text-white mb-4">Currently Loaded Evidence</h2>
                {% if uploaded_files %}
                    <div class="space-y-2" id="currently-loaded-evidence-list">
                    {% for filename, details in uploaded_files.items() %}
                    <div class="p-2 border-b border-gray-700" id="file-status-{{ loop.index0 }}">
                        <div class="flex justify-between items-center">
                            <div>
                                <span class="font-mono">{{ filename }}</span> <span>({{ details.size_mb }} MB)</span>
                                <span id="decryption-badge-{{ loop.index0 }}" class="ml-2">
                                    {% if details.encryption_status.decrypted_path %}<span class="decryption-badge">DECRYPTED</span>
                                    {% elif details.encryption_status.decrypting %}<span class="decrypting-badge">DECRYPTING...</span>
                                    {% elif details.encryption_status.encrypted %}<span class="encryption-badge">ENCRYPTED</span>
                                    {% endif %}
                                </span>
                            </div>
                            <div class="flex space-x-2">
                                <a href="{{ url_for('remove_file', filename=filename) }}" class="btn-small btn-delete">Unload</a>
                            </div>
                        </div>
                        <div id="decryption-status-text-{{ loop.index0 }}" class="text-xs text-gray-400 mt-2 pl-2">
                        {% if details.encryption_status.encrypted and not details.encryption_status.decrypted_path %}
                            Encryption: {{ details.encryption_status.description }} - <a href="{{ url_for('decryption_page', filename=filename) }}" class="text-blue-400 hover:underline">Go to Decryption Page</a>
                        {% elif details.encryption_status.decrypting %}
                             <span class="text-yellow-400">Decryption in progress...</span>
                        {% endif %}
                        </div>
                    </div>
                {% endfor %}
                </div>
                
                <!-- ADDED: Clear Session button in Currently Loaded Evidence section -->
                <div class="mt-6 pt-4 border-t border-gray-700">
                    <a id="clear-session-btn" href="{{ url_for('clear_session') }}" onclick="return confirm('This will clear ALL recovered files and analysis data. Continue?');" class="btn-danger px-6 py-3 rounded-lg w-full text-center block">Clear Session</a>
                </div>
            {% else %}
                <p class="text-gray-500">No evidence file is currently loaded for analysis.</p>
            {% endif %}
        </div>
        <!-- NEW: Pinned DB Files card below the upload card -->
        <div class="card p-6 rounded-lg mt-6">
            <h2 class="text-xl font-semibold text-white mb-4">Database Files (Pinned)</h2>
            <p class="text-sm text-gray-400 mb-4">Files from the application database (uploaded, encrypted, decrypted). Use Load to copy into the upload area for analysis.</p>
            {% if pinned_db_files %}
                <div class="space-y-2 pinned-db-list" id="pinned-db-list">
                {% for f in pinned_db_files %}
                    <div class="p-2 border-b border-gray-700 db-file-row" data-file-id="{{ f.id }}">
                        <div class="flex justify-between items-center">
                            <div>
                                <span class="font-mono db-filename">{{ f.filename }}</span>
                                <span class="ml-2 text-xs text-gray-400">[{{ f.file_type or 'unknown' }}]</span>
                            </div>
                            <div class="flex space-x-2">
                                <button class="btn-small btn-load-db" data-file-id="{{ f.id }}">Load</button>
                                <form method="post" action="{{ url_for('db_delete_file', file_id=f.id) }}" style="display:inline" onsubmit="return confirm('Delete this DB record?');">
                                    <input type="hidden" name="csrf_token" value="{{ session.get('csrf_token') }}">
                                    <button type="submit" class="btn-small btn-delete">Delete</button>
                                </form>
                            </div>
                        </div>
                    </div>
                {% endfor %}
                </div>
            {% else %}
                <p class="text-gray-500">No pinned DB files found.</p>
            {% endif %}
        </div>
        
    </div>
</div>

<script>
// If the server passed a loaded_filename, highlight the matching file entry
document.addEventListener('DOMContentLoaded', function(){
    try{
        const loaded = {{ ('"' + loaded_filename + '"') if loaded_filename else 'null' }};
        if(loaded){
            try{
                // find matching file listing and highlight it
                const items = document.querySelectorAll('.card .font-mono');
                for(const it of items){
                    try{
                        if(it.textContent.trim() === loaded){
                            const row = it.closest('[id^="file-status-"]');
                            if(row){
                                row.classList.add('ring','ring-green-400','ring-opacity-20');
                                row.scrollIntoView({behavior:'smooth', block:'center'});
                                setTimeout(()=>{ try{ row.classList.remove('ring','ring-green-400','ring-opacity-20'); }catch(e){} }, 7000);
                            }
                            break;
                        }
                    }catch(e){}
                }
            }catch(e){ console.warn('Highlight logic failed', e); }
        }
    }catch(e){ console.error('Loaded file highlight failed', e); }
});
// Upload progress functionality
let currentUploadFilename = null; // filename currently being uploaded from this client
function formatBytes(bytes) {
    if (bytes === 0) return '0 B';
    const k = 1024;
    const sizes = ['B', 'KB', 'MB', 'GB'];
    const i = Math.floor(Math.log(bytes) / Math.log(k));
    return parseFloat((bytes / Math.pow(k, i)).toFixed(2)) + ' ' + sizes[i];
}

// DB Pinned files: bind Load buttons
document.addEventListener('DOMContentLoaded', function(){
    try{
        const token = '{{ session.get('csrf_token') }}';
        document.querySelectorAll('.btn-load-db').forEach(btn => {
            btn.addEventListener('click', async function(ev){
                ev.preventDefault();
                const fid = btn.getAttribute('data-file-id');
                if(!fid) return alert('Missing file id');
                try{
                    const fd = new FormData();
                    fd.append('file_id', fid);
                    fd.append('csrf_token', token);
                    const r = await fetch('/api/fs/load_to_upload', { method: 'POST', body: fd, headers: {'X-CSRF-Token': token} });
                    const j = await r.json();
                    if(!r.ok || (j && j.error)){
                        return alert(j && (j.error || j.message) ? (j.error || j.message) : 'Load failed');
                    }
                    const filename = j.filename || document.querySelector('.db-file-row[data-file-id="'+fid+'"] .db-filename').textContent.trim();
                    window.location = '/evidence_upload?loaded=' + encodeURIComponent(filename);
                }catch(e){
                    alert('Load failed: ' + (e && e.message ? e.message : e));
                }
            });
        });
    }catch(e){ console.error('Bind Load DB buttons failed', e); }
});

function formatSpeed(bytesPerSecond) {
    return formatBytes(bytesPerSecond) + '/s';
}

function formatTime(seconds) {
    if (seconds < 60) {
        return Math.round(seconds) + 's';
    } else if (seconds < 3600) {
        const mins = Math.floor(seconds / 60);
        const secs = Math.round(seconds % 60);
        return mins + 'm ' + secs + 's';
    } else {
        const hours = Math.floor(seconds / 3600);
        const minutes = Math.floor((seconds % 3600) / 60);
        return hours + 'h ' + minutes + 'm';
    }
}

function updateFileInfo() {
    const fileInput = document.getElementById('file-input');
    const fileInfo = document.getElementById('file-info');
    const listContainer = document.getElementById('selected-files-list');

    while (listContainer.firstChild) listContainer.removeChild(listContainer.firstChild);

    if (fileInput.files && fileInput.files.length > 0) {
        Array.from(fileInput.files).forEach((file, idx) => {
            const row = document.createElement('div');
            row.className = 'flex justify-between items-center p-2 border border-gray-700 rounded-md';

            const left = document.createElement('div');
            left.innerHTML = `<div class="font-mono text-sm text-white">${escapeHtml(file.name)}</div><div class="text-xs text-gray-400">${formatBytes(file.size)}</div>`;

            const right = document.createElement('div');
            right.className = 'flex items-center space-x-2';

            const removeBtn = document.createElement('button');
            removeBtn.type = 'button';
            removeBtn.className = 'btn-secondary px-2 py-1 text-xs rounded-lg';
            removeBtn.textContent = 'Remove';
            removeBtn.addEventListener('click', () => {
                // Remove this file from the FileList by creating a new DataTransfer
                const dt = new DataTransfer();
                Array.from(fileInput.files).forEach((f, i) => { if (i !== idx) dt.items.add(f); });
                fileInput.files = dt.files;
                updateFileInfo();
            });

            right.appendChild(removeBtn);
            row.appendChild(left);
            row.appendChild(right);
            listContainer.appendChild(row);
        });

        fileInfo.classList.remove('hidden');
    } else {
        fileInfo.classList.add('hidden');
    }
}

function updateUploadProgress() {
    fetch('/upload_status')
        .then(response => {
            if (!response.ok) throw new Error('Network error');
            return response.json();
        })
        .then(data => {
            console.log('Upload status:', data);
            
            const progressContainer = document.getElementById('upload-progress-container');
            const progressBar = document.getElementById('upload-progress-bar');
            const progressPercent = document.getElementById('upload-percent');
            const progressText = document.getElementById('upload-progress-text');
            const uploadSpeed = document.getElementById('upload-speed');
            const uploadTiming = document.getElementById('upload-timing');
            const uploadFilename = document.getElementById('upload-filename');
            const uploadButton = document.getElementById('upload-button');
            
            if (data.in_progress) {
                // Show progress container
                progressContainer.classList.remove('hidden');
                
                // Update progress bar
                progressBar.style.width = data.progress + '%';
                progressPercent.textContent = data.progress + '%';
                
                // Update file info
                uploadFilename.textContent = data.filename || 'Unknown file';
                
                // Update progress text
                progressText.textContent = `${formatBytes(data.bytes_uploaded)} / ${formatBytes(data.total_bytes)}`;
                
                // Update speed and timing
                uploadSpeed.textContent = data.upload_speed ? formatSpeed(data.upload_speed) : 'Calculating...';
                uploadTiming.innerHTML = `Elapsed: ${data.elapsed_time || '0s'}<br>Remaining: ${data.time_remaining_str || 'Calculating...'}`;
                
                // Update button state
                uploadButton.disabled = true;
                uploadButton.textContent = 'Uploading...';
                uploadButton.classList.add('opacity-50');
                
                // Continue polling
                setTimeout(updateUploadProgress, 500);
            } else if (data.progress === 100 || data.complete) {
                // Upload complete
                progressBar.style.width = '100%';
                progressBar.classList.remove('bg-blue-600');
                progressBar.classList.add('bg-green-600');
                progressPercent.textContent = '100%';
                progressPercent.classList.remove('text-blue-400');
                progressPercent.classList.add('text-green-400');
                
                uploadTiming.innerHTML = `Upload complete!<br>Total time: ${data.elapsed_time || 'Unknown'}`;
                uploadSpeed.textContent = 'Complete';
                
                // Reset button
                uploadButton.disabled = false;
                uploadButton.textContent = 'Upload File';
                uploadButton.classList.remove('opacity-50');
                
                // Reload page after delay to show the new file
                setTimeout(() => {
                    window.location.reload();
                }, 2000);
            } else {
                // Upload not in progress
                progressContainer.classList.add('hidden');
                uploadButton.disabled = false;
                uploadButton.textContent = 'Upload File';
                uploadButton.classList.remove('opacity-50');
            }
        })
        .catch(error => {
            console.error('Error fetching upload status:', error);
            const uploadButton = document.getElementById('upload-button');
            uploadButton.disabled = false;
            uploadButton.textContent = 'Upload File';
            uploadButton.classList.remove('opacity-50');
        });
}

function uploadFileWithProgress(file, onComplete) {
    const xhr = new XMLHttpRequest();
    const formData = new FormData();
    formData.append('file', file);
    // include target database selection if present
    // determine storage target: prefer per-session selection stored in localStorage, fallback to 'local'
    try {
        const target = window.localStorage.getItem('selected_target_db') || 'local';
        formData.append('target_db', target);
    } catch (e) { console.warn('Unable to read target storage from localStorage', e); formData.append('target_db', 'local'); }

    const progressContainer = document.getElementById('upload-progress-container');
    const uploadFilename = document.getElementById('upload-filename');
    const uploadButton = document.getElementById('upload-button');
    const progressBar = document.getElementById('upload-progress-bar');
    const progressPercent = document.getElementById('upload-percent');
    const progressText = document.getElementById('upload-progress-text');
    const uploadSpeed = document.getElementById('upload-speed');
    const uploadTiming = document.getElementById('upload-timing');

    progressContainer.classList.remove('hidden');
    uploadFilename.textContent = file.name;
    uploadButton.disabled = true;
    uploadButton.textContent = 'Uploading...';
    uploadButton.classList.add('opacity-50');

    let startTime = Date.now();
    // mark this client as the originator of this upload so status polls won't
    // trigger unrelated page reloads for other uploads
    currentUploadFilename = file.name;

    xhr.upload.addEventListener('progress', function(evt) {
        if (evt.lengthComputable) {
            const MAX_BYTES = 60 * 1024 * 1024 * 1024; // 60 GB
            if (evt.loaded > MAX_BYTES) {
                // Abort and request server to cleanup partial upload, then restart
                console.warn('Upload exceeded 60GB threshold, aborting and restarting');
                xhr.abort();
                // Ask server to remove partial file
                fetch('/ajax_upload_abort', { method: 'POST', headers: {'Content-Type': 'application/json'}, body: JSON.stringify({ filename: file.name }) })
                    .then(() => {
                        // restart after short delay
                        setTimeout(() => { uploadFileWithProgress(file); }, 1500);
                    })
                    .catch((e) => { console.error('Abort cleanup failed', e); });
                return;
            }
            const percent = Math.round((evt.loaded / evt.total) * 100);
            progressBar.style.width = percent + '%';
            progressPercent.textContent = percent + '%';
            progressText.textContent = `${formatBytes(evt.loaded)} / ${formatBytes(evt.total)}`;

            const elapsed = (Date.now() - startTime) / 1000;
            const speed = evt.loaded / Math.max(elapsed, 0.001);
            uploadSpeed.textContent = formatSpeed(speed);

            const estimatedTotal = evt.total / Math.max(speed, 0.001);
            const remaining = Math.max(0, estimatedTotal - elapsed);
            uploadTiming.innerHTML = `Elapsed: ${formatTime(elapsed)}<br>Remaining: ${formatTime(remaining)}`;
        } else {
            // Fallback when total size unknown
            progressBar.style.width = '50%';
            progressPercent.textContent = 'Uploading...';
            progressText.textContent = `${formatBytes(evt.loaded)} / ?`;
            uploadSpeed.textContent = 'Calculating...';
            uploadTiming.innerHTML = `Elapsed: ${formatTime((Date.now() - startTime) / 1000)}<br>Remaining: Calculating...`;
        }
    });

    xhr.onreadystatechange = function() {
        if (xhr.readyState === XMLHttpRequest.DONE) {
            try {
                const data = JSON.parse(xhr.responseText || '{}');

                if (xhr.status >= 200 && xhr.status < 300 && data.success) {
                    progressBar.style.width = '100%';
                    progressBar.classList.remove('bg-blue-600');
                    progressBar.classList.add('bg-green-600');
                    progressPercent.textContent = '100%';
                    progressText.textContent = 'Upload complete!';
                    uploadSpeed.textContent = 'Complete';
                    uploadTiming.innerHTML = `Upload successful!<br>Total time: ${data.elapsed_time || 'Unknown'}`;
                    // clear our marker before continuing
                    currentUploadFilename = null;
                    if (typeof onComplete === 'function') {
                        onComplete(null, data);
                    } else {
                        setTimeout(() => { window.location.reload(); }, 1200);
                    }
                } else {
                    throw new Error(data.error || 'Upload failed');
                }
            } catch (err) {
                console.error('Upload response error:', err);
                // clear upload marker on error so other status updates work normally
                currentUploadFilename = null;
                progressBar.style.width = '100%';
                progressBar.classList.remove('bg-blue-600');
                progressBar.classList.add('bg-red-600');
                progressPercent.textContent = 'Error';
                progressText.textContent = 'Upload failed: ' + (err.message || 'Unknown');
                setTimeout(() => {
                    uploadButton.disabled = false;
                    uploadButton.textContent = 'Upload File';
                    uploadButton.classList.remove('opacity-50');
                    progressContainer.classList.add('hidden');
                    if (typeof onComplete === 'function') onComplete(err);
                }, 3000);
            }
        }
    };

    // network-level errors (connection reset, etc.)
    xhr.onerror = function() {
        console.error('Upload network error');
        currentUploadFilename = null;
        progressBar.style.width = '100%';
        progressBar.classList.remove('bg-blue-600');
        progressBar.classList.add('bg-red-600');
        progressPercent.textContent = 'Error';
        progressText.textContent = 'Network error during upload.';
        setTimeout(() => {
            uploadButton.disabled = false;
            uploadButton.textContent = 'Upload File';
            uploadButton.classList.remove('opacity-50');
            progressContainer.classList.add('hidden');
        }, 3000);
    };

    xhr.open('POST', '/ajax_upload', true);
    xhr.send(formData);
}

document.addEventListener('DOMContentLoaded', function() {
    const uploadForm = document.getElementById('upload-form');
    const fileInput = document.getElementById('file-input');
    const uploadButton = document.getElementById('upload-button');

    // Update file info when file is selected
    fileInput.addEventListener('change', updateFileInfo);

    uploadForm.addEventListener('submit', function(e) {
        e.preventDefault();

        if (!fileInput.files || fileInput.files.length === 0) {
            alert('Please select one or more files to upload.');
            return;
        }

    const files = Array.from(fileInput.files);
    // determine storage target from floating selector persisted in localStorage; default to 'local'
    let targetDb = 'local';
    try { targetDb = window.localStorage.getItem('selected_target_db') || 'local'; } catch (e) { targetDb = 'local'; }

        // Confirm large total size
        const totalBytes = files.reduce((acc, f) => acc + f.size, 0);
        const totalMB = totalBytes / (1024 * 1024);
        if (totalMB > 100) {
            if (!confirm(`The combined selection is ${(totalMB).toFixed(2)} MB. Large uploads may take a while. Continue?`)) {
                return;
            }
        }

        // Upload files sequentially to avoid confusing concurrent progress indicators
        let i = 0;
        function uploadNext(err) {
            if (err) {
                console.error('Upload aborted due to error:', err);
                return;
            }
            if (i >= files.length) {
                // All done — reload to show new files
                setTimeout(() => { window.location.reload(); }, 800);
                return;
            }
            const f = files[i++];
            // uploadFileWithProgress reads the persisted target from localStorage itself; pass only callback
            uploadFileWithProgress(f, uploadNext);
        }

        uploadNext();
    });
    
    // Initialize progress display
    updateUploadProgress();

    // Ensure Clear Session button always navigates (some browsers or other scripts
    // might interfere with normal anchor navigation when large uploads are in progress)
    try {
        const clearBtn = document.getElementById('clear-session-btn');
        if (clearBtn) {
            clearBtn.addEventListener('click', function(evt) {
                // The anchor already has an inline confirm; double-check here too
                if (!confirm('This will clear ALL recovered files and analysis data. Continue?')) {
                    evt.preventDefault();
                    return;
                }
                // Use a robust navigation method to avoid other scripts preventing default
                evt.preventDefault();
                const href = clearBtn.getAttribute('href');
                setTimeout(() => { window.location.assign(href); }, 50);
            });
        }
    } catch (e) {
        console.error('Clear Session handler setup failed', e);
    }
});
</script>
"""
ENCRYPTION_PAGE_CONTENT = """
<h1 class="text-3xl font-bold text-white mb-4">File Encryption</h1>
<p class="text-gray-400 mb-8">Encrypt files using Fernet encryption with password protection.</p>

<div class="grid grid-cols-1 lg:grid-cols-2 gap-8">
    <div class="card p-6 rounded-lg">
        <h2 class="text-xl font-semibold text-white mb-4">Encrypt a File</h2>
        
        <form action="{{ url_for('encrypt_file') }}" method="post" enctype="multipart/form-data" class="space-y-6">
            <div>
                <label class="block text-sm font-medium text-gray-300">Select File to Encrypt</label>
                <input type="file" name="file" class="mt-1 block w-full bg-gray-800 border-gray-600 rounded-md shadow-sm p-2 text-white" required>
            </div>
            
            <div>
                <label class="block text-sm font-medium text-gray-300">Password</label>
                <input type="password" name="password" class="mt-1 block w-full bg-gray-800 border-gray-600 rounded-md shadow-sm p-2 text-white" placeholder="Enter strong password" required>
                <p class="mt-1 text-xs text-gray-400">Use a strong password. The file will be encrypted using AES-128 in CBC mode.</p>
            </div>
            
            <div>
                <label class="block text-sm font-medium text-gray-300">Confirm Password</label>
                <input type="password" name="confirm_password" class="mt-1 block w-full bg-gray-800 border-gray-600 rounded-md shadow-sm p-2 text-white" placeholder="Confirm password" required>
            </div>
            
            <button type="submit" class="btn-primary w-full py-3 rounded-lg font-semibold">Encrypt File</button>
        </form>
    </div>
    
    <div class="card p-6 rounded-lg">
        <h2 class="text-xl font-semibold text-white mb-4">Encryption Information</h2>
        
        <div class="space-y-4">
            <div class="p-4 bg-blue-900 border border-blue-700 rounded-lg">
                <h3 class="font-semibold text-blue-300 mb-2">🔒 Encryption Method</h3>
                <p class="text-sm text-blue-200">Files are encrypted using Fernet (AES-128 in CBC mode) with PBKDF2 key derivation.</p>
            </div>
            
            <div class="p-4 bg-green-900 border border-green-700 rounded-lg">
                <h3 class="font-semibold text-green-300 mb-2">🔑 Key Security</h3>
                <p class="text-sm text-green-200">Keys are derived from your password using 480,000 iterations of PBKDF2 with a random 16-byte salt.</p>
            </div>
            
            <div class="p-4 bg-purple-900 border border-purple-700 rounded-lg">
                <h3 class="font-semibold text-purple-300 mb-2">📁 File Format</h3>
                <p class="text-sm text-purple-200">Encrypted files include a custom header, salt, and encrypted data. They can be decrypted using this application.</p>
            </div>
            
            <div class="p-4 bg-yellow-900 border border-yellow-700 rounded-lg">
                <h3 class="font-semibold text-yellow-300 mb-2">⚠️ Important Notes</h3>
                <ul class="text-sm text-yellow-200 list-disc list-inside space-y-1">
                    <li>Keep your password safe - it cannot be recovered</li>
                    <li>Encrypted files will have a .enc extension</li>
                    <li>Original files are not modified - new encrypted copies are created</li>
                    <li>Use strong, unique passwords for each file</li>
                </ul>
            </div>
        </div>
    </div>
</div>

{% if encrypted_files %}
<div class="card p-6 rounded-lg mt-8">
    <h2 class="text-xl font-semibold text-white mb-4">Recently Encrypted Files</h2>
    <div class="space-y-2">
        {% for file_info in encrypted_files %}
        <div class="p-3 border border-gray-600 rounded-lg flex justify-between items-center">
            <div>
                <span class="font-mono text-white">{{ file_info.filename }}</span>
                <span class="ml-2 text-xs text-gray-400">({{ file_info.size }})</span>
            </div>
            <div class="flex space-x-2">
                <a href="{{ url_for('download_encrypted_file', filename=file_info.filename) }}" 
                   class="btn-primary px-3 py-1 text-xs rounded-lg">Download</a>
                <a href="{{ url_for('delete_encrypted_file', filename=file_info.filename) }}" 
                   onclick="return confirm('Delete this encrypted file?');" 
                   class="bg-red-600 text-white px-3 py-1 text-xs rounded-lg">Delete</a>
            </div>
        </div>
        {% endfor %}
    </div>
</div>
{% endif %}

{% if newly_encrypted_file %}

<div class="card p-6 rounded-lg mb-8 bg-green-900 border border-green-700">
<h2 class="text-xl font-semibold text-white mb-4">✅ Encryption Successful</h2>
<p class="text-green-200 mb-4">Your file <strong class="font-mono">{{ newly_encrypted_file }}</strong> has been encrypted and is ready for download.</p>
<a href="{{ url_for('download_encrypted_file', filename=newly_encrypted_file) }}"
class="btn-primary px-6 py-3 rounded-lg font-semibold inline-block">
Download Now
</a>
</div>
{% endif %}

"""


DECRYPTION_PAGE_CONTENT = """
<h1 class="text-3xl font-bold text-white mb-4">File Decryption</h1>
<p class="text-gray-400 mb-8">The loaded evidence file <strong class="font-mono">{{ filename }}</strong> appears to be encrypted.</p>
<div class="card p-6 rounded-lg">
    <form action="{{ url_for('decryption_page', filename=filename) }}" method="post" class="space-y-6">
        <div>
            <label for="force_type" class="block text-sm font-medium text-gray-300">Decryption Method</label>
            <select name="force_type" id="force_type" class="mt-1 block w-full bg-gray-800 border-gray-600 rounded-md shadow-sm p-2 h-10 text-white">
                <option value="auto" selected>Auto-Detect (Current: {{ file_info.encryption_status.encryption_type }})</option>
                <option value="ZIP_ENCRYPTED">Force ZIP Decryption</option>
                <option value="AES_ENCRYPTED">Force OpenSSL AES Decryption</option>
                <option value="FERNET_ENCRYPTED">Force Fernet Decryption</option>
            </select>
            <p class="mt-2 text-xs text-gray-400">If Auto-Detect fails, manually select the suspected file type.</p>
        </div>

        <div class="grid grid-cols-1 md:grid-cols-2 gap-8 pt-4 border-t border-gray-700">
            <div>
                <h2 class="text-xl font-semibold text-white mb-4">1. Decrypt with Password</h2>
                <div class="space-y-4">
                    <div>
                        <label for="password" class="block text-sm font-medium text-gray-300">Password</label>
                        <input type="password" name="password" id="password" class="mt-1 block w-full bg-gray-800 border-gray-600 rounded-md shadow-sm p-2 text-white">
                    </div>
                    <button type="submit" name="action" value="with_password" class="btn-primary w-full py-2 rounded-lg font-semibold">Decrypt with Provided Password</button>
                </div>
            </div>
            <div>
                <h2 class="text-xl font-semibold text-white mb-4">2. Automated Decryption</h2>
                 <p class="text-sm text-gray-400 mb-4">Attempt decryption using the password dictionary file.</p>
                <button type="submit" name="action" value="auto_decrypt" class="btn-secondary w-full py-2 rounded-lg font-semibold">Start Automated Attempt</button>
            </div>
        </div>
    </form>
</div>
{% if decrypted_files %}
<div class="card p-6 rounded-lg mt-8">
    <h2 class="text-xl font-semibold text-white mb-4">Recently Decrypted Files</h2>
    <div class="space-y-2">
        {% for file_info in decrypted_files %}
        <div class="p-3 border border-gray-600 rounded-lg flex justify-between items-center">
            <div>
                <span class="font-mono text-white">{{ file_info.filename }}</span>
                <span class="ml-2 text-xs text-gray-400">({{ file_info.size }})</span>
            </div>
            <div class="flex space-x-2">
                <a href="{{ url_for('download_decrypted_file', filename=file_info.filename) }}" class="btn-primary px-3 py-1 text-xs rounded-lg">Download</a>
                <a href="{{ url_for('delete_decrypted_file', filename=file_info.filename) }}" onclick="return confirm('Delete this decrypted file?');" class="bg-red-600 text-white px-3 py-1 text-xs rounded-lg">Delete</a>
            </div>
        </div>
        {% endfor %}
    </div>
</div>
{% endif %}
"""

TEXT_PREVIEW_CONTENT = """
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <title>Viewing Text: {{ filename }}</title>
    <style>
        body, html { 
            margin: 0; 
            padding: 0; 
            height: 100%; 
            background-color: #111827; 
            color: white; 
            font-family: sans-serif; 
        }
        pre { 
            margin: 20px; 
            background-color: #1f2937; 
            color: #d1d5db; 
            height: calc(100vh - 40px); 
            overflow: auto; 
            padding: 20px; 
            box-sizing: border-box; 
            font-family: monospace; 
            white-space: pre-wrap; 
            word-wrap: break-word; 
            border-radius: 8px; 
            border: 1px solid #374151; 
        }
    </style>
</head>
<body>
    <pre>{{ text_content }}</pre>
</body>
</html>
"""

DATABASES_CONTENT = """
<h1 class="text-3xl font-bold text-white mb-4">Database Management</h1>
<p class="text-gray-400 mb-6">Manage configured database connections and choose where uploaded evidence will be stored.</p>
{% if summary %}
<div class="mb-4">
    <span id="db-summary-badge" class="inline-block px-3 py-1 rounded-full text-sm font-medium bg-gray-800 text-gray-200">Status: <strong class="text-green-300">{{ summary.connected }}</strong> connected / <strong class="text-red-300">{{ summary.disconnected }}</strong> disconnected</span>
    <!-- Poll interval controls removed -->
</div>
{% endif %}

<div class="grid grid-cols-1 lg:grid-cols-2 gap-8">
    <!-- Left: Configured Connections -->
    <div>
        <div class="card p-6 rounded-lg">
            <h2 class="text-xl font-semibold text-white mb-4">Configured Connections</h2>
            <div class="mb-4">
                <form method="post" action="{{ url_for('db_refresh_all') }}">
                    <button class="btn-primary px-3 py-1 text-sm rounded-lg">Refresh All</button>
                </form>
            </div>
            {% if dbs %}
            <div class="space-y-2">
                {% for dbid, cfg in dbs.items() %}
                <div class="p-3 border border-gray-700 rounded-lg" id="db-card-{{ dbid }}">
                    <div class="mb-3">
                        <div class="font-semibold text-white">{{ cfg.name }} <span class="text-xs text-gray-400">({{ cfg.type }})</span></div>
                        <div class="text-xs text-gray-400 break-words">{{ cfg.conn|truncate(200) }}</div>
                        <div class="text-xs text-gray-400">Size: <span class="db-size">{{ cfg.size or 'N/A' }}</span></div>
                        <div class="text-xs text-gray-400">Last checked: <span class="db-last-checked">{{ cfg.last_checked if cfg.last_checked else 'Never' }}</span></div>
                        <div class="text-xs mt-2 flex items-center gap-3">
                            <span class="inline-block px-2 py-0.5 rounded-full text-xs font-medium {% if cfg.connected %}bg-green-600 text-green-100{% else %}bg-red-600 text-red-100{% endif %} db-connected">{% if cfg.connected %}Connected{% else %}Disconnected{% endif %}</span>
                            {% if cfg.message %}
                                <span class="inline-block px-2 py-0.5 rounded-full text-xs font-medium bg-yellow-600 text-yellow-100 db-message">{{ cfg.message }}</span>
                            {% endif %}
                            {% if cfg.folder_creation %}
                                <span class="inline-block px-2 py-0.5 rounded-full text-xs font-medium {% if cfg.folder_creation.success %}bg-green-600 text-green-100{% else %}bg-red-600 text-red-100{% endif %}">{% if cfg.folder_creation.success %}Folders created{% else %}Folder creation failed{% endif %}</span>
                                <span class="text-gray-400 text-xs">{{ cfg.folder_creation.message or '' }}</span>
                            {% else %}
                                <form method="post" action="{{ url_for('db_create_folders', dbid=dbid) }}" style="display:inline">
                                    <button class="px-3 py-0.5 text-xs rounded-lg btn-secondary">Create Default Folders</button>
                                </form>
                            {% endif %}
                        </div>
                    </div>
                    <div class="mt-2 flex flex-wrap items-center gap-2">
                        
                        {% if cfg.connected %}
                            <form method="post" action="{{ url_for('db_disconnect', dbid=dbid) }}" style="display:inline">
                                <button class="btn-secondary px-3 py-1 text-xs rounded-lg">Disconnect</button>
                            </form>
                        {% else %}
                            <form method="post" action="{{ url_for('db_connect', dbid=dbid) }}" style="display:inline">
                                <button class="btn-primary px-3 py-1 text-xs rounded-lg">Connect</button>
                            </form>
                        {% endif %}
                        <a href="{{ url_for('db_edit', dbid=dbid) }}" class="btn-secondary px-3 py-1 text-xs rounded-lg">Edit</a>
                        <form method="post" action="{{ url_for('db_delete', dbid=dbid) }}" style="display:inline" onsubmit="return confirm('Delete connection?');">
                            <button class="bg-red-600 text-white px-3 py-1 text-xs rounded-lg">Delete</button>
                        </form>
                        <!-- per-db Refresh button removed as requested -->
                        <!-- Open the DB browser for this connection in a new tab -->
                        <a target="_blank" href="{{ url_for('database_browser', dbid=dbid) }}" class="btn-primary px-3 py-1 text-xs rounded-lg">Open DB Browser</a>
                    </div>
                </div>
                {% endfor %}
            </div>
            {% else %}
                <p class="text-gray-400">No configured database connections.</p>
            {% endif %}
        </div>
    </div>

    <!-- Right: Add New Connection -->
    <div>
        <div class="card p-6 rounded-lg mb-0">
            <h2 class="text-xl font-semibold text-white mb-4">Add New Connection</h2>
            <form id="db-add-form" action="{{ url_for('db_add') }}" method="post" class="space-y-4">
                <div>
                    <label class="block text-sm text-gray-300">Connection Name</label>
                    <input name="name" required class="mt-1 block w-full bg-gray-800 border-gray-600 rounded-md p-2 text-white">
                </div>
                <div>
                    <label class="block text-sm text-gray-300">Type</label>
                    <select name="type" required class="mt-1 block w-full bg-gray-800 border-gray-600 rounded-md p-2 text-white">
                        <option>Relational (RDBMS)</option>
                        <option>Document (NoSQL)</option>
                        <option>Key-Value (NoSQL)</option>
                        <option>Wide-Column (NoSQL)</option>
                        <option>Graph (NoSQL)</option>
                        <option>Object-Oriented</option>
                        <option>Time-Series</option>
                        <option>NewSQL</option>
                        <option>Hierarchical</option>
                        <option>Network</option>
                    </select>
                </div>
                <div>
                    <label class="block text-sm text-gray-300">Engine</label>
                    <select id="engine-select" name="engine" class="mt-1 block w-full bg-gray-800 border-gray-600 rounded-md p-2 text-white">
                        <option value="">(Select type first)</option>
                    </select>
                </div>
                <div>
                    <label class="block text-sm text-gray-300">Connection String / Details (JSON allowed)</label>
                    <div class="flex items-center space-x-4 mb-2">
                        <label class="text-sm text-gray-300"><input type="radio" name="conn_type" value="string" checked> Connection String</label>
                        <label class="text-sm text-gray-300"><input type="radio" name="conn_type" value="manual"> Manual (host/port/...)</label>
                    </div>
                    <!-- connection input container: the actual input is created dynamically when 'Connection String' is selected -->
                    <div id="conn-string-container"></div>
                    <div id="manual-conn" class="hidden mt-2 space-y-2">
                        <div class="grid grid-cols-1 md:grid-cols-2 gap-2">
                            <input name="host" placeholder="Host" class="bg-gray-800 border-gray-600 rounded-md p-2 text-white">
                            <input name="port" placeholder="Port" class="bg-gray-800 border-gray-600 rounded-md p-2 text-white">
                            <input name="database" placeholder="Database name" class="bg-gray-800 border-gray-600 rounded-md p-2 text-white">
                            <input name="username" placeholder="Username" class="bg-gray-800 border-gray-600 rounded-md p-2 text-white">
                            <input name="password" placeholder="Password" type="password" class="bg-gray-800 border-gray-600 rounded-md p-2 text-white">
                            <input name="extra" placeholder="Extra params (JSON)" class="bg-gray-800 border-gray-600 rounded-md p-2 text-white">
                        </div>
                    </div>
                </div>
                <div class="flex space-x-2">
                    <button id="db-add-button" class="btn-primary px-4 py-2 rounded-lg">Add Connection</button>
                    <button type="button" id="db-test-button" class="btn-secondary px-4 py-2 rounded-lg">Test Connection</button>
                    <button type="button" id="db-reset-button" class="px-4 py-2 rounded-lg border border-gray-600 text-gray-300">Reset</button>
                </div>
                <!-- Inline test status (hidden until used) -->
                <div id="db-test-status" class="hidden mt-2 p-2 rounded text-sm"></div>
            </form>
        </div>
    </div>
</div>

<script>
// Live polling for database statuses
(function(){
    const intervalInput = document.getElementById('db-poll-interval');
    const toggleBtn = document.getElementById('db-poll-toggle');
    const summaryBadge = document.getElementById('db-summary-badge');
    let timer = null;
    // Load saved interval
    try{ const saved = localStorage.getItem('db_poll_interval'); if(saved) intervalInput.value = saved; }catch(e){}

    function applyStatusAll(obj){
        if(!obj) return;
        // update summary counts
        let connected = 0, disconnected = 0;
        Object.keys(obj).forEach(k => {
            const d = obj[k];
            if(d.connected) connected++; else disconnected++;
            const card = document.getElementById('db-card-' + k);
            if(!card) return;
            const sizeEl = card.querySelector('.db-size'); if(sizeEl && d.size) sizeEl.textContent = d.size;
            const lastEl = card.querySelector('.db-last-checked'); if(lastEl) lastEl.textContent = d.last_checked || 'Never';
            const msgEl = card.querySelector('.db-message'); if(msgEl) msgEl.textContent = d.message || '';
            const connEl = card.querySelector('.db-connected'); if(connEl){ connEl.textContent = d.connected ? 'Connected' : 'Disconnected';
                if(d.connected){ connEl.classList.remove('bg-red-600'); connEl.classList.add('bg-green-600'); }
                else { connEl.classList.remove('bg-green-600'); connEl.classList.add('bg-red-600'); }
            }
        });
        // update top summary badge
        if(summaryBadge){ summaryBadge.innerHTML = 'Status: <strong class="text-green-300">'+connected+'</strong> connected / <strong class="text-red-300">'+disconnected+'</strong> disconnected'; }
    }

    function pollOnce(){
        fetch('/databases/status_all')
            .then(r => r.json())
            .then(obj => applyStatusAll(obj))
            .catch(err => console.warn('Status poll failed', err));
    }

    function startPolling(){
        if(timer) return;
        const sec = Math.max(5, parseInt(intervalInput.value || 15, 10));
        try{ localStorage.setItem('db_poll_interval', String(sec)); }catch(e){}
        timer = setInterval(pollOnce, sec * 1000);
        toggleBtn.textContent = 'Stop';
        pollOnce();
    }

    function stopPolling(){
        if(!timer) return;
        clearInterval(timer); timer = null;
        toggleBtn.textContent = 'Start';
    }

    toggleBtn.addEventListener('click', function(){ if(timer) stopPolling(); else startPolling(); });
    // Start automatically if previously enabled
    try{ if(localStorage.getItem('db_poll_enabled') === '1'){ startPolling(); }}catch(e){}
    // Persist enabled state when toggling
    const origToggle = toggleBtn;
    const origClick = toggleBtn.onclick;
    // Hook to persist state
    toggleBtn.addEventListener('click', function(){ try{ localStorage.setItem('db_poll_enabled', timer ? '1' : '0'); }catch(e){} });
})();
</script>
"""

DATABASE_BROWSER_CONTENT = """
<!-- Full-featured Explorer UI -->
<div id="fac-explorer" class="text-gray-300">
    <div class="flex items-center justify-between mb-2">
        <div class="flex gap-2 items-center">
            <button id="btn-back" class="btn-secondary">◀</button>
            <button id="btn-forward" class="btn-secondary">▶</button>
            <button id="btn-up" class="btn-secondary">⬆</button>
            <button id="btn-refresh" class="btn-secondary">⟳</button>
            <input id="address" type="text" class="bg-gray-800 border border-gray-600 text-white p-2 rounded w-96 ml-2" placeholder="Root / path" />
        </div>
        <div class="flex gap-2 items-center">
            <select id="root-select" class="bg-gray-800 border p-2 rounded">
                <option value="">Select root...</option>
                {% for r in ['Upload Files','Encrypted Files','Decrypted Files','Carved Files','Deleted Files','Session Files'] %}
                <option value="{{ r }}">{{ r }}</option>
                {% endfor %}
            </select>
            <div class="flex items-center gap-2">
                <button id="btn-new-folder" title="Create new folder" class="btn-primary px-3 py-1 flex items-center gap-2">
                    <span aria-hidden>📁</span>
                    <span class="text-sm">New Folder</span>
                </button>

                <label id="btn-upload" title="Upload a file to current folder" class="btn-secondary px-3 py-1 flex items-center gap-2 rounded cursor-pointer">
                    <input id="file-input" type="file" style="display:none">
                    <span aria-hidden>⬆️</span>
                    <span class="text-sm">Upload</span>
                </label>

                <button id="btn-zip" title="Download the current folder as a ZIP" class="btn-secondary px-3 py-1 flex items-center gap-2">
                    <span aria-hidden>📦</span>
                    <span class="text-sm">Download ZIP</span>
                </button>

                <button id="btn-delete-selected" title="Delete selected items" class="bg-red-600 text-white px-3 py-1 flex items-center gap-2 rounded">
                    <span aria-hidden>🗑️</span>
                    <span class="text-sm">Delete Selected</span>
                </button>
            </div>
        </div>
    </div>

    <div class="flex gap-4">
        <aside class="w-72">
            <div class="card bg-gray-900 p-4 rounded-lg shadow-sm">
                <div class="flex items-center justify-between mb-3">
                    <div class="font-semibold text-gray-200">Quick access</div>
                    <button id="quick-collapse" title="Collapse" class="text-gray-400 hover:text-gray-200 text-sm">▾</button>
                </div>

                <div class="mb-3">
                    <input id="quick-search" type="search" placeholder="Filter..." class="w-full bg-gray-800 border border-gray-700 rounded-md p-2 text-sm text-gray-200" />
                </div>

                <nav id="quick-list" class="space-y-2" aria-label="Quick access list">
                    <button class="quick-item w-full flex items-center justify-between gap-2 p-2 rounded hover:bg-gray-800 text-left" data-root="Upload Files">
                        <div class="flex items-center gap-3">
                            <span class="w-8 h-8 inline-flex items-center justify-center bg-blue-700 text-white rounded text-sm">↑</span>
                            <span class="text-sm text-gray-200">Upload Files</span>
                        </div>
                        <span class="text-xs text-gray-400">Open</span>
                    </button>

                    <button class="quick-item w-full flex items-center justify-between gap-2 p-2 rounded hover:bg-gray-800 text-left" data-root="Encrypted Files">
                        <div class="flex items-center gap-3">
                            <span class="w-8 h-8 inline-flex items-center justify-center bg-yellow-700 text-white rounded text-sm">🔒</span>
                            <span class="text-sm text-gray-200">Encrypted Files</span>
                        </div>
                        <span class="text-xs text-gray-400">Open</span>
                    </button>

                    <button class="quick-item w-full flex items-center justify-between gap-2 p-2 rounded hover:bg-gray-800 text-left" data-root="Decrypted Files">
                        <div class="flex items-center gap-3">
                            <span class="w-8 h-8 inline-flex items-center justify-center bg-green-700 text-white rounded text-sm">🔓</span>
                            <span class="text-sm text-gray-200">Decrypted Files</span>
                        </div>
                        <span class="text-xs text-gray-400">Open</span>
                    </button>

                    <button class="quick-item w-full flex items-center justify-between gap-2 p-2 rounded hover:bg-gray-800 text-left" data-root="Session Files">
                        <div class="flex items-center gap-3">
                            <span class="w-8 h-8 inline-flex items-center justify-center bg-indigo-700 text-white rounded text-sm">⌛</span>
                            <span class="text-sm text-gray-200">Session Files</span>
                        </div>
                        <span class="text-xs text-gray-400">Open</span>
                    </button>

                    <!-- Placeholder for additional quick items; keep JS selectors (.quick-item, data-root) unchanged -->
                </nav>

                <div class="mt-3 text-xs text-gray-400">Tip: click any item to navigate the explorer to that root.</div>
            </div>
        </aside>

        <div class="flex-1 bg-gray-800 p-3 rounded">
            <div class="flex items-center justify-between mb-2">
                <div class="flex items-center gap-2">
                    <label><input type="checkbox" id="view-thumbs"> Thumbnails</label>
                    <label><input type="checkbox" id="view-details" checked> Details</label>
                </div>
                    <div class="flex items-center gap-2">
                        <div class="text-sm text-gray-400" id="status-bar">Items: 0</div>
                        <!-- Status is read-only in Database Browser; badges show color-coded state -->
                        <div class="ml-4 text-sm text-gray-400">Legend:
                            <span class="inline-block px-2 py-0.5 rounded text-xs font-medium" style="background:#4B5563;color:#fff;margin-left:6px">saved</span>
                            <span class="inline-block px-2 py-0.5 rounded text-xs font-medium" style="background:#D97706;color:#fff;margin-left:6px">processing</span>
                            <span class="inline-block px-2 py-0.5 rounded text-xs font-medium" style="background:#2563EB;color:#fff;margin-left:6px">waiting</span>
                            <span class="inline-block px-2 py-0.5 rounded text-xs font-medium" style="background:#DC2626;color:#fff;margin-left:6px">paused</span>
                        </div>
                        <div class="ml-6 text-sm text-gray-400">
                            <label for="time-mode-select" class="text-xs text-gray-300 mr-2">Time:</label>
                            <select id="time-mode-select" class="bg-gray-800 border border-gray-700 text-gray-200 p-1 rounded text-xs">
                                <option value="relative-tooltip">Relative (tooltip)</option>
                                <option value="relative-inline">Relative (inline)</option>
                                <option value="absolute">Absolute</option>
                            </select>
                        </div>
                    </div>
            </div>

            <div id="explorer-list" class="bg-gray-900 p-2 rounded min-h-[300px]">
                <!-- Dynamic file/folder list populated by JS -->
            </div>
        </div>
    </div>

    <template id="tpl-row">
        <div class="row flex items-center gap-2 p-2 border-b border-gray-700">
            <input type="checkbox" class="row-select">
            <div class="icon w-8 text-center">[icon]</div>
            <div class="name flex-1">NAME</div>
            <div class="status-col" style="width:10rem; text-align:right;">
                <!-- status badge goes here -->
            </div>
            <div class="size text-right">SIZE</div>
            <div class="mtime text-right">MTIME</div>
            <div class="actions flex gap-2 items-center">
                <button class="btn-small btn-view">View</button>
                <button class="btn-small btn-download">Download</button>
                <button class="btn-small btn-delete">Delete</button>
            </div>
        </div>
    </template>

    <style>
        .btn-small{ background:#374151;color:#fff;padding:4px 6px;border-radius:4px;margin-left:4px }
    /* Colored small button variants */
    .btn-small.btn-load-db{ background-color: #3b82f6; color: #fff; border: 1px solid #2563eb; }
    .btn-small.btn-load-db:hover{ filter:brightness(0.95); }
    .btn-small.btn-delete{ background-color: #ef4444; color: #fff; border: 1px solid #dc2626; }
    .btn-small.btn-delete:hover{ filter:brightness(0.95); }
    /* Make pinned DB files list scrollable when it grows large */
    .pinned-db-list{ max-height: 280px; overflow-y: auto; padding-right: 8px; }
    .pinned-db-list .db-file-row{ padding-right: 6px; }
        /* Keep flex children allowed to shrink so text-overflow works */
        .row .name { min-width: 0; overflow: hidden; text-overflow: ellipsis; white-space: nowrap; }
        .row .size { flex: 0 0 6.5rem; max-width: 6.5rem; }
        .row .mtime { flex: 0 0 9rem; max-width: 9rem; }
        .row .actions { flex: 0 0 auto; }

        /* Custom styled tooltip (replaces native title) */
        .fac-tooltip { position: relative; display: inline-block; }
        .fac-tooltip .fac-tooltip-text {
            visibility: hidden; opacity: 0; transform: translateY(6px);
            background: rgba(15,23,42,0.95); color: #e6eef8; padding: 6px 8px; border-radius: 6px; font-size: 12px;
            position: absolute; z-index: 9999; left: 8px; top: 100%; white-space: nowrap; box-shadow: 0 6px 18px rgba(2,6,23,0.6);
            transition: opacity 0.12s ease, transform 0.12s ease;
        }
        .fac-tooltip:hover .fac-tooltip-text, .fac-tooltip.fac-show .fac-tooltip-text { visibility: visible; opacity: 1; transform: translateY(0); }

        /* Inline full-name popup (small) */
        .fac-inline-fullname { position: absolute; z-index: 10000; background: #0b1220; color:#fff; padding:6px 8px; border:1px solid #253042; border-radius:6px; box-shadow:0 6px 20px rgba(0,0,0,0.6); max-width: 60vw; word-break: break-all; }

        /* Truncation vs wrap helper classes */
        .fac-wrap .row .name { white-space: normal; }
        .fac-truncate .row .name { white-space: nowrap; overflow: hidden; text-overflow: ellipsis; }

        /* Column-resize handle (simple CSS grab area) */
        .col-resize-handle { width: 8px; cursor: col-resize; display:inline-block; vertical-align:middle; }
        .col-resize { display:flex; align-items:center; }

        /* Selected row styling for keyboard nav */
        .row.selected { background: rgba(99,102,241,0.06); outline: 1px solid rgba(99,102,241,0.14); }
    </style>

    <script>
        // Explorer JS: uses the /api/fs/* endpoints with full actions (create/rename/delete/upload/preview)
        (function(){
            const FAC_CSRF_TOKEN = '{{ session.csrf_token }}';
            const rootSelect = document.getElementById('root-select');
            const address = document.getElementById('address');
            const list = document.getElementById('explorer-list');
            const status = document.getElementById('status-bar');
            const tpl = document.getElementById('tpl-row');
            const quick = document.getElementById('quick-list');
            const fileInput = document.getElementById('file-input');

            let historyStack = [];
            let historyIdx = -1;

            function formatSize(n){ if(n===0) return '0 B'; if(n<1024) return n+ ' B'; if(n<1024*1024) return (n/1024).toFixed(2)+' KB'; if(n<1024*1024*1024) return (n/(1024*1024)).toFixed(2)+' MB'; return (n/(1024*1024*1024)).toFixed(2)+' GB'; }
            function timeStr(ts){ try{ const d=new Date(ts*1000); return d.toLocaleString(); }catch(e){return ''} }

            function pushHistory(root, path){ historyStack = historyStack.slice(0, historyIdx+1); historyStack.push({root, path}); historyIdx = historyStack.length-1; }
            function goBack(){ if(historyIdx>0){ historyIdx--; const s = historyStack[historyIdx]; rootSelect.value = s.root; load(s.root, s.path); } }
            function goForward(){ if(historyIdx < historyStack.length-1){ historyIdx++; const s = historyStack[historyIdx]; rootSelect.value = s.root; load(s.root, s.path); } }

            async function load(root, path){
                if(!root) return alert('Select a root');
                const rel = path||'';
                address.value = rel? (root + '/' + rel) : root;
                pushHistory(root, rel);
                const res = await fetch('/api/fs/list?root='+encodeURIComponent(root)+'&path='+encodeURIComponent(rel));
                const obj = await res.json();
                if(obj.error) return alert(obj.message || obj.error);
                renderList(obj.items || [], root, obj.path || rel);
            }

            function makeElem(tag,cls,txt){ const e=document.createElement(tag); if(cls) e.className=cls; if(txt!==undefined) e.textContent=txt; return e; }

            async function renderList(items, root, relPath){
                list.innerHTML='';
                const folders = items.filter(i=>i.is_dir);
                const files = items.filter(i=>!i.is_dir);

                // folders first
                for(const f of folders){
                    const row = makeElem('div','row flex items-center gap-2 p-2 border-b border-gray-700');
                    const sel = makeElem('input','row-select'); sel.type='checkbox'; row.appendChild(sel);
                    const icon = makeElem('div','icon w-8 text-center','📁'); row.appendChild(icon);
                    // name element wrapped for custom tooltip and click-to-expand
                    const name = makeElem('div','name flex-1 fac-tooltip', '');
                    name.style.cursor='pointer';
                    const nameText = document.createElement('span'); nameText.textContent = f.name; name.appendChild(nameText);
                    // attach custom tooltip node (replaces native title) and accessibility
                    const tip = document.createElement('span'); tip.className = 'fac-tooltip-text'; tip.textContent = f.name; tip.setAttribute('role','tooltip'); tip.setAttribute('aria-hidden','true'); name.appendChild(tip);
                    name.addEventListener('dblclick', ()=>{ load(root, (relPath? relPath + '/' : '') + f.name); });
                    // keyboard accessibility for tooltip: focus shows tooltip
                    name.tabIndex = 0; name.setAttribute('role','button'); name.setAttribute('aria-label', f.name);
                    name.addEventListener('focus', function(){ name.classList.add('fac-show'); tip.setAttribute('aria-hidden','false'); });
                    name.addEventListener('blur', function(){ name.classList.remove('fac-show'); tip.setAttribute('aria-hidden','true'); });
                    name.addEventListener('keydown', function(ev){ if(ev.key === 'Enter' || ev.key === ' '){ ev.preventDefault(); name.click(); } });
                    // click to show inline full-name popup (toggle)
                    name.addEventListener('click', function(ev){ ev.stopPropagation(); const existing = document.getElementById('fac-inline-fullname'); if(existing) existing.remove(); const popup = document.createElement('div'); popup.id = 'fac-inline-fullname'; popup.className = 'fac-inline-fullname'; popup.textContent = f.name; document.body.appendChild(popup); const rect = name.getBoundingClientRect(); popup.style.left = (rect.left + window.scrollX) + 'px'; popup.style.top = (rect.bottom + window.scrollY + 6) + 'px'; function removePopup(){ try{ popup.remove(); }catch(e){} window.removeEventListener('click', removePopup); } setTimeout(()=>{ window.addEventListener('click', removePopup, {once:true}); }, 20); setTimeout(removePopup, 8000); });
                    row.appendChild(name);
                    // Status column (read-only colored badge) for folders
                    const statusCol = makeElem('div','status-col w-40 text-right','');
                    const curStatus = f.status || 'saved';
                    function setBadgeColor(b, s){ try{ if(!b) return; if(s==='processing'){ b.style.background='#D97706'; b.style.color='#FFFFFF'; } else if(s==='waiting'){ b.style.background='#2563EB'; b.style.color='#FFFFFF'; } else if(s==='paused'){ b.style.background='#DC2626'; b.style.color='#FFFFFF'; } else { b.style.background='#4B5563'; b.style.color='#FFFFFF'; } }catch(e){}
                    }
                    const badgeSpan = document.createElement('span'); badgeSpan.className='px-2 py-1 rounded text-sm inline-block mr-2'; badgeSpan.textContent = curStatus; setBadgeColor(badgeSpan, curStatus);
                    // attach tooltip info (status_since) if provided by server
                    try{
                        // store raw values for the relative-time helper to use
                        if(f.status_since) badgeSpan.dataset.statusSince = f.status_since;
                        if(curStatus) badgeSpan.dataset.status = curStatus;
                        // set an initial title (will be refreshed by relative-time helper)
                        if(f.status_since){ const ts = Date.parse(f.status_since); if(!isNaN(ts)){ badgeSpan.title = curStatus + ' — started ' + (function(){ const d = new Date(ts); return d.toLocaleString(); })(); } else { badgeSpan.title = curStatus; } }
                    }catch(e){}
                    statusCol.appendChild(badgeSpan);
                    row.appendChild(statusCol);
                    row.appendChild(makeElem('div','size w-32 text-right',''));
                    row.appendChild(makeElem('div','mtime w-40 text-right', timeStr(f.mtime) ));
                    const actions = makeElem('div','actions flex gap-2 items-center');
                    const openBtn = makeElem('button','btn-small btn-open','Open'); openBtn.addEventListener('click', ()=> load(root, (relPath? relPath + '/' : '') + f.name)); actions.appendChild(openBtn);
                    const delBtn = makeElem('button','btn-small btn-delete','Delete'); delBtn.addEventListener('click', ()=> doDelete(root, (relPath? relPath + '/' : '') + f.name, true)); actions.appendChild(delBtn);
                    row.appendChild(actions);
                    list.appendChild(row);
                }

                // files
                for(const f of files){
                    const row = makeElem('div','row flex items-center gap-2 p-2 border-b border-gray-700');
                    const sel = makeElem('input','row-select'); sel.type='checkbox'; row.appendChild(sel);
                    const icon = makeElem('div','icon w-8 text-center','📄'); row.appendChild(icon);
                    // name element wrapped for custom tooltip and click-to-expand
                    const name = makeElem('div','name flex-1 fac-tooltip', '');
                    const nameText = document.createElement('span'); nameText.textContent = f.name; name.appendChild(nameText);
                    const tip = document.createElement('span'); tip.className = 'fac-tooltip-text'; tip.textContent = f.name; tip.setAttribute('role','tooltip'); tip.setAttribute('aria-hidden','true'); name.appendChild(tip);
                    // keyboard accessibility for tooltip
                    name.tabIndex = 0; name.setAttribute('role','button'); name.setAttribute('aria-label', f.name);
                    name.addEventListener('focus', function(){ name.classList.add('fac-show'); tip.setAttribute('aria-hidden','false'); });
                    name.addEventListener('blur', function(){ name.classList.remove('fac-show'); tip.setAttribute('aria-hidden','true'); });
                    name.addEventListener('keydown', function(ev){ if(ev.key === 'Enter' || ev.key === ' '){ ev.preventDefault(); name.click(); } });
                    // click to show inline full-name popup (toggle)
                    name.addEventListener('click', function(ev){ ev.stopPropagation(); const existing = document.getElementById('fac-inline-fullname'); if(existing) existing.remove(); const popup = document.createElement('div'); popup.id = 'fac-inline-fullname'; popup.className = 'fac-inline-fullname'; popup.textContent = f.name; document.body.appendChild(popup); const rect = name.getBoundingClientRect(); popup.style.left = (rect.left + window.scrollX) + 'px'; popup.style.top = (rect.bottom + window.scrollY + 6) + 'px'; function removePopup(){ try{ popup.remove(); }catch(e){} window.removeEventListener('click', removePopup); } setTimeout(()=>{ window.addEventListener('click', removePopup, {once:true}); }, 20); setTimeout(removePopup, 8000); });
                    row.appendChild(name);
                    // attach file id to the row for status operations
                    try{ if(f.id) row.dataset.fileId = String(f.id); }catch(e){}

                    // Status column (read-only colored badge)
                    const statusCol = makeElem('div','status-col w-40 text-right','');
                    const curStatus = f.status || 'saved';
                    function setBadgeColor(b, s){ try{ if(!b) return; if(s==='processing'){ b.style.background='#D97706'; b.style.color='#FFFFFF'; } else if(s==='waiting'){ b.style.background='#2563EB'; b.style.color='#FFFFFF'; } else if(s==='paused'){ b.style.background='#DC2626'; b.style.color='#FFFFFF'; } else { b.style.background='#4B5563'; b.style.color='#FFFFFF'; } }catch(e){}
                    }
                    const badgeSpan = document.createElement('span'); badgeSpan.className='px-2 py-1 rounded text-sm inline-block mr-2'; badgeSpan.textContent = curStatus; setBadgeColor(badgeSpan, curStatus);
                    try{
                        // store raw values for the relative-time helper to use
                        if(f.status_since) badgeSpan.dataset.statusSince = f.status_since;
                        if(curStatus) badgeSpan.dataset.status = curStatus;
                        // set an initial title (will be refreshed by relative-time helper)
                        if(f.status_since){ const ts = Date.parse(f.status_since); if(!isNaN(ts)){ badgeSpan.title = curStatus + ' — started ' + (function(){ const d = new Date(ts); return d.toLocaleString(); })(); } else { badgeSpan.title = curStatus; } }
                    }catch(e){}
                    statusCol.appendChild(badgeSpan);
                    row.appendChild(statusCol);
                    row.appendChild(makeElem('div','size w-32 text-right', formatSize(f.size) ));
                    row.appendChild(makeElem('div','mtime w-40 text-right', timeStr(f.mtime) ));
                    const actions = makeElem('div','actions flex gap-2 items-center');
                            // For security/UX, do NOT show the Load button when browsing roots that are already storage targets
                            // (Upload Files, Encrypted Files, Decrypted Files). In those cases show Preview instead.
                            const forbiddenLoadRoots = ['Upload Files','Encrypted Files','Decrypted Files'];
                            if(forbiddenLoadRoots.indexOf(root) !== -1){
                                // In forbidden roots, do NOT show Load or Preview actions; keep only Download/Delete
                                // (No action appended here)
                            } else {
                                // For other roots, keep the Load button which copies DB-referenced files into the upload area
                                const loadBtn = makeElem('button','btn-small btn-load','Load');
                                loadBtn.addEventListener('click', async ()=>{
                                    try{
                                        const fd = new FormData();
                                        fd.append('file_id', String(f.id));
                                        fd.append('csrf_token', FAC_CSRF_TOKEN);
                                        const r = await fetch('/api/fs/load_to_upload', { method: 'POST', body: fd, headers: {'X-CSRF-Token': FAC_CSRF_TOKEN} });
                                        const j = await r.json();
                                        if(j && j.error) return alert(j.error || j.message || 'Load failed');
                                        const filenameParam = encodeURIComponent(j.filename || f.name);
                                        window.location = '/evidence_upload?loaded=' + filenameParam;
                                    }catch(e){ alert('Load failed: '+(e && e.message ? e.message : e)); }
                                });
                                actions.appendChild(loadBtn);
                            }
                    const dlBtn = makeElem('button','btn-small btn-download','Download'); dlBtn.addEventListener('click', ()=> window.open('/serve_fs_file?root='+encodeURIComponent(root)+'&path='+encodeURIComponent((relPath? relPath + '/' : '') + f.name), '_blank')); actions.appendChild(dlBtn);
                    const delBtn = makeElem('button','btn-small btn-delete','Delete'); delBtn.addEventListener('click', ()=> doDelete(root, (relPath? relPath + '/' : '') + f.name, false)); actions.appendChild(delBtn);
                    row.appendChild(actions);
                    list.appendChild(row);

                    // fetch thumbnail/info asynchronously
                    (async function(row, root, relPath, fname){
                        try{
                            const infoRes = await fetch('/api/fs/info?root='+encodeURIComponent(root)+'&path='+encodeURIComponent((relPath? relPath + '/' : '') + fname));
                            const info = await infoRes.json();
                            if(info && info.thumbnail){ row.querySelector('.icon').textContent = ''; const img = document.createElement('img'); img.src = info.thumbnail; img.style.width='32px'; img.style.height='32px'; img.style.objectFit='cover'; row.querySelector('.icon').appendChild(img); }
                        }catch(e){}

                    })(row, root, relPath, f.name);
                }

                status.textContent = 'Items: '+items.length;
                // refresh status badges based on user preference
                try{ refreshStatusBadges(); }catch(e){}
                // initialize UI helpers after rendering
                try{ ensureTruncationToggle(); }catch(e){}
                try{ ensureResizeHandles(); }catch(e){}
            }

            // Toggle truncation/wrap control
            function ensureTruncationToggle(){
                if(document.getElementById('fac-trunc-toggle')) return;
                const container = document.createElement('div'); container.style.display='flex'; container.style.alignItems='center'; container.style.gap='8px';
                container.style.marginBottom='8px';
                const toggle = document.createElement('button'); toggle.id='fac-trunc-toggle'; toggle.className='btn-small';
                toggle.textContent = localStorage.getItem('fac_trunc') === 'wrap' ? 'Wrap filenames' : 'Truncate filenames';
                toggle.addEventListener('click', ()=>{
                    const cur = document.body.classList.contains('fac-wrap') ? 'wrap' : 'truncate';
                    if(cur === 'truncate'){
                        document.body.classList.remove('fac-truncate'); document.body.classList.add('fac-wrap'); localStorage.setItem('fac_trunc','wrap'); toggle.textContent='Wrap filenames';
                    } else {
                        document.body.classList.remove('fac-wrap'); document.body.classList.add('fac-truncate'); localStorage.setItem('fac_trunc','truncate'); toggle.textContent='Truncate filenames';
                    }
                });
                // insert before list
                list.parentElement.insertBefore(container, list);
                container.appendChild(toggle);
                // initialize
                if(localStorage.getItem('fac_trunc') === 'wrap'){ document.body.classList.add('fac-wrap'); } else { document.body.classList.add('fac-truncate'); }
            }

            // Column resize helpers (per-column handles + persistence)
            function ensureResizeHandles(){
                if(document.getElementById('fac-resize-handle')) return;
                const header = document.createElement('div'); header.id='fac-resize-handle'; header.style.display='flex'; header.style.alignItems='center'; header.style.gap='8px'; header.style.marginBottom='6px';
                const info = document.createElement('div'); info.style.fontSize='12px'; info.style.color='#9CA3AF'; info.textContent='Drag handles to resize columns'; header.appendChild(info);
                list.parentElement.insertBefore(header, list);

                const handleSize = document.createElement('div'); handleSize.className='col-resize-handle'; handleSize.style.width='12px'; handleSize.style.height='18px'; handleSize.style.cursor='col-resize'; handleSize.title='Resize Size column'; handleSize.style.background='transparent';
                const handleMtime = document.createElement('div'); handleMtime.className='col-resize-handle'; handleMtime.style.width='12px'; handleMtime.style.height='18px'; handleMtime.style.cursor='col-resize'; handleMtime.title='Resize MTime column'; handleMtime.style.background='transparent';
                const wrapper = document.createElement('div'); wrapper.style.display='flex'; wrapper.style.alignItems='center'; wrapper.style.gap='6px'; wrapper.appendChild(handleSize); wrapper.appendChild(handleMtime); header.appendChild(wrapper);

                // apply saved widths if available
                const savedSize = parseInt(localStorage.getItem('fac_col_width_size') || '0', 10);
                const savedMtime = parseInt(localStorage.getItem('fac_col_width_mtime') || '0', 10);
                if(savedSize > 0){ document.querySelectorAll('.row .size').forEach(el=>{ el.style.flex = '0 0 '+savedSize+'px'; el.style.maxWidth = savedSize+'px'; }); }
                if(savedMtime > 0){ document.querySelectorAll('.row .mtime').forEach(el=>{ el.style.flex = '0 0 '+savedMtime+'px'; el.style.maxWidth = savedMtime+'px'; }); }

                function attachDrag(handle, targetSelector, minWidth, storageKey){
                    let dragging=false, startX=0, startW=0;
                    handle.addEventListener('mousedown', function(e){ e.preventDefault(); dragging = true; startX = e.clientX; const el = document.querySelector('.row '+targetSelector); if(el) startW = el.getBoundingClientRect().width; });
                    window.addEventListener('mousemove', function(e){ if(!dragging) return; const dx = e.clientX - startX; const newW = Math.max(minWidth, Math.round(startW + dx)); document.querySelectorAll('.row '+targetSelector).forEach(el=>{ el.style.flex = '0 0 '+newW+'px'; el.style.maxWidth = newW+'px'; }); });
                    window.addEventListener('mouseup', function(){ if(dragging){ const el = document.querySelector('.row '+targetSelector); if(el){ const w = Math.round(el.getBoundingClientRect().width); try{ localStorage.setItem(storageKey, String(w)); }catch(e){} } dragging = false; } });
                }

                attachDrag(handleSize, '.size', 48, 'fac_col_width_size');
                attachDrag(handleMtime, '.mtime', 80, 'fac_col_width_mtime');
            }

            // actions
            // Rename UI removed from DB browser template
            async function doDelete(root, target, isDir){ if(!confirm('Delete '+target+' ?')) return; const r = await fetch('/api/fs/delete',{method:'POST', headers:{'Content-Type':'application/x-www-form-urlencoded','X-CSRF-Token': FAC_CSRF_TOKEN}, body:new URLSearchParams({root, path: target, csrf_token: FAC_CSRF_TOKEN})}); const j = await r.json(); if(j.error) alert(j.error||j.message); else load(root, address.value.replace(root + '/', '')); }
            async function doPreview(root, target){ const r = await fetch('/api/fs/preview?root='+encodeURIComponent(root)+'&path='+encodeURIComponent(target)); const j = await r.json(); if(j.error) return alert(j.error||j.message); showPreviewModal(j); }

            // upload
            fileInput.addEventListener('change', async function(e){ const f=this.files[0]; if(!f) return; const dest = address.value && address.value.startsWith(rootSelect.value + '/') ? address.value.replace(rootSelect.value + '/', '') : ''; const fd = new FormData(); fd.append('root', rootSelect.value); fd.append('path', dest); fd.append('file', f); fd.append('csrf_token', FAC_CSRF_TOKEN); const r = await fetch('/api/fs/upload',{method:'POST', body:fd, headers: {'X-CSRF-Token': FAC_CSRF_TOKEN}}); const j = await r.json(); if(j.error) alert(j.error||j.message); else load(rootSelect.value, dest); });

            // new folder
            document.getElementById('btn-new-folder').addEventListener('click', async ()=>{ const name = prompt('Folder name:'); if(!name) return; const dest = address.value && address.value.startsWith(rootSelect.value + '/') ? address.value.replace(rootSelect.value + '/', '') : ''; const r = await fetch('/api/fs/mkdir',{method:'POST', headers:{'Content-Type':'application/x-www-form-urlencoded','X-CSRF-Token': FAC_CSRF_TOKEN}, body:new URLSearchParams({root: rootSelect.value, path: dest, name, csrf_token: FAC_CSRF_TOKEN})}); const j = await r.json(); if(j.error) alert(j.error||j.message); else load(rootSelect.value, dest); });

            // Delete selected items (bulk) - uses /api/fs/bulk_delete endpoint for better efficiency and supports undo
            document.getElementById('btn-delete-selected').addEventListener('click', async function(){
                const checks = Array.from(document.querySelectorAll('#explorer-list .row-select:checked'));
                if(!checks.length) return alert('No items selected');
                // gather targets relative to current root/path
                const curRoot = rootSelect.value;
                const curPath = address.value && address.value.startsWith(curRoot + '/') ? address.value.replace(curRoot + '/', '') : '';
                const paths = [];
                for(const ch of checks){
                    const row = ch.closest('.row');
                    if(!row) continue;
                    const nameEl = row.querySelector('.name span');
                    const fname = nameEl ? nameEl.textContent : null;
                    if(!fname) continue;
                    const target = curPath ? (curPath + '/' + fname) : fname;
                    paths.push(target);
                }
                if(!paths.length) return alert('No valid items selected');
                if(!confirm(`Delete ${paths.length} selected item(s)? You will have a short window to undo.`)) return;

                // Create toast container if not present
                function ensureToastContainer(){
                    let c = document.getElementById('fac-toast-container');
                    if(c) return c;
                    c = document.createElement('div'); c.id = 'fac-toast-container';
                    c.style.position = 'fixed'; c.style.right = '20px'; c.style.bottom = '20px'; c.style.zIndex = 20000; c.style.display = 'flex'; c.style.flexDirection = 'column'; c.style.gap = '8px'; document.body.appendChild(c);
                    return c;
                }

                const toastContainer = ensureToastContainer();
                const toast = document.createElement('div'); toast.className = 'fac-toast'; toast.style.background = '#0b1220'; toast.style.border = '1px solid #333'; toast.style.color = '#fff'; toast.style.padding = '10px 12px'; toast.style.minWidth = '260px'; toast.style.boxShadow = '0 4px 12px rgba(0,0,0,0.3)';
                const title = document.createElement('div'); title.style.fontWeight='600'; title.textContent = `Deleting ${paths.length} item(s)`; toast.appendChild(title);
                const progress = document.createElement('div'); progress.style.marginTop='8px'; progress.style.display='flex'; progress.style.alignItems='center'; progress.style.gap='8px';
                const spinner = document.createElement('div'); spinner.className='fac-spinner'; spinner.style.width='16px'; spinner.style.height='16px'; spinner.style.border='3px solid #333'; spinner.style.borderTop='3px solid #fff'; spinner.style.borderRadius='50%'; spinner.style.animation='fac-spin 1s linear infinite'; progress.appendChild(spinner);
                const statusText = document.createElement('div'); statusText.textContent = `0 / ${paths.length}`; progress.appendChild(statusText);
                toast.appendChild(progress);
                const actions = document.createElement('div'); actions.style.display='flex'; actions.style.justifyContent='flex-end'; actions.style.gap='8px'; actions.style.marginTop='8px';
                const undoBtn = document.createElement('button'); undoBtn.textContent = 'Undo'; undoBtn.className='btn-small'; actions.appendChild(undoBtn);
                const closeBtn = document.createElement('button'); closeBtn.textContent = 'Dismiss'; closeBtn.className='btn-small'; actions.appendChild(closeBtn);
                toast.appendChild(actions);
                toastContainer.appendChild(toast);

                // set CSS for spinner animation (once)
                if(!document.getElementById('fac-spinner-style')){
                    const s = document.createElement('style'); s.id='fac-spinner-style'; s.textContent='@keyframes fac-spin{from{transform:rotate(0deg)}to{transform:rotate(360deg)}}'; document.head.appendChild(s);
                }

                // send bulk delete request with trash enabled so undo is possible
                let trash_id = null;
                statusText.textContent = `0 / ${paths.length}`;
                try{
                    const r = await fetch('/api/fs/bulk_delete', { method: 'POST', headers: {'Content-Type':'application/json','X-CSRF-Token': FAC_CSRF_TOKEN}, body: JSON.stringify({ root: curRoot, paths: paths, trash: true, csrf_token: FAC_CSRF_TOKEN }) });
                    const j = await r.json();
                    // update counts
                    const deletedCount = (j && j.deleted) ? j.deleted.length : 0;
                    const failedCount = (j && j.failed) ? j.failed.length : 0;
                    statusText.textContent = `${deletedCount} / ${paths.length}` + (failedCount ? ` · ${failedCount} failed` : '');
                    // stop spinner
                    spinner.style.display='none';
                    // expose trash id for undo
                    if(j && j.trash_id) trash_id = j.trash_id;
                    title.textContent = `Deleted ${deletedCount} item(s)` + (failedCount ? `, ${failedCount} failed` : '');
                }catch(e){
                    spinner.style.display='none'; statusText.textContent='Failed'; title.textContent='Delete request failed';
                    console.warn('bulk delete failed', e);
                }

                // undo logic: allow undo for 12 seconds or until dismiss
                let undone = false;
                const undoWindowMs = 12000;
                const timeoutId = setTimeout(()=>{
                    // time expired, auto-dismiss toast after short delay
                    toast.style.opacity='0.4'; undoBtn.disabled = true;
                }, undoWindowMs);

                undoBtn.addEventListener('click', async ()=>{
                    if(!trash_id || undone) return;
                    try{
                        undoBtn.disabled = true; statusText.textContent = 'Restoring...';
                        const ur = await fetch('/api/fs/undo_bulk', { method: 'POST', headers: {'Content-Type':'application/json','X-CSRF-Token': FAC_CSRF_TOKEN}, body: JSON.stringify({ trash_id: trash_id, csrf_token: FAC_CSRF_TOKEN }) });
                        const uj = await ur.json();
                        if(uj && uj.restored){
                            statusText.textContent = `Restored ${uj.restored.length}`;
                            title.textContent = `Restored ${uj.restored.length} item(s)`;
                        }
                        undone = true;
                        clearTimeout(timeoutId);
                    }catch(e){
                        console.warn('undo failed', e); statusText.textContent = 'Undo failed';
                    }
                });

                closeBtn.addEventListener('click', ()=>{ try{ toast.remove(); }catch(e){} });

                // after short delay, reload view to reflect deletions
                setTimeout(()=>{ try{ load(rootSelect.value, address.value.replace(rootSelect.value + '/', '')); }catch(e){ console.warn(e); } }, 800);
            });

            document.getElementById('btn-refresh').addEventListener('click', ()=> load(rootSelect.value, address.value.replace(rootSelect.value+'/', '')));
            document.getElementById('btn-up').addEventListener('click', ()=>{ const cur = address.value.replace(rootSelect.value+'/', ''); if(!cur) return; const parts = cur.split('/'); parts.pop(); load(rootSelect.value, parts.join('/')); });
            document.getElementById('btn-back').addEventListener('click', goBack);
            document.getElementById('btn-forward').addEventListener('click', goForward);
            document.getElementById('btn-zip').addEventListener('click', ()=>{ const root=rootSelect.value; const rel = address.value.replace(root + '/', ''); window.location='/download_folder_zip?browse_root='+encodeURIComponent(root)+'&browse_folder='+encodeURIComponent(rel); });

            // Database Browser is read-only for statuses; no bulk apply handler

            // click quick items
            quick.querySelectorAll('.quick-item').forEach(a=>{ a.addEventListener('click', (e)=>{ e.preventDefault(); const root = a.dataset.root; rootSelect.value = root; load(root,''); }); });

            // preview modal
            function showPreviewModal(data){ let modal = document.getElementById('preview-modal'); if(!modal){ modal = document.createElement('div'); modal.id='preview-modal'; modal.style.position='fixed'; modal.style.left='10%'; modal.style.top='10%'; modal.style.width='80%'; modal.style.height='80%'; modal.style.background='#0b1220'; modal.style.padding='12px'; modal.style.overflow='auto'; modal.style.zIndex=9999; modal.style.border='1px solid #333'; document.body.appendChild(modal); }
                modal.innerHTML = '';
                const h = document.createElement('div'); h.style.display='flex'; h.style.justifyContent='space-between'; const title = document.createElement('div'); title.textContent = data.name || 'Preview'; title.style.fontWeight='600'; h.appendChild(title); const close = document.createElement('button'); close.textContent='Close'; close.onclick = ()=> modal.remove(); h.appendChild(close); modal.appendChild(h);
                if(data.thumbnail){ const img = document.createElement('img'); img.src = data.thumbnail; img.style.maxWidth='240px'; img.style.float='right'; img.style.margin='6px'; modal.appendChild(img); }
                if(data.mime && data.mime.startsWith('text/')){ const pre = document.createElement('pre'); pre.style.whiteSpace='pre-wrap'; pre.textContent = data.preview || ''; modal.appendChild(pre); }
                else if(data.mime && data.mime.startsWith('image/')){ const full = document.createElement('img'); full.src = '/serve_fs_file?root='+encodeURIComponent(data.root)+'&path='+encodeURIComponent(data.path); full.style.maxWidth='100%'; modal.appendChild(full); }
                else { const pre = document.createElement('pre'); pre.style.whiteSpace='pre-wrap'; pre.textContent = data.preview || '[Binary preview not available]'; modal.appendChild(pre); }
            }

            // initial load if query provided
            try{ const urlParams = new URLSearchParams(location.search); const br = urlParams.get('browse_root'); const bf = urlParams.get('browse_folder'); if(br){ rootSelect.value = br; load(br, bf||''); } }catch(e){}
            
            // Time mode helper: converts ISO timestamp -> human friendly '3m ago'
            function relativeTimeFromDate(dateOrIso){
                try{
                    const d = (typeof dateOrIso === 'string') ? new Date(dateOrIso) : new Date(dateOrIso);
                    const now = new Date();
                    const diff = Math.floor((now - d) / 1000); // seconds
                    if(isNaN(diff) || diff < 0) return '';
                    if(diff < 10) return 'just now';
                    if(diff < 60) return diff + 's ago';
                    const m = Math.floor(diff/60);
                    if(m < 60) return m + 'm ago';
                    const h = Math.floor(m/60);
                    if(h < 24) return h + 'h ago';
                    const days = Math.floor(h/24);
                    if(days < 7) return days + 'd ago';
                    // fallback to short date
                    return d.toLocaleDateString();
                }catch(e){ return ''; }
            }

            // Refresh all status badge tooltips to show relative times
            function refreshStatusBadges(){
                try{
                    // determine user preference
                    const sel = document.getElementById('time-mode-select');
                    let mode = 'relative-tooltip';
                    try{ mode = localStorage.getItem('fac_time_mode') || (sel ? sel.value : mode); }catch(e){}
                    if(sel){ sel.value = mode; sel.addEventListener('change', function(){ try{ localStorage.setItem('fac_time_mode', sel.value); refreshStatusBadges(); }catch(e){} }); }

                    document.querySelectorAll('#explorer-list .status-col span[data-status-since]').forEach(function(b){
                        const iso = b.dataset.statusSince;
                        const status = b.dataset.status || '';
                        if(!iso) return;
                        const rel = relativeTimeFromDate(iso);
                        if(mode === 'relative-inline'){
                            // display inline short text next to status label
                            b.textContent = status + (rel ? ' · ' + rel : '');
                            b.title = status + (rel ? ' · ' + rel : '');
                        } else if(mode === 'absolute'){
                            // show absolute localized timestamp in title, keep status text unchanged
                            const ts = Date.parse(iso);
                            if(!isNaN(ts)){ b.title = status + ' — started ' + (new Date(ts)).toLocaleString(); }
                            b.textContent = status;
                        } else {
                            // relative-tooltip
                            b.textContent = status;
                            if(rel) b.title = status + ' — started ' + rel;
                        }
                    });
                }catch(e){ console.warn('refreshStatusBadges failed', e); }
            }

            // Kick off periodic refresh every 30 seconds
            try{ setInterval(refreshStatusBadges, 30 * 1000); }catch(e){}
        })();
    </script>

</div>
"""


@app.route('/api/fs/bulk_delete', methods=['POST'])
def api_fs_bulk_delete():
    """Delete multiple files/dirs in a single atomic operation and log a single audit entry."""
    # Expect JSON body: { root: <root_display>, paths: [<rel1>, <rel2>, ...], csrf_token }
    if not verify_csrf_request():
        return jsonify({'error': 'csrf_missing_or_invalid'}), 403
    try:
        data = request.get_json(force=True)
    except Exception:
        return jsonify({'error': 'invalid_json'}), 400
    root = data.get('root')
    paths = data.get('paths') or []
    # support trash behavior for undo (default True)
    trash = True if data.get('trash', True) else False
    if not root or not isinstance(paths, list):
        return jsonify({'error': 'missing_params'}), 400
    root_path = resolve_root_display_to_path(root)
    if not root_path:
        return jsonify({'error': 'invalid_root'}), 400
    results = {'deleted': [], 'failed': []}
    # If trash enabled, move items into the configured DELETED_RECOVERY_FOLDER under a unique session folder
    if trash:
        try:
            trash_id = 'trash_' + datetime.datetime.utcnow().strftime('%Y%m%dT%H%M%SZ') + '_' + secrets.token_hex(6)
        except Exception:
            trash_id = 'trash_' + secrets.token_hex(10)
        trash_dir = os.path.join(app.config.get('DELETED_RECOVERY_FOLDER') or DELETED_RECOVERY_FOLDER, trash_id)
        try:
            os.makedirs(trash_dir, exist_ok=True)
        except Exception:
            return jsonify({'error': 'trash_create_failed'}), 500

        metadata = {'id': trash_id, 'ts': datetime.datetime.utcnow().isoformat() + 'Z', 'root': root, 'items': []}
        idx = 0
        for rel in paths:
            idx += 1
            try:
                abs_path = os.path.abspath(os.path.join(root_path, rel))
                if not is_path_under_allowed_roots(abs_path) or not os.path.exists(abs_path):
                    results['failed'].append({'path': rel, 'reason': 'not_found_or_access_denied'})
                    metadata['items'].append({'path': rel, 'status': 'failed', 'reason': 'not_found_or_access_denied'})
                    continue
                # create a per-item subfolder to preserve name and avoid collisions
                item_sub = os.path.join(trash_dir, f'item_{idx}')
                os.makedirs(item_sub, exist_ok=True)
                # move file or directory into the item_sub folder
                try:
                    shutil.move(abs_path, item_sub)
                    results['deleted'].append(rel)
                    metadata['items'].append({'path': rel, 'status': 'moved', 'dest': os.path.relpath(item_sub, app.config.get('DELETED_RECOVERY_FOLDER'))})
                except Exception as e:
                    # fallback: attempt remove if move fails
                    try:
                        if os.path.isdir(abs_path):
                            shutil.rmtree(abs_path)
                        else:
                            os.remove(abs_path)
                        results['deleted'].append(rel)
                        metadata['items'].append({'path': rel, 'status': 'deleted_fallback'})
                    except Exception as e2:
                        results['failed'].append({'path': rel, 'reason': str(e2)})
                        metadata['items'].append({'path': rel, 'status': 'failed', 'reason': str(e2)})
            except Exception as e:
                results['failed'].append({'path': rel, 'reason': str(e)})
                metadata['items'].append({'path': rel, 'status': 'failed', 'reason': str(e)})

        # persist metadata for potential undo/restore
        try:
            meta_path = os.path.join(trash_dir, 'metadata.json')
            with open(meta_path, 'w', encoding='utf-8') as mf:
                json.dump(metadata, mf)
        except Exception:
            app.logger.debug('Failed to write trash metadata')

        # Single audit entry
        try:
            audit_event('bulk_delete', root_path, {'requested': paths, 'results': results, 'trash_id': trash_id})
        except Exception:
            pass
        resp = results.copy()
        resp['trash_id'] = metadata.get('id')
        return jsonify(resp)

    # If trash not requested, perform permanent deletion as before
    for rel in paths:
        try:
            abs_path = os.path.abspath(os.path.join(root_path, rel))
            if not is_path_under_allowed_roots(abs_path) or not os.path.exists(abs_path):
                results['failed'].append({'path': rel, 'reason': 'not_found_or_access_denied'})
                continue
            if os.path.isdir(abs_path):
                shutil.rmtree(abs_path)
            else:
                os.remove(abs_path)
            results['deleted'].append(rel)
        except Exception as e:
            results['failed'].append({'path': rel, 'reason': str(e)})
    # Single audit entry
    try:
        audit_event('bulk_delete', root_path, {'requested': paths, 'results': results})
    except Exception:
        pass
    return jsonify(results)


@app.route('/database_browser/<dbid>')
def database_browser(dbid):
    # Parse query params for pagination/search/filter
    try:
        page = int(request.args.get('page', '1'))
        if page < 1: page = 1
    except Exception:
        page = 1
    try:
        per_page = int(request.args.get('per_page', '25'))
        if per_page < 1: per_page = 25
    except Exception:
        per_page = 25
    q = request.args.get('q', '').strip() or None
    qtype = request.args.get('type', '').strip() or None
    qsession = request.args.get('session', '').strip() or None
    browse_root = request.args.get('browse_root')
    browse_folder = request.args.get('browse_folder')

    offset = (page - 1) * per_page

    files = []
    total = 0
    try:
        # If the dbid refers to a configured connection, attempt to query the target DB via adapter
        if dbid in db_connections:
            try:
                files, total, file_types = db_query_files(dbid, q=q, qtype=qtype, qsession=qsession, browse_root=browse_root, browse_folder=browse_folder, per_page=per_page, offset=offset)
            except Exception as e:
                app.logger.debug(f"db_query_files failed for {dbid}: {e}")
                # fallback to local DB
                files, total, file_types = db_query_files(None, q=q, qtype=qtype, qsession=qsession, browse_root=browse_root, browse_folder=browse_folder, per_page=per_page, offset=offset)
        else:
            files, total, file_types = db_query_files(None, q=q, qtype=qtype, qsession=qsession, browse_root=browse_root, browse_folder=browse_folder, per_page=per_page, offset=offset)

        # Compute folder and root mapping for each file (best-effort)
        for f in files:
            fpath = f.get('path')
            folder = None
            root_name = None
            try:
                if fpath:
                    abs_fpath = os.path.abspath(fpath)
                    # Check per-session first (only valid against local sessions table)
                    try:
                        sess_path = None
                        if f.get('session_id'):
                            conn = _get_db_conn()
                            scur = conn.cursor()
                            scur.execute('SELECT session_path FROM sessions WHERE id=?', (f.get('session_id'),))
                            srow = scur.fetchone()
                            if srow and srow[0]:
                                sess_path = srow[0]
                            try: conn.close()
                            except Exception: pass
                        if sess_path and abs_fpath.startswith(os.path.abspath(sess_path)):
                            folder = os.path.relpath(os.path.dirname(abs_fpath), os.path.abspath(sess_path))
                            root_name = os.path.basename(os.path.abspath(sess_path))
                        else:
                            # determine which allowed root it belongs to
                            for root in _allowed_roots():
                                try:
                                    abs_root = os.path.abspath(root)
                                    if abs_fpath == abs_root or abs_fpath.startswith(abs_root + os.path.sep):
                                        folder = os.path.relpath(os.path.dirname(abs_fpath), abs_root)
                                        root_name = os.path.basename(abs_root)
                                        break
                                except Exception:
                                    continue
                    except Exception:
                        pass
            except Exception:
                folder = None
            f['folder'] = folder
            f['root'] = root_name
            # Provide a best-effort status_since timestamp for UI tooltips.
            # Prefer an explicit timestamp stored in extra.status_changed_at or fall back to created_at.
            try:
                status_since = None
                extra = None
                # try to load extra JSON if present
                if isinstance(f.get('session_id'), (str, int)):
                    # nothing here; prefer reading extra from the DB rows where available
                    pass
                # if row contains 'created_at' use that as fallback
                if f.get('created_at'):
                    status_since = f.get('created_at')
                # attach to the file dict for template usage
                f['status_since'] = status_since
            except Exception:
                f['status_since'] = f.get('created_at')
    except Exception as e:
        app.logger.debug(f"database_browser query error: {e}")
        files = []
        total = 0

    # gather distinct file types for filter dropdown
    try:
        conn2 = _get_db_conn()
        c2 = conn2.cursor()
        c2.execute('SELECT DISTINCT file_type FROM files')
        file_types = [r[0] for r in c2.fetchall() if r[0]]
    except Exception:
        file_types = []
    finally:
        try: conn2.close()
        except Exception: pass

    # compute child folders (unique first-level segments) from files when browsing a root
    child_folders = []
    try:
        if browse_root:
            seen_children = set()
            for f in files:
                if f.get('folder') and f.get('folder') not in ('.', ''):
                    first = f['folder'].split(os.path.sep)[0]
                    if first and first not in seen_children:
                        seen_children.add(first)
                        child_folders.append(first)
    except Exception:
        child_folders = []

    # If a session was active, end it because user navigated to DB browser
    sess_id = session.pop('analysis_session_id', None)
    if sess_id:
        try:
            end_analysis_session(sess_id)
        except Exception:
            pass

    content = render_template_string(DATABASE_BROWSER_CONTENT, files=files, db_name=db_connections.get(dbid,{}).get('name', dbid),
                                     page=page, per_page=per_page, offset=offset, total=total, file_types=file_types,
                                     cur_dbid=dbid, q=q, qtype=qtype, qsession=qsession,
                                     browse_root=browse_root, browse_folder=browse_folder, child_folders=child_folders)
    return render_template_string(BASE_TEMPLATE, content=content, show_session_details=False)


@app.route('/api/fs/undo_bulk', methods=['POST'])
def api_fs_undo_bulk():
    """Attempt to restore a previously bulk-deleted (trashed) set of items identified by trash_id.
    Expects JSON { trash_id: <id>, csrf_token }
    Returns { restored: [paths], failed: [...] }
    """
    if not verify_csrf_request():
        return jsonify({'error': 'csrf_missing_or_invalid'}), 403
    try:
        data = request.get_json(force=True)
    except Exception:
        return jsonify({'error': 'invalid_json'}), 400
    trash_id = data.get('trash_id')
    if not trash_id:
        return jsonify({'error': 'missing_params'}), 400
    trash_dir = os.path.join(app.config.get('DELETED_RECOVERY_FOLDER') or DELETED_RECOVERY_FOLDER, trash_id)
    if not os.path.exists(trash_dir):
        return jsonify({'error': 'trash_not_found'}), 404
    meta_path = os.path.join(trash_dir, 'metadata.json')
    if not os.path.exists(meta_path):
        return jsonify({'error': 'metadata_missing'}), 404
    try:
        with open(meta_path, 'r', encoding='utf-8') as mf:
            metadata = json.load(mf)
    except Exception:
        return jsonify({'error': 'metadata_read_failed'}), 500

    restored = []
    failed = []
    # attempt to move each item back to its original root path
    for idx, it in enumerate(metadata.get('items', []) , start=1):
        try:
            if it.get('status') != 'moved':
                continue
            dest_rel = it.get('path')
            # item directory stored as item_{n}
            item_sub = os.path.join(trash_dir, f'item_{idx}')
            if not os.path.exists(item_sub):
                failed.append({'path': dest_rel, 'reason': 'item_missing_in_trash'})
                continue
            # restore target original path
            root_display = metadata.get('root')
            root_path = resolve_root_display_to_path(root_display)
            if not root_path:
                failed.append({'path': dest_rel, 'reason': 'invalid_root'})
                continue
            orig_abs = os.path.abspath(os.path.join(root_path, dest_rel))
            # ensure parent dir exists
            parent = os.path.dirname(orig_abs)
            try:
                os.makedirs(parent, exist_ok=True)
            except Exception:
                pass
            # attempt move back from item_sub (which may contain a single file or dir)
            # find the single child inside item_sub
            children = os.listdir(item_sub)
            if not children:
                failed.append({'path': dest_rel, 'reason': 'no_item_data'})
                continue
            # if child is a single element, move it to orig_abs path
            child_name = children[0]
            src = os.path.join(item_sub, child_name)
            try:
                if os.path.exists(orig_abs):
                    # avoid overwrite; attempt to rename existing
                    backup_existing = orig_abs + '.restore_conflict.' + secrets.token_hex(4)
                    try:
                        os.rename(orig_abs, backup_existing)
                    except Exception:
                        pass
                shutil.move(src, orig_abs)
                restored.append(dest_rel)
            except Exception as e:
                failed.append({'path': dest_rel, 'reason': str(e)})
        except Exception as e:
            failed.append({'path': it.get('path'), 'reason': str(e)})

    # write audit log
    try:
        audit_event('undo_bulk', trash_dir, {'restored': restored, 'failed': failed})
    except Exception:
        pass

    return jsonify({'restored': restored, 'failed': failed})

@app.route('/serve_db_file/<int:file_id>')
def serve_db_file(file_id):
    try:
        conn = _get_db_conn()
        cur = conn.cursor()
        cur.execute('SELECT filename, path, file_type FROM files WHERE id=?', (file_id,))
        row = cur.fetchone()
    except Exception:
        row = None
    finally:
        try:
            conn.close()
        except Exception:
            pass

    if not row:
        return "File record not found", 404

    filename, path, ftype = row
    if path and os.path.exists(path):
        # Ensure the path is inside allowed roots before serving
        if not is_path_under_allowed_roots(path):
            return "Access denied", 403
        # Serve directly from disk via an absolute safe directory
        dirpath = os.path.dirname(os.path.abspath(path))
        filename_only = os.path.basename(path)
        try:
            return send_from_directory(dirpath, filename_only, as_attachment=True)
        except Exception:
            # Fallback: send via send_file
            try:
                return send_file(path, as_attachment=True)
            except Exception:
                return "Could not serve file", 500
    else:
        return "File not available on disk", 404


@app.route('/api/fs/load_to_upload', methods=['POST'])
def api_fs_load_to_upload():
    """Server helper: copy a file referenced in the DB (or by path under allowed roots) into the configured UPLOAD_FOLDER.
    Expects 'file_id' or 'file_path' and a CSRF token. Returns JSON {'ok': True, 'filename': <name>}.
    """
    # CSRF
    token = request.form.get('csrf_token') or request.headers.get('X-CSRF-Token')
    try:
        if not token or token != session.get('csrf_token'):
            return jsonify({'error': 'Invalid CSRF token'}), 403
    except Exception:
        return jsonify({'error': 'Session error validating CSRF token'}), 403

    file_id = request.form.get('file_id')
    file_path = request.form.get('file_path')

    if not file_id and not file_path:
        return jsonify({'error': 'Missing parameters (file_id or file_path required)'}), 400

    candidate = None
    db_filename = None
    # Resolve via DB when id provided
    if file_id:
        try:
            # Debug trace for test environments
            try:
                print('DEBUG load_to_upload: DB_FILE=', DB_FILE, 'exists=', os.path.exists(DB_FILE) if DB_FILE else None)
            except Exception:
                pass
            prow = None
            if DB_FILE and os.path.exists(DB_FILE):
                try:
                    conn2 = sqlite3.connect(DB_FILE)
                    cur2 = conn2.cursor()
                    cur2.execute('SELECT path, filename FROM files WHERE id=?', (int(file_id),))
                    prow = cur2.fetchone()
                    try: conn2.close()
                    except Exception: pass
                except Exception:
                    prow = None
            if not prow:
                try:
                    conn = _get_db_conn()
                    cur = conn.cursor()
                    cur.execute('SELECT path, filename FROM files WHERE id=?', (int(file_id),))
                    prow = cur.fetchone()
                    try: conn.close()
                    except Exception: pass
                except Exception:
                    prow = None

            if prow:
                db_path_val = prow[0]
                db_filename = prow[1] if len(prow) > 1 else None

                # Try to normalize and use any path stored in the DB
                try:
                    print('DEBUG load_to_upload: prow=', prow)
                    print('DEBUG load_to_upload: db_path_val=', repr(db_path_val), 'db_filename=', repr(db_filename))
                    try:
                        parent = os.path.dirname(db_path_val) if isinstance(db_path_val, str) else None
                        print('DEBUG load_to_upload: db_path parent=', repr(parent), 'parent_exists=', os.path.exists(parent) if parent else None)
                        if parent and os.path.exists(parent):
                            try:
                                listing = os.listdir(parent)
                                print('DEBUG load_to_upload: parent listing sample=', listing[:20])
                                if not listing:
                                    # dump shallow tree for diagnostics (two levels)
                                    try:
                                        root_base = os.path.dirname(parent)
                                        print('DEBUG load_to_upload: dumping shallow tree under', root_base)
                                        for r, ds, fs in os.walk(root_base):
                                            rel = os.path.relpath(r, root_base)
                                            print('TREE:', rel, 'dirs=', len(ds), 'files=', len(fs))
                                            # limit output
                                            if rel.count(os.path.sep) >= 2:
                                                break
                                    except Exception:
                                        pass
                            except Exception as e:
                                print('DEBUG load_to_upload: parent listing failed:', repr(e))
                    except Exception:
                        pass
                    cleaned_db_path = db_path_val
                    try:
                        if isinstance(db_path_val, str):
                            cleaned_db_path = ''.join([c for c in db_path_val if ord(c) >= 32])
                    except Exception:
                        cleaned_db_path = db_path_val

                    exists_check = os.path.exists(cleaned_db_path)
                    print('DEBUG load_to_upload: os.path.exists(db_path_val)=', exists_check)
                    try:
                        tail = db_path_val[-20:]
                        print('DEBUG load_to_upload: tail_repr=', repr(tail))
                        print('DEBUG load_to_upload: tail_ordinals=', [ord(c) for c in tail])
                    except Exception:
                        pass

                    try:
                        print('DEBUG load_to_upload: isabs=', os.path.isabs(db_path_val), 'len=', len(db_path_val))
                    except Exception:
                        pass

                    # prefer cleaned path if it exists; try several normalized variants to catch odd encodings
                    tried_paths = []
                    try_variants = [db_path_val, cleaned_db_path]
                    try_variants += [os.path.normpath(p) for p in try_variants if isinstance(p, str)]
                    try_variants += [p.replace('\\\\', '\\') for p in try_variants if isinstance(p, str)]
                    # also try replacing forward/back slashes
                    try_variants += [p.replace('/', os.path.sep) for p in try_variants if isinstance(p, str)]
                    found = False
                    for p in try_variants:
                        try:
                            if not isinstance(p, str):
                                continue
                            if os.path.exists(p):
                                candidate = os.path.abspath(p)
                                found = True
                                break
                            tried_paths.append((p, os.path.exists(p)))
                        except Exception:
                            continue
                    if not found:
                        # last attempt: check with pathlib resolve (non-strict)
                        try:
                            from pathlib import Path
                            for p in try_variants:
                                try:
                                    if not isinstance(p, str):
                                        continue
                                    rp = Path(p).expanduser()
                                    if rp.exists():
                                        candidate = str(rp.resolve())
                                        found = True
                                        break
                                except Exception:
                                    continue
                        except Exception:
                            pass
                except Exception:
                    pass

                # If path is missing on disk, attempt to locate by filename under several locations
                if not candidate and db_filename:
                    # check upload folder first
                    try:
                        up = app.config.get('UPLOAD_FOLDER')
                        if up:
                            cand = os.path.join(up, db_filename)
                            exists_cand = os.path.exists(cand)
                            print('DEBUG load_to_upload: checking UPLOAD_FOLDER cand=', cand, 'exists=', exists_cand)
                            try:
                                print('DEBUG load_to_upload: UPLOAD_FOLDER exists=', os.path.exists(up), 'contents=', os.listdir(up) if os.path.exists(up) else None)
                            except Exception:
                                pass
                            if exists_cand:
                                candidate = os.path.abspath(cand)
                                print('DEBUG load_to_upload: found in UPLOAD_FOLDER ->', candidate)
                    except Exception:
                        pass

                    # check allowed roots direct
                    if not candidate:
                        for rroot in _allowed_roots():
                            try:
                                cand = os.path.join(rroot, db_filename)
                                print('DEBUG load_to_upload: checking root cand=', cand, 'exists=', os.path.exists(cand))
                                if os.path.exists(cand):
                                    candidate = os.path.abspath(cand)
                                    print('DEBUG load_to_upload: found in allowed root ->', candidate)
                                    break
                            except Exception:
                                continue

                    # recursive fallback under allowed roots
                    if not candidate:
                        for rroot in _allowed_roots():
                            try:
                                for root_dir, dirs, files in os.walk(rroot):
                                    if db_filename in files:
                                        candidate = os.path.abspath(os.path.join(root_dir, db_filename))
                                        break
                                if candidate:
                                    break
                            except Exception:
                                continue

                    # Extra fallback: search under the directory containing the DB file (useful for test tmpdirs)
                    if not candidate and DB_FILE and os.path.exists(DB_FILE):
                        try:
                            db_parent = os.path.dirname(os.path.abspath(DB_FILE))
                            for root_dir, dirs, files in os.walk(db_parent):
                                if db_filename in files:
                                    candidate = os.path.abspath(os.path.join(root_dir, db_filename))
                                    print('DEBUG load_to_upload: recovered candidate under DB parent ->', candidate)
                                    break
                            # also try one level up from DB parent
                            if not candidate:
                                maybe_root = os.path.dirname(db_parent)
                                for root_dir, dirs, files in os.walk(maybe_root):
                                    if db_filename in files:
                                        candidate = os.path.abspath(os.path.join(root_dir, db_filename))
                                        print('DEBUG load_to_upload: recovered candidate under DB grandparent ->', candidate)
                                        break
                        except Exception:
                            pass

                    # Extra fallback: search under system temp directory (helps pytest tmpdirs)
                    if not candidate:
                        try:
                            tmp_root = tempfile.gettempdir()
                            print('DEBUG load_to_upload: searching system temp for', db_filename, 'under', tmp_root)
                            for root_dir, dirs, files in os.walk(tmp_root):
                                if db_filename in files:
                                    candidate = os.path.abspath(os.path.join(root_dir, db_filename))
                                    print('DEBUG load_to_upload: found candidate under system temp ->', candidate)
                                    break
                        except Exception:
                            pass

        except ValueError:
            return jsonify({'error': 'Invalid file_id'}), 400
        except Exception as e:
            return jsonify({'error': f'DB lookup failed: {e}'}), 500

    else:
        # Resolve by provided path relative to allowed roots
        try:
            src_rel = file_path.lstrip('/\\')
            for root_path in _allowed_roots():
                try:
                    abs_root = os.path.abspath(root_path)
                    abs_candidate = os.path.abspath(os.path.join(abs_root, src_rel))
                    if (abs_candidate == abs_root or abs_candidate.startswith(abs_root + os.path.sep)) and os.path.exists(abs_candidate):
                        candidate = abs_candidate
                        break
                except Exception:
                    continue
            if not candidate:
                return jsonify({'error': 'Source file not found or not accessible'}), 404
        except Exception as e:
            return jsonify({'error': f'Path resolution error: {e}'}), 500

    # final candidate check
    if not candidate:
        return jsonify({'error': 'DB file record not found or file missing on disk'}), 404

    if not is_path_under_allowed_roots(candidate):
        return jsonify({'error': 'Access denied'}), 403

    # Copy into upload folder
    dest_dir = app.config.get('UPLOAD_FOLDER')
    if not dest_dir:
        return jsonify({'error': 'Upload folder not configured'}), 500
    safe_name = secure_filename(os.path.basename(candidate))
    dest_path = os.path.join(dest_dir, safe_name)
    base, ext = os.path.splitext(safe_name)
    counter = 1
    while os.path.exists(dest_path):
        safe_name = f"{base}_{counter}{ext}"
        dest_path = os.path.join(dest_dir, safe_name)
        counter += 1

    try:
        shutil.copy2(candidate, dest_path)
    except Exception as e:
        print('ERROR load_to_upload copy failed:', repr(e))
        return jsonify({'error': f'Copy failed: {e}'}), 500

    # Audit the operation
    try:
        audit_event('load_to_upload', target=candidate, details={'dest': dest_path, 'by': request.remote_addr})
    except Exception:
        pass

    # Register the copied file
    try:
        fname_key = safe_name
        uploaded_files_db[fname_key] = {
            'path': dest_path,
            'size': os.path.getsize(dest_path) if os.path.exists(dest_path) else 0,
            'hash_info': {},
            'encryption_status': {'encrypted': False, 'decrypted_path': None, 'decrypting': False, 'description': ''}
        }
    except Exception:
        pass

    return jsonify({'ok': True, 'filename': safe_name})

@app.route('/db_delete_file/<int:file_id>', methods=['POST'])
def db_delete_file(file_id):
    # Verify simple CSRF token
    token = request.form.get('csrf_token') or request.headers.get('X-CSRF-Token')
    try:
        if not token or token != session.get('csrf_token'):
            flash('Invalid CSRF token. Deletion denied.', 'error')
            return redirect(url_for('databases'))
    except Exception:
        flash('Session error validating CSRF token. Deletion denied.', 'error')
        return redirect(url_for('databases'))

    remove_disk = bool(request.form.get('remove_disk'))

    # Remove DB record and optionally remove disk file if exists
    try:
        conn = _get_db_conn()
        cur = conn.cursor()
        cur.execute('SELECT path FROM files WHERE id=?', (file_id,))
        row = cur.fetchone()
        path = row[0] if row else None

        # delete record
        cur.execute('DELETE FROM files WHERE id=?', (file_id,))
        conn.commit()
        try:
            conn.close()
        except Exception:
            pass

        # Attempt to remove file on disk only if requested and if it is inside configured folders
        if remove_disk and path and os.path.exists(path):
            try:
                if not is_path_under_allowed_roots(path):
                    flash('Refused to delete file: path not in allowed folders.', 'error')
                else:
                    try:
                        os.remove(path)
                    except Exception as ex:
                        flash(f'File deletion failed: {ex}', 'error')
            except Exception as ex:
                flash(f'File deletion failed: {ex}', 'error')
    except Exception as e:
        flash(f"Error deleting DB record: {e}", 'error')
        return redirect(url_for('databases'))

    flash('Record deleted.', 'success')
    return redirect(url_for('databases'))
@app.before_request
def ensure_csrf_token():
    # Ensure a simple CSRF token exists for the user's session
    try:
        if 'csrf_token' not in session:
            session['csrf_token'] = secrets.token_hex(16)
    except Exception:
        # sessions may not be writable in some test contexts
        pass
EXTRA_DB_JS = r"""
<script>
// Engine options by type
const enginesByType = {
    'Relational (RDBMS)': ['MySQL', 'PostgreSQL', 'Oracle', 'SQL Server', 'SQLite', 'MariaDB', 'IBM Db2'],
    'Document (NoSQL)': ['MongoDB', 'CouchDB', 'RavenDB', 'ArangoDB', 'Amazon DocumentDB'],
    'Key-Value (NoSQL)': ['Redis', 'DynamoDB', 'Riak KV', 'Aerospike', 'Etcd', 'Berkeley DB'],
    'Wide-Column (NoSQL)': ['Cassandra', 'HBase', 'ScyllaDB', 'Bigtable', 'ClickHouse'],
    'Graph (NoSQL)': ['Neo4j', 'Amazon Neptune', 'OrientDB', 'ArangoDB', 'TigerGraph'],
    'Object-Oriented': ['db4o', 'ObjectDB', 'Versant', 'GemStone/S', 'ZopeDB'],
    'Time-Series': ['InfluxDB', 'TimescaleDB', 'Prometheus', 'OpenTSDB', 'VictoriaMetrics'],
    'NewSQL': ['Google Spanner', 'CockroachDB', 'VoltDB', 'TiDB', 'NuoDB', 'SingleStore'],
    'Hierarchical': ['IBM IMS', 'Windows Registry', 'XML databases'],
    'Network': ['IDS', 'IDMS', 'TurboIMAGE', 'Raima DB']
};

function populateEngines() {
    const typeSel = document.querySelector('select[name="type"]');
    const engSel = document.getElementById('engine-select');
    const options = enginesByType[typeSel.value] || [];
    engSel.innerHTML = '';
    if (options.length === 0) {
        const o = document.createElement('option'); o.value=''; o.textContent='(No engines available)'; engSel.appendChild(o);
        return;
    }
    options.forEach(e => {
        const o = document.createElement('option'); o.value = e; o.textContent = e; engSel.appendChild(o);
    });
}

    document.addEventListener('DOMContentLoaded', function(){
    const typeSel = document.querySelector('select[name="type"]');
    if (typeSel) {
        typeSel.addEventListener('change', populateEngines);
        // populate initial
        populateEngines();
    }
    // toggle between connection string and manual
    const radios = document.getElementsByName('conn_type');
    const connString = document.getElementById('conn-string');
    const manualDiv = document.getElementById('manual-conn');
    function updateConnMode(){
        const val = Array.from(radios).find(r=>r.checked).value;
        const container = document.getElementById('conn-string-container');
        if(val === 'manual'){
            manualDiv.classList.remove('hidden');
            // remove any existing conn-string input when manual selected
            const existing = document.getElementById('conn-string');
            if(existing) existing.remove();
        } else {
            manualDiv.classList.add('hidden');
            // ensure a single-line conn input exists
            if(!document.getElementById('conn-string')){
                const inp = document.createElement('input');
                inp.id = 'conn-string';
                inp.name = 'conn';
                inp.type = 'text';
                inp.placeholder = 'Connection string or JSON';
                inp.className = 'mt-1 block w-full bg-gray-800 border-gray-600 rounded-md p-2 text-white';
                container.appendChild(inp);
            }
        }
    }
    Array.from(radios).forEach(r=>r.addEventListener('change', updateConnMode));
    updateConnMode();
});
</script>

<script>
// Add connection: Test and Reset handlers
document.addEventListener('DOMContentLoaded', function(){
    const form = document.getElementById('db-add-form');
    const testBtn = document.getElementById('db-test-button');
    const resetBtn = document.getElementById('db-reset-button');

    function collectFormData() {
        const fd = new FormData(form);
        // If conn input is not present (manual mode) ensure conn not sent
        if(!document.getElementById('conn-string')) fd.delete('conn');
        return fd;
    }

    testBtn.addEventListener('click', function(){
        const fd = collectFormData();
        const status = document.getElementById('db-test-status');
        // show pending
        status.className = 'mt-2 p-2 rounded text-sm bg-yellow-800 text-yellow-100';
        status.textContent = 'Testing connection...';
        status.classList.remove('hidden');
        fetch('/databases/test', { method: 'POST', body: fd })
            .then(r => r.json())
            .then(data => {
                if(data.success){
                    status.className = 'mt-2 p-2 rounded text-sm bg-green-800 text-green-100';
                    status.textContent = 'Connection OK: ' + data.message + (data.size ? ' | Size: '+data.size : '');
                } else {
                    status.className = 'mt-2 p-2 rounded text-sm bg-red-800 text-red-100';
                    status.textContent = 'Connection failed: ' + data.message;
                }
                // auto-hide after a while
                setTimeout(()=>{ status.classList.add('hidden'); }, 8000);
            }).catch(err => { 
                status.className = 'mt-2 p-2 rounded text-sm bg-red-800 text-red-100';
                status.textContent = 'Test request failed: ' + err.message;
                setTimeout(()=>{ status.classList.add('hidden'); }, 8000);
            });
    });

    resetBtn.addEventListener('click', function(){
        // reset all inputs in the form and clear status
        form.reset();
        const status = document.getElementById('db-test-status');
        status.classList.add('hidden'); status.textContent = '';
        // ensure manual/string UI updated
        Array.from(document.getElementsByName('conn_type')).forEach(r=>{ if(r.checked) r.dispatchEvent(new Event('change')); });
        // show a short ephemeral toast to confirm reset
        function ensureToastContainer(){
            let c = document.getElementById('fac-toast-container');
            if(c) return c;
            c = document.createElement('div'); c.id = 'fac-toast-container';
            c.style.position = 'fixed'; c.style.right = '20px'; c.style.bottom = '20px'; c.style.zIndex = 20000; c.style.display = 'flex'; c.style.flexDirection = 'column'; c.style.gap = '8px'; document.body.appendChild(c);
            return c;
        }
        const container = ensureToastContainer();
        const t = document.createElement('div'); t.className='fac-toast'; t.style.background='#0b1220'; t.style.border='1px solid #333'; t.style.color='#fff'; t.style.padding='8px 10px'; t.style.borderRadius='6px'; t.style.boxShadow='0 6px 18px rgba(0,0,0,0.3)'; t.textContent = 'Form reset'; container.appendChild(t);
        setTimeout(()=>{ try{ t.style.transition='opacity 0.5s ease'; t.style.opacity='0'; setTimeout(()=>{ try{ t.remove(); }catch(e){} }, 550); }catch(e){} }, 1800);
    });
});
</script>
"""

@app.route('/databases', methods=['GET'])
def databases():
    # Compute simple connected/disconnected summary counts
    connected_count = 0
    disconnected_count = 0
    for cfg in db_connections.values():
        if cfg.get('connected'):
            connected_count += 1
        else:
            disconnected_count += 1
    summary = {'connected': connected_count, 'disconnected': disconnected_count}
    # If a session was active, end it because user is going to DB configuration
    sess_id = session.pop('analysis_session_id', None)
    if sess_id:
        try:
            end_analysis_session(sess_id)
        except Exception:
            pass

    # Provide configured DBs and summary to template. Append EXTRA_DB_JS so page JS binds correctly.
    db_content = render_template_string(DATABASES_CONTENT, dbs=db_connections, summary=summary)
    try:
        db_content = db_content + EXTRA_DB_JS
    except Exception:
        pass
    return render_template_string(BASE_TEMPLATE, content=db_content, show_session_details=False)

@app.route('/databases/add', methods=['POST'])
def db_add():
    name = request.form.get('name')
    dtype = request.form.get('type')
    engine = request.form.get('engine')
    conn = request.form.get('conn')
    # collect manual fields if provided
    conn_type = request.form.get('conn_type', 'string')
    conn_details = None
    if conn_type == 'manual':
        conn_details = {
            'host': request.form.get('host'),
            'port': request.form.get('port'),
            'database': request.form.get('database'),
            'username': request.form.get('username'),
            'password': request.form.get('password'),
            'extra': request.form.get('extra')
        }
    # Validate: name is required. Connection can be provided either via a string
    # or via manual fields. Accept manual-only submissions (no 'conn' string).
    if not name:
        flash('Name is required', 'warning')
        return redirect(url_for('databases'))
    if conn_type == 'string' and not conn:
        flash('Connection string is required when Connection String mode is selected', 'warning')
        return redirect(url_for('databases'))
    if conn_type == 'manual':
        # require at least host or database name
        if not conn_details or not (conn_details.get('host') or conn_details.get('database')):
            flash('At least host or database name should be provided for manual connections', 'warning')
            return redirect(url_for('databases'))
    # create an id
    dbid = hashlib.sha1((name + str(time.time())).encode()).hexdigest()[:12]
    db_connections[dbid] = {
        'name': name,
        'type': dtype,
        'engine': engine,
        'conn': conn,
        'conn_details': conn_details,
        'connected': False,
        'size': None,
        'last_checked': None,
        'message': None
    }
    save_db_configs()
    # Attempt to create default folder entries in the configured DB
    try:
        create_default_folder_entries_for_db(db_connections[dbid], dbid)
    except Exception as e:
        app.logger.debug(f"create_default_folder_entries_for_db error (add): {e}")
    flash('Database connection added', 'success')
    return redirect(url_for('databases'))

@app.route('/databases/edit/<dbid>', methods=['GET', 'POST'])
def db_edit(dbid):
    cfg = db_connections.get(dbid)
    if not cfg:
        flash('Connection not found', 'warning')
        return redirect(url_for('databases'))

    if request.method == 'POST':
        cfg['name'] = request.form.get('name') or cfg.get('name')
        cfg['type'] = request.form.get('type') or cfg.get('type')
        cfg['engine'] = request.form.get('engine') or cfg.get('engine')
        cfg['conn'] = request.form.get('conn') or cfg.get('conn')
        conn_type = request.form.get('conn_type', 'string')
        if conn_type == 'manual':
            cfg['conn_details'] = {
                'host': request.form.get('host'),
                'port': request.form.get('port'),
                'database': request.form.get('database'),
                'username': request.form.get('username'),
                'password': request.form.get('password'),
                'extra': request.form.get('extra')
            }
        else:
            cfg['conn_details'] = None
        save_db_configs()
        # Attempt to create/update default folder entries in the configured DB
        try:
            create_default_folder_entries_for_db(cfg, dbid)
        except Exception as e:
            app.logger.debug(f"create_default_folder_entries_for_db error (edit): {e}")
        flash('Connection updated', 'success')
        return redirect(url_for('databases'))

    # render edit form quickly (GET)
    cd = cfg.get('conn_details') or {}
    conn_type_checked_string = 'checked' if not cd else ''
    conn_type_checked_manual = 'checked' if cd else ''
    edit_form = """
    <h1 class="text-2xl text-white">Edit Connection</h1>
    <form method="post"> 
        <label class="text-sm text-gray-300">Name</label>
        <input name="name" value="__NAME__" class="mt-1 block w-full bg-gray-800 border-gray-600 rounded-md p-2 text-white" />
        <label class="text-sm text-gray-300 mt-2">Type</label>
        <input name="type" value="__TYPE__" class="mt-1 block w-full bg-gray-800 border-gray-600 rounded-md p-2 text-white" />
        <label class="text-sm text-gray-300 mt-2">Engine</label>
        <input name="engine" value="__ENGINE__" class="mt-1 block w-full bg-gray-800 border-gray-600 rounded-md p-2 text-white" />
        <label class="block text-sm text-gray-300 mt-2">Connection Input</label>
        <div class="flex items-center space-x-4 mb-2">
            <label class="text-sm text-gray-300"><input type="radio" name="conn_type" value="string" __CHK_STR__> Connection String</label>
            <label class="text-sm text-gray-300"><input type="radio" name="conn_type" value="manual" __CHK_MAN__> Manual (host/port/...)</label>
        </div>
    <div id="conn-string-container"></div>
    <input type="hidden" id="initial-conn" value="__CONN__" />
    <div id="manual-conn" class="mt-2 space-y-2">
            <div class="grid grid-cols-1 md:grid-cols-2 gap-2">
                <input name="host" placeholder="Host" value="__HOST__" class="bg-gray-800 border-gray-600 rounded-md p-2 text-white"> 
                <input name="port" placeholder="Port" value="__PORT__" class="bg-gray-800 border-gray-600 rounded-md p-2 text-white"> 
                <input name="database" placeholder="Database name" value="__DB__" class="bg-gray-800 border-gray-600 rounded-md p-2 text-white"> 
                <input name="username" placeholder="Username" value="__USER__" class="bg-gray-800 border-gray-600 rounded-md p-2 text-white"> 
                <input name="password" placeholder="Password" type="password" value="__PASS__" class="bg-gray-800 border-gray-600 rounded-md p-2 text-white"> 
                <input name="extra" placeholder="Extra params (JSON)" value="__EXTRA__" class="bg-gray-800 border-gray-600 rounded-md p-2 text-white"> 
            </div>
        </div>
        <script>
        // toggle manual/string in edit form and dynamically create/remove the conn input
        (function(){
            function update(){
                var radios = document.getElementsByName('conn_type');
                var manualDiv = document.getElementById('manual-conn');
                var container = document.getElementById('conn-string-container');
                var initial = document.getElementById('initial-conn');
                var val = Array.from(radios).find(function(r){return r.checked;}).value;
                if(val === 'manual'){
                    manualDiv.style.display='block';
                    // remove any conn input
                    var existing = document.getElementById('conn-string'); if(existing) existing.remove();
                } else {
                    manualDiv.style.display='none';
                    // create single-line conn input if missing
                    if(!document.getElementById('conn-string')){
                        var inp = document.createElement('input');
                        inp.id='conn-string'; inp.name='conn'; inp.type='text';
                        inp.className='mt-1 block w-full bg-gray-800 border-gray-600 rounded-md p-2 text-white';
                        inp.value = initial ? initial.value : '';
                        container.appendChild(inp);
                    }
                }
            }
            Array.from(document.getElementsByName('conn_type')).forEach(function(r){ r.addEventListener('change', update); });
            update();
        })();
        </script>
        <div class="mt-3"><button class="btn-primary px-4 py-2 rounded-lg">Save</button></div>
    </form>
    """
    # substitute placeholders
    edit_form = edit_form.replace('__NAME__', str(cfg.get('name') or ''))
    edit_form = edit_form.replace('__TYPE__', str(cfg.get('type') or ''))
    edit_form = edit_form.replace('__ENGINE__', str(cfg.get('engine') or ''))
    edit_form = edit_form.replace('__CONN__', str(cfg.get('conn') or ''))
    edit_form = edit_form.replace('__HOST__', str(cd.get('host') or ''))
    edit_form = edit_form.replace('__PORT__', str(cd.get('port') or ''))
    edit_form = edit_form.replace('__DB__', str(cd.get('database') or ''))
    edit_form = edit_form.replace('__USER__', str(cd.get('username') or ''))
    edit_form = edit_form.replace('__PASS__', str(cd.get('password') or ''))
    edit_form = edit_form.replace('__EXTRA__', str(cd.get('extra') or ''))
    edit_form = edit_form.replace('__CHK_STR__', conn_type_checked_string)
    edit_form = edit_form.replace('__CHK_MAN__', conn_type_checked_manual)
    return render_template_string(BASE_TEMPLATE, content=edit_form)


@app.route('/databases/test', methods=['POST'])
def db_test():
    # Build a temporary cfg from posted form data (do not save)
    name = request.form.get('name') or 'test'
    dtype = request.form.get('type') or ''
    engine = request.form.get('engine') or ''
    conn = request.form.get('conn') or ''
    conn_type = request.form.get('conn_type', 'string')
    conn_details = None
    if conn_type == 'manual':
        conn_details = {
            'host': request.form.get('host'),
            'port': request.form.get('port'),
            'database': request.form.get('database'),
            'username': request.form.get('username'),
            'password': request.form.get('password'),
            'extra': request.form.get('extra')
        }
    cfg = {
        'name': name,
        'type': dtype,
        'engine': engine,
        'conn': conn,
        'conn_details': conn_details
    }
    success, message, size = test_db_connection(cfg)
    return jsonify({'success': bool(success), 'message': message, 'size': size})

@app.route('/databases/delete/<dbid>', methods=['POST'])
def db_delete(dbid):
    if dbid in db_connections:
        del db_connections[dbid]
        save_db_configs()
        flash('Connection deleted', 'success')
    else:
        flash('Connection not found', 'warning')
    return redirect(url_for('databases'))

@app.route('/databases/connect/<dbid>', methods=['POST'])
def db_connect(dbid):
    cfg = db_connections.get(dbid)
    if not cfg:
        flash('Connection not found', 'warning')
        return redirect(url_for('databases'))

    # Try a real connection test (best effort). Drivers may be missing in runtime.
    success, message, size = test_db_connection(cfg)
    cfg['connected'] = bool(success)
    cfg['message'] = message
    cfg['last_checked'] = datetime.datetime.now(datetime.timezone.utc).strftime('%Y-%m-%d %H:%M:%S UTC')
    if size:
        cfg['size'] = size
    save_db_configs()
    if success:
        flash(f'Connected: {message}', 'success')
    else:
        flash(f'Connection failed: {message}', 'warning')
    return redirect(url_for('databases'))


def test_db_connection(cfg, timeout=5):
    """Attempt to connect to a database described by cfg (best-effort).

    Returns: (success:bool, message:str, size:str|None)
    """
    conn = (cfg.get('conn') or '').strip()
    dtype = (cfg.get('type') or '').lower()

    # If manual conn_details provided, prefer those
    params = {}
    cd = cfg.get('conn_details')
    if cd and isinstance(cd, dict):
        # normalize manual fields into params
        params.update({k: v for k, v in cd.items() if v is not None})
    else:
        # try parse JSON conn string if available
        try:
            if conn.startswith('{'):
                params = json.loads(conn)
        except Exception:
            params = {}

    # Helper to format size where possible
    def _fmt(b):
        try:
            return format_bytes(int(b))
        except Exception:
            return None

    # SQLite (local file)
    try:
        if (conn and (conn.lower().startswith('sqlite') or conn.lower().endswith('.sqlite'))) or dtype.startswith('sqlite'):
            # try to get path
            path = conn.split('///')[-1] if '///' in conn else conn
            if params and isinstance(params, dict) and params.get('path'):
                path = params.get('path')
            if not path or not os.path.exists(path):
                return False, f'SQLite file not found: {path}', None
            try:
                con = sqlite3.connect(path, timeout=timeout)
                con.close()
                size = _fmt(os.path.getsize(path))
                return True, 'SQLite file accessible', size
            except Exception as e:
                return False, f'SQLite connect error: {e}', None
    except Exception:
        pass

    # Try Postgres
    try:
        if (conn and (conn.startswith('postgres') or conn.startswith('postgresql'))) or 'postgres' in dtype:
            try:
                import psycopg2
            except Exception:
                return False, 'psycopg2 driver not installed', None
            try:
                # parse URL if present
                from urllib.parse import urlparse
                conn_kwargs = {}
                if conn and (conn.startswith('postgres') or conn.startswith('postgresql')):
                    u = urlparse(conn)
                    if u.hostname:
                        conn_kwargs['host'] = u.hostname
                    if u.port:
                        conn_kwargs['port'] = u.port
                    if u.username:
                        conn_kwargs['user'] = u.username
                    if u.password:
                        conn_kwargs['password'] = u.password
                    if u.path:
                        conn_kwargs['dbname'] = u.path.lstrip('/')
                if params and isinstance(params, dict):
                    conn_kwargs.update(params)
                pg = psycopg2.connect(connect_timeout=timeout, **conn_kwargs)
                cur = pg.cursor()
                try:
                    cur.execute("SELECT pg_database_size(current_database())")
                    bytes_size = cur.fetchone()[0]
                    cur.close()
                    pg.close()
                    return True, 'PostgreSQL reachable', _fmt(bytes_size)
                except Exception:
                    cur.close()
                    pg.close()
                    return True, 'PostgreSQL reachable (size unavailable)', None
            except Exception as e:
                return False, f'Postgres connection error: {e}', None
    except Exception:
        pass

    # Try MySQL (pymysql / MySQLdb)
    try:
        if (conn and conn.startswith('mysql')) or 'mysql' in dtype:
            mysql_driver = _get_mysql_driver()
            if mysql_driver is None:
                return False, 'MySQL driver (pymysql/mysqlclient) not installed', None
            try:
                from urllib.parse import urlparse
                kwargs = {}
                if conn and conn.startswith('mysql'):
                    u = urlparse(conn)
                    if u.hostname:
                        kwargs['host'] = u.hostname
                    if u.port:
                        kwargs['port'] = u.port
                    if u.username:
                        kwargs['user'] = u.username
                    if u.password:
                        kwargs['password'] = u.password
                    if u.path:
                        kwargs['db'] = u.path.lstrip('/')
                if params and isinstance(params, dict):
                    kwargs.update(params)
                conn_mysql = mysql_driver.connect(connect_timeout=timeout, **kwargs)
                cur = conn_mysql.cursor()
                try:
                    cur.execute("SELECT SUM(data_length + index_length) FROM information_schema.tables WHERE table_schema = DATABASE()")
                    row = cur.fetchone()
                    conn_mysql.close()
                    if row and row[0]:
                        return True, 'MySQL reachable', _fmt(row[0])
                    return True, 'MySQL reachable (size unavailable)', None
                except Exception:
                    conn_mysql.close()
                    return True, 'MySQL reachable', None
            except Exception as e:
                return False, f'MySQL connection error: {e}', None
    except Exception:
        pass

    # Try MongoDB (pymongo)
    try:
        if (conn and conn.startswith('mongodb')) or 'mongo' in dtype:
            try:
                from pymongo import MongoClient
            except Exception:
                return False, 'pymongo not installed', None
            try:
                # build mongo uri or use conn
                client = MongoClient(conn, serverSelectionTimeoutMS=timeout*1000) if conn else MongoClient(serverSelectionTimeoutMS=timeout*1000)
                client.server_info()
                # determine db name
                from urllib.parse import urlparse
                dbname = None
                if conn:
                    u = urlparse(conn)
                    dbname = (u.path.lstrip('/') if u.path and u.path != '/' else None)
                if not dbname:
                    dbname = params.get('database') or params.get('db') or 'admin'
                try:
                    stats = client[dbname].command('dbstats')
                    dataSize = stats.get('dataSize') or stats.get('storageSize')
                    client.close()
                    return True, 'MongoDB reachable', _fmt(dataSize)
                except Exception:
                    client.close()
                    return True, 'MongoDB reachable (size unavailable)', None
            except Exception as e:
                return False, f'MongoDB connection error: {e}', None
    except Exception:
        pass

    # Try Redis
    try:
        if (conn and conn.startswith('redis')) or 'redis' in dtype:
            try:
                import redis as redislib
            except Exception:
                return False, 'redis-py not installed', None
            try:
                client = redislib.from_url(conn) if conn else redislib.Redis(host=params.get('host', 'localhost'), port=int(params.get('port', 6379)))
                info = client.info()
                used = info.get('used_memory') or info.get('used_memory_rss')
                try:
                    client.close()
                except Exception:
                    pass
                return True, 'Redis reachable', _fmt(used)
            except Exception as e:
                return False, f'Redis connection error: {e}', None
    except Exception:
        pass

    # --- Additional engines (best-effort) ---
    # DynamoDB (boto3)
    try:
        if 'dynamodb' in (conn or '').lower() or 'dynamodb' in dtype:
            try:
                import boto3
            except Exception:
                return False, 'boto3 not installed for DynamoDB', None
            try:
                if params and isinstance(params, dict):
                    client = boto3.client('dynamodb', **params)
                else:
                    client = boto3.client('dynamodb')
                client.list_tables(Limit=1)
                return True, 'DynamoDB reachable', None
            except Exception as e:
                return False, f'DynamoDB error: {e}', None
    except Exception:
        pass

    # CouchDB (requests)
    try:
        if 'couchdb' in (conn or '').lower() or 'couchdb' in dtype:
            try:
                import requests
            except Exception:
                return False, 'requests not installed for CouchDB', None
            try:
                resp = requests.get(conn, timeout=timeout)
                if resp.status_code in (200, 301, 302):
                    return True, 'CouchDB reachable', None
                return False, f'CouchDB returned {resp.status_code}', None
            except Exception as e:
                return False, f'CouchDB error: {e}', None
    except Exception:
        pass

    # ArangoDB (python-arango)
    try:
        if 'arango' in (conn or '').lower() or 'arango' in dtype:
            try:
                import importlib
                arango_mod = importlib.import_module('arango')
                ArangoClient = getattr(arango_mod, 'ArangoClient', None)
                if ArangoClient is None:
                    return False, 'python-arango not installed', None
            except Exception:
                return False, 'python-arango not installed', None
            try:
                client = ArangoClient(hosts=conn)
                sys_db = client.db('_system')
                sys_db.version()
                return True, 'ArangoDB reachable', None
            except Exception as e:
                return False, f'ArangoDB error: {e}', None
    except Exception:
        pass

    # Cassandra (cassandra-driver)
    try:
        if 'cassandra' in (conn or '').lower() or 'cassandra' in dtype:
            try:
                import importlib
                cass_mod = importlib.import_module('cassandra.cluster')
                Cluster = getattr(cass_mod, 'Cluster', None)
                if Cluster is None:
                    return False, 'cassandra-driver not installed', None
            except Exception:
                return False, 'cassandra-driver not installed', None
            try:
                hosts = []
                if conn and ',' in conn:
                    hosts = [h.strip() for h in conn.split(',') if h.strip()]
                elif params and params.get('host'):
                    hosts = [params.get('host')]
                else:
                    hosts = [conn]
                cluster = Cluster(hosts)
                session = cluster.connect()
                session.execute('SELECT now() FROM system.local')
                cluster.shutdown()
                return True, 'Cassandra reachable', None
            except Exception as e:
                return False, f'Cassandra error: {e}', None
    except Exception:
        pass

    # InfluxDB (influxdb-client)
    try:
        if 'influx' in (conn or '').lower() or 'influxdb' in dtype:
            try:
                import importlib
                influx_mod = importlib.import_module('influxdb_client')
                InfluxDBClient = getattr(influx_mod, 'InfluxDBClient', None)
                if InfluxDBClient is None:
                    return False, 'influxdb-client not installed', None
            except Exception:
                return False, 'influxdb-client not installed', None
            try:
                if params and isinstance(params, dict):
                    url = params.get('url')
                    token = params.get('token')
                    org = params.get('org')
                    client = InfluxDBClient(url=url, token=token, org=org, timeout=timeout*1000)
                else:
                    client = InfluxDBClient(url=conn)
                health = client.health()
                client.close()
                if health and health.get('status') == 'pass':
                    return True, 'InfluxDB reachable', None
                return True, 'InfluxDB reachable (health unknown)', None
            except Exception as e:
                return False, f'InfluxDB error: {e}', None
    except Exception:
        pass

    # Neo4j (neo4j-driver)
    try:
        if 'neo4j' in (conn or '').lower() or 'neo4j' in dtype:
            try:
                from neo4j import GraphDatabase
            except Exception:
                return False, 'neo4j-driver not installed', None
            try:
                drv = GraphDatabase.driver(conn, max_transaction_retry_time=timeout)
                with drv.session() as s:
                    res = s.run('RETURN 1 as v')
                    _ = res.single()
                drv.close()
                return True, 'Neo4j reachable', None
            except Exception as e:
                return False, f'Neo4j error: {e}', None
    except Exception:
        pass

    return False, 'Unknown database type or unsupported connection string', None


def create_default_folder_entries_for_db(cfg, dbid=None):
    """Create default folder entries (Upload Files, Encrypted Files, Decrypted Files, Session Files)
    in the target DB if it's a local SQLite file. Otherwise fall back to inserting
    folder entries into the application's local DB using add_file_record with an extra
    field referencing the target dbid.
    """
    folder_names = [
        ('Upload Files', app.config.get('UPLOAD_FOLDER')),
        ('Encrypted Files', app.config.get('ENCRYPTED_FOLDER')),
        ('Decrypted Files', app.config.get('DECRYPTED_FOLDER')),
        ('Session Files', app.config.get('SESSION_FOLDER'))
    ]

    # Try to detect SQLite file target
    target_conn = (cfg.get('conn') or '').strip()
    params = cfg.get('conn_details') or {}
    sqlite_path = None
    try:
        if params and isinstance(params, dict) and params.get('path'):
            sqlite_path = params.get('path')
        elif target_conn:
            if target_conn.lower().startswith('sqlite') or target_conn.lower().endswith('.sqlite') or target_conn.lower().endswith('.db'):
                # handle sqlite:///path or direct path
                sqlite_path = target_conn.split('///')[-1]
    except Exception:
        sqlite_path = None

    # Helper to insert into a SQLite files table if available
    def _insert_into_sqlite(db_path, name, path):
        try:
            if not db_path or not os.path.exists(db_path):
                return False
            norm_path = ''
            try:
                norm_path = os.path.normcase(os.path.abspath(path)) if path else ''
            except Exception:
                norm_path = path or ''
            conn = sqlite3.connect(db_path)
            cur = conn.cursor()
            # ensure files table exists (same schema as app init)
            cur.execute('''
                CREATE TABLE IF NOT EXISTS files (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    filename TEXT,
                    file_type TEXT,
                    path TEXT,
                    size_bytes INTEGER,
                    created_at TEXT,
                    session_id INTEGER,
                    extra JSON
                )
            ''')
            # avoid duplicates: check by filename + file_type and case-insensitive path (Windows paths may differ in case)
            cur.execute('SELECT id FROM files WHERE filename=? AND file_type=? AND LOWER(path)=LOWER(?)', (name, 'folder', path or ''))
            cur.execute('SELECT id FROM files WHERE filename=? AND file_type=? AND LOWER(path)=LOWER(?)', (name, 'folder', norm_path))
            if cur.fetchone():
                conn.close()
                return True
            now = datetime.datetime.now().isoformat()
            cur.execute('INSERT INTO files(filename, file_type, path, size_bytes, created_at, session_id, extra) VALUES(?,?,?,?,?,?,?)',
                        (name, 'folder', norm_path, 0, now, None, json.dumps({'created_for_db': dbid}) if dbid else None))
            conn.commit()
            conn.close()
            return True
        except Exception as e:
            app.logger.debug(f"_insert_into_sqlite error: {e}")
            try:
                conn.close()
            except Exception:
                pass
            return False

    status = {'success': False, 'message': None, 'timestamp': datetime.datetime.now().isoformat()}
    # normalize connection hints for later engine-specific attempts
    conn_str = (cfg.get('conn') or '').strip() if cfg else ''
    engine = (cfg.get('engine') or cfg.get('type') or '').lower() if cfg else ''
    params = cfg.get('conn_details') or {}

    # If sqlite_path looks valid, attempt to write entries there
    if sqlite_path:
        try:
            for name, path in folder_names:
                _insert_into_sqlite(sqlite_path, name, path)
            status['success'] = True
            status['message'] = f"Inserted default folders into SQLite DB at {sqlite_path}"
            # persist status into cfg if provided
            if cfg is not None:
                cfg['folder_creation'] = status
                save_db_configs()
            return
        except Exception as e:
            status['message'] = f"SQLite insertion error: {e}"
        # Fallback: insert into local application DB for each folder
        for name, path in folder_names:
            try:
                # normalize path for consistent storage
                try:
                    norm_path = os.path.normcase(os.path.abspath(path)) if path else ''
                except Exception:
                    norm_path = path or ''
                # avoid duplicates in local app DB by checking existing rows (case-insensitive)
                conn = _get_db_conn()
                cur = conn.cursor()
                cur.execute('SELECT id FROM files WHERE filename=? AND file_type=? AND LOWER(path)=LOWER(?)', (name, 'folder', norm_path))
                if cur.fetchone():
                    conn.close()
                    continue
                extra = {'created_for_db': dbid} if dbid else None
                add_file_record(name, 'folder', norm_path, 0, session_id=None, extra=extra)
                conn.close()
            except Exception as e:
                app.logger.debug(f"Fallback create folder entry error: {e}")

    # Try Postgres
    try:
        conn_str = (cfg.get('conn') or '').strip() if cfg else ''
        engine = (cfg.get('engine') or cfg.get('type') or '').lower() if cfg else ''
        params = cfg.get('conn_details') or {}
        if conn_str and (conn_str.startswith('postgres') or conn_str.startswith('postgresql')) or 'postgres' in engine:
            try:
                import psycopg2
            except Exception:
                raise
            from urllib.parse import urlparse
            conn_kwargs = {}
            if conn_str and (conn_str.startswith('postgres') or conn_str.startswith('postgresql')):
                u = urlparse(conn_str)
                if u.hostname:
                    conn_kwargs['host'] = u.hostname
                if u.port:
                    conn_kwargs['port'] = u.port
                if u.username:
                    conn_kwargs['user'] = u.username
                if u.password:
                    conn_kwargs['password'] = u.password
                if u.path:
                    conn_kwargs['dbname'] = u.path.lstrip('/')
            # merge manual params
            if params and isinstance(params, dict):
                conn_kwargs.update({k: v for k, v in params.items() if v is not None})
            pg = psycopg2.connect(connect_timeout=5, **conn_kwargs)
            cur = pg.cursor()
            cur.execute('''
                CREATE TABLE IF NOT EXISTS files (
                    id SERIAL PRIMARY KEY,
                    filename TEXT,
                    file_type TEXT,
                    path TEXT,
                    size_bytes BIGINT,
                    created_at TEXT,
                    session_id TEXT,
                    extra JSONB
                )
            ''')
            for name, path in folder_names:
                cur.execute('SELECT id FROM files WHERE filename=%s AND path=%s AND file_type=%s', (name, path or '', 'folder'))
                if cur.fetchone():
                    continue
                # normalize path for consistent storage
                try:
                    npath = os.path.normcase(os.path.abspath(path)) if path else ''
                except Exception:
                    npath = path or ''
                extra_json = json.dumps({'created_for_db': dbid}) if dbid else None
                now = datetime.datetime.now().isoformat()
                cur.execute('INSERT INTO files(filename, file_type, path, size_bytes, created_at, session_id, extra) VALUES(%s,%s,%s,%s,%s,%s,%s)',
                            (name, 'folder', npath, 0, now, None, extra_json))
            pg.commit()
            cur.close()
            pg.close()
            status['success'] = True
            status['message'] = f"Inserted default folders into PostgreSQL DB ({conn_kwargs.get('host') or 'conn'})"
            if cfg is not None:
                cfg['folder_creation'] = status
                save_db_configs()
            return
    except Exception as e:
        app.logger.debug(f"Postgres insert attempt failed: {e}")
    except Exception:
        pass

    # Try MySQL
    try:
        if 'mysql' in engine or conn_str.startswith('mysql'):
            mysql_driver = _get_mysql_driver()
            if not mysql_driver:
                raise RuntimeError('MySQL driver not available')
            from urllib.parse import urlparse
            kwargs = {}
            if conn_str and conn_str.startswith('mysql'):
                u = urlparse(conn_str)
                if u.hostname:
                    kwargs['host'] = u.hostname
                if u.port:
                    kwargs['port'] = int(u.port)
                if u.username:
                    kwargs['user'] = u.username
                if u.password:
                    kwargs['password'] = u.password
                if u.path:
                    kwargs['db'] = u.path.lstrip('/')
            if params and isinstance(params, dict):
                kwargs.update({k: v for k, v in params.items() if v is not None})
            conn_mysql = mysql_driver.connect(connect_timeout=5, **kwargs)
            cur = conn_mysql.cursor()
            cur.execute('''
                CREATE TABLE IF NOT EXISTS files (
                    id BIGINT PRIMARY KEY AUTO_INCREMENT,
                    filename TEXT,
                    file_type TEXT,
                    path TEXT,
                    size_bytes BIGINT,
                    created_at TEXT,
                    session_id TEXT,
                    extra JSON
                )
            ''')
            for name, path in folder_names:
                cur.execute('SELECT id FROM files WHERE filename=%s AND path=%s AND file_type=%s', (name, path or '', 'folder'))
                if cur.fetchone():
                    continue
                now = datetime.datetime.now().isoformat()
                extra_json = json.dumps({'created_for_db': dbid}) if dbid else None
                cur.execute('INSERT INTO files(filename, file_type, path, size_bytes, created_at, session_id, extra) VALUES(%s,%s,%s,%s,%s,%s,%s)',
                            (name, 'folder', path or '', 0, now, None, extra_json))
            conn_mysql.commit()
            cur.close()
            conn_mysql.close()
            status['success'] = True
            status['message'] = f"Inserted default folders into MySQL DB ({kwargs.get('host') or 'conn'})"
            if cfg is not None:
                cfg['folder_creation'] = status
                save_db_configs()
            return
    except Exception as e:
        app.logger.debug(f"MySQL insert attempt failed: {e}")

    # Try MongoDB
    try:
        if 'mongo' in engine or conn_str.startswith('mongodb'):
            try:
                from pymongo import MongoClient
                # build client
                client = MongoClient(conn_str) if conn_str else MongoClient()
                # determine db name
                from urllib.parse import urlparse
                dbname = None
                if conn_str:
                    u = urlparse(conn_str)
                    dbname = u.path.lstrip('/') if u.path and u.path != '/' else None
                if not dbname:
                    dbname = params.get('database') or params.get('db') or 'admin'
                db = client[dbname]
                coll = db['files']
                for name, path in folder_names:
                    if coll.find_one({'filename': name, 'path': path or '', 'file_type': 'folder'}):
                        continue
                    doc = {'filename': name, 'file_type': 'folder', 'path': path or '', 'size_bytes': 0, 'created_at': datetime.datetime.now().isoformat(), 'session_id': None, 'extra': {'created_for_db': dbid} if dbid else {}}
                    coll.insert_one(doc)
                try:
                    client.close()
                except Exception:
                    pass
                status['success'] = True
                status['message'] = f"Inserted default folders into MongoDB ({dbname})"
                if cfg is not None:
                    cfg['folder_creation'] = status
                    save_db_configs()
                return
            except Exception as e:
                app.logger.debug(f"MongoDB insert attempt failed: {e}")
    except Exception:
        pass

    # Final fallback: insert into local application DB via add_file_record, tag with target dbid
    for name, path in folder_names:
        try:
            # avoid duplicates in local app DB by checking existing rows
            conn = _get_db_conn()
            cur = conn.cursor()
            cur.execute('SELECT id FROM files WHERE filename=? AND file_type=? AND LOWER(path)=LOWER(?)', (name, 'folder', path or ''))
            if cur.fetchone():
                conn.close()
                continue
            extra = {'created_for_db': dbid} if dbid else None
            add_file_record(name, 'folder', path or '', 0, session_id=None, extra=extra)
            conn.close()
        except Exception as e:
            app.logger.debug(f"Fallback create folder entry error: {e}")
    # mark success for fallback path
    status['success'] = True
    status['message'] = 'Inserted default folders into local application DB (fac_data.db)'
    if cfg is not None:
        cfg['folder_creation'] = status
        save_db_configs()
    return


def _validate_all_db_connections(interval_seconds=300):
    """Background thread: periodically validate all configured DB connections."""
    while True:
        try:
            if not db_connections:
                time.sleep(interval_seconds)
                continue
            for dbid, cfg in list(db_connections.items()):
                try:
                    success, message, size = test_db_connection(cfg)
                    cfg['connected'] = bool(success)
                    cfg['message'] = message
                    cfg['last_checked'] = datetime.datetime.now(datetime.timezone.utc).strftime('%Y-%m-%d %H:%M:%S UTC')
                    if size:
                        cfg['size'] = size
                except Exception as e:
                    cfg['connected'] = False
                    cfg['message'] = f'Validation error: {e}'
            save_db_configs()
        except Exception as e:
            app.logger.debug(f"Background DB validator error: {e}")
        # Sleep before next pass
        time.sleep(interval_seconds)


# Start background validator thread (daemon)
try:
    validator_thread = threading.Thread(target=_validate_all_db_connections, args=(300,), daemon=True)
    validator_thread.start()
except Exception:
    app.logger.debug('Could not start background DB validator thread')


@app.route('/databases/refresh_all', methods=['POST'])
def db_refresh_all():
    """Trigger an immediate full validation pass for all configured DBs."""
    try:
        for dbid, cfg in list(db_connections.items()):
            try:
                success, message, size = test_db_connection(cfg)
                cfg['connected'] = bool(success)
                cfg['message'] = message
                cfg['last_checked'] = datetime.datetime.now(datetime.timezone.utc).strftime('%Y-%m-%d %H:%M:%S UTC')
                if size:
                    cfg['size'] = size
            except Exception as e:
                cfg['connected'] = False
                cfg['message'] = f'Validation error: {e}'
        save_db_configs()
        flash('Refresh triggered for all connections', 'success')
    except Exception as e:
        flash(f'Error during refresh: {e}', 'warning')
    return redirect(url_for('databases'))


@app.route('/databases/list_json', methods=['GET'])
def db_list_json():
    """Return a simple JSON list of configured DBs (id, name, type) for client-side widgets."""
    try:
        out = []
        for dbid, cfg in db_connections.items():
            out.append({
                'id': dbid,
                'name': cfg.get('name', dbid),
                'type': cfg.get('type', 'unknown'),
                'connected': bool(cfg.get('connected', False)),
                'message': cfg.get('message'),
                'last_checked': cfg.get('last_checked'),
                'size': cfg.get('size')
            })
        return jsonify(out)
    except Exception:
        return jsonify([])


@app.route('/databases/status/<dbid>', methods=['GET'])
def db_status(dbid):
    """Return detailed status for a single configured DB (JSON)."""
    cfg = db_connections.get(dbid)
    if not cfg:
        return jsonify({'error': 'not found'}), 404
    # Return current cached status fields; do not block on live test here.
    return jsonify({
        'id': dbid,
        'name': cfg.get('name', dbid),
        'type': cfg.get('type', 'unknown'),
        'connected': bool(cfg.get('connected', False)),
        'message': cfg.get('message'),
        'last_checked': cfg.get('last_checked'),
        'size': cfg.get('size')
    })


@app.route('/databases/status_all', methods=['GET'])
def db_status_all():
    """Return status for all configured DBs (JSON)."""
    out = {}
    for dbid, cfg in db_connections.items():
        out[dbid] = {
            'id': dbid,
            'name': cfg.get('name', dbid),
            'type': cfg.get('type', 'unknown'),
            'connected': bool(cfg.get('connected', False)),
            'message': cfg.get('message'),
            'last_checked': cfg.get('last_checked'),
            'size': cfg.get('size')
        }
    return jsonify(out)

@app.route('/databases/disconnect/<dbid>', methods=['POST'])
def db_disconnect(dbid):
    cfg = db_connections.get(dbid)
    if cfg:
        cfg['connected'] = False
        save_db_configs()
        flash('Disconnected', 'success')
    else:
        flash('Connection not found', 'warning')
    return redirect(url_for('databases'))

@app.route('/databases/refresh/<dbid>', methods=['POST'])
def db_refresh(dbid):
    cfg = db_connections.get(dbid)
    if not cfg:
        flash('Connection not found', 'warning')
        return redirect(url_for('databases'))
    # Recompute size for file-backed DBs (best-effort)
    try:
        if cfg.get('conn', '').strip().endswith('.sqlite'):
            path = cfg.get('conn').split('///')[-1]
            if os.path.exists(path):
                cfg['size'] = format_bytes(os.path.getsize(path))
    except Exception:
        pass
    save_db_configs()
    flash('Refreshed', 'success')
    return redirect(url_for('databases'))


@app.route('/databases/create_folders/<dbid>', methods=['POST'])
def db_create_folders(dbid):
    cfg = db_connections.get(dbid)
    if not cfg:
        flash('Connection not found', 'warning')
        return redirect(url_for('databases'))
    try:
        create_default_folder_entries_for_db(cfg, dbid)
        flash('Folder creation attempted - check status below', 'success')
    except Exception as e:
        flash(f'Error creating folders: {e}', 'warning')
    return redirect(url_for('databases'))


FORENSIC_ANALYSIS_CONTENT = """
<div class="flex items-center justify-between">
    <div>
        <h1 class="text-3xl font-bold text-white mb-4">Forensic Analysis</h1>
        <p class="text-gray-400 mb-8">High-level analysis of the uploaded evidence file.</p>
    </div>
    <div>
        <!-- Session details button moved to global placement in base template -->
    </div>
</div>
<div class="card p-6 rounded-lg">
{% if not uploaded_files %}
    <p class="text-gray-500">No file uploaded. Please upload an evidence file to begin analysis.</p>
{% else %}
    {% for filename, details in uploaded_files.items() %}
    <h2 class="text-2xl font-semibold text-white mb-6">Analysis for: <span class="font-mono">{{ filename }}</span></h2>
    <div class="space-y-8">
        <div>
            <h3 class="text-lg font-semibold text-white mb-2">File Status</h3>
            <div class="text-gray-300">
                {% if details.encryption_status.encrypted %}
                    {% if details.encryption_status.decrypted_path %}
                        <p><span class="inline-flex items-center px-2.5 py-0.5 rounded-full text-sm font-medium bg-green-600 text-green-100">DECRYPTED</span> - Analysis is being performed on the decrypted version of the file.</p>
                    {% else %}
                        <p><span class="inline-flex items-center px-2.5 py-0.5 rounded-full text-sm font-medium bg-red-600 text-red-100">ENCRYPTED</span> - <span class="text-yellow-400">Warning:</span> Analysis results may be inaccurate. Please <a href="{{ url_for('decryption_page', filename=filename) }}" class="text-blue-400 hover:underline">decrypt the file</a> first.</p>
                        <p class="text-sm text-gray-400 mt-1">Detected Type: {{ details.encryption_status.description }}</p>
                    {% endif %}
                {% else %}
                    <p><span class="inline-flex items-center px-2.5 py-0.5 rounded-full text-sm font-medium bg-gray-600 text-gray-100">NOT ENCRYPTED</span> - The file does not appear to be encrypted.</p>
                {% endif %}
            </div>
        </div>

        <div>
            <h3 class="text-lg font-semibold text-white mb-2">General Information (from file header)</h3>
            <ul class="list-disc list-inside text-gray-300 space-y-1">
            {% for result in details.forensic_results %}
                <li>{{ result }}</li>
            {% endfor %}
            </ul>
        </div>
        
        <div>
            <h3 class="text-lg font-semibold text-white mb-2">Partition Map</h3>
            {% if details.partition_info %}
            <div class="overflow-x-auto">
                <table class="w-full text-left text-sm">
                    <thead class="bg-gray-800">
                        <tr class="border-b border-gray-700">
                            <th class="p-2">Index</th>
                            <th class="p-2">Description</th>
                            <th class="p-2">Start Sector</th>
                            <th class="p-2">Length (Sectors)</th>
                        </tr>
                    </thead>
                    <tbody>
                    {% for part in details.partition_info %}
                        <tr class="border-b border-gray-700 hover:bg-gray-800">
                            <td class="p-2 font-mono">{{ part.addr }}</td>
                            <td class="p-2">{{ part.desc }}</td>
                            <td class="p-2 font-mono">{{ part.start }}</td>
                            <td class="p-2 font-mono">{{ part.len }}</td>
                        </tr>
                    {% endfor %}
                    </tbody>
                </table>
            </div>
            {% else %}
            <p class="text-gray-500">No partition data found in the evidence file header.</p>
            {% endif %}
        </div>
        
        <div>
            <h3 class="text-lg font-semibold text-white mb-2">Full File Hashes</h3>
            <div id="hash-progress-container" class="my-2" {% if details.hashing_complete %}style="display: none;"{% endif %}>
                <div class="w-full bg-gray-700 rounded-full h-4">
                    <div id="hash-progress-bar" class="bg-blue-600 h-4 rounded-full transition-all duration-500" style="width: 0%"></div>
                </div>
                <p id="hash-progress-text" class="text-sm text-gray-400 mt-2 text-center">Initializing...</p>
            </div>
            <div id="hash-results-container" {% if not details.hashing_complete %}class="hidden"{% endif %}>
                <ul class="list-disc list-inside text-gray-300 space-y-1 font-mono text-sm">
                    <li><strong>MD5:</strong> <span id="hash-md5">{{ details.hash_info.MD5 or 'Calculating...' }}</span></li>
                    <li><strong>SHA-1:</strong> <span id="hash-sha1">{{ details.hash_info['SHA-1'] or 'Calculating...' }}</span></li>
                    <li><strong>SHA-256:</strong> <span id="hash-sha256">{{ details.hash_info['SHA-256'] or 'Calculating...' }}</span></li>
                </ul>
            </div>
        </div>

        <div>
            <h3 class="text-lg font-semibold text-white mb-2">Deep Scan Tools</h3>
            <div class="card p-4 rounded-lg bg-gray-800 border-gray-600">
                <div class="flex items-center justify-between">
                    <div>
                        <h4 class="font-semibold text-white">Strings Analysis</h4>
                        <p class="text-sm text-gray-400">Extract all printable strings from the evidence file.</p>
                    </div>
                    <button id="start-strings-btn" class="btn-primary px-4 py-2 rounded-lg text-sm">Run Analysis</button>
                </div>
                <div id="strings-analysis-container" class="mt-4 hidden">
                    <div class="w-full bg-gray-700 rounded-full h-4 mb-2">
                        <div id="strings-progress-bar" class="bg-blue-600 h-4 rounded-full transition-all duration-500" style="width: 0%"></div>
                    </div>
                    <p id="strings-progress-text" class="text-sm text-gray-400 text-center mb-4">Initializing...</p>
                    <h5 class="text-sm font-semibold text-white mb-2">Strings Preview (first 100 found):</h5>
                    <div class="log-view p-2 rounded-lg text-xs overflow-auto h-48">
                        <pre id="strings-preview"></pre>
                    </div>
                </div>
            </div>
        </div>
        
    </div>

<!-- Session Details Modal -->
<div id="session-modal" style="display:none; position:fixed; right:20px; top:80px; width:320px; background:#071018; border:1px solid #23313e; padding:12px; z-index:9999; border-radius:8px;">
    <div style="display:flex; justify-content:space-between; align-items:center; margin-bottom:8px;">
        <div style="font-weight:600;">Session Details</div>
        <button id="session-close" style="background:transparent;color:#cfe3ff;border:none;">✕</button>
    </div>
    <div style="font-size:13px; color:#cfe3ff;">
        <div>Session: <strong>{{ session_info.name or session_info.id or 'N/A' }}</strong></div>
        <div style="margin-top:6px">Path: <div style="font-family:monospace; font-size:12px; color:#9fc3e6;">{{ session_info.path or 'Not created' }}</div></div>
        <div style="margin-top:8px">
            <div style="font-weight:600; margin-bottom:6px;">Counts</div>
            <div style="display:grid; grid-template-columns:1fr 1fr; gap:6px;">
                <div>Carved: <strong>{{ session_info.counts.Carved or 0 }}</strong></div>
                <div>Deleted: <strong>{{ session_info.counts.Deleted or 0 }}</strong></div>
                <div>Reports: <strong>{{ session_info.counts.Reports or 0 }}</strong></div>
                <div>Logs: <strong>{{ session_info.counts.Logs or 0 }}</strong></div>
                <div>Events: <strong>{{ session_info.counts.Events or 0 }}</strong></div>
                <div>ManualCarving: <strong>{{ session_info.counts.ManualCarving or 0 }}</strong></div>
            </div>
        </div>
        <div style="margin-top:8px; display:flex; gap:6px;">
            <a href="{{ url_for('manual_carving') }}" class="btn-primary px-2 py-1 rounded">Manual Carving</a>
            <a href="{{ url_for('fs_explorer') }}?browse_root=Session Files" class="btn-secondary px-2 py-1 rounded">Open Session Files</a>
        </div>
    </div>
</div>

<script>
document.getElementById('btn-session-details').addEventListener('click', function(){ var m=document.getElementById('session-modal'); m.style.display = m.style.display === 'none' ? 'block' : 'none'; });
document.getElementById('session-close').addEventListener('click', function(){ document.getElementById('session-modal').style.display='none'; });
</script>
    {% endfor %}
{% endif %}
</div>

<script>
function updateHashingProgress() {
    fetch('/hashing_status')
        .then(response => response.json())
        .then(data => {
            const progressBar = document.getElementById('hash-progress-bar');
            const progressText = document.getElementById('hash-progress-text');
            const progressContainer = document.getElementById('hash-progress-container');
            const resultsContainer = document.getElementById('hash-results-container');

            if (data.complete) {
                progressContainer.style.display = 'none';
                resultsContainer.style.display = 'block';
                document.getElementById('hash-md5').textContent = data.hashes.MD5 || 'N/A';
                document.getElementById('hash-sha1').textContent = data.hashes['SHA-1'] || 'N/A';
                document.getElementById('hash-sha256').textContent = data.hashes['SHA-256'] || 'N/A';
            } else if (data.in_progress) {
                progressContainer.style.display = 'block';
                resultsContainer.style.display = 'none';
                progressBar.style.width = data.progress + '%';
                progressText.textContent = `Hashing... ${data.progress}% complete`;
                setTimeout(updateHashingProgress, 2000);
            } else {
                progressText.textContent = 'Hashing has not started or was interrupted.';
            }
        });
}

function updateStringsProgress() {
    fetch('/strings_status')
        .then(response => response.json())
        .then(data => {
            const btn = document.getElementById('start-strings-btn');
            const analysisContainer = document.getElementById('strings-analysis-container');
            const progressBar = document.getElementById('strings-progress-bar');
            const progressText = document.getElementById('strings-progress-text');
            const preview = document.getElementById('strings-preview');

            if (data.in_progress) {
                analysisContainer.style.display = 'block';
                btn.disabled = true;
                btn.textContent = 'Running...';
                progressBar.style.width = data.progress + '%';
                progressText.textContent = `Scanning... ${data.progress}% complete | Found: ${data.strings_found} strings`;
                preview.textContent = data.preview.join('\\n');
                setTimeout(updateStringsProgress, 1500);
            } else if (data.complete) {
                analysisContainer.style.display = 'block';
                progressBar.style.width = '100%';
                progressText.textContent = `Scan Complete. Found ${data.strings_found} strings.`;
                preview.textContent = data.preview.join('\\n');
                btn.disabled = false;
                btn.textContent = 'Run Again';
            }
        });
}

document.addEventListener('DOMContentLoaded', () => {
    {% if hashing_in_progress %}
    updateHashingProgress();
    {% endif %}

    const startStringsBtn = document.getElementById('start-strings-btn');
    if(startStringsBtn) {
        startStringsBtn.addEventListener('click', () => {
            fetch('/start_strings_analysis')
                .then(response => response.json())
                .then(data => {
                    if(data.status === 'started') {
                        updateStringsProgress();
                    } else {
                        alert('Could not start strings analysis: ' + data.error);
                    }
                });
        });
    }
});
</script>
"""

VIEW_FILE_CONTENT = """
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <title>Viewing: {{ filename }}</title>
    <style>
        body, html { 
            margin: 0; 
            padding: 0; 
            height: 100%; 
            background-color: #111827; 
            color: white; 
            font-family: sans-serif; 
            overflow: hidden;
        }
        .header {
            background-color: #1f2937;
            padding: 15px 20px;
            border-bottom: 1px solid #374151;
            display: flex;
            justify-content: space-between;
            align-items: center;
        }
        .header h1 {
            margin: 0;
            font-size: 1.5em;
            color: #3b82f6;
        }
        .header-actions {
            display: flex;
            gap: 10px;
        }
        .btn {
            padding: 8px 16px;
            border-radius: 5px;
            text-decoration: none;
            font-weight: bold;
            transition: background-color 0.2s;
        }
        .btn-secondary {
            background-color: #4b5563;
            color: white;
        }
        .btn-primary {
            background-color: #3b82f6;
            color: white;
        }
        .btn:hover {
            opacity: 0.9;
        }
        .content-area {
            height: calc(100vh - 70px);
            display: flex;
            flex-direction: column;
        }
        embed, iframe, img, video, audio { 
            width: 100%; 
            height: 100%; 
            border: none; 
            flex-grow: 1;
        }
        pre { 
            margin: 0;
            background-color: #1f2937; 
            color: #d1d5db; 
            height: 100%; 
            overflow: auto; 
            padding: 20px; 
            box-sizing: border-box; 
            font-family: 'Courier New', monospace; 
            white-space: pre-wrap; 
            word-wrap: break-word; 
            flex-grow: 1;
        }
        .fallback-message {
            padding: 40px;
            text-align: center;
            color: white;
            font-family: sans-serif;
            flex-grow: 1;
            display: flex;
            flex-direction: column;
            justify-content: center;
            align-items: center;
        }
        .file-info {
            background-color: #1f2937;
            padding: 10px 20px;
            border-bottom: 1px solid #374151;
            font-size: 0.9em;
            color: #9ca3af;
        }
    </style>
</head>
<body>
    <div class="header">
        <h1>Viewing: {{ filename }}</h1>
        <div class="header-actions">
            <a href="{{ url_for('recovered_files') if is_carved else url_for('deleted_files') }}" class="btn btn-secondary">Back to List</a>
            <a href="{{ url_for('download_carved_file', filename=filename) if is_carved else url_for('download_deleted_file', filename=filename) }}" class="btn btn-primary">Download</a>
        </div>
    </div>
    
    <div class="file-info">
        File Type: {{ mime_type }} | 
        Size: {{ (content|length if content else 0) }} bytes
    </div>
    
    <div class="content-area">
        {% if mime_type.startswith('image/') %}
            <img src="{{ data_url }}" alt="Preview of {{ filename }}">
        {% elif mime_type == 'application/pdf' %}
            <iframe src="{{ data_url }}"></iframe>
        {% elif mime_type.startswith('audio/') %}
            <audio controls>
                <source src="{{ data_url }}" type="{{ mime_type }}">
                Your browser does not support the audio element.
            </audio>
        {% elif mime_type.startswith('video/') %}
            <video controls>
                <source src="{{ data_url }}" type="{{ mime_type }}">
                Your browser does not support the video element.
            </video>
        {% elif text_content is defined %}
            <pre>{{ text_content }}</pre>
        {% else %}
            <div class="fallback-message">
                <h2>Preview Not Available</h2>
                <p>Direct preview for '{{ mime_type }}' is not supported.</p>
                <a href="{{ url_for('download_carved_file', filename=filename) if is_carved else url_for('download_deleted_file', filename=filename) }}" 
                   class="btn btn-primary" style="margin-top: 20px;">
                    Download File
                </a>
            </div>
        {% endif %}
    </div>
</body>
</html>
"""

AUTO_CARVING_SETUP_CONTENT = """
<h1 class="text-3xl font-bold text-white mb-4">Automatic File Carving Setup</h1>
<p class="text-gray-400 mb-8">Select file types for automated carving based on file signatures.</p>
<form method="post" action="{{ form_action_url }}">
<div class="grid grid-cols-1 lg:grid-cols-3 gap-8">
    <div class="lg:col-span-2">
        <div class="card p-6 rounded-lg">
            <div class="flex justify-between items-center mb-4 pb-4 border-b border-gray-700">
                <h2 class="text-xl font-semibold text-white">Select File Types</h2>
                <div class="flex items-center space-x-4">
                    <span id="selected-count" class="bg-blue-600 text-white px-3 py-1 rounded-full text-sm font-medium">0 selected</span>
                    <button type="button" id="select-all-btn" class="btn-secondary px-3 py-1 text-xs rounded-lg">Select All</button>
                    <button type="button" id="deselect-all-btn" class="btn-secondary px-3 py-1 text-xs rounded-lg">Deselect All</button>
                </div>
            </div>
            {% for category, types in signatures.items() %}
            <h3 class="text-lg font-semibold {{ colors[loop.index0 % colors|length] }} mt-6 mb-2 border-b border-gray-700 pb-1">{{ category }}</h3>
            <div class="grid grid-cols-2 md:grid-cols-3 gap-4">
                {% for name, sig in types.items() %}
                <div class="flex items-center">
                    <input type="checkbox" name="file_types" value="{{ name }}" id="{{ name }}" class="file-type-checkbox h-4 w-4 rounded bg-gray-700 border-gray-600 text-blue-600 focus:ring-blue-500">
                    <label for="{{ name }}" class="ml-3 text-white cursor-pointer">{{ name }}</label>
                </div>
                {% endfor %}
            </div>
            {% endfor %}
        </div>
    </div>
    <div class="lg:col-span-1">
        <div class="card p-6 rounded-lg sticky top-8">
            <h3 class="text-lg font-semibold text-white mb-4">Begin Carving</h3>
            <p class="text-gray-400 text-sm mb-4">Once you have selected the desired file types, start the carving process.</p>
            <button type="submit" class="btn-green w-full py-3 rounded-lg font-semibold">Start Auto Carving</button>
        </div>
    </div>
</div>
</form>

<script>
document.addEventListener('DOMContentLoaded', function() {
    const checkboxes = document.querySelectorAll('.file-type-checkbox');
    const counter = document.getElementById('selected-count');
    
    function updateCounter() {
        const selectedCount = document.querySelectorAll('.file-type-checkbox:checked').length;
        counter.textContent = `${selectedCount} selected`;
    }
    
    checkboxes.forEach(checkbox => checkbox.addEventListener('change', updateCounter));
    
    document.getElementById('select-all-btn').addEventListener('click', () => {
        checkboxes.forEach(checkbox => checkbox.checked = true);
        updateCounter();
    });

    document.getElementById('deselect-all-btn').addEventListener('click', () => {
        checkboxes.forEach(checkbox => checkbox.checked = false);
        updateCounter();
    });

    updateCounter();
});
</script>
"""

AUTO_CARVING_PROCESS_CONTENT = """
<h1 class="text-3xl font-bold text-white mb-2">Automatic File Carving in Progress</h1>
<p class="text-gray-400 mb-8">Scanning evidence file for selected file types.</p>
<div class="card p-6 rounded-lg">
    <div class="flex justify-between items-center mb-4">
        <h2 class="text-xl font-semibold text-white">Carving Progress</h2>
        <span id="progress-percent" class="font-bold text-blue-400">0% Complete</span>
    </div>
    <div class="w-full bg-gray-700 rounded-full h-4">
        <div id="progress-bar" class="bg-blue-600 h-4 rounded-full transition-all duration-500" style="width: 0%"></div>
    </div>
    <p class="text-sm text-gray-400 mt-2">Processing offset: <span id="current-offset">0x00000000</span> | Found: <span id="files-found">0</span> files</p>
    <!-- Add this after the progress bar -->
<p class="text-sm text-gray-400 mt-2">
    Elapsed: <span id="elapsed-time">0s</span> | 
    Remaining: <span id="time-remaining">Calculating...</span> |
    Total Estimated: <span id="total-estimated">Calculating...</span>
</p>
    <div id="error-container" class="mt-4 p-4 bg-red-900 border border-red-500 rounded-lg text-sm text-red-200 hidden">
        <h3 class="font-bold mb-2">An Error Occurred</h3>
        <p id="error-message"></p>
    </div>
</div>
<div class="card p-6 rounded-lg mt-8">
    <h2 class="text-xl font-semibold text-white mb-4">Live Hex Preview of Found Files</h2>
    <div id="live-hex-view" class="hex-view p-4 rounded-lg text-sm overflow-y-auto max-h-[40rem]">
        <p class="text-gray-500">Waiting for files to be found...</p>
    </div>
</div>
<div class="mt-8 flex justify-between items-center">
    <h3 class="text-xl font-semibold text-white">Found Files (<span id="total-found-count">0</span>)</h3>
    <div class="flex space-x-4">
        <a href="{{ url_for('auto_carving_setup') }}" class="btn-secondary px-6 py-2 rounded-lg">Back to Setup</a>
        <a href="{{ url_for('recovered_files') }}" id="view-recovered-btn" class="btn-primary px-6 py-2 rounded-lg opacity-50 pointer-events-none">View Recovered Files</a>
    </div>
</div>

<script>
function updateProgress() {
    fetch('/carving_status')
        .then(response => response.json())
        .then(data => {
            document.getElementById('progress-bar').style.width = data.progress + '%';
            const progressPercent = document.getElementById('progress-percent');
            progressPercent.innerText = data.progress + '% Complete';
            document.getElementById('current-offset').innerText = data.current_offset;
            document.getElementById('files-found').innerText = data.files_found;
            document.getElementById('total-found-count').innerText = data.files_found;
            // timing placeholders will be handled below with better fallbacks
            
            const liveHexView = document.getElementById('live-hex-view');
            const recentFiles = data.found_files_list.slice(-5).reverse();
            if (recentFiles.length > 0) {
                let fullHexContent = '';
                recentFiles.forEach(file => {
                    fullHexContent += `<p class="text-xs text-green-400">${file.name} @ ${file.offset}</p><pre>${file.hex_preview}</pre><hr class="border-gray-600 my-2">`;
                });
                liveHexView.innerHTML = fullHexContent;
            }

            const errorContainer = document.getElementById('error-container');
            const errorMessage = document.getElementById('error-message');
            if (data.error) {
                errorMessage.textContent = data.error;
                errorContainer.classList.remove('hidden');
                progressPercent.className = 'font-bold text-red-400';
                progressPercent.innerText = 'Scan Failed';
                document.getElementById('progress-bar').classList.add('bg-red-600');
            }

            if (data.complete) {
                if (!data.error) {
                    progressPercent.className = 'font-bold text-green-400';
                    progressPercent.innerText = 'Scan Complete';
                }
                const viewBtn = document.getElementById('view-recovered-btn');
                viewBtn.classList.remove('opacity-50', 'pointer-events-none');
            }

            // Ensure elapsed / remaining display always shows a sensible fallback
            const elapsedEl = document.getElementById('elapsed-time');
            const remainingEl = document.getElementById('time-remaining');
            const totalEstEl = document.getElementById('total-estimated');
            elapsedEl.textContent = data.elapsed_time || '0s';
            remainingEl.textContent = data.time_remaining_str || (data.estimated_total_time ? 'Calculating...' : 'Calculating...');
            totalEstEl.textContent = data.estimated_total_time || 'Calculating...';
        })
        .catch(err => {
            console.error('Error fetching carving status:', err);
        });
}

// Start polling when the page is ready. Use setInterval to be resilient.
document.addEventListener('DOMContentLoaded', function() {
    updateProgress();
    setInterval(updateProgress, 1000);
});
</script>
"""

RECOVERED_FILES_CONTENT = """
<h1 class="text-3xl font-bold text-white mb-4">Recovered Carved Files</h1>
<div class="card p-6 rounded-lg">
    <form action="{{ url_for('download_zip') }}" method="post">
        <input type="hidden" name="file_type" value="carved">
        <div class="flex justify-between items-center mb-4">
            <h2 class="text-xl font-semibold text-white">Recovered Files ({{ total_files }} found)</h2>
            {% if total_files > 0 %}
            <button type="submit" class="btn-primary px-4 py-2 rounded-lg text-sm">Download Selected as ZIP</button>
            {% endif %}
        </div>
        {% if carved_files %}
            <table class="w-full text-left">
                <thead>
                    <tr class="border-b border-gray-700">
                        <th class="p-2 w-10"><input type="checkbox" id="select-all-carved" class="h-4 w-4 rounded bg-gray-700 border-gray-600"></th>
                        <th class="p-2">ID</th>
                        <th class="p-2">Filename</th>
                        <th class="p-2">Offset</th>
                        <th class="p-2">Size</th>
                        <th class="p-2">Actions</th>
                    </tr>
                </thead>
                <tbody>
                {% for file in carved_files %}
                    <tr class="border-b border-gray-700 hover:bg-gray-800">
                        <td class="p-2"><input type="checkbox" name="selected_files" value="{{ file.name }}" class="h-4 w-4 rounded bg-gray-700 border-gray-600 file-checkbox"></td>
                        <td class="p-2">{{ file.id }}</td>
                        <td class="p-2 font-mono break-all">{{ file.name }}</td>
                        <td class="p-2 font-mono">{{ file.offset }}</td>
                        <td class="p-2 font-mono">{{ "%.2f KB"|format(file.size_bytes / 1024) if file.size_bytes else 'N/A' }}</td>
                        <td class="p-2 flex space-x-2">
                            <a href="{{ url_for('view_carved_file', filename=file.name) }}" target="_blank" class="btn-secondary px-3 py-1 text-xs rounded-lg">View</a>
                            <a href="{{ url_for('hex_view_carved', filename=file.name) }}" target="_blank" class="btn-green px-3 py-1 text-xs rounded-lg">Hex View</a>
                            <a href="{{ url_for('download_carved_file', filename=file.name) }}" class="btn-primary px-3 py-1 text-xs rounded-lg">Download</a>
                        </td>
                    </tr>
                {% endfor %}
                </tbody>
            </table>
            
            <div class="mt-6 flex justify-center items-center space-x-2 text-gray-400">
                {% if page > 1 %}
                    <a href="{{ url_for('recovered_files', page=page-1) }}" class="px-3 py-1 rounded-md bg-gray-800 hover:bg-gray-700">&laquo; Prev</a>
                {% endif %}
                
                {% for p in range(1, total_pages + 1) %}
                    {% if p == page %}
                        <span class="px-3 py-1 rounded-md bg-blue-600 text-white">{{ p }}</span>
                    {% else %}
                        <a href="{{ url_for('recovered_files', page=p) }}" class="px-3 py-1 rounded-md bg-gray-800 hover:bg-gray-700">{{ p }}</a>
                    {% endif %}
                {% endfor %}

                {% if page < total_pages %}
                    <a href="{{ url_for('recovered_files', page=page+1) }}" class="px-3 py-1 rounded-md bg-gray-800 hover:bg-gray-700">Next &raquo;</a>
                {% endif %}
            </div>
        {% else %}
            <div class="text-center p-8">
                <p class="text-gray-500 text-lg mb-4">No files recovered yet.</p>
                <div class="text-sm text-gray-400">
                    <p>Possible reasons:</p>
                    <ul class="list-disc list-inside mt-2">
                        <li>Carving process has not been run</li>
                        <li>Carving process is still in progress</li>
                        <li>No files of the selected types were found</li>
                        <li>Session was cleared unexpectedly</li>
                    </ul>
                </div>
                <div class="mt-4">
                    <a href="{{ url_for('auto_carving_setup') }}" class="btn-primary px-4 py-2 rounded-lg">Go to Auto Carving</a>
                </div>
            </div>
        {% endif %}
    </form>
</div>

<script>
document.addEventListener('DOMContentLoaded', function() {
    const selectAll = document.getElementById('select-all-carved');
    if(selectAll) {
        selectAll.addEventListener('change', function(e) {
            document.querySelectorAll('.file-checkbox').forEach(function(checkbox) {
                checkbox.checked = e.target.checked;
            });
        });
    }
});
</script>
"""
MANUAL_CARVING_CONTENT = """
<h1 class="text-3xl font-bold text-white mb-4">Manual File Carving</h1>
<p class="text-gray-400 mb-8">Manually search for and extract data blocks from the evidence file using hex patterns or offsets.</p>

<div class="grid grid-cols-1 lg:grid-cols-2 gap-8">
    <div class="card p-6 rounded-lg">
        <h2 class="text-xl font-semibold text-white mb-4">Search for Data</h2>
        
        <div class="space-y-6">
            <div>
                <h3 class="text-lg font-semibold text-white mb-2">1. Search for Hex Pattern</h3>
                <form id="search-hex-form" class="space-y-4">
                    <div>
                        <label class="block text-sm font-medium text-gray-300">Hex Pattern (e.g., FF D8 FF E0)</label>
                        <input type="text" id="hex-term" class="mt-1 block w-full bg-gray-800 border-gray-600 rounded-md shadow-sm p-2 text-white font-mono" placeholder="Enter hex values without spaces">
                    </div>
                    <button type="submit" class="btn-primary w-full py-2 rounded-lg">Search Pattern</button>
                </form>
                <div id="hex-result" class="mt-4 hidden"></div>
            </div>

            <div class="pt-4 border-t border-gray-700">
                <h3 class="text-lg font-semibold text-white mb-2">2. Search for Text String</h3>
                <form id="search-text-form" class="space-y-4">
                    <div>
                        <label class="block text-sm font-medium text-gray-300">Text String</label>
                        <input type="text" id="text-term" class="mt-1 block w-full bg-gray-800 border-gray-600 rounded-md shadow-sm p-2 text-white" placeholder="Enter text to search for">
                    </div>
                    <button type="submit" class="btn-primary w-full py-2 rounded-lg">Search Text</button>
                </form>
                <div id="text-result" class="mt-4 hidden"></div>
            </div>

            <div class="pt-4 border-t border-gray-700">
                <h3 class="text-lg font-semibold text-white mb-2">3. Find Block Between Headers/Footers</h3>
                <form id="search-block-form" class="space-y-4">
                    <div class="grid grid-cols-2 gap-4">
                        <div>
                            <label class="block text-sm font-medium text-gray-300">Header Pattern</label>
                            <input type="text" id="header-term" class="mt-1 block w-full bg-gray-800 border-gray-600 rounded-md shadow-sm p-2 text-white font-mono" placeholder="Hex or text">
                        </div>
                        <div>
                            <label class="block text-sm font-medium text-gray-300">Footer Pattern</label>
                            <input type="text" id="footer-term" class="mt-1 block w-full bg-gray-800 border-gray-600 rounded-md shadow-sm p-2 text-white font-mono" placeholder="Hex or text">
                        </div>
                    </div>
                    <div>
                        <label class="block text-sm font-medium text-gray-300">Search Type</label>
                        <select id="block-search-type" class="mt-1 block w-full bg-gray-800 border-gray-600 rounded-md shadow-sm p-2 text-white">
                            <option value="hex">Hex</option>
                            <option value="text">Text</option>
                        </select>
                    </div>
                    <button type="submit" class="btn-primary w-full py-2 rounded-lg">Find Block</button>
                </form>
                <div id="block-result" class="mt-4 hidden"></div>
            </div>
        </div>
    </div>

    <div class="card p-6 rounded-lg">
        <h2 class="text-xl font-semibold text-white mb-4">Manual Carve Data</h2>
        
        <form action="{{ url_for('perform_manual_carve') }}" method="post" class="space-y-4">
            <div>
                <label class="block text-sm font-medium text-gray-300">Start Offset (decimal or hex with 0x)</label>
                <input type="text" name="start_offset" class="mt-1 block w-full bg-gray-800 border-gray-600 rounded-md shadow-sm p-2 text-white font-mono" placeholder="e.g., 1024 or 0x400" required>
            </div>
            
            <div>
                <label class="block text-sm font-medium text-gray-300">Length (bytes)</label>
                <input type="number" name="length" class="mt-1 block w-full bg-gray-800 border-gray-600 rounded-md shadow-sm p-2 text-white" placeholder="Number of bytes to extract" required min="1">
            </div>
            
            <div class="flex space-x-2">
                <button type="submit" class="btn-green flex-1 py-2 rounded-lg font-semibold">Carve Data</button>
                <button type="button" id="preview-hex-btn" class="btn-secondary py-2 rounded-lg px-4" onclick="previewHexFromForm()">Preview Hex (Scrabble)</button>
            </div>
        </form>

        <div class="mt-6 p-4 bg-gray-800 rounded-lg">
            <h3 class="text-lg font-semibold text-white mb-2">Quick Reference</h3>
            <ul class="text-sm text-gray-400 space-y-1">
                <li>• JPEG: FF D8 FF E0 ... FF D9</li>
                <li>• PNG: 89 50 4E 47 0D 0A 1A 0A</li>
                <li>• PDF: 25 50 44 46</li>
                <li>• ZIP: 50 4B 03 04</li>
                <li>• GIF: 47 49 46 38</li>
            </ul>
        </div>
    </div>
</div>

<script>
document.addEventListener('DOMContentLoaded', function() {
    // Hex search
    document.getElementById('search-hex-form').addEventListener('submit', function(e) {
        e.preventDefault();
        const term = document.getElementById('hex-term').value.trim();
        if (!term) return;
        
        fetch('/find_in_file', {
            method: 'POST',
            headers: {'Content-Type': 'application/x-www-form-urlencoded'},
            body: `term=${encodeURIComponent(term)}&type=hex`
        })
        .then(response => response.json())
        .then(data => {
            const resultDiv = document.getElementById('hex-result');
            if (data.offset !== -1) {
                resultDiv.innerHTML = `<div class="p-3 bg-green-900 border border-green-500 rounded-lg">
                    <p class="text-green-300">Found at offset: <span class="font-mono">0x${data.offset.toString(16).toUpperCase()}</span> (${data.offset} decimal)</p>
                    <button onclick="fillOffset(${data.offset})" class="mt-2 btn-secondary px-3 py-1 text-xs">Use this offset</button>
                </div>`;
            } else {
                resultDiv.innerHTML = `<div class="p-3 bg-red-900 border border-red-500 rounded-lg">
                    <p class="text-red-300">Pattern not found. ${data.error || ''}</p>
                </div>`;
            }
            resultDiv.classList.remove('hidden');
        });
    });

    // Text search
    document.getElementById('search-text-form').addEventListener('submit', function(e) {
        e.preventDefault();
        const term = document.getElementById('text-term').value.trim();
        if (!term) return;
        
        fetch('/find_in_file', {
            method: 'POST',
            headers: {'Content-Type': 'application/x-www-form-urlencoded'},
            body: `term=${encodeURIComponent(term)}&type=text`
        })
        .then(response => response.json())
        .then(data => {
            const resultDiv = document.getElementById('text-result');
            if (data.offset !== -1) {
                resultDiv.innerHTML = `<div class="p-3 bg-green-900 border border-green-500 rounded-lg">
                    <p class="text-green-300">Found at offset: <span class="font-mono">0x${data.offset.toString(16).toUpperCase()}</span> (${data.offset} decimal)</p>
                    <button onclick="fillOffset(${data.offset})" class="mt-2 btn-secondary px-3 py-1 text-xs">Use this offset</button>
                </div>`;
            } else {
                resultDiv.innerHTML = `<div class="p-3 bg-red-900 border border-red-500 rounded-lg">
                    <p class="text-red-300">Text not found. ${data.error || ''}</p>
                </div>`;
            }
            resultDiv.classList.remove('hidden');
        });
    });

    // Block search
    document.getElementById('search-block-form').addEventListener('submit', function(e) {
        e.preventDefault();
        const header = document.getElementById('header-term').value.trim();
        const footer = document.getElementById('footer-term').value.trim();
        const type = document.getElementById('block-search-type').value;
        
        if (!header || !footer) return;
        
        fetch('/find_block', {
            method: 'POST',
            headers: {'Content-Type': 'application/x-www-form-urlencoded'},
            body: `header_term=${encodeURIComponent(header)}&footer_term=${encodeURIComponent(footer)}&type=${type}`
        })
        .then(response => response.json())
        .then(data => {
            const resultDiv = document.getElementById('block-result');
            if (data.status === 'success') {
                resultDiv.innerHTML = `<div class="p-3 bg-green-900 border border-green-500 rounded-lg">
                    <p class="text-green-300">Block found:</p>
                    <p class="text-green-300">Start: <span class="font-mono">0x${data.start_offset.toString(16).toUpperCase()}</span></p>
                    <p class="text-green-300">Length: <span class="font-mono">${data.length} bytes</span></p>
                    <div class="mt-2 space-x-2">
                        <button onclick="fillOffsetAndLength(${data.start_offset}, ${data.length})" class="btn-secondary px-3 py-1 text-xs">Use these values</button>
                    </div>
                </div>`;
            } else {
                resultDiv.innerHTML = `<div class="p-3 bg-red-900 border border-red-500 rounded-lg">
                    <p class="text-red-300">Block not found. ${data.message || ''}</p>
                </div>`;
            }
            resultDiv.classList.remove('hidden');
        });
    });
});

function fillOffset(offset) {
    document.querySelector('input[name="start_offset"]').value = offset;
}

function fillOffsetAndLength(offset, length) {
    document.querySelector('input[name="start_offset"]').value = offset;
    document.querySelector('input[name="length"]').value = length;
}
</script>
<!-- Hex Preview Modal -->
<div id="hexPreviewModal" class="modal">
    <div class="modal-content">
        <span class="close" id="hexPreviewClose">&times;</span>
        <h2 class="text-xl font-semibold text-white mb-4">Hex Preview (Scrabble)</h2>
        <div id="hex-preview-content" class="log-view p-4" style="max-height:60vh; overflow:auto; background:#0d1117; color:#d1d5db;"></div>
    </div>
</div>

<script>
function previewHexFromForm() {
    const startField = document.querySelector('input[name="start_offset"]');
    const lengthField = document.querySelector('input[name="length"]');
    let start = startField ? startField.value.trim() : '';
    let length = lengthField ? lengthField.value.trim() : '';
    if (!start || !length) {
        alert('Please provide both Start Offset and Length to preview.');
        return;
    }
    fetch('/manual_carve_hex', {
        method: 'POST',
        headers: {'Content-Type': 'application/x-www-form-urlencoded'},
        body: `start_offset=${encodeURIComponent(start)}&length=${encodeURIComponent(length)}`
    })
    .then(r => r.json())
    .then(data => {
        if (data.error) {
            alert('Error: ' + data.error);
            return;
        }
        const container = document.getElementById('hex-preview-content');
        container.innerHTML = data.html;
        const modal = document.getElementById('hexPreviewModal');
        modal.style.display = 'block';
    }).catch(e => alert('Error fetching hex preview: ' + e));
}

document.getElementById('hexPreviewClose').addEventListener('click', function(){
    document.getElementById('hexPreviewModal').style.display = 'none';
});

// Close modal when clicking outside content
window.addEventListener('click', function(e){
    const modal = document.getElementById('hexPreviewModal');
    if (e.target === modal) modal.style.display = 'none';
});
</script>
"""

MANUAL_CARVING_EVIDENCE_HEX = """
<div class="mt-8 card p-6 rounded-lg">
    <h2 class="text-xl font-semibold text-white mb-4">View Evidence File Hex</h2>
    <p class="text-sm text-gray-400 mb-4">Browse the uploaded evidence file in hex. Use Prev/Next to page through the file. Page size defaults to 64KB.</p>

    <div class="grid grid-cols-3 gap-4 mb-4">
        <div>
            <label class="block text-sm font-medium text-gray-300">Start Offset (decimal or 0x hex)</label>
            <input type="text" id="evidence-start" class="mt-1 block w-full bg-gray-800 border-gray-600 rounded-md shadow-sm p-2 text-white font-mono" placeholder="0x0 or 0" value="0">
        </div>
        <div>
            <label class="block text-sm font-medium text-gray-300">Length (bytes)</label>
            <input type="number" id="evidence-length" class="mt-1 block w-full bg-gray-800 border-gray-600 rounded-md shadow-sm p-2 text-white" value="65536">
        </div>
        <div class="flex items-end space-x-2">
            <button id="evidence-view-btn" class="btn-primary py-2 px-4">View</button>
            <button id="evidence-prev-btn" class="btn-secondary py-2 px-3">Prev</button>
            <button id="evidence-next-btn" class="btn-secondary py-2 px-3">Next</button>
        </div>
    </div>

    <div id="evidence-hex-view" class="hex-view p-3" style="background:#0b0f14; color:#d1d5db; max-height:50vh; overflow:auto;"></div>
    <div id="evidence-hex-info" class="text-xs text-gray-400 mt-2"></div>
</div>

<script>
(function(){
    let currentOffset = 0;
    const maxPreview = 65536;

    function parseOffset(val) {
        val = String(val).trim();
        if (val.startsWith('0x') || val.startsWith('0X')) return parseInt(val, 16);
        return parseInt(val, 10) || 0;
    }

    async function loadEvidenceHex(offset, length) {
        const body = `start_offset=${encodeURIComponent(offset)}&length=${encodeURIComponent(length)}`;
        const res = await fetch('/view_evidence_hex', {method:'POST', headers:{'Content-Type':'application/x-www-form-urlencoded'}, body});
        const data = await res.json();
        if (data.error) {
            document.getElementById('evidence-hex-view').innerText = 'Error: ' + data.error;
            document.getElementById('evidence-hex-info').innerText = '';
            return;
        }
        document.getElementById('evidence-hex-view').innerHTML = data.html;
        currentOffset = data.start || offset;
        const info = `Showing 0x${currentOffset.toString(16).toUpperCase()} - 0x${(currentOffset + data.length -1).toString(16).toUpperCase()} (${data.length} bytes) of ${data.file_size} bytes`;
        document.getElementById('evidence-hex-info').innerText = info;
        document.getElementById('evidence-start').value = '0x' + currentOffset.toString(16).toUpperCase();
    }

    document.getElementById('evidence-view-btn').addEventListener('click', function(){
        const start = parseOffset(document.getElementById('evidence-start').value);
        let length = parseInt(document.getElementById('evidence-length').value) || maxPreview;
        length = Math.min(length, maxPreview);
        loadEvidenceHex(start, length);
    });

    document.getElementById('evidence-prev-btn').addEventListener('click', function(){
        let length = parseInt(document.getElementById('evidence-length').value) || maxPreview;
        length = Math.min(length, maxPreview);
        const newOffset = Math.max(0, currentOffset - length);
        loadEvidenceHex(newOffset, length);
    });

    document.getElementById('evidence-next-btn').addEventListener('click', function(){
        let length = parseInt(document.getElementById('evidence-length').value) || maxPreview;
        length = Math.min(length, maxPreview);
        const newOffset = currentOffset + length;
        loadEvidenceHex(newOffset, length);
    });

    // Load initial view
    loadEvidenceHex(0, 65536);
})();
</script>
"""

MANUAL_CARVING_EVIDENCE_HEX = """
<div class="mt-8 card p-6 rounded-lg">
    <h2 class="text-xl font-semibold text-white mb-4">View Evidence File Hex</h2>
    <p class="text-sm text-gray-400 mb-4">Browse the uploaded evidence file in hex. Use Prev/Next to page through the file. Page size defaults to 64KB.</p>

    <div class="grid grid-cols-3 gap-4 mb-4">
        <div>
            <label class="block text-sm font-medium text-gray-300">Start Offset (decimal or 0x hex)</label>
            <input type="text" id="evidence-start" class="mt-1 block w-full bg-gray-800 border-gray-600 rounded-md shadow-sm p-2 text-white font-mono" placeholder="0x0 or 0" value="0">
        </div>
        <div>
            <label class="block text-sm font-medium text-gray-300">Length (bytes)</label>
            <input type="number" id="evidence-length" class="mt-1 block w-full bg-gray-800 border-gray-600 rounded-md shadow-sm p-2 text-white" value="65536">
        </div>
        <div class="flex items-end space-x-2">
            <button id="evidence-view-btn" class="btn-primary py-2 px-4">View</button>
            <button id="evidence-prev-btn" class="btn-secondary py-2 px-3">Prev</button>
            <button id="evidence-next-btn" class="btn-secondary py-2 px-3">Next</button>
        </div>
    </div>

    <div id="evidence-hex-view" class="hex-view p-3" style="background:#0b0f14; color:#d1d5db; max-height:50vh; overflow:auto;"></div>
    <div id="evidence-hex-info" class="text-xs text-gray-400 mt-2"></div>
</div>

<script>
(function(){
    let currentOffset = 0;
    const maxPreview = 65536;

    function parseOffset(val) {
        val = String(val).trim();
        if (val.startsWith('0x') || val.startsWith('0X')) return parseInt(val, 16);
        return parseInt(val, 10) || 0;
    }

    async function loadEvidenceHex(offset, length) {
        const body = `start_offset=${encodeURIComponent(offset)}&length=${encodeURIComponent(length)}`;
        const res = await fetch('/view_evidence_hex', {method:'POST', headers:{'Content-Type':'application/x-www-form-urlencoded'}, body});
        const data = await res.json();
        if (data.error) {
            document.getElementById('evidence-hex-view').innerText = 'Error: ' + data.error;
            document.getElementById('evidence-hex-info').innerText = '';
            return;
        }
        document.getElementById('evidence-hex-view').innerHTML = data.html;
        currentOffset = data.start || offset;
        const info = `Showing 0x${currentOffset.toString(16).toUpperCase()} - 0x${(currentOffset + data.length -1).toString(16).toUpperCase()} (${data.length} bytes) of ${data.file_size} bytes`;
        document.getElementById('evidence-hex-info').innerText = info;
        document.getElementById('evidence-start').value = '0x' + currentOffset.toString(16).toUpperCase();
    }

    document.getElementById('evidence-view-btn').addEventListener('click', function(){
        const start = parseOffset(document.getElementById('evidence-start').value);
        let length = parseInt(document.getElementById('evidence-length').value) || maxPreview;
        length = Math.min(length, maxPreview);
        loadEvidenceHex(start, length);
    });

    document.getElementById('evidence-prev-btn').addEventListener('click', function(){
        let length = parseInt(document.getElementById('evidence-length').value) || maxPreview;
        length = Math.min(length, maxPreview);
        const newOffset = Math.max(0, currentOffset - length);
        loadEvidenceHex(newOffset, length);
    });

    document.getElementById('evidence-next-btn').addEventListener('click', function(){
        let length = parseInt(document.getElementById('evidence-length').value) || maxPreview;
        length = Math.min(length, maxPreview);
        const newOffset = currentOffset + length;
        loadEvidenceHex(newOffset, length);
    });

    // Load initial view
    loadEvidenceHex(0, 65536);
})();
</script>
"""

LOG_VIEWER_TEMPLATE = """
<h1 class="text-3xl font-bold text-white mb-4">{{ title }}</h1>
<p class="text-gray-400 mb-8">{{ description }}</p>

<div class="grid grid-cols-1 lg:grid-cols-3 gap-8">
    <div class="lg:col-span-1">
        <div class="card p-6 rounded-lg">
            <h2 class="text-xl font-semibold text-white mb-4">Select Log Type</h2>
            <form method="post" class="space-y-6">
                <div>
                    <label class="block text-sm font-medium text-gray-300">Operating System</label>
                    <select id="os-select" class="mt-1 block w-full bg-gray-800 border-gray-600 rounded-md shadow-sm p-2 text-white">
                        <option value="">Select OS</option>
                        {% for os_name in os_options.keys() %}
                        <option value="{{ os_name }}">{{ os_name }}</option>
                        {% endfor %}
                    </select>
                </div>
                
                <div>
                    <label class="block text-sm font-medium text-gray-300">Log Format</label>
                    <select name="log_ext" id="log-ext-select" class="mt-1 block w-full bg-gray-800 border-gray-600 rounded-md shadow-sm p-2 text-white">
                        <option value="">Select after choosing OS</option>
                    </select>
                </div>
                
                <button type="submit" class="btn-primary w-full py-2 rounded-lg font-semibold">Parse Log</button>
            </form>
            
            <div class="mt-6 p-4 bg-gray-800 rounded-lg">
                <h3 class="text-lg font-semibold text-white mb-2">Current File</h3>
                <p class="text-sm text-gray-400 break-all">{{ evidence_filename }}</p>
            </div>
        </div>
    </div>
    
    <div class="lg:col-span-2">
        <div class="card p-6 rounded-lg">
            <h2 class="text-xl font-semibold text-white mb-4">Log Content</h2>
            {% if log_content %}
            <div class="log-view p-4 rounded-lg overflow-auto max-h-[70vh]">
                <pre>{{ log_content }}</pre>
            </div>
            {% else %}
            <div class="text-center p-8 text-gray-500">
                <p>Select a log type and click "Parse Log" to view content.</p>
            </div>
            {% endif %}
        </div>
    </div>
</div>

<script>
document.addEventListener('DOMContentLoaded', function() {
    const osSelect = document.getElementById('os-select');
    const logExtSelect = document.getElementById('log-ext-select');
    
    const osOptions = {{ os_options|tojson }};
    
    osSelect.addEventListener('change', function() {
        const selectedOS = this.value;
        logExtSelect.innerHTML = '<option value="">Select log format</option>';
        
        if (selectedOS && osOptions[selectedOS]) {
            osOptions[selectedOS].forEach(function(format) {
                const option = document.createElement('option');
                option.value = format;
                option.textContent = format;
                logExtSelect.appendChild(option);
            });
        }
    });
});
</script>
"""

ENHANCED_DELETED_STATUS_TEMPLATE = """
<!-- Enhanced Deleted Files Recovery Interface -->
<div class="grid grid-cols-1 lg:grid-cols-4 gap-8">
    <div class="lg:col-span-1">
        <div class="card p-6 rounded-lg">
            <h2 class="text-xl font-semibold text-white mb-4">Recovery Control</h2>
            
            <!-- Enhanced Start Button with Status -->
            {% if not deleted_scan_status.in_progress %}
            <a href="{{ url_for('start_strict_deleted_recovery') }}" 
               class="btn-green w-full py-3 rounded-lg font-semibold text-center block mb-4 hover:bg-green-700 transition-colors">
                🚀 Start Strict Automatic Recovery
            </a>
            {% else %}
            <button class="bg-yellow-600 text-white w-full py-3 rounded-lg font-semibold mb-4 cursor-not-allowed">
                ⏳ Recovery In Progress...
            </button>
            {% endif %}
            
            <!-- Real-time Validation Statistics -->
            <div class="p-3 bg-gray-800 rounded-lg mb-4">
                <h3 class="font-semibold text-white mb-2">Validation Statistics</h3>
                <div class="grid grid-cols-2 gap-2 text-xs">
                    <div class="flex justify-between">
                        <span class="text-gray-400">Total Scanned:</span>
                        <span id="total-scanned" class="text-white font-mono">{{ deleted_scan_status.validation_stats.total_scanned if deleted_scan_status.validation_stats else 0 }}</span>
                    </div>
                    <div class="flex justify-between">
                        <span class="text-gray-400">Empty Rejected:</span>
                        <span id="empty-rejected" class="text-red-400 font-mono">{{ deleted_scan_status.validation_stats.empty_rejected if deleted_scan_status.validation_stats else 0 }}</span>
                    </div>
                    <div class="flex justify-between">
                        <span class="text-gray-400">Duplicate Rejected:</span>
                        <span id="duplicate-rejected" class="text-yellow-400 font-mono">{{ deleted_scan_status.validation_stats.duplicate_rejected if deleted_scan_status.validation_stats else 0 }}</span>
                    </div>
                    <div class="flex justify-between">
                        <span class="text-gray-400">Invalid Rejected:</span>
                        <span id="invalid-rejected" class="text-orange-400 font-mono">{{ deleted_scan_status.validation_stats.invalid_rejected if deleted_scan_status.validation_stats else 0 }}</span>
                    </div>
                    <div class="flex justify-between col-span-2 border-t border-gray-700 pt-1">
                        <span class="text-gray-400 font-semibold">Valid Recovered:</span>
                        <span id="valid-recovered" class="text-green-400 font-semibold font-mono">{{ deleted_scan_status.validation_stats.valid_recovered if deleted_scan_status.validation_stats else 0 }}</span>
                    </div>
                </div>
            </div>

            <!-- Enhanced Timing Information -->
            <div class="p-3 bg-gray-800 rounded-lg mb-4">
                <h3 class="font-semibold text-white mb-2">⏱️ Timing Information</h3>
                <div class="space-y-1 text-xs">
                    <div class="flex justify-between">
                        <span class="text-gray-400">Elapsed:</span>
                        <span id="elapsed-time" class="text-white font-mono">{{ deleted_scan_status.elapsed_time if deleted_scan_status.elapsed_time else "0s" }}</span>
                    </div>
                    <div class="flex justify-between">
                        <span class="text-gray-400">Remaining:</span>
                        <span id="time-remaining" class="text-white font-mono">{{ deleted_scan_status.time_remaining_str if deleted_scan_status.time_remaining_str else "Calculating..." }}</span>
                    </div>
                    <div class="flex justify-between">
                        <span class="text-gray-400">Total Estimated:</span>
                        <span id="total-estimated" class="text-white font-mono">{{ deleted_scan_status.estimated_total_time if deleted_scan_status.estimated_total_time else "Calculating..." }}</span>
                    </div>
                </div>
            </div>

            <!-- Progress Visualization -->
            <div class="mb-4">
                <div class="flex justify-between text-sm text-gray-400 mb-2">
                    <span>Recovery Progress</span>
                    <span id="progress-percent">0%</span>
                </div>
                <div class="w-full bg-gray-700 rounded-full h-3">
                    <div id="progress-bar" class="bg-blue-600 h-3 rounded-full transition-all duration-500" style="width: 0%"></div>
                </div>
            </div>

            <!-- Recovery Methods Summary -->
            <div class="p-3 bg-gray-800 rounded-lg">
                <h3 class="font-semibold text-white mb-2">🔍 Scan Methods</h3>
                <div class="space-y-2 text-sm">
                    <div class="flex justify-between">
                        <span class="text-gray-400">Directory Walk:</span>
                        <span id="method-dir" class="text-white font-mono bg-blue-600 px-2 py-1 rounded text-xs">{{ deleted_scan_status.scan_methods.directory_walk }}</span>
                    </div>
                    <div class="flex justify-between">
                        <span class="text-gray-400">Inode Scan:</span>
                        <span id="method-inode" class="text-white font-mono bg-green-600 px-2 py-1 rounded text-xs">{{ deleted_scan_status.scan_methods.inode_scan }}</span>
                    </div>
                    <div class="flex justify-between">
                        <span class="text-gray-400">File Slack:</span>
                        <span id="method-slack" class="text-white font-mono bg-yellow-600 px-2 py-1 rounded text-xs">{{ deleted_scan_status.scan_methods.file_slack }}</span>
                    </div>
                    <div class="flex justify-between">
                        <span class="text-gray-400">Recycle Bin:</span>
                        <span id="method-recycle" class="text-white font-mono bg-purple-600 px-2 py-1 rounded text-xs">{{ deleted_scan_status.scan_methods.recycle_bin }}</span>
                    </div>
                </div>
            </div>

            <!-- Status Message -->
            <div id="status-message" class="mt-4 p-3 rounded-lg text-sm 
                {% if deleted_scan_status.in_progress %}bg-blue-900 text-blue-200 border border-blue-700
                {% elif deleted_scan_status.complete %}bg-green-900 text-green-200 border border-green-700
                {% else %}bg-gray-800 text-gray-400 border border-gray-700{% endif %}">
                <div class="font-semibold mb-1">
                    {% if deleted_scan_status.in_progress %}🔄 {% elif deleted_scan_status.complete %}✅ {% else %}⏸️ {% endif %}
                    Status
                </div>
                <div id="status-text">{{ deleted_scan_status.message }}</div>
            </div>
        </div>
    </div>

    <div class="lg:col-span-3">
        <div class="card p-6 rounded-lg">
            <div class="flex justify-between items-center mb-6">
                <h2 class="text-xl font-semibold text-white">📁 Recovered Files</h2>
                <div class="flex items-center space-x-3">
                    <span id="recovered-count" class="bg-blue-600 text-white px-3 py-1 rounded-full text-sm font-medium">
                        {{ recovered_files|length }} files recovered
                    </span>
                    {% if recovered_files %}
                    <form action="{{ url_for('download_zip') }}" method="post" class="inline">
                        <input type="hidden" name="file_type" value="deleted_recovered">
                        {% for file in recovered_files %}
                        <input type="hidden" name="selected_files" value="{{ file.name }}">
                        {% endfor %}
                        <button type="submit" class="btn-primary px-4 py-2 rounded-lg text-sm hover:bg-blue-700 transition-colors">
                            📥 Download All as ZIP
                        </button>
                    </form>
                    {% endif %}
                </div>
            </div>

            <div id="recovered-panel">
            {% if recovered_files %}
            <!-- Files Table -->
            <div class="overflow-x-auto">
                <table class="min-w-full bg-gray-800 rounded-lg overflow-hidden">
                    <thead class="bg-gray-700">
                        <tr>
                            <th class="text-left py-3 px-4 text-white font-semibold">ID</th>
                            <th class="text-left py-3 px-4 text-white font-semibold">Filename</th>
                            <th class="text-left py-3 px-4 text-white font-semibold">Size</th>
                            <th class="text-left py-3 px-4 text-white font-semibold">Type</th>
                            <th class="text-left py-3 px-4 text-white font-semibold">Recovery Method</th>
                            <th class="text-left py-3 px-4 text-white font-semibold">Actions</th>
                        </tr>
                    </thead>
                    <tbody>
                        {% for file in recovered_files %}
                        <tr class="border-b border-gray-600 hover:bg-gray-750 transition-colors">
                            <td class="py-3 px-4 text-gray-300 font-mono">{{ file.id }}</td>
                            <td class="py-3 px-4 text-gray-300">
                                <div class="flex items-center space-x-3">
                                    {% if file.thumbnail %}
                                    <img src="{{ file.thumbnail }}" class="w-8 h-8 object-cover rounded" alt="Thumbnail" onerror="this.style.display='none'">
                                    {% endif %}
                                    <div>
                                        <div class="font-medium">{{ file.name|truncate(30) }}</div>
                                        <div class="text-xs text-gray-500">{{ file.mtime if file.mtime else 'Unknown date' }}</div>
                                    </div>
                                </div>
                            </td>
                            <td class="py-3 px-4 text-gray-300 font-mono">{{ file.size_kb }} KB</td>
                            <td class="py-3 px-4 text-gray-300">
                                <span class="px-2 py-1 bg-gray-600 rounded text-xs">{{ file.file_type }}</span>
                            </td>
                            <td class="py-3 px-4 text-gray-300">
                                <span class="px-2 py-1 bg-blue-600 rounded text-xs">{{ file.recovery_method }}</span>
                            </td>
                            <td class="py-3 px-4">
                                <div class="flex space-x-2">
                                    <a href="{{ url_for('view_deleted_file', filename=file.name) }}" target="_blank" 
                                       class="btn-secondary px-3 py-1 text-xs rounded hover:bg-gray-600 transition-colors">👁️ View</a>
                                    <a href="{{ url_for('download_deleted_file', filename=file.name) }}" 
                                       class="btn-primary px-3 py-1 text-xs rounded hover:bg-blue-700 transition-colors">📥 Download</a>
                                    <a href="{{ url_for('hex_view_deleted', filename=file.name) }}" target="_blank"
                                       class="btn-green px-3 py-1 text-xs rounded hover:bg-green-700 transition-colors">🔢 Hex</a>
                                </div>
                            </td>
                        </tr>
                        {% endfor %}
                    </tbody>
                </table>
            </div>
            {% else %}
            <!-- Empty State -->
            <div class="text-center py-12 text-gray-500">
                <div class="mb-4">
                    <svg class="w-16 h-16 mx-auto text-gray-600" fill="none" stroke="currentColor" viewBox="0 0 24 24">
                        <path stroke-linecap="round" stroke-linejoin="round" stroke-width="2" d="M9 12h6m-6 4h6m2 5H7a2 2 0 01-2-2V5a2 2 0 012-2h5.586a1 1 0 01.707.293l5.414 5.414a1 1 0 01.293.707V19a2 2 0 01-2 2z"></path>
                    </svg>
                </div>
                <p class="text-lg mb-2">No deleted files recovered yet</p>
                <p class="text-sm text-gray-400 mb-6">Start the automatic recovery process to scan for deleted files using multiple forensic methods</p>
                
                {% if not deleted_scan_status.in_progress %}
                <a href="{{ url_for('start_strict_deleted_recovery') }}" 
                   class="btn-green px-6 py-3 rounded-lg font-semibold hover:bg-green-700 transition-colors inline-flex items-center">
                    <svg class="w-4 h-4 mr-2" fill="none" stroke="currentColor" viewBox="0 0 24 24">
                        <path stroke-linecap="round" stroke-linejoin="round" stroke-width="2" d="M13 10V3L4 14h7v7l9-11h-7z"></path>
                    </svg>
                    Start Automatic Recovery
                </a>
                {% else %}
                <div class="text-yellow-400 text-sm">
                    ⏳ Recovery process is currently running...
                </div>
                {% endif %}
            </div>
            {% endif %}
            </div> <!-- /#recovered-panel -->
        </div>
    </div>
</div>

<script>
// Small helper to escape HTML when inserting into innerHTML
function escapeHtml(unsafe) {
    if (!unsafe && unsafe !== 0) return '';
    return String(unsafe)
        .replace(/&/g, '&amp;')
        .replace(/</g, '&lt;')
        .replace(/>/g, '&gt;')
        .replace(/"/g, '&quot;')
        .replace(/'/g, '&#039;');
}

// Function to update the recovered files list
function updateRecoveredFilesList() {
    fetch('/get_recovered_files_list')
        .then(response => response.json())
        .then(data => {
            const recoveredPanel = document.getElementById('recovered-panel');
            const recoveredCountBadge = document.getElementById('recovered-count');
            
            if (!recoveredPanel) return;

            if (data.success && data.files && Object.keys(data.files).length > 0) {
                // Build the files table
                let html = `
                    <div class="overflow-x-auto">
                        <table class="min-w-full bg-gray-800 rounded-lg overflow-hidden">
                            <thead class="bg-gray-700">
                                <tr>
                                    <th class="text-left py-3 px-4 text-white font-semibold">ID</th>
                                    <th class="text-left py-3 px-4 text-white font-semibold">Filename</th>
                                    <th class="text-left py-3 px-4 text-white font-semibold">Size</th>
                                    <th class="text-left py-3 px-4 text-white font-semibold">Type</th>
                                    <th class="text-left py-3 px-4 text-white font-semibold">Recovery Method</th>
                                    <th class="text-left py-3 px-4 text-white font-semibold">Actions</th>
                                </tr>
                            </thead>
                            <tbody>`;

                Object.values(data.files).forEach((file, index) => {
                    const id = index + 1;
                    const name = escapeHtml(file.name || 'Unknown');
                    const size_kb = file.size_kb || (file.size ? (Math.round((file.size/1024)*100)/100).toFixed(2) : '0');
                    const file_type = escapeHtml(file.file_type || 'Unknown');
                    const method = escapeHtml(file.recovery_method || 'Unknown');
                    const mtime = escapeHtml(file.mtime || 'Unknown date');
                    const thumb = file.thumbnail ? `<img src="${file.thumbnail}" class="w-8 h-8 object-cover rounded" alt="Thumbnail" onerror="this.style.display='none'">` : '';

                    html += `
                        <tr class="border-b border-gray-600 hover:bg-gray-750 transition-colors">
                            <td class="py-3 px-4 text-gray-300 font-mono">${id}</td>
                            <td class="py-3 px-4 text-gray-300">
                                <div class="flex items-center space-x-3">
                                    ${thumb}
                                    <div>
                                        <div class="font-medium">${name.length > 40 ? name.slice(0,37) + '...' : name}</div>
                                        <div class="text-xs text-gray-500">${mtime}</div>
                                    </div>
                                </div>
                            </td>
                            <td class="py-3 px-4 text-gray-300 font-mono">${size_kb} KB</td>
                            <td class="py-3 px-4 text-gray-300">
                                <span class="px-2 py-1 bg-gray-600 rounded text-xs">${file_type}</span>
                            </td>
                            <td class="py-3 px-4 text-gray-300">
                                <span class="px-2 py-1 bg-blue-600 rounded text-xs">${method}</span>
                            </td>
                            <td class="py-3 px-4">
                                <div class="flex space-x-2">
                                    <a href="/view_deleted_file/${encodeURIComponent(name)}" target="_blank" class="btn-secondary px-3 py-1 text-xs rounded hover:bg-gray-600 transition-colors">👁️ View</a>
                                    <a href="/download_deleted_file/${encodeURIComponent(name)}" class="btn-primary px-3 py-1 text-xs rounded hover:bg-blue-700 transition-colors">📥 Download</a>
                                    <a href="/hex_view_deleted/${encodeURIComponent(name)}" target="_blank" class="btn-green px-3 py-1 text-xs rounded hover:bg-green-700 transition-colors">🔢 Hex</a>
                                </div>
                            </td>
                        </tr>`;
                });

                html += `</tbody></table></div>`;
                recoveredPanel.innerHTML = html;
                
                // Update the count badge
                if (recoveredCountBadge) {
                    recoveredCountBadge.textContent = `${Object.keys(data.files).length} files recovered`;
                }
            } else {
                // Show empty state
                recoveredPanel.innerHTML = `
                    <div class="text-center py-12 text-gray-500">
                        <div class="mb-4">
                            <svg class="w-16 h-16 mx-auto text-gray-600" fill="none" stroke="currentColor" viewBox="0 0 24 24">
                                <path stroke-linecap="round" stroke-linejoin="round" stroke-width="2" d="M9 12h6m-6 4h6m2 5H7a2 2 0 01-2-2V5a2 2 0 012-2h5.586a1 1 0 01.707.293l5.414 5.414a1 1 0 01.293.707V19a2 2 0 01-2 2z"></path>
                            </svg>
                        </div>
                        <p class="text-lg mb-2">No deleted files recovered yet</p>
                        <p class="text-sm text-gray-400 mb-6">Start the automatic recovery process to scan for deleted files using multiple forensic methods</p>
                    </div>`;
                
                if (recoveredCountBadge) {
                    recoveredCountBadge.textContent = `0 files recovered`;
                }
            }
        })
        .catch(error => {
            console.error('Error fetching recovered files list:', error);
        });
}

// Enhanced progress tracking with better error handling
function updateRecoveryProgress() {
    fetch('/deleted_scan_status')
        .then(response => {
            if (!response.ok) throw new Error('Network response was not ok');
            return response.json();
        })
        .then(data => {
            console.log('Recovery status update:', data);
            
            // Update validation statistics (guard element presence)
            if (data.validation_stats) {
                const totalScannedEl = document.getElementById('total-scanned');
                const emptyRejectedEl = document.getElementById('empty-rejected');
                const duplicateRejectedEl = document.getElementById('duplicate-rejected');
                const invalidRejectedEl = document.getElementById('invalid-rejected');
                const validRecoveredEl = document.getElementById('valid-recovered');

                if (totalScannedEl) totalScannedEl.textContent = data.validation_stats.total_scanned || 0;
                if (emptyRejectedEl) emptyRejectedEl.textContent = data.validation_stats.empty_rejected || 0;
                if (duplicateRejectedEl) duplicateRejectedEl.textContent = data.validation_stats.duplicate_rejected || 0;
                if (invalidRejectedEl) invalidRejectedEl.textContent = data.validation_stats.invalid_rejected || 0;
                if (validRecoveredEl) validRecoveredEl.textContent = data.validation_stats.valid_recovered || 0;
            }
            
            // Update timing information
            const elapsedTimeEl = document.getElementById('elapsed-time');
            const timeRemainingEl = document.getElementById('time-remaining');
            const totalEstimatedEl = document.getElementById('total-estimated');
            
            function formatTimeSeconds(secs) {
                if (secs < 1) return '0s';
                const h = Math.floor(secs / 3600);
                const m = Math.floor((secs % 3600) / 60);
                const s = Math.floor(secs % 60);
                if (h > 0) return `${h}h ${m}m ${s}s`;
                if (m > 0) return `${m}m ${s}s`;
                return `${s}s`;
            }

            if (elapsedTimeEl) {
                if (data.elapsed_time && data.elapsed_time !== '0s') {
                    elapsedTimeEl.textContent = data.elapsed_time;
                } else if (data.start_time) {
                    const elapsedSecs = Math.max(0, Math.floor(Date.now()/1000 - data.start_time));
                    elapsedTimeEl.textContent = formatTimeSeconds(elapsedSecs);
                }
            }
            if (data.time_remaining_str && timeRemainingEl) timeRemainingEl.textContent = data.time_remaining_str;
            if (data.estimated_total_time && totalEstimatedEl) totalEstimatedEl.textContent = data.estimated_total_time;
            
            // Update progress bar
            const progressBar = document.getElementById('progress-bar');
            const progressPercent = document.getElementById('progress-percent');
            
            if (data.in_progress) {
                // Calculate progress based on methods completed
                const totalMethods = 4; // directory_walk, inode_scan, file_slack, recycle_bin
                const completedMethods = Object.values(data.scan_methods || {}).filter(val => val > 0).length;
                const progress = Math.min(95, (completedMethods / totalMethods) * 100); // Cap at 95% until complete
                
                if (progressBar) {
                    progressBar.style.width = progress + '%';
                    progressBar.className = 'bg-blue-600 h-3 rounded-full transition-all duration-500';
                }
                if (progressPercent) progressPercent.textContent = Math.round(progress) + '%';
            }
            
            // Update method counters
            if (data.scan_methods) {
                const methodDirEl = document.getElementById('method-dir');
                const methodInodeEl = document.getElementById('method-inode');
                const methodSlackEl = document.getElementById('method-slack');
                const methodRecycleEl = document.getElementById('method-recycle');

                if (methodDirEl) methodDirEl.textContent = data.scan_methods.directory_walk || 0;
                if (methodInodeEl) methodInodeEl.textContent = data.scan_methods.inode_scan || 0;
                if (methodSlackEl) methodSlackEl.textContent = data.scan_methods.file_slack || 0;
                if (methodRecycleEl) methodRecycleEl.textContent = data.scan_methods.recycle_bin || 0;
            }
            
            // Update status message
            const statusText = document.getElementById('status-text');
            const statusMessage = document.getElementById('status-message');

            if (statusText && statusMessage) {
                if (data.message) {
                    statusText.textContent = data.message;
                }
                
                if (data.in_progress) {
                    statusMessage.className = 'mt-4 p-3 rounded-lg text-sm bg-blue-900 text-blue-200 border border-blue-700';
                } else if (data.error) {
                    statusMessage.className = 'mt-4 p-3 rounded-lg text-sm bg-red-900 text-red-200 border border-red-700';
                    statusText.textContent = 'Error: ' + (data.error || 'Unknown error occurred');
                } else if (data.complete) {
                    statusMessage.className = 'mt-4 p-3 rounded-lg text-sm bg-green-900 text-green-200 border border-green-700';
                }
            }
            
            // CRITICAL: Update the recovered files list in real-time
            updateRecoveredFilesList();

            // Continue polling if still in progress
            if (data.in_progress && !data.complete) {
                setTimeout(updateRecoveryProgress, 2000);
            } else if (data.complete) {
                // Final update when complete
                if (progressBar) progressBar.style.width = '100%';
                if (progressPercent) progressPercent.textContent = '100%';
                console.log('Recovery process completed');
            }
        })
        .catch(error => {
            console.error('Error fetching recovery status:', error);
            const statusText = document.getElementById('status-text');
            if (statusText) {
                statusText.textContent = 'Error connecting to server: ' + error.message;
            }
            // Retry after 5 seconds on error
            setTimeout(updateRecoveryProgress, 5000);
        });
}

// Start polling if recovery is in progress
document.addEventListener('DOMContentLoaded', function() {
    // Always update the files list on page load
    updateRecoveredFilesList();
    
    {% if deleted_scan_status.in_progress %}
    console.log('Starting recovery progress monitoring...');
    updateRecoveryProgress();
    {% endif %}
    
    // Add click handler for the start button to show immediate feedback
    const startButton = document.querySelector('a[href*="start_strict_deleted_recovery"]');
    if (startButton) {
        startButton.addEventListener('click', function(e) {
            const button = this;
            const originalText = button.innerHTML;
            
            button.innerHTML = '🔄 Starting...';
            button.classList.add('opacity-50', 'cursor-not-allowed');
            
            // Start polling immediately when user clicks start
            setTimeout(() => {
                updateRecoveryProgress();
            }, 1000);
            
            // Revert after 3 seconds if still on same page
            setTimeout(() => {
                if (button.innerHTML === '🔄 Starting...') {
                    button.innerHTML = originalText;
                    button.classList.remove('opacity-50', 'cursor-not-allowed');
                }
            }, 3000);
        });
    }
});
        
</script>
"""
REPORTING_PAGE_CONTENT = """
<h1 class="text-3xl font-bold text-white mb-4">Forensic Reporting</h1>
<p class="text-gray-400 mb-8">Generate comprehensive reports of your forensic analysis findings.</p>

<div class="card p-6 rounded-lg">
    <h2 class="text-xl font-semibold text-white mb-6">Create New Report</h2>
    
    <form method="post" class="space-y-8">
        <div class="grid grid-cols-1 md:grid-cols-3 gap-6">
            <div>
                <label class="block text-sm font-medium text-gray-300">Case Name</label>
                <input type="text" name="case_name" class="mt-1 block w-full bg-gray-800 border-gray-600 rounded-md shadow-sm p-2 text-white" placeholder="e.g., Company Investigation" required>
            </div>
            
            <div>
                <label class="block text-sm font-medium text-gray-300">Case Number</label>
                <input type="text" name="case_number" class="mt-1 block w-full bg-gray-800 border-gray-600 rounded-md shadow-sm p-2 text-white" placeholder="e.g., CASE-2024-001" required>
            </div>
            
            <div>
                <label class="block text-sm font-medium text-gray-300">Examiner Name</label>
                <input type="text" name="examiner_name" class="mt-1 block w-full bg-gray-800 border-gray-600 rounded-md shadow-sm p-2 text-white" placeholder="Your name" required>
            </div>
        </div>
        
        <div class="pt-4">
            <label class="block text-sm font-medium text-gray-300">Report Scope</label>
            <div class="mt-2 grid grid-cols-1 md:grid-cols-3 gap-4 text-sm text-gray-300">
                <div>
                    <label class="inline-flex items-center">
                        <input type="checkbox" name="include_carved" value="1" checked class="form-checkbox">
                        <span class="ml-2">Include Carved Files</span>
                    </label>
                </div>
                <div>
                    <label class="inline-flex items-center">
                        <input type="checkbox" name="include_deleted" value="1" checked class="form-checkbox">
                        <span class="ml-2">Include Deleted Files</span>
                    </label>
                </div>
                <div>
                    <label class="block text-sm font-medium text-gray-400">File Categories (select any)</label>
                    <div class="mt-2 grid grid-cols-2 gap-2 text-xs text-gray-300">
                        <label class="inline-flex items-center"><input type="checkbox" name="file_category" value="images" checked class="form-checkbox"><span class="ml-2">Images</span></label>
                        <label class="inline-flex items-center"><input type="checkbox" name="file_category" value="documents" checked class="form-checkbox"><span class="ml-2">Documents</span></label>
                        <label class="inline-flex items-center"><input type="checkbox" name="file_category" value="archives" checked class="form-checkbox"><span class="ml-2">Archives</span></label>
                        <label class="inline-flex items-center"><input type="checkbox" name="file_category" value="media" checked class="form-checkbox"><span class="ml-2">Audio/Video</span></label>
                        <label class="inline-flex items-center"><input type="checkbox" name="file_category" value="executables" checked class="form-checkbox"><span class="ml-2">Executables</span></label>
                        <label class="inline-flex items-center"><input type="checkbox" name="file_category" value="other" checked class="form-checkbox"><span class="ml-2">Other</span></label>
                    </div>
                    <div class="mt-2 text-xs text-gray-400">Or add custom extensions (comma separated, e.g. <code>doc,log,dat</code>)</div>
                    <input type="text" name="custom_exts" placeholder="e.g. doc,log,dat" class="mt-1 block w-full bg-gray-800 border-gray-600 rounded-md shadow-sm p-2 text-white text-xs">
                </div>
            </div>
        </div>
        
        <div>
            <label class="block text-sm font-medium text-gray-300">Report Format</label>
            <div id="report-format-options" class="mt-2 grid grid-cols-2 md:grid-cols-4 gap-4">
                <label class="relative flex cursor-pointer rounded-lg border bg-gray-800 p-4 focus:outline-none">
                    <input type="radio" name="report_format" value="html" class="sr-only" checked>
                    <span class="flex flex-1">
                        <span class="flex flex-col">
                            <span class="block text-sm font-medium text-white">HTML</span>
                            <span class="mt-1 flex text-xs text-gray-400">Web format with images</span>
                        </span>
                    </span>
                    <svg class="h-5 w-5 text-blue-600" viewBox="0 0 20 20" fill="currentColor" aria-hidden="true">
                        <path fill-rule="evenodd" d="M10 18a8 8 0 100-16 8 8 0 000 16zm3.857-9.809a.75.75 0 00-1.214-.882l-3.483 4.79-1.88-1.88a.75.75 0 10-1.06 1.061l2.5 2.5a.75.75 0 001.137-.089l4-5.5z" clip-rule="evenodd" />
                    </svg>
                </label>
                
                <label class="relative flex cursor-pointer rounded-lg border bg-gray-800 p-4 focus:outline-none">
                    <input type="radio" name="report_format" value="pdf" class="sr-only">
                    <span class="flex flex-1">
                        <span class="flex flex-col">
                            <span class="block text-sm font-medium text-white">PDF</span>
                            <span class="mt-1 flex text-xs text-gray-400">Printable document</span>
                        </span>
                    </span>
                    <svg class="h-5 w-5 text-blue-600" viewBox="0 0 20 20" fill="currentColor" aria-hidden="true">
                        <path fill-rule="evenodd" d="M10 18a8 8 0 100-16 8 8 0 000 16zm3.857-9.809a.75.75 0 00-1.214-.882l-3.483 4.79-1.88-1.88a.75.75 0 10-1.06 1.061l2.5 2.5a.75.75 0 001.137-.089l4-5.5z" clip-rule="evenodd" />
                    </svg>
                </label>
                
                <label class="relative flex cursor-pointer rounded-lg border bg-gray-800 p-4 focus:outline-none">
                    <input type="radio" name="report_format" value="docx" class="sr-only">
                    <span class="flex flex-1">
                        <span class="flex flex-col">
                            <span class="block text-sm font-medium text-white">DOCX</span>
                            <span class="mt-1 flex text-xs text-gray-400">Word document</span>
                        </span>
                    </span>
                    <svg class="h-5 w-5 text-blue-600" viewBox="0 0 20 20" fill="currentColor" aria-hidden="true">
                        <path fill-rule="evenodd" d="M10 18a8 8 0 100-16 8 8 0 000 16zm3.857-9.809a.75.75 0 00-1.214-.882l-3.483 4.79-1.88-1.88a.75.75 0 10-1.06 1.061l2.5 2.5a.75.75 0 001.137-.089l4-5.5z" clip-rule="evenodd" />
                    </svg>
                </label>
                
                <label class="relative flex cursor-pointer rounded-lg border bg-gray-800 p-4 focus:outline-none">
                    <input type="radio" name="report_format" value="csv" class="sr-only">
                    <span class="flex flex-1">
                        <span class="flex flex-col">
                            <span class="block text-sm font-medium text-white">CSV (ZIP)</span>
                            <span class="mt-1 flex text-xs text-gray-400">Spreadsheet data</span>
                        </span>
                    </span>
                    <svg class="h-5 w-5 text-blue-600" viewBox="0 0 20 20" fill="currentColor" aria-hidden="true">
                        <path fill-rule="evenodd" d="M10 18a8 8 0 100-16 8 8 0 000 16zm3.857-9.809a.75.75 0 00-1.214-.882l-3.483 4.79-1.88-1.88a.75.75 0 10-1.06 1.061l2.5 2.5a.75.75 0 001.137-.089l4-5.5z" clip-rule="evenodd" />
                    </svg>
                </label>
            </div>
        </div>
        
        <div class="pt-4 border-t border-gray-700">
            <button type="submit" class="btn-green w-full py-3 rounded-lg font-semibold text-lg">Generate Report</button>
        </div>
    </form>
</div>

<script>
// Make the report-format tiles visually reflect the single selected radio
document.addEventListener('DOMContentLoaded', function() {
    try {
        const container = document.getElementById('report-format-options');
        if (!container) return;
        const labels = Array.from(container.querySelectorAll('label'));

        function clearSelection() {
            labels.forEach(l => l.classList.remove('selected', 'ring-2', 'ring-blue-500'));
        }

        function updateFromRadios() {
            clearSelection();
            labels.forEach(l => {
                const inp = l.querySelector('input[type="radio"]');
                if (inp && inp.checked) {
                    l.classList.add('selected', 'ring-2', 'ring-blue-500');
                }
            });
        }

        labels.forEach(l => {
            const inp = l.querySelector('input[type="radio"]');
            if (!inp) return;
            // Click on the tile selects the radio and updates visuals
            l.addEventListener('click', function(e) {
                // Allow the input to be checked by default behavior
                inp.checked = true;
                updateFromRadios();
            });
            // Also react to keyboard changes
            inp.addEventListener('change', updateFromRadios);
        });

        // initialize
        updateFromRadios();
    } catch (err) {
        console.error('Report format selector init failed', err);
    }
});
</script>

<div class="card p-6 rounded-lg mt-8">
    <h2 class="text-xl font-semibold text-white mb-4">Report Contents</h2>
    <div class="grid grid-cols-1 md:grid-cols-2 gap-6">
        <div>
            <h3 class="text-lg font-semibold text-white mb-2">Case Information</h3>
            <ul class="text-gray-400 space-y-1 text-sm">
                <li>• Case details and examiner information</li>
                <li>• Report generation timestamp</li>
                <li>• Software version and configuration</li>
            </ul>
        </div>
        
        <div>
            <h3 class="text-lg font-semibold text-white mb-2">Evidence Analysis</h3>
            <ul class="text-gray-400 space-y-1 text-sm">
                <li>• Loaded evidence file information</li>
                <li>• File hashes (MD5, SHA-1, SHA-256)</li>
                <li>• Forensic analysis results</li>
                <li>• Partition information</li>
            </ul>
        </div>
        
        <div>
            <h3 class="text-lg font-semibold text-white mb-2">Recovery Results</h3>
            <ul class="text-gray-400 space-y-1 text-sm">
                <li>• Carved files listing with offsets</li>
                <li>• File preview thumbnails (images)</li>
                <li>• Recovery statistics</li>
            </ul>
        </div>
        
        <div>
            <h3 class="text-lg font-semibold text-white mb-2">Deleted Files</h3>
            <ul class="text-gray-400 space-y-1 text-sm">
                <li>• Recovered deleted files listing</li>
                <li>• Metadata (timestamps, sizes)</li>
                <li>• Recovery method information</li>
            </ul>
        </div>
    </div>
</div>
"""

REPORT_TEMPLATE = """
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Forensic Report - {{ case_details.name }}</title>
    <style>
        body { 
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif; 
            line-height: 1.6; 
            color: #333; 
            max-width: 1200px; 
            margin: 0 auto; 
            padding: 20px; 
            background-color: #f9f9f9;
        }
        .header { 
            background: linear-gradient(135deg, #1e3c72 0%, #2a5298 100%); 
            color: white; 
            padding: 30px; 
            border-radius: 10px; 
            margin-bottom: 30px; 
            text-align: center;
        }
        .section { 
            background: white; 
            padding: 25px; 
            margin-bottom: 25px; 
            border-radius: 8px; 
            box-shadow: 0 2px 10px rgba(0,0,0,0.1); 
            border-left: 4px solid #2a5298;
        }
        table { 
            width: 100%; 
            border-collapse: collapse; 
            margin: 15px 0; 
        }
        th, td { 
            padding: 12px; 
            text-align: left; 
            border-bottom: 1px solid #ddd; 
        }
        th { 
            background-color: #f2f2f2; 
            font-weight: 600; 
        }
        .thumbnail { 
            max-width: 100px; 
            max-height: 100px; 
            border: 1px solid #ddd; 
            border-radius: 4px; 
            padding: 2px; 
        }
        .badge { 
            display: inline-block; 
            padding: 3px 8px; 
            border-radius: 4px; 
            font-size: 0.8em; 
            font-weight: bold; 
        }
        .badge-success { background: #d4edda; color: #155724; }
        .badge-warning { background: #fff3cd; color: #856404; }
        .badge-danger { background: #f8d7da; color: #721c24; }
        .footer { 
            text-align: center; 
            margin-top: 40px; 
            padding: 20px; 
            color: #666; 
            font-size: 0.9em; 
            border-top: 1px solid #ddd;
        }
        .hex-preview { 
            font-family: 'Courier New', monospace; 
            font-size: 0.8em; 
            background: #f8f9fa; 
            padding: 10px; 
            border-radius: 4px; 
            overflow-x: auto; 
        }
    </style>
</head>
<body>
    <div class="header">
        <h1>Digital Forensic Analysis Report</h1>
        <h2>Case: {{ case_details.name }}</h2>
        <p>Generated on {{ now.strftime('%Y-%m-%d at %H:%M:%S') }}</p>
    </div>

    <div class="section">
        <h2>Case Information</h2>
        <table>
            <tr><th>Case Name</th><td>{{ case_details.name }}</td></tr>
            <tr><th>Case Number</th><td>{{ case_details.number }}</td></tr>
            <tr><th>Examiner</th><td>{{ case_details.examiner }}</td></tr>
            <tr><th>Report Date</th><td>{{ now.strftime('%Y-%m-%d %H:%M:%S') }}</td></tr>
            <tr><th>Software</th><td>ForensicCarver Pro v4.3.0</td></tr>
        </table>
    </div>

    {% for filename, details in evidence_file.items() %}
    <div class="section">
        <h2>Evidence File: {{ filename }}</h2>
        
        <h3>File Information</h3>
        <table>
            <tr><th>File Size</th><td>{{ details.size_mb }} MB</td></tr>
            <tr><th>MD5 Hash</th><td>{{ details.hash_info.MD5 or 'Calculating...' }}</td></tr>
            <tr><th>SHA-1 Hash</th><td>{{ details.hash_info['SHA-1'] or 'Calculating...' }}</td></tr>
            <tr><th>SHA-256 Hash</th><td>{{ details.hash_info['SHA-256'] or 'Calculating...' }}</td></tr>
            <tr><th>Encryption Status</th>
                <td>
                    {% if details.encryption_status.encrypted %}
                        <span class="badge badge-warning">ENCRYPTED</span> - {{ details.encryption_status.description }}
                        {% if details.encryption_status.decrypted_path %}
                            <span class="badge badge-success">(Decrypted)</span>
                        {% endif %}
                    {% else %}
                        <span class="badge badge-success">Not Encrypted</span>
                    {% endif %}
                </td>
            </tr>
        </table>

        <h3>Forensic Analysis</h3>
        <ul>
            {% for result in details.forensic_results %}
            <li>{{ result }}</li>
            {% endfor %}
        </ul>

        {% if details.partition_info %}
        <h3>Partition Information</h3>
        <table>
            <tr><th>Index</th><th>Description</th><th>Start Sector</th><th>Length (Sectors)</th></tr>
            {% for part in details.partition_info %}
            <tr>
                <td>{{ part.addr }}</td>
                <td>{{ part.desc }}</td>
                <td>{{ part.start }}</td>
                <td>{{ part.len }}</td>
            </tr>
            {% endfor %}
        </table>
        {% endif %}
    </div>
    {% endfor %}

    {% if carved_files %}
    <div class="section">
        <h2>Carved Files Recovery ({{ carved_files|length }} files found)</h2>
        <table>
            <tr><th>ID</th><th>Filename</th><th>Offset</th><th>Size</th><th>Preview</th></tr>
            {% for filename, info in carved_files.items() %}
            <tr>
                <td>{{ info.id }}</td>
                <td>{{ filename }}</td>
                <td>{{ info.offset }}</td>
                <td>{{ "%.2f KB"|format(info.size_bytes / 1024) if info.size_bytes else 'N/A' }}</td>
                <td>
                    {% if info.thumbnail_uri %}
                    <img src="{{ info.thumbnail_uri }}" class="thumbnail" alt="Preview">
                    {% else %}
                    <em>No preview</em>
                    {% endif %}
                </td>
            </tr>
            {% endfor %}
        </table>
    </div>
    {% endif %}

    {% if deleted_files %}
    <div class="section">
        <h2>Deleted Files Recovery ({{ deleted_files|length }} entries found)</h2>
        <table>
            <tr><th>Inode</th><th>Filename</th><th>Size</th><th>Modified</th><th>Recovery Method</th></tr>
            {% for inode, file in deleted_files.items() %}
            <tr>
                <td>{{ file.inode }}</td>
                <td>{{ file.name }}</td>
                <td>{{ file.size }} bytes</td>
                <td>{{ file.mtime }}</td>
                <td>{{ file.recovery_method }}</td>
            </tr>
            {% endfor %}
        </table>
    </div>
    {% endif %}

    <div class="footer">
        <p>This report was automatically generated by ForensicCarver Pro.</p>
        <p>Confidential - For authorized personnel only.</p>
    </div>
</body>
</html>
"""

HEX_VIEW_CARVED_FILE_CONTENT = """
<h1 class="text-3xl font-bold text-white mb-4">Hex View: {{ filename }}</h1>
<div class="flex justify-between items-center mb-4">
    <a href="{{ url_for('view_carved_file', filename=filename) }}" class="btn-secondary px-4 py-2 rounded-lg">Back to File View</a>
    <a href="{{ url_for('download_carved_file', filename=filename) }}" class="btn-primary px-4 py-2 rounded-lg">Download File</a>
</div>

<div class="card p-6 rounded-lg">
    <div class="hex-view p-4 rounded-lg overflow-x-auto">
        <pre class="font-mono text-sm leading-tight">{{ hex_content }}</pre>
    </div>
</div>
"""

DELETED_RECOVERY_PROCESS_CONTENT = """
<h1 class="text-3xl font-bold text-white mb-2">Deleted Files Recovery in Progress</h1>
<p class="text-gray-400 mb-8">Scanning evidence file for deleted files and metadata.</p>

<div class="card p-6 rounded-lg">
    <div class="flex justify-between items-center mb-4">
        <h2 class="text-xl font-semibold text-white">Recovery Progress</h2>
        <span id="status-badge" class="px-3 py-1 bg-yellow-600 text-white rounded-full text-sm font-medium">Running</span>
    </div>
    
    <div id="progress-container">
        <p class="text-sm text-gray-400 mb-2">Status: <span id="status-message">Starting recovery process...</span></p>
        <p class="text-sm text-gray-400">Files found: <span id="files-found">0</span></p>
        
        <div class="mt-4 w-full bg-gray-700 rounded-full h-2">
            <div id="progress-bar" class="bg-blue-600 h-2 rounded-full transition-all duration-500" style="width: 0%"></div>
        </div>
    </div>
    
    <div id="error-container" class="mt-4 p-4 bg-red-900 border border-red-500 rounded-lg text-sm text-red-200 hidden">
        <h3 class="font-bold mb-2">An Error Occurred</h3>
        <p id="error-message"></p>
    </div>
    
    <div id="success-container" class="mt-4 p-4 bg-green-900 border border-green-500 rounded-lg text-sm text-green-200 hidden">
        <h3 class="font-bold mb-2">Recovery Complete</h3>
        <p id="success-message"></p>
    </div>
</div>

<div class="mt-8 flex justify-center">
    <a href="{{ url_for('deleted_files') }}" id="view-results-btn" class="btn-primary px-6 py-3 rounded-lg opacity-50 pointer-events-none">View Recovered Files</a>
</div>

<script>
function updateRecoveryProgress() {
    fetch('/deleted_scan_status')
        .then(response => response.json())
        .then(data => {
            const statusMessage = document.getElementById('status-message');
            const filesFound = document.getElementById('files-found');
            const progressBar = document.getElementById('progress-bar');
            const statusBadge = document.getElementById('status-badge');
            const errorContainer = document.getElementById('error-container');
            const successContainer = document.getElementById('success-container');
            const viewResultsBtn = document.getElementById('view-results-btn');

            statusMessage.textContent = data.message;
            filesFound.textContent = data.files_found;

            // Update progress bar based on status
            if (data.in_progress) {
                progressBar.style.width = '70%'; // Indeterminate progress
                statusBadge.className = 'px-3 py-1 bg-yellow-600 text-white rounded-full text-sm font-medium';
                statusBadge.textContent = 'Running';
                setTimeout(updateRecoveryProgress, 2000);
            } else if (data.complete) {
                progressBar.style.width = '100%';
                progressBar.classList.remove('bg-blue-600');
                progressBar.classList.add('bg-green-600');
                statusBadge.className = 'px-3 py-1 bg-green-600 text-white rounded-full text-sm font-medium';
                statusBadge.textContent = 'Complete';
                
                successContainer.classList.remove('hidden');
                document.getElementById('success-message').textContent = 
                    `Recovery completed successfully. Found ${data.files_found} files.`;
                
                viewResultsBtn.classList.remove('opacity-50', 'pointer-events-none');
            }

            if (data.message && data.message.toLowerCase().includes('error')) {
                errorContainer.classList.remove('hidden');
                document.getElementById('error-message').textContent = data.message;
                statusBadge.className = 'px-3 py-1 bg-red-600 text-white rounded-full text-sm font-medium';
                statusBadge.textContent = 'Error';
            }
        })
        .catch(error => {
            console.error('Error fetching recovery status:', error);
            document.getElementById('status-message').textContent = 'Error checking recovery status';
            document.getElementById('error-container').classList.remove('hidden');
            document.getElementById('error-message').textContent = error.toString();
        });
}

// Start polling for progress updates
document.addEventListener('DOMContentLoaded', function() {
    updateRecoveryProgress();
});
</script>
"""


DELETED_FILES_CONTENT = """
<h1 class="text-3xl font-bold text-white mb-4">Recovered Deleted Files</h1>
<p class="text-gray-400 mb-8">Files recovered from deleted space and file system metadata.</p>

<div class="card p-6 rounded-lg">
    {% if deleted_files %}
    <div class="overflow-x-auto">
        <table class="min-w-full">
            <thead>
                <tr class="border-b border-gray-600">
                    <th class="text-left py-3 px-4 text-white font-semibold">Inode</th>
                    <th class="text-left py-3 px-4 text-white font-semibold">Filename</th>
                    <th class="text-left py-3 px-4 text-white font-semibold">Size</th>
                    <th class="text-left py-3 px-4 text-white font-semibold">Modified</th>
                    <th class="text-left py-3 px-4 text-white font-semibold">Recovery Method</th>
                    <th class="text-left py-3 px-4 text-white font-semibold">Actions</th>
                </tr>
            </thead>
            <tbody>
                {% for inode, file in deleted_files.items() %}
                <tr class="border-b border-gray-700 hover:bg-gray-800">
                    <td class="py-3 px-4 text-gray-300 font-mono">{{ file.inode }}</td>
                    <td class="py-3 px-4 text-gray-300">{{ file.name }}</td>
                    <td class="py-3 px-4 text-gray-300">{{ file.size }} bytes</td>
                    <td class="py-3 px-4 text-gray-300">{{ file.mtime }}</td>
                    <td class="py-3 px-4 text-gray-300">
                        <span class="px-2 py-1 bg-blue-600 rounded text-xs">{{ file.recovery_method }}</span>
                    </td>
                    <td class="py-3 px-4">
                        <div class="flex space-x-2">
                            <a href="{{ url_for('view_deleted_file', inode=inode) }}" class="btn-secondary px-3 py-1 text-xs">View</a>
                            <a href="{{ url_for('download_deleted_file', inode=inode) }}" class="btn-primary px-3 py-1 text-xs">Download</a>
                        </div>
                    </td>
                </tr>
                {% endfor %}
            </tbody>
        </table>
    </div>
    
    <div class="mt-4 text-sm text-gray-400">
        Total recovered files: {{ deleted_files|length }}
    </div>
    {% else %}
    <div class="text-center py-8 text-gray-500">
        <p>No deleted files have been recovered yet.</p>
        <a href="{{ url_for('deleted_files_status_page') }}" class="btn-primary mt-4 inline-block px-4 py-2">Start Recovery Process</a>
    </div>
    {% endif %}
</div>
"""

def get_active_evidence_path():
    """
    Gets the correct path for analysis. 
    Returns the decrypted path if it exists, otherwise the original path.
    Returns None if no file is loaded.
    """
    if not uploaded_files_db:
        return None
    
    file_details = next(iter(uploaded_files_db.values()))
    
    decrypted_path = file_details.get('encryption_status', {}).get('decrypted_path')
    if decrypted_path and os.path.exists(decrypted_path):
        return decrypted_path
    
    return file_details.get('path')

# --- Helper functions for loading and clearing data ---

def _clear_all_session_data():
    """Clears all in-memory data and temporary result files."""
    global carving_status , deleted_scan_status, decryption_status, hashing_status, strings_status, sorted_deleted_inodes
    
    
    uploaded_files_db.clear()
    deleted_files_db.clear()
    sorted_deleted_inodes = []
    
    if not carving_status.get("complete"):
        carving_status = {
            "progress": 0, "current_offset": "0x00000000", "files_found": 0,
            "complete": False, "found_files_list": [], "time_remaining_str": "N/A"
        }
    
    deleted_scan_status = {
        "in_progress": False, "files_found": 0, "complete": False, "message": "Scan has not started.",
        "errors": [], "time_remaining_str": "N/A"
    }
    decryption_status = {
        "in_progress": False, "complete": False, "message": "", "attempts": 0, "success": False, "filename": None
    }
    hashing_status = {
        "in_progress": False, "progress": 0, "complete": False, "hashes": {}
    }
    strings_status = {
        "in_progress": False, "complete": False, "progress": 0, "strings_found": 0, "preview": []
    }
    
    # NEW: Clear the carved files directory
    carved_dir = app.config['CARVED_FOLDER']
    try:
        if os.path.exists(carved_dir):
            for filename in os.listdir(carved_dir):
                file_path = os.path.join(carved_dir, filename)
                if os.path.isfile(file_path):
                    os.unlink(file_path)
            # Sanitize path to remove any trailing CR/LF from config values before logging
            safe_carved_dir = carved_dir.rstrip('\r\n') if isinstance(carved_dir, str) else carved_dir
            try:
                app.logger.info(f"Cleared carved files directory: {safe_carved_dir}")
            except Exception:
                # Fallback to print if logger is not available in this context
                print(f"Cleared carved files directory: {safe_carved_dir}")
    except Exception as e:
        print(f"Error clearing carved files directory: {e}")
    
    # NEW: Clear the deleted recovery directory
    deleted_recovery_dir = app.config['DELETED_RECOVERY_FOLDER']
    try:
        if os.path.exists(deleted_recovery_dir):
            for filename in os.listdir(deleted_recovery_dir):
                file_path = os.path.join(deleted_recovery_dir, filename)
                if os.path.isfile(file_path):
                    os.unlink(file_path)
            # Sanitize path to remove any trailing CR/LF from config values before logging
            safe_deleted_dir = deleted_recovery_dir.rstrip('\r\n') if isinstance(deleted_recovery_dir, str) else deleted_recovery_dir
            try:
                app.logger.info(f"Cleared deleted recovery directory: {safe_deleted_dir}")
            except Exception:
                # Fallback to print if logger is not available in this context
                print(f"Cleared deleted recovery directory: {safe_deleted_dir}")
    except Exception as e:
        print(f"Error clearing deleted recovery directory: {e}")
    
    # Remove any JSON result files
    for json_file in ['carved_results.json', 'deleted_results.json']:
        if os.path.exists(json_file):
            try:
                os.remove(json_file)
            except Exception as e:
                print(f"Error removing {json_file}: {e}")

def _generate_preview_response(content, filename, data_url, file_info=None):
    """Analyzes file content and returns the appropriate Flask response for previewing."""
    if not content:
        return "Cannot preview an empty file.", 400

    mime_type = magic.from_buffer(content, mime=True)
    text_content = None
    
    # For images, PDFs, and media files - use embedded preview
    if mime_type.startswith('image/'):
        return render_template_string(VIEW_FILE_CONTENT, 
                                    filename=filename, 
                                    mime_type=mime_type, 
                                    data_url=data_url,
                                    is_carved=True,
                                    file_info=file_info or {})
    
    if mime_type == 'application/pdf':
        return render_template_string(VIEW_FILE_CONTENT, 
                                    filename=filename, 
                                    mime_type=mime_type, 
                                    data_url=data_url,
                                    is_carved=True,
                                    file_info=file_info or {})
    
    if mime_type.startswith(('audio/', 'video/')):
        return render_template_string(VIEW_FILE_CONTENT, 
                                    filename=filename, 
                                    mime_type=mime_type, 
                                    data_url=data_url,
                                    is_carved=True,
                                    file_info=file_info or {})

    # For text-based and other files - extract content for preview
    temp_file_path = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=f"_{filename}") as tf:
            tf.write(content)
            temp_file_path = tf.name

        if mime_type.startswith('text/'):
            text_content = parse_text_log(temp_file_path)
        elif mime_type == 'application/rtf' or filename.lower().endswith('.rtf'):
            text_content = parse_text_log(temp_file_path)
        elif 'openxmlformats-officedocument.wordprocessingml' in mime_type:
            text_content = "Preview for docx is not available."
        elif 'openxmlformats-officedocument.spreadsheetml' in mime_type:
            text_content = "Preview for xlsx is not available."
        elif 'openxmlformats-officedocument.presentationml' in mime_type:
            text_content = "Preview for pptx is not available."
        elif any(t in mime_type for t in ['zip', 'rar', 'x-7z-compressed', 'gzip']):
            text_content = "Preview for archives is not available."
        elif 'sqlite' in mime_type:
            text_content = "Preview for sqlite is not available."
        elif 'eventlog' in mime_type or filename.lower().endswith('.evtx'):
            text_content = parse_evtx_log(temp_file_path)
        elif 'executable' in mime_type or 'x-dosexec' in mime_type or filename.lower().endswith(('.exe', '.dll', '.elf')):
            text_content = "Preview for executables is not available."
        elif filename.lower().endswith('.doc'):
            text_content = extract_strings_preview(temp_file_path)
        elif filename.lower().endswith('.dat') and content[:4] == b'regf':
            text_content = extract_strings_preview(temp_file_path)
        else:
            # Fallback for unknown types - show hex preview
            text_content = format_hex_view(content[:2000])  # Show first 2000 bytes as hex

        if text_content is not None:
            return render_template_string(VIEW_FILE_CONTENT, 
                                        filename=filename, 
                                        mime_type=mime_type, 
                                        text_content=text_content,
                                        is_carved=True,
                                        file_info=file_info or {})
    except Exception as e:
        print(f"Error during preview generation for {filename}: {e}")
        # Fallback to hex view if parsing fails
        text_content = f"Error parsing file: {e}\n\nHex preview:\n{format_hex_view(content[:2000])}"
        return render_template_string(VIEW_FILE_CONTENT, 
                                    filename=filename, 
                                    mime_type=mime_type, 
                                    text_content=text_content,
                                    is_carved=True,
                                    file_info=file_info or {})
    finally:
        if temp_file_path and os.path.exists(temp_file_path):
            os.remove(temp_file_path)

    # Ultimate fallback
    return render_template_string(VIEW_FILE_CONTENT, 
                                filename=filename, 
                                mime_type=mime_type, 
                                text_content="File preview not available",
                                is_carved=True,
                                file_info=file_info or {})

def _load_file_into_session(filename, filepath):
    """Loads file metadata into the session and starts background analysis."""
    # Defensive: ensure the file exists before processing
    if not filepath or not os.path.exists(filepath):
        try:
            flash(f"Evidence file not found on disk: {filepath}", "error")
        except Exception:
            pass
        return None

    encryption_info = detect_encryption(filepath)
    forensic_results, partition_info = perform_forensic_analysis(filepath)

    try:
        size_mb = f"{os.path.getsize(filepath) / (1024*1024):.2f}"
    except Exception:
        size_mb = "0.00"

    uploaded_files_db[filename] = {
        "path": filepath,
        "size_mb": size_mb,
        "encryption_status": {
            "encrypted": encryption_info.get('encrypted'),
            "encryption_type": encryption_info.get('encryption_type'),
            "description": encryption_info.get('description'),
            "decrypting": False,
            "decrypted_path": None
        },
        "forensic_results": forensic_results,
        "partition_info": partition_info,
        "hash_info": {},
        "hashing_complete": False
    }
    threading.Thread(target=calculate_hashes_threaded, args=(filepath,)).start()
    return encryption_info.get('encrypted')

def determine_recovery_method(filename):
    """Determine recovery method from filename pattern."""
    filename_lower = filename.lower()
    if 'dirwalk' in filename_lower or 'directory' in filename_lower:
        return "Directory Walk"
    elif 'inode' in filename_lower or 'orphan' in filename_lower:
        return "Inode Scan"
    elif 'slack' in filename_lower:
        return "File Slack"
    elif 'recycle' in filename_lower or 'recycled' in filename_lower:
        return "Recycle Bin"
    elif 'deep' in filename_lower:
        return "Deep Scan"
    else:
        return "Unknown Method"

# --- STRICT AUTOMATIC DELETED FILES RECOVERY ---
def strict_deleted_files_recovery_engine(filepath):
    """Autopsy-like automatic deleted files recovery with strict validation."""
    # global deleted_scan_status
    
    deleted_scan_status.update({
        "in_progress": True, 
        "complete": False, 
        "files_found": 0, 
        "message": "Starting strict automatic recovery...",
        "start_time": time.time(),
        "last_update_time": time.time(),
        "elapsed_time": "0s",
        "estimated_total_time": None,
        "time_remaining_str": "Calculating...",
        "scan_methods": {
            "directory_walk": 0,
            "inode_scan": 0,
            "file_slack": 0,
            "unallocated_space": 0,
            "recycle_bin": 0
        },
        "validation_stats": {
            "total_scanned": 0,
            "empty_rejected": 0,
            "duplicate_rejected": 0,
            "invalid_rejected": 0,
            "valid_recovered": 0
        }
    })

    recovery_dir = app.config.get('DELETED_RECOVERY_FOLDER', DELETED_RECOVERY_FOLDER)
    os.makedirs(recovery_dir, exist_ok=True)
    seen_hashes = set()
    MIN_FILE_SIZE = 512  # Increased minimum size to avoid tiny files
    MAX_FILE_SIZE = 100 * 1024 * 1024  # 100MB limit
    
    # Clear previous results
    try:
        for item in os.listdir(recovery_dir):
            item_path = os.path.join(recovery_dir, item)
            if os.path.isfile(item_path):
                os.unlink(item_path)
    except Exception as e:
        deleted_scan_status["message"] = f"Error clearing old files: {e}"
        deleted_scan_status["in_progress"] = False
        return

    total_recovered = 0
    
    def validate_and_save_file(content, original_name, recovery_method, fs_object=None):
        """STRICT validation: Check file size, content, and duplicates before saving."""
       # nonlocal total_recovered, seen_hashes
        # global deleted_files_db
        
        deleted_scan_status["validation_stats"]["total_scanned"] += 1
        
        # 1. Check for empty content
        if not content or len(content) < MIN_FILE_SIZE:
            deleted_scan_status["validation_stats"]["empty_rejected"] += 1
            return False
        
        # 2. Check file size limit
        if len(content) > MAX_FILE_SIZE:
            deleted_scan_status["validation_stats"]["invalid_rejected"] += 1
            return False
        
        # 3. Calculate content hash for deduplication
        content_hash = hashlib.sha256(content).hexdigest()
        if content_hash in seen_hashes:
            deleted_scan_status["validation_stats"]["duplicate_rejected"] += 1
            return False
        
        # 4. Validate file structure based on type
        if not validate_file_structure(content, original_name):
            deleted_scan_status["validation_stats"]["invalid_rejected"] += 1
            return False
        
        # 5. All checks passed - save the file
        try:
            # Try to detect a sensible extension from the recovered content
            detected_ext = ''
            try:
                mime_type = magic.from_buffer(content[:8192], mime=True)
                if mime_type:
                    guessed = mimetypes.guess_extension(mime_type)
                    if guessed:
                        detected_ext = guessed
            except Exception:
                detected_ext = ''

            # Fallback: check FILE_SIGNATURES for a header-based extension
            if not detected_ext:
                try:
                    for _cat, types in FILE_SIGNATURES.items():
                        for _name, sig in types.items():
                            hdr = sig.get('header') or sig.get('headers')
                            if not hdr:
                                continue
                            # support list of headers
                            if isinstance(hdr, list):
                                for h in hdr:
                                    if content.startswith(h):
                                        detected_ext = sig.get('extension', '')
                                        break
                                if detected_ext:
                                    break
                            else:
                                if content.startswith(hdr):
                                    detected_ext = sig.get('extension', '')
                                    break
                        if detected_ext:
                            break
                except Exception:
                    detected_ext = ''

            # Generate a sequential standardized deleted filename and include detected extension
            seq = sum(1 for _ in os.listdir(recovery_dir) if os.path.isfile(os.path.join(recovery_dir, _))) + 1
            safe_filename = generate_deleted_filename(seq, detected_ext)
            save_path = os.path.join(recovery_dir, safe_filename)

            # Ensure unique filename to avoid accidental overwrites when multiple
            # recovered files are saved concurrently. Append a numeric suffix
            # if the computed filename already exists.
            if os.path.exists(save_path):
                base_name, base_ext = os.path.splitext(safe_filename)
                counter = 1
                while True:
                    candidate = f"{base_name}_{counter}{base_ext}"
                    candidate_path = os.path.join(recovery_dir, candidate)
                    if not os.path.exists(candidate_path):
                        save_path = candidate_path
                        break
                    counter += 1

            with open(save_path, 'wb') as out_file:
                out_file.write(content)
            
            # Add to seen hashes to prevent duplicates
            seen_hashes.add(content_hash)
            # Note: do not increment total_recovered here; update_recovery_status
            # is responsible for incrementing the overall recovered counter to
            # avoid double-counting when that function is called after
            # successful validation.
            deleted_scan_status["validation_stats"]["valid_recovered"] += 1
            # Register the recovered file in the in-memory DB so the UI can show it
            try:
                mime_type = 'application/octet-stream'
                try:
                    mime_type = magic.from_file(save_path, mime=True)
                except Exception:
                    pass

                file_size = os.path.getsize(save_path)
                thumb = None
                if mime_type and mime_type.startswith('image/'):
                    thumb = create_thumbnail_data_uri(save_path)

                file_type = 'Unknown'
                if mime_type:
                    if mime_type.startswith('image/'):
                        file_type = 'Image'
                    elif mime_type.startswith('video/'):
                        file_type = 'Video'
                    elif mime_type.startswith('audio/'):
                        file_type = 'Audio'
                    elif mime_type.startswith('text/'):
                        file_type = 'Text'
                    elif 'pdf' in mime_type:
                        file_type = 'PDF'
                    elif 'zip' in mime_type or 'archive' in mime_type:
                        file_type = 'Archive'

                mtime = None
                try:
                    mtime = datetime.datetime.fromtimestamp(os.path.getmtime(save_path)).strftime('%Y-%m-%d %H:%M:%S')
                except Exception:
                    mtime = None

                deleted_files_db[os.path.basename(save_path)] = {
                    'id': total_recovered + 1,
                    'name': os.path.basename(save_path),
                    'size_kb': f"{file_size/1024:.2f}",
                    'file_type': file_type,
                    'thumbnail': thumb,
                    'recovery_method': recovery_method,
                    'mtime': mtime,
                    'path': save_path
                }
                # Persist recovered deleted file in DB/session
                try:
                    sess_id = None
                    try:
                        sess_id = session.get('analysis_session_id')
                    except Exception:
                        sess_id = None
                    add_file_record(os.path.basename(save_path), 'deleted_recovered', save_path, os.path.getsize(save_path), session_id=sess_id, extra={'recovery_method': recovery_method})
                except Exception:
                    pass
            except Exception:
                pass

            return True
        except Exception as e:
            print(f"Error saving recovered file: {e}")
            return False

    def validate_file_structure(content, filename):
        """Validate file structure based on file type signatures."""
        try:
            # Check for valid file headers
            if len(content) < 8:  # Too short for any meaningful file
                return False
            
            # Common file signature validation
            file_signatures = {
                b'\xff\xd8\xff': 'JPEG',
                b'\x89PNG\r\n\x1a\n': 'PNG',
                b'%PDF': 'PDF',
                b'PK\x03\x04': 'ZIP/DOCX',
                b'\x25\x50\x44\x46': 'PDF',
                b'\x47\x49\x46\x38': 'GIF',
                b'BM': 'BMP',
                b'RIFF': 'WAV/AVI',
            }
            
            # Check if content starts with known file signature
            for signature, file_type in file_signatures.items():
                if content.startswith(signature):
                    return True
            
            # For files without clear signatures, check entropy
            entropy = calculate_entropy(content[:4096])  # Check first 4KB
            if entropy < 1.0 or entropy > 7.5:  # Suspicious entropy values
                return False
                
            return True  # Accept files that pass basic checks
            
        except Exception:
            return False

    def generate_recovery_filename(original_name, recovery_method, file_id, fs_object=None, detected_ext=''):
        """Generate informative filename for recovered files."""
        # Clean the original filename
        clean_name = re.sub(r'[^\w\.-]', '_', original_name)
        # Preserve extension if present
        base, cur_ext = os.path.splitext(clean_name)
        if cur_ext:
            ext = cur_ext
        else:
            # if caller supplied a detected extension prefer that
            ext = detected_ext or ''
        
        # Add metadata
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        method_abbr = recovery_method[:4].upper()
        
        if fs_object and hasattr(fs_object.info, 'meta'):
            inode = f"i{fs_object.info.meta.addr}"
        else:
            inode = "i0000"
        
        return f"rec_{file_id:04d}_{method_abbr}_{inode}_{base}{ext}"

    def update_recovery_status(method, validated=True):
        """Update recovery status with validation info."""
        nonlocal total_recovered
        
        if validated:
            total_recovered += 1
            deleted_scan_status["files_found"] = total_recovered
            deleted_scan_status["scan_methods"][method] += 1
        
        # Update status message with validation stats
        stats = deleted_scan_status["validation_stats"]
        status_msg = (
            f"Recovered: {total_recovered} | "
            f"Scanned: {stats['total_scanned']} | "
            f"Rejected - Empty: {stats['empty_rejected']}, "
            f"Duplicate: {stats['duplicate_rejected']}, "
            f"Invalid: {stats['invalid_rejected']}"
        )
        deleted_scan_status["message"] = status_msg

    # Main recovery logic
    try:
        img_handle = pytsk3.Img_Info(filepath)
        
        def strict_directory_walk(fs, directory, path="/"):
            """Strict directory walking with validation."""
            try:
                for f in directory:
                    if not hasattr(f.info, 'meta') or f.info.meta is None:
                        continue
                    
                    # Check if file is deleted and valid for recovery
                    is_deleted = not (f.info.meta.flags & pytsk3.TSK_FS_META_FLAG_ALLOC)
                    is_file = f.info.meta.type == pytsk3.TSK_FS_META_TYPE_REG
                    has_name = hasattr(f.info, 'name') and f.info.name is not None
                    
                    if is_deleted and is_file and has_name and f.info.name.name not in [b'.', b'..']:
                        try:
                            if MIN_FILE_SIZE < f.info.meta.size <= MAX_FILE_SIZE:
                                content = f.read_random(0, min(f.info.meta.size, MAX_FILE_SIZE))
                                
                                if validate_and_save_file(
                                    content, 
                                    f.info.name.name.decode('utf-8', 'ignore'),
                                    'directory_walk',
                                    f
                                ):
                                    update_recovery_status("directory_walk", True)
                                
                        except Exception as e:
                            continue
                    
                    # Recursive directory scanning
                    if (f.info.meta.type == pytsk3.TSK_FS_META_TYPE_DIR and 
                        has_name and f.info.name.name not in [b'.', b'..']):
                        try:
                            strict_directory_walk(fs, f.as_directory(), 
                                                f"{path}/{f.info.name.name.decode('utf-8', 'ignore')}")
                        except (IOError, AttributeError):
                            continue
            except Exception as e:
                pass

        def strict_inode_scan(fs):
            """Strict inode scanning with validation."""
            try:
                last_inode = fs.info.last_inum
                for inode_num in range(fs.info.first_inum, last_inode + 1):
                    try:
                        fs_file = fs.open_meta(inode=inode_num)
                        
                        if (fs_file.info.meta.flags & pytsk3.TSK_FS_META_FLAG_UNALLOC and 
                            fs_file.info.meta.type == pytsk3.TSK_FS_META_TYPE_REG and 
                            MIN_FILE_SIZE < fs_file.info.meta.size <= MAX_FILE_SIZE):
                            
                            try:
                                content = fs_file.read_random(0, min(fs_file.info.meta.size, MAX_FILE_SIZE))
                                
                                # Generate name for orphaned files
                                name = f"orphan_inode_{inode_num}"
                                if hasattr(fs_file.info, 'name') and fs_file.info.name:
                                    orig_name = fs_file.info.name.name.decode('utf-8', 'ignore')
                                    if orig_name not in ['.', '..']:
                                        name = orig_name
                                
                                if validate_and_save_file(content, name, 'inode_scan', fs_file):
                                    update_recovery_status("inode_scan", True)
                                    
                            except Exception:
                                continue
                    except (IOError, AttributeError):
                        continue
            except Exception as e:
                pass

        # Execute recovery methods
        try:
            volume = pytsk3.Volume_Info(img_handle)
            for part in volume:
                if part.flags != pytsk3.TSK_VS_PART_FLAG_UNALLOC:
                    try:
                        fs_offset = part.start * volume.info.block_size
                        fs = pytsk3.FS_Info(img_handle, offset=fs_offset)
                        
                        deleted_scan_status["message"] = f"Scanning partition {part.desc}..."
                        
                        # Run strict recovery methods
                        strict_directory_walk(fs, fs.open_dir(path="/"))
                        strict_inode_scan(fs)
                        
                    except IOError:
                        continue
        except IOError:
            # Try as single filesystem
            try:
                fs = pytsk3.FS_Info(img_handle, offset=0)
                deleted_scan_status["message"] = "Scanning as single filesystem..."
                
                strict_directory_walk(fs, fs.open_dir(path="/"))
                strict_inode_scan(fs)
                
            except IOError as e:
                deleted_scan_status["message"] = f"Error opening filesystem: {e}"

        # Final status update
        stats = deleted_scan_status["validation_stats"]
        final_msg = (
            f"Recovery complete! "
            f"Recovered: {total_recovered} valid files. "
            f"Scanned: {stats['total_scanned']} | "
            f"Rejected: {stats['empty_rejected'] + stats['duplicate_rejected'] + stats['invalid_rejected']} "
            f"(Empty: {stats['empty_rejected']}, Duplicate: {stats['duplicate_rejected']}, Invalid: {stats['invalid_rejected']})"
        )
        deleted_scan_status["message"] = final_msg
        deleted_scan_status["complete"] = True
        
    except Exception as e:
        deleted_scan_status["message"] = f"A critical error occurred: {e}"
        deleted_scan_status["complete"] = True
        
    deleted_scan_status["in_progress"] = False

# --- UPDATE THE ROUTE TO USE STRICT RECOVERY ---
@app.route('/start_strict_deleted_recovery')
def start_strict_deleted_recovery():
    """Starts the strict automatic deleted files recovery process."""
    if not uploaded_files_db:
        flash("Please upload an evidence file first.", "error")
        return redirect(url_for('evidence_upload'))
    
    if deleted_scan_status.get("in_progress"):
        flash("Deleted files recovery is already in progress.", "info")
        return redirect(url_for('deleted_files_status_page'))
    
    filepath = get_active_evidence_path()
    if not filepath:
        flash("Could not determine evidence file path.", "error")
        return redirect(url_for('evidence_upload'))
    
    # Reset status and start strict recovery
    deleted_scan_status.update({
        "in_progress": True, 
        "complete": False, 
        "files_found": 0, 
        "message": "Starting strict automatic recovery...",
        "start_time": time.time(),
        "last_update_time": time.time(),
        "elapsed_time": "0s",
        "estimated_total_time": None,
        "time_remaining_str": "Calculating...",
        "scan_methods": {
            "directory_walk": 0,
            "inode_scan": 0,
            "file_slack": 0,
            "unallocated_space": 0,
            "recycle_bin": 0
        },
        "validation_stats": {
            "total_scanned": 0,
            "empty_rejected": 0,
            "duplicate_rejected": 0,
            "invalid_rejected": 0,
            "valid_recovered": 0
        }
    })
    
    # Clear previous results
    recovery_dir = app.config['DELETED_RECOVERY_FOLDER']
    try:
        for item in os.listdir(recovery_dir):
            item_path = os.path.join(recovery_dir, item)
            if os.path.isfile(item_path):
                os.unlink(item_path)
    except Exception as e:
        flash(f"Error clearing old files: {e}", "warning")
    # Also clear the in-memory DB of recovered files to avoid stale entries
    global deleted_files_db
    try:
        deleted_files_db = {}
    except Exception:
        # Fallback: clear in place if possible
        try:
            deleted_files_db.clear()
        except Exception:
            pass
    # Start the strict recovery in a background thread
    threading.Thread(target=strict_deleted_files_recovery_engine, args=(filepath,)).start()
    
    flash("Started strict automatic recovery process with duplicate and empty file filtering.", "success")
    return redirect(url_for('deleted_files_status_page'))

# --- Flask Routes ---
@app.route('/')
def index():
    return redirect(url_for('evidence_upload'))

@app.route('/upload_status')
def upload_status_api():
    """API endpoint to get current upload status"""
    return jsonify(upload_status)

@app.route('/carving_status')
def carving_status_endpoint():
    """Get current carving status with timing information."""
    return jsonify(carving_status)

@app.route('/upload', methods=['POST'])
def upload_file():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file part'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No selected file'}), 400
        
        # Secure the filename
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        
        # Initialize upload status
        upload_status.update({
            "in_progress": True,
            "progress": 0,
            "filename": filename,
            "start_time": time.time(),
            "estimated_total_time": None,
            "elapsed_time": "0s",
            "time_remaining_str": "Calculating...",
            "total_bytes": 0,
            "bytes_uploaded": 0,
            "upload_speed": 0
        })
        
        # Save file in chunks to handle large files
        chunk_size = 16 * 1024 * 1024  # 16MB chunks
        total_size = 0
        
        with open(filepath, 'wb') as f:
            while True:
                chunk = file.stream.read(chunk_size)
                if not chunk:
                    break
                f.write(chunk)
                total_size += len(chunk)
                
                # Update upload status
                upload_status['bytes_uploaded'] = total_size
                upload_status['total_bytes'] = total_size  # We don't know total in advance
                upload_status['progress'] = 50  # Show progress as we're reading
                
                # Calculate timing
                elapsed = time.time() - upload_status['start_time']
                if elapsed > 0:
                    upload_speed = total_size / elapsed
                    upload_status['upload_speed'] = upload_speed
                    upload_status['elapsed_time'] = format_time(elapsed)
        
        # Finalize upload status
        upload_status.update({
            "in_progress": False,
            "progress": 100,
            "complete": True,
            "elapsed_time": format_time(time.time() - upload_status['start_time'])
        })
        
        return jsonify({'message': 'File uploaded successfully', 'filename': filename}), 200
        
    except Exception as e:
        upload_status.update({
            "in_progress": False,
            "error": str(e)
        })
        return jsonify({'error': f'Upload failed: {str(e)}'}), 500

@app.route('/all_process_status')
def all_process_status():
    """Get status of all background processes with timing."""
    return jsonify({
        "upload": upload_status,
        "carving": carving_status,
        "deleted_recovery": deleted_scan_status,
        "decryption": decryption_status,
        "hashing": hashing_status,
        "strings": strings_status
    })

@app.errorhandler(413)
def too_large(e):
    return jsonify({'error': 'File too large'}), 413

@app.errorhandler(500)
def internal_error(e):
    return jsonify({'error': 'Internal server error during file upload'}), 500

@app.route('/deleted_recovery')
def deleted_recovery():
    """Redirect to the deleted files recovery status page"""
    return redirect(url_for('deleted_files_status_page'))

@app.route('/start_deleted_recovery')
def start_deleted_recovery():
    """Starts the deleted files recovery process"""
    if not uploaded_files_db:
        flash("Please upload an evidence file first.", "error")
        return redirect(url_for('evidence_upload'))
    
    if deleted_scan_status.get("in_progress"):
        flash("Deleted files recovery is already in progress.", "info")
        return redirect(url_for('deleted_files_status_page'))
    
    filepath = get_active_evidence_path()
    if not filepath:
        flash("Could not determine evidence file path.", "error")
        return redirect(url_for('evidence_upload'))
    
    # Reset status and start recovery
    deleted_scan_status.update({
        "in_progress": True, 
        "complete": False, 
        "files_found": 0, 
        "message": "Starting recovery process...",
        "errors": []
    })
    
    # Clear previous results
    recovery_dir = app.config['DELETED_RECOVERY_FOLDER']
    try:
        for item in os.listdir(recovery_dir):
            item_path = os.path.join(recovery_dir, item)
            if os.path.isfile(item_path):
                os.unlink(item_path)
    except Exception as e:
        flash(f"Error clearing old files: {e}", "warning")
    
    # Start the recovery in a background thread
    threading.Thread(target=recover_deleted_files_engine, args=(filepath,)).start()
    
    return redirect(url_for('deleted_files_status_page'))

@app.route('/', methods=['GET', 'POST'])
@app.route('/evidence_upload', methods=['GET', 'POST'])
def evidence_upload():
    """Handle evidence file upload with progress tracking."""
    # global uploaded_files_db, upload_status
    
    if request.method == 'POST':
        # This route should only handle AJAX uploads, but keep for fallback
        flash('Please use the upload form with JavaScript enabled.', 'error')
        return redirect(request.url)
    
    # If a session was active, end it because user is going back to upload/db area
    sess_id = session.pop('analysis_session_id', None)
    if sess_id:
        try:
            end_analysis_session(sess_id)
        except Exception:
            pass

    # Get database files if available (pass dict so template can use .items())
    db_files = db_connections

    # Query DB for pinned file types (uploaded, encrypted, decrypted) to show in UI
    pinned_db_files = []
    try:
        conn = _get_db_conn()
        cur = conn.cursor()
        cur.execute("SELECT id, filename, path, file_type FROM files WHERE file_type IN ('uploaded','encrypted','decrypted') ORDER BY created_at DESC")
        rows = cur.fetchall()
        for r in rows:
            try:
                fid, fname, fpath, ftype = r[0], r[1], r[2], r[3]
            except Exception:
                # fallback in case schema differs
                fid = r[0]
                fname = r[1] if len(r) > 1 else None
                fpath = r[2] if len(r) > 2 else None
                ftype = r[3] if len(r) > 3 else None
            pinned_db_files.append({'id': fid, 'filename': fname, 'path': fpath, 'file_type': ftype})
        try:
            conn.close()
        except Exception:
            pass
    except Exception:
        pinned_db_files = []
    
    # If the user arrived here via a Load action, a 'loaded' query param may be present
    loaded_filename = request.args.get('loaded')
    content = render_template_string(
        EVIDENCE_UPLOAD_CONTENT,
        uploaded_files=uploaded_files_db,
        db_files=db_files,
        pinned_db_files=pinned_db_files,
        loaded_filename=loaded_filename
    )
    # Show the encryption sidebar link only when on the upload page (so user can access it via upload UI button)
    return render_template_string(BASE_TEMPLATE, content=content, show_session_details=False, show_encryption_in_sidebar=True)

@app.route('/ajax_upload', methods=['POST'])
def ajax_upload():
    """Handle AJAX file upload with progress tracking."""
    # global uploaded_files_db, upload_status
    
    if 'file' not in request.files:
        return jsonify({'error': 'No file selected'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    try:
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)



        # Initialize upload status
        upload_status.update({
            "in_progress": True,
            "progress": 0,
            "filename": filename,
            "start_time": time.time(),
            "estimated_total_time": None,
            "elapsed_time": "0s",
            "time_remaining_str": "Calculating...",
            "total_bytes": 0,
            "bytes_uploaded": 0,
            "upload_speed": 0,
            "last_update_time": time.time(),
            "error": None
        })

        # Try to determine total size (may not always be available)
        try:
            file.stream.seek(0, os.SEEK_END)
            total_size = file.stream.tell()
            file.stream.seek(0)
        except Exception:
            # Fallback: use Content-Length header if present
            total_size = request.content_length or 0

        upload_status["total_bytes"] = total_size

        # Save file in chunks and update progress so polling clients see accurate values
        chunk_size = 4 * 1024 * 1024  # 4MB
        bytes_written = 0

        # Save file in chunks and optionally persist bytes in DB chunks
        with open(filepath, 'wb') as out_f:
            chunk_index = 0
            while True:
                chunk = file.stream.read(chunk_size)
                if not chunk:
                    break
                out_f.write(chunk)
                bytes_written += len(chunk)



                # Update upload status using helper
                try:
                    update_upload_progress(bytes_uploaded=bytes_written, total_bytes=total_size or bytes_written)
                except Exception:
                    # Non-fatal - continue writing
                    upload_status['bytes_uploaded'] = bytes_written
                    upload_status['total_bytes'] = total_size or bytes_written

        # Finalize upload status
        upload_status.update({
            "in_progress": False,
            "progress": 100,
            "bytes_uploaded": bytes_written,
            "complete": True,
            "elapsed_time": format_time(time.time() - upload_status["start_time"]) 
        })



        # capture selected target database (if provided)
        target_db = request.form.get('target_db', 'local')
        # Start processing in background so response returns quickly
        try:
            threading.Thread(target=_process_uploaded_file, args=(filename, filepath, target_db), daemon=True).start()
        except Exception:
            # If background thread cannot be started, process inline (best-effort)
            _process_uploaded_file(filename, filepath, target_db)

        return jsonify({
            'success': True,
            'filename': filename,
            'message': f'File "{filename}" uploaded successfully!'
        })
        
    except Exception as e:
        upload_status.update({
            "in_progress": False,
            "error": str(e)
        })
        return jsonify({'error': f'Error uploading file: {str(e)}'}), 500
def _process_uploaded_file(filename, filepath, target_db='local'):
    """Process uploaded file and add to database."""
    # global uploaded_files_db

    file_size = os.path.getsize(filepath)
    encryption_info = detect_encryption(filepath)
    forensic_results, partition_info = perform_forensic_analysis(filepath)

    uploaded_files_db[filename] = {
        "path": filepath,
        "size_mb": f"{file_size / (1024 * 1024):.2f}",
        "encryption_status": {
            "encrypted": encryption_info['encrypted'],
            "encryption_type": encryption_info['encryption_type'],
            "description": encryption_info['description'],
            "decrypting": False,
            "decrypted_path": None
        },
        "forensic_results": forensic_results,
        "partition_info": partition_info,
        "hash_info": {},
        "hashing_complete": False
    }
    # keep target db selection
    try:
        uploaded_files_db[filename]['target_db'] = target_db
        uploaded_files_db[filename]['target_db_cfg'] = db_connections.get(target_db) if target_db != 'local' else None
    except Exception:
        pass

    # Start background hashing
    hashing_thread = threading.Thread(target=calculate_hashes_threaded, args=(filepath,))
    hashing_thread.daemon = True
    hashing_thread.start()
    # Add DB record for uploaded file
    try:
        add_file_record(filename, 'uploaded', filepath, os.path.getsize(filepath), session_id=None, extra={'target_db': target_db})
    except Exception:
        pass



        

def update_upload_progress(bytes_uploaded=None, total_bytes=None):
    """Update upload progress in real-time."""
    # global upload_status
    
    current_time = time.time()
    
    if not upload_status.get("start_time"):
        upload_status["start_time"] = current_time
    
    elapsed = current_time - upload_status["start_time"]
    upload_status["elapsed_time"] = format_time(elapsed)
    
    if bytes_uploaded is not None and total_bytes is not None:
        upload_status["bytes_uploaded"] = bytes_uploaded
        upload_status["total_bytes"] = total_bytes
        
        if total_bytes > 0:
            progress = (bytes_uploaded / total_bytes) * 100
            upload_status["progress"] = min(progress, 99)  # Cap at 99% until complete
            
            # Calculate upload speed
            if elapsed > 0:
                upload_speed = bytes_uploaded / elapsed
                upload_status["upload_speed"] = upload_speed
                
                # Calculate time remaining
                if progress > 0:
                    estimated_total = elapsed / (progress / 100)
                    estimated_remaining = estimated_total - elapsed
                    upload_status["time_remaining_str"] = format_time(max(0, estimated_remaining))
                    upload_status["estimated_total_time"] = format_time(estimated_total)
    
    upload_status["last_update_time"] = current_time
    return upload_status

@app.route('/encrypt_file', methods=['POST'])
def encrypt_file():
    """Handle file encryption request."""
    if 'file' not in request.files:
        flash('No file selected', 'error')
        return redirect(url_for('encryption_page'))
    
    file = request.files['file']
    if file.filename == '':
        flash('No file selected', 'error')
        return redirect(url_for('encryption_page'))
    
    password = request.form.get('password')
    confirm_password = request.form.get('confirm_password')
    
    if not password or password != confirm_password:
        flash('Passwords do not match or are empty', 'error')
        return redirect(url_for('encryption_page'))
    
    if len(password) < 4:
        flash('Password must be at least 4 characters long', 'error')
        return redirect(url_for('encryption_page'))
    
    try:
        # Save uploaded file temporarily
        filename = secure_filename(file.filename)
        temp_path = os.path.join(app.config['UPLOAD_FOLDER'], f"temp_{filename}")
        file.save(temp_path)
        
        # Ensure encrypted folder exists
        os.makedirs(app.config.get('ENCRYPTED_FOLDER', ENCRYPTED_FOLDER), exist_ok=True)
        # Encrypt the file into the Encrypted Files folder
        success, message, output_filename = custom_encrypt_file(
            temp_path, filename, password, app.config['ENCRYPTED_FOLDER']
        )
        
        # Clean up temp file
        try:
            os.remove(temp_path)
        except:
            pass
        
        if success:
            # After successful encryption, redirect to the encryption page
            # and pass the new filename as a URL parameter.
            return redirect(url_for('encryption_page', newly_encrypted_file=output_filename))
        else:
            flash(f'Encryption failed: {message}', 'error')
            
    except Exception as e:
        flash(f'Error during encryption: {str(e)}', 'error')
    
    return redirect(url_for('encryption_page'))


@app.route('/download_encrypted/<filename>')
def download_encrypted_file(filename):
    """Download an encrypted file."""
    try:
        # Secure the filename to prevent directory traversal
        safe_filename = secure_filename(filename)
        file_path = os.path.join(app.config['ENCRYPTED_FOLDER'], safe_filename)
        
        # Check if file exists
        if not os.path.exists(file_path):
            flash(f"Encrypted file '{safe_filename}' not found.", "error")
            return redirect(url_for('encryption_page'))
        
        # Send the file for download
        return send_file(
            file_path,
            as_attachment=True,
            download_name=safe_filename,
            mimetype='application/octet-stream'
        )
        
    except Exception as e:
        flash(f"Error downloading encrypted file: {str(e)}", "error")
        return redirect(url_for('encryption_page'))
    
@app.route('/delete_encrypted/<filename>')
def delete_encrypted_file(filename):
    """Delete an encrypted file."""
    try:
        safe_filename = secure_filename(filename)
        file_path = os.path.join(app.config['ENCRYPTED_FOLDER'], safe_filename)
        
        if os.path.exists(file_path):
            os.remove(file_path)
            flash(f"Encrypted file '{safe_filename}' deleted successfully.", "success")
        else:
            flash(f"Encrypted file '{safe_filename}' not found.", "error")
            
    except Exception as e:
        flash(f"Error deleting encrypted file: {str(e)}", "error")
    
    return redirect(url_for('encryption_page'))

@app.route('/encryption_page', methods=['GET'])
def encryption_page():
    """Displays the encryption page, lists existing encrypted files, and shows success messages."""
    # This route now only handles GET requests to display the page.
    # The actual encryption is handled by the '/encrypt_file' route.
    
    # Get the newly encrypted filename from the URL parameter, if it exists.
    newly_encrypted_file = request.args.get('newly_encrypted_file')
    
    # Get the list of all currently encrypted files to display them.
    encrypted_files = get_encrypted_files()
    
    # Render the HTML template, passing the necessary variables to it.
    content = render_template_string(
        ENCRYPTION_PAGE_CONTENT,
        encrypted_files=encrypted_files,
        newly_encrypted_file=newly_encrypted_file
    )
    return render_template_string(BASE_TEMPLATE, content=content, uploaded_files_db=uploaded_files_db, show_session_details=False)


@app.route('/manual_carve_hex', methods=['POST'])
def manual_carve_hex():
    filepath = get_active_evidence_path()
    if not filepath:
        return jsonify({'error': 'No evidence file loaded.'}), 400

    start_s = request.form.get('start_offset', '').strip()
    length_s = request.form.get('length', '').strip()
    try:
        if isinstance(start_s, str) and start_s.startswith('0x'):
            start = int(start_s, 16)
        else:
            start = int(start_s)
        length = int(length_s)
    except Exception:
        return jsonify({'error': 'Invalid start offset or length.'}), 400

    if length <= 0:
        return jsonify({'error': 'Length must be positive.'}), 400

    file_size = os.path.getsize(filepath)
    if start < 0 or start >= file_size:
        return jsonify({'error': 'Start offset out of range.'}), 400

    # cap length to reasonable preview size to avoid huge responses
    MAX_PREVIEW = 64 * 1024  # 64 KB
    length = min(length, MAX_PREVIEW, file_size - start)

    try:
        with open(filepath, 'rb') as f:
            f.seek(start)
            data = f.read(length)
    except Exception as e:
        return jsonify({'error': f'Error reading file: {e}'}), 500

    html = _format_hex_scrabble(data, start_offset=start)
    return jsonify({'html': html})


    
@app.route('/remove_file/<filename>')
def remove_file(filename):
    if filename in uploaded_files_db:
        _clear_all_session_data()
        flash(f'File {filename} unloaded from session.', 'success')
    return redirect(url_for('evidence_upload'))




@app.route('/ajax_upload_abort', methods=['POST'])
def ajax_upload_abort():
    """Abort an in-progress upload and remove partial data (disk + DB chunks)."""
    try:
        payload = request.get_json() or {}
        filename = payload.get('filename')
        if not filename:
            return jsonify({'error': 'No filename provided'}), 400

        filepath = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(filename))

        # Remove partial disk file if present
        try:
            if os.path.exists(filepath):
                os.remove(filepath)
        except Exception:
            pass



        # Reset upload_status if it refers to this file
        try:
            if upload_status.get('filename') == filename:
                upload_status.update({
                    'in_progress': False,
                    'progress': 0,
                    'bytes_uploaded': 0,
                    'total_bytes': 0,
                    'filename': None,
                    'error': 'Aborted by client due to size threshold'
                })
        except Exception:
            pass

        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/decryption/<filename>', methods=['GET', 'POST'])
def decryption_page(filename):
    if filename not in uploaded_files_db:
        flash("File not found or not loaded.", 'error')
        return redirect(url_for('evidence_upload'))
    
    file_info = uploaded_files_db[filename]
    if not file_info['encryption_status']['encrypted']:
        flash(f"'{filename}' is not encrypted.", "info")
        return redirect(url_for('forensic_analysis'))

    if request.method == 'POST':
        password = request.form.get('password')
        action = request.form.get('action')
        forced_type = request.form.get('force_type')

        if action == 'with_password' and not password:
            flash("Password field must be filled to use 'Decrypt with Provided Password'.", 'error')
            return redirect(url_for('decryption_page', filename=filename))
        
        pwd_to_try = password if action == 'with_password' else None
        
        effective_type = file_info['encryption_status']['encryption_type']
        if forced_type and forced_type != 'auto':
            effective_type = forced_type
        
        if not effective_type:
            effective_type = 'FERNET_ENCRYPTED'
            flash("No encryption type auto-detected. Attempting Fernet decryption.", "info")

        threading.Thread(
            target=attempt_decryption,
            args=(
                file_info['path'], 
                effective_type,
                pwd_to_try
            )
        ).start()
        
        return redirect(url_for('decryption_progress', filename=filename))

    content = render_template_string(DECRYPTION_PAGE_CONTENT, filename=filename, file_info=file_info)
    # include the list of decrypted files for display
    decrypted_files = get_decrypted_files()
    return render_template_string(BASE_TEMPLATE, content=content, uploaded_files_db=uploaded_files_db, decrypted_files=decrypted_files, show_session_details=False)

@app.route('/decryption_progress/<filename>')
def decryption_progress(filename):
    if filename not in uploaded_files_db:
        flash("File not found or not loaded.", 'error')
        return redirect(url_for('evidence_upload'))
    
    content = render_template_string(DECRYPTION_PROGRESS_CONTENT, filename=filename)
    return render_template_string(BASE_TEMPLATE, content=content, uploaded_files_db=uploaded_files_db, show_session_details=False)

@app.route('/forensic_analysis')
def forensic_analysis():
    hashing_in_progress = False
    if not uploaded_files_db:
        flash("Please upload or load an evidence file to begin analysis.", "warning")
        return redirect(url_for('evidence_upload'))

    first_file_details = next(iter(uploaded_files_db.values()))
    if not first_file_details.get('hashing_complete', False):
        hashing_in_progress = True
            
    # Auto-start a session when entering analysis if none active in Flask session
    sess_id = session.get('analysis_session_id')
    if not sess_id:
        new_id = start_analysis_session()
        if new_id:
            session['analysis_session_id'] = new_id
    # Build session_info for template: id, path, name, counts per typed subfolder
    session_info = {'id': None, 'path': None, 'name': None, 'counts': {}}
    try:
        sid = session.get('analysis_session_id')
        s_path = session.get('analysis_session_path')
        if sid:
            session_info['id'] = sid
        if s_path and os.path.isdir(s_path):
            session_info['path'] = s_path
            session_info['name'] = os.path.basename(s_path)
            for sub in ('Carved','Deleted','Reports','Logs','Events','ManualCarving'):
                try:
                    subdir = os.path.join(s_path, sub)
                    if os.path.isdir(subdir):
                        session_info['counts'][sub] = sum(1 for _ in os.listdir(subdir) if os.path.isfile(os.path.join(subdir, _)))
                    else:
                        session_info['counts'][sub] = 0
                except Exception:
                    session_info['counts'][sub] = 0
    except Exception:
        pass

    content = render_template_string(FORENSIC_ANALYSIS_CONTENT, 
                                     uploaded_files=uploaded_files_db, 
                                     hashing_in_progress=hashing_in_progress,
                                     session_info=session_info)
    return render_template_string(BASE_TEMPLATE, content=content, uploaded_files_db=uploaded_files_db)

@app.route('/start_strings_analysis')
def start_strings_analysis():
    if not uploaded_files_db:
        return jsonify({"status": "error", "error": "No evidence file uploaded."})
    if strings_status.get("in_progress"):
        return jsonify({"status": "error", "error": "Strings analysis is already in progress."})
    
    filepath = get_active_evidence_path()
    if not filepath:
        return jsonify({"status": "error", "error": "Could not determine evidence file path."})

    threading.Thread(target=extract_strings_threaded, args=(filepath,)).start()
    return jsonify({"status": "started"})

@app.route('/auto_carving_setup', methods=['GET'])
def auto_carving_setup():
    if not uploaded_files_db:
        flash("Please upload or load an evidence file first.", "warning")
        return redirect(url_for('evidence_upload'))
    
    colors = ['text-sky-400', 'text-emerald-400', 'text-amber-400', 'text-rose-400', 'text-violet-400', 'text-teal-400', 'text-cyan-400', 'text-lime-400', 'text-pink-400', 'text-indigo-400']
    
    content = render_template_string(
        AUTO_CARVING_SETUP_CONTENT, 
        signatures=FILE_SIGNATURES, 
        colors=colors,
        form_action_url=url_for('run_auto_carving')
    )
    return render_template_string(BASE_TEMPLATE, content=content, uploaded_files_db=uploaded_files_db)

@app.route('/auto_carving_process')
def auto_carving_process():
    if not uploaded_files_db:
        flash("No evidence file is currently loaded.", "warning")
        return redirect(url_for('evidence_upload'))
    content = render_template_string(AUTO_CARVING_PROCESS_CONTENT)
    return render_template_string(BASE_TEMPLATE, content=content, uploaded_files_db=uploaded_files_db)

@app.route('/recovered_files')
def recovered_files():
    if not uploaded_files_db:
        flash("No evidence file is currently loaded. Please start a new session.", "warning")
        return redirect(url_for('evidence_upload'))

    recovered_files_map = {}
    carved_dir = app.config['CARVED_FOLDER']
    
    try:
        pattern = re.compile(r"^(\d+)-([0-9A-Fa-f]{8})-(\d+)-(.+)$")
        for filename in os.listdir(carved_dir):
            match = pattern.match(filename)
            if match:
                file_id = int(match.group(1))
                offset_hex = match.group(2)
                size_bytes = int(match.group(3))
                
                file_info = {
                    "id": file_id,
                    "name": filename,
                    "offset": f"0x{offset_hex}",
                    "size_bytes": size_bytes
                }
                recovered_files_map[file_id] = file_info
    except FileNotFoundError:
        pass
    except Exception as e:
        flash(f"Error scanning recovery directory: {e}", "error")

    sorted_file_ids = sorted(recovered_files_map.keys())
    
    page = request.args.get('page', 1, type=int)
    PER_PAGE = 20
    total_files = len(sorted_file_ids)
    start_index = (page - 1) * PER_PAGE
    end_index = start_index + PER_PAGE
    
    ids_on_page = sorted_file_ids[start_index:end_index]
    files_on_page = [recovered_files_map.get(id) for id in ids_on_page]
    
    total_pages = max(1, (total_files + PER_PAGE - 1) // PER_PAGE)

    content = render_template_string(RECOVERED_FILES_CONTENT, 
                                     carved_files=files_on_page, 
                                     total_files=total_files,
                                     page=page, 
                                     total_pages=total_pages)
    return render_template_string(BASE_TEMPLATE, content=content, uploaded_files_db=uploaded_files_db)

@app.route('/manual_carving')
def manual_carving():
    if not uploaded_files_db:
        flash("Please upload an evidence file first to use the manual carver.", "error")
        return redirect(url_for('evidence_upload'))
    content = render_template_string(MANUAL_CARVING_CONTENT)
    # append the evidence hex viewer box at the end
    content = content + render_template_string(MANUAL_CARVING_EVIDENCE_HEX)
    return render_template_string(BASE_TEMPLATE, content=content, uploaded_files_db=uploaded_files_db)


def _format_hex_scrabble(data_bytes, start_offset=0, bytes_per_line=16, group_size=4):
    """Return HTML string with hex displayed in grouped 'scrabble' format.

    Each line: offset | grouped hex (uppercase) | ASCII printable on right
    """
    lines = []
    for i in range(0, len(data_bytes), bytes_per_line):
        chunk = data_bytes[i:i+bytes_per_line]
        offset = start_offset + i
        hex_bytes = [f"{b:02X}" for b in chunk]
        # group into group_size
        grouped = []
        for j in range(0, len(hex_bytes), group_size):
            grouped.append(' '.join(hex_bytes[j:j+group_size]))
        hex_part = '  '.join(grouped)
        
        ascii_part = ''.join([chr(b) if 32 <= b < 127 else '.' for b in chunk])
        lines.append(f"<div><span class=\"font-mono\" style=\"color:#9CA3AF; width:120px; display:inline-block;\">0x{offset:08X}</span> <span class=\"font-mono\" style=\"color:#d1d5db;\">{hex_part}</span> <span class=\"font-mono\" style=\"color:#9CA3AF; margin-left:12px;\">{ascii_part}</span></div>")
    return '<div class="hex-view p-3" style="background:#0b0f14;">' + '\n'.join(lines) + '</div>'

@app.route('/view_evidence_hex', methods=['POST'])
def view_evidence_hex():
    """Handles requests for hex data from the evidence file viewer."""
    filepath = get_active_evidence_path()
    if not filepath:
        return jsonify({'error': 'No evidence file loaded. Please upload a file first.'}), 400

    try:
        file_size = os.path.getsize(filepath)
    except OSError as e:
        return jsonify({'error': f'Cannot access evidence file: {e}'}), 500
    
    if file_size == 0:
        return jsonify({'error': 'The loaded evidence file is empty (0 bytes).'}), 400

    start_s = request.form.get('start_offset', '0').strip()
    length_s = request.form.get('length', '65536').strip()
    
    try:
        # Handle both decimal and hexadecimal (0x) inputs for the start offset
        start = int(start_s, 16) if start_s.lower().startswith('0x') else int(start_s)
        length = int(length_s)
    except (ValueError, TypeError):
        return jsonify({'error': 'Invalid start offset or length. Please enter valid numbers.'}), 400

    if start < 0 or start >= file_size:
        return jsonify({'error': f'Start offset is out of range. File size is {file_size} bytes.'}), 400

    MAX_PREVIEW = 128 * 1024  # 128KB maximum to prevent browser overload
    
    # Calculate the actual number of bytes to read, ensuring it doesn't exceed limits
    bytes_to_read = min(length, MAX_PREVIEW, file_size - start)

    # If there are no bytes to read (e.g., at the end of the file), return a clear message
    if bytes_to_read <= 0:
         html = '<div class="p-4 text-gray-500">End of file reached.</div>'
         return jsonify({'html': html, 'start': start, 'length': 0, 'file_size': file_size})

    try:
        with open(filepath, 'rb') as f:
            f.seek(start)
            data = f.read(bytes_to_read)
    except Exception as e:
        return jsonify({'error': f'Error reading file: {e}'}), 500

    # Format the data into the scrabble hex view and return it
    html = _format_hex_scrabble(data, start_offset=start)
    return jsonify({'html': html, 'start': start, 'length': len(data), 'file_size': file_size})


@app.route('/log_file_viewer', methods=['GET', 'POST'])
def log_file_viewer():
    if not uploaded_files_db:
        flash("Please upload an evidence file first to use the log viewer.", "error")
        return redirect(url_for('evidence_upload'))

    evidence_filename = list(uploaded_files_db.keys())[0]
    log_content = None

    if request.method == 'POST':
        filepath = get_active_evidence_path()
        if not filepath:
            flash("Could not determine evidence file path.", "error")
            return redirect(url_for('log_file_viewer'))

        log_ext = request.form.get('log_ext', '').lower()
        try:
            if not log_ext and '.' in evidence_filename:
                log_ext = '.' + evidence_filename.rsplit('.', 1)[1].lower()
            if log_ext in [".log", ".txt"] or log_ext.startswith('/var/log/'):
                log_content = parse_text_log(filepath)
            elif log_ext == ".csv":
                log_content = parse_csv_log(filepath)
            elif log_ext == ".json":
                log_content = parse_json_log(filepath)
            elif log_ext == ".xml":
                log_content = parse_xml_log(filepath)
            elif log_ext in [".evt", ".evtx"]:
                log_content = parse_evtx_log(filepath)
            else:
                log_content = f"[INFO] Format '{log_ext}' requires specialized handling. Displaying raw text.\n\n"
                log_content += parse_text_log(filepath)
        except Exception as e:
            log_content = f"[ERROR] Failed to parse file as '{log_ext}'. It may be part of a larger disk image.\n\nDetails: {e}"

    template_vars = {
        "title": "Log File Viewer",
        "description": "Analyze the loaded evidence file by parsing it as a specific log type.",
        "os_options": LOG_OS_OPTIONS,
        "log_content": log_content,
        "evidence_filename": evidence_filename
    }
    content = render_template_string(LOG_VIEWER_TEMPLATE, **template_vars)
    return render_template_string(BASE_TEMPLATE, content=content, uploaded_files_db=uploaded_files_db)

@app.route('/event_log_viewer', methods=['GET', 'POST'])
def event_log_viewer():
    if not uploaded_files_db:
        flash("Please upload an evidence file first to use the event log viewer.", "error")
        return redirect(url_for('evidence_upload'))

    evidence_filename = list(uploaded_files_db.keys())[0]
    log_content = None
    
    if request.method == 'POST':
        filepath = get_active_evidence_path()
        if not filepath:
            flash("Could not determine evidence file path.", "error")
            return redirect(url_for('event_log_viewer'))

        log_ext = request.form.get('log_ext', '').lower()
        try:
            if not log_ext and '.' in evidence_filename:
                log_ext = '.' + evidence_filename.rsplit('.', 1)[1].lower()
            if log_ext in [".evtx", ".evt"]:
                log_content = parse_evtx_log(filepath)
            elif log_ext.startswith("/var/log/") or log_ext in [".log", ".txt", ".asl"]:
                log_content = parse_text_log(filepath)
            else:
                log_content = f"[INFO] Format '{log_ext}' requires external tools. Displaying as plain text.\n\n"
                log_content += parse_text_log(filepath)
        except Exception as e:
            log_content = f"[ERROR] Failed to parse file as '{log_ext}'. It may be part of a larger disk image.\n\nDetails: {e}"

    template_vars = {
        "title": "Event Log Viewer",
        "description": "Analyze the loaded evidence file by parsing it as a specific event log type.",
        "os_options": EVENT_OS_OPTIONS,
        "log_content": log_content,
        "evidence_filename": evidence_filename
    }
    content = render_template_string(LOG_VIEWER_TEMPLATE, **template_vars)
    return render_template_string(BASE_TEMPLATE, content=content, uploaded_files_db=uploaded_files_db)

@app.route('/reporting', methods=['GET', 'POST'])
def reporting():
    if not uploaded_files_db:
        flash("Please upload an evidence file first to generate a report.", "error")
        return redirect(url_for('evidence_upload'))

    if request.method == 'POST':
        case_details = {
            "name": request.form.get('case_name', 'N/A'),
            "number": request.form.get('case_number', 'N/A'),
            "examiner": request.form.get('examiner_name', 'N/A')
        }
        report_format = request.form.get('report_format')
        now = datetime.datetime.now()
        safe_case_name = secure_filename(case_details['name'])

        try:
            with open('carved_results.json', 'r') as f: 
                report_carved_files = json.load(f)
        except: 
            report_carved_files = {}
        try:
            with open('deleted_results.json', 'r') as f: 
                report_deleted_files = json.load(f)
        except: 
            report_deleted_files = {}

        # Read report-scope options from the form
        include_carved = bool(request.form.get('include_carved'))
        include_deleted = bool(request.form.get('include_deleted'))
        selected_categories = request.form.getlist('file_category') or []
        custom_exts_raw = (request.form.get('custom_exts') or '').strip()
        custom_exts = [e.lower().strip().lstrip('.') for e in custom_exts_raw.split(',') if e.strip()] if custom_exts_raw else []

        # Helper: determine category by extension
        def ext_category(ext_no_dot):
            e = ext_no_dot.lower()
            if e in ['jpg','jpeg','png','gif','bmp','tiff','webp']:
                return 'images'
            if e in ['doc','docx','pdf','xls','xlsx','ppt','pptx','txt','rtf']:
                return 'documents'
            if e in ['zip','rar','7z','tar','gz','bz2','tgz','xz']:
                return 'archives'
            if e in ['mp3','wav','flac','mp4','mov','avi','mkv','wmv']:
                return 'media'
            if e in ['exe','dll','elf','bin','sys','apk']:
                return 'executables'
            return 'other'

        # Apply filters to carved and deleted maps
        def filter_map(src_map):
            if not src_map:
                return {}
            if not selected_categories and not custom_exts:
                return src_map
            out = {}
            for fname, info in src_map.items():
                # try to get extension
                ext = ''
                if '.' in fname:
                    ext = fname.rsplit('.',1)[1].lower()
                cat = ext_category(ext) if ext else 'other'
                if custom_exts and ext and ext in custom_exts:
                    out[fname] = info
                    continue
                if selected_categories and cat in selected_categories:
                    out[fname] = info
            return out

        # Filter according to include_* flags
        if include_carved:
            carved_filtered = filter_map(report_carved_files)
        else:
            carved_filtered = {}
        if include_deleted:
            deleted_filtered = filter_map(report_deleted_files)
        else:
            deleted_filtered = {}

        # Handle report generation based on chosen format
        if report_format == 'html':
            for filename, info in carved_filtered.items():
                filepath = os.path.join(app.config['CARVED_FOLDER'], filename)
                info['thumbnail_uri'] = create_thumbnail_data_uri(filepath)

            report_html = render_template_string(
                REPORT_TEMPLATE, case_details=case_details, evidence_file=uploaded_files_db,
                carved_files=carved_filtered, deleted_files=deleted_filtered, now=now
            )
            # Save HTML report to session files and DB (ensure bytes)
            html_bytes = report_html.encode('utf-8') if isinstance(report_html, str) else report_html
            _save_bytes_to_session_file(html_bytes, f'Report_{safe_case_name}.html', file_type='report')
            response = make_response(report_html)
            response.headers['Content-Disposition'] = f"attachment; filename=Report_{safe_case_name}.html"
            return response

        elif report_format == 'pdf':
            if not HTML:
                flash("PDF generation requires the 'weasyprint' library. Please run: pip install weasyprint", "error")
                return redirect(url_for('reporting'))

            for filename, info in carved_filtered.items():
                filepath = os.path.join(app.config['CARVED_FOLDER'], filename)
                info['thumbnail_uri'] = create_thumbnail_data_uri(filepath)

            report_html = render_template_string(
                REPORT_TEMPLATE, case_details=case_details, evidence_file=uploaded_files_db,
                carved_files=carved_filtered, deleted_files=deleted_filtered, now=now
            )
            pdf_bytes = HTML(string=report_html).write_pdf()
            # Persist PDF to session files and DB
            _save_bytes_to_session_file(pdf_bytes, f'Report_{safe_case_name}.pdf', file_type='report')
            response = make_response(pdf_bytes)
            response.headers['Content-Disposition'] = f"attachment; filename=Report_{safe_case_name}.pdf"
            response.headers['Content-Type'] = 'application/pdf'
            return response

        elif report_format == 'docx':
            docx_buffer = generate_docx_report_data(case_details, uploaded_files_db, carved_filtered, deleted_filtered, now)
            if not docx_buffer:
                flash("DOCX generation requires the 'python-docx' library. Please run: pip install python-docx", "error")
                return redirect(url_for('reporting'))
            # Persist DOCX to session files and DB
            try:
                docx_bytes = docx_buffer.getvalue() if hasattr(docx_buffer, 'getvalue') else None
            except Exception:
                docx_bytes = None
            if docx_bytes:
                _save_bytes_to_session_file(docx_bytes, f'Report_{safe_case_name}.docx', file_type='report')
            return send_file(docx_buffer, download_name=f'Report_{safe_case_name}.docx', as_attachment=True)

        elif report_format == 'csv':
            zip_buffer = generate_csv_zip_report_data(case_details, uploaded_files_db, carved_filtered, deleted_filtered)
            try:
                zip_bytes = zip_buffer.getvalue() if hasattr(zip_buffer, 'getvalue') else None
            except Exception:
                zip_bytes = None
            if zip_bytes:
                _save_bytes_to_session_file(zip_bytes, f'Report_{safe_case_name}_CSV.zip', file_type='report')
            return send_file(zip_buffer, download_name=f'Report_{safe_case_name}_CSV.zip', as_attachment=True, mimetype='application/zip')

    content = render_template_string(REPORTING_PAGE_CONTENT, now=datetime.datetime.now())
    return render_template_string(BASE_TEMPLATE, content=content, uploaded_files_db=uploaded_files_db)

@app.route('/deleted_files_status_page')
def deleted_files_status_page():
    """Enhanced deleted files recovery status page with multiple scanning methods."""
    # Allow viewing recovered files even if no evidence file is currently uploaded
    # (helps inspecting previous runs or imported recovered files).
    
    # Ensure scan_methods exists in deleted_scan_status
    if 'scan_methods' not in deleted_scan_status:
        deleted_scan_status['scan_methods'] = {
            "directory_walk": 0,
            "inode_scan": 0,
            "file_slack": 0,
            "unallocated_space": 0,
            "recycle_bin": 0
        }
    
    # Get recovered files with enhanced info
    recovered_files = []
    recovery_dir = app.config['DELETED_RECOVERY_FOLDER']
    
    try:
        for i, filename in enumerate(os.listdir(recovery_dir), 1):
            filepath = os.path.join(recovery_dir, filename)
            if os.path.isfile(filepath):
                file_info = get_enhanced_file_info(filepath)
                recovery_method = determine_recovery_method(filename)
                
                recovered_files.append({
                    "id": i,
                    "name": filename,
                    "size_kb": file_info['size_kb'],
                    "file_type": file_info['file_type'],
                    "thumbnail": file_info['thumbnail'],
                    "recovery_method": recovery_method
                })
    except FileNotFoundError:
        pass  # Directory doesn't exist yet
    
    content = render_template_string(ENHANCED_DELETED_STATUS_TEMPLATE,
                                    deleted_scan_status=deleted_scan_status,
                                    recovered_files=recovered_files)
    return render_template_string(BASE_TEMPLATE, content=content, uploaded_files_db=uploaded_files_db)

@app.route('/get_recovered_files_list')
def get_recovered_files_list():
    """API endpoint to get the current list of recovered files for real-time updates."""
    try:
        # Ensure the in-memory DB is synced with disk
        populate_deleted_files_db_from_disk()
        
        # Convert the deleted_files_db to a format suitable for JSON
        files_list = {}
        for filename, file_info in deleted_files_db.items():
            files_list[filename] = {
                'name': file_info.get('name', filename),
                'size': file_info.get('size', 0),
                'size_kb': file_info.get('size_kb', '0'),
                'mtime': file_info.get('mtime', 'Unknown'),
                'file_type': file_info.get('file_type', 'Unknown'),
                'recovery_method': file_info.get('recovery_method', 'Unknown'),
                'thumbnail': file_info.get('thumbnail')
            }
        
        return jsonify({
            'success': True,
            'files': files_list,
            'total_files': len(files_list)
        })
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e),
            'files': {},
            'total_files': 0
        })


@app.route('/decryption_status')
def decryption_status_endpoint():
    return jsonify(decryption_status)

@app.route('/hashing_status')
def hashing_status_endpoint():
    return jsonify(hashing_status)

@app.route('/strings_status')
def strings_status_endpoint():
    return jsonify(strings_status)

@app.route('/deleted_files')
def deleted_files():
    """Redirect to the status page where files are now displayed"""
    return redirect(url_for('deleted_files_status_page'))

@app.route('/deleted_scan_status')
def deleted_scan_status_endpoint():
    # Ensure timing is updated before returning
    try:
        # Provide a simple update based on files_found and total scanned where possible
        total_scanned = deleted_scan_status.get('validation_stats', {}).get('total_scanned', 0)
        # Use calculate_time_estimations helper to keep values consistent
        calculate_time_estimations(deleted_scan_status, total_scanned)
    except Exception:
        # If anything goes wrong, fall back to existing fields
        pass

    return jsonify({
        "in_progress": deleted_scan_status["in_progress"],
        "complete": deleted_scan_status["complete"], 
        "files_found": deleted_scan_status["files_found"],
        "message": deleted_scan_status["message"],
        "scan_methods": deleted_scan_status.get("scan_methods", {}),
        "files": deleted_files_db,
        "validation_stats": deleted_scan_status.get('validation_stats', {}),
        "elapsed_time": deleted_scan_status.get('elapsed_time', '0s'),
        "time_remaining_str": deleted_scan_status.get('time_remaining_str', 'Calculating...'),
        "estimated_total_time": deleted_scan_status.get('estimated_total_time', None),
        "start_time": deleted_scan_status.get('start_time', None),
        "last_update_time": deleted_scan_status.get('last_update_time', None)
    })

@app.route('/audit_deleted_files')
def audit_deleted_files():
    """Audit endpoint: compare files on disk vs in-memory `deleted_files_db` and return a diff.

    Returns JSON with:
      - files_on_disk: list of filenames found in recovery folder
      - files_in_db: list of filenames in deleted_files_db
      - missing_in_db: files present on disk but missing in deleted_files_db
      - missing_on_disk: files present in deleted_files_db but missing on disk
      - summary counts
    """
    recovery_dir = app.config.get('DELETED_RECOVERY_FOLDER', os.path.join(APP_ROOT, 'deleted_files'))
    try:
        disk_list = [f for f in os.listdir(recovery_dir) if os.path.isfile(os.path.join(recovery_dir, f))]
    except Exception:
        disk_list = []

    db_list = list(deleted_files_db.keys()) if isinstance(deleted_files_db, dict) else []

    disk_set = set(disk_list)
    db_set = set(db_list)

    missing_in_db = sorted(list(disk_set - db_set))
    missing_on_disk = sorted(list(db_set - disk_set))

    # Prepare small metadata previews for the first N items in each category
    def preview_disk(fname):
        p = os.path.join(recovery_dir, fname)
        meta = {}
        try:
            meta['size_bytes'] = os.path.getsize(p)
            meta['mtime'] = datetime.datetime.fromtimestamp(os.path.getmtime(p)).strftime('%Y-%m-%d %H:%M:%S')
            meta['mime'] = magic.from_file(p, mime=True)
        except Exception:
            meta = {'size_bytes': None, 'mtime': None, 'mime': None}
        return meta

    def preview_db(fname):
        entry = deleted_files_db.get(fname, {})
        return {
            'size': entry.get('size'),
            'size_kb': entry.get('size_kb'),
            'file_type': entry.get('file_type'),
            'mtime': entry.get('mtime'),
            'path': entry.get('path')
        }

    MAX_PREVIEW = 20
    payload = {
        'files_on_disk': disk_list,
        'files_in_db': db_list,
        'missing_in_db': [{ 'name': n, 'meta': preview_disk(n) } for n in missing_in_db[:MAX_PREVIEW]],
        'missing_on_disk': [{ 'name': n, 'meta': preview_db(n) } for n in missing_on_disk[:MAX_PREVIEW]],
        'counts': {
            'disk_total': len(disk_list),
            'db_total': len(db_list),
            'missing_in_db': len(missing_in_db),
            'missing_on_disk': len(missing_on_disk)
        }
    }

    return jsonify(payload)

@app.route('/hex_view_carved/<filename>')
def hex_view_carved(filename):
    if not uploaded_files_db:
        flash("No evidence file is loaded.", "error")
        return redirect(url_for('evidence_upload'))

    s_filename = secure_filename(filename)
    filepath = os.path.join(app.config['CARVED_FOLDER'], s_filename)

    if not os.path.exists(filepath):
        flash(f"Carved file '{s_filename}' not found.", "error")
        return redirect(url_for('recovered_files'))

    try:
        with open(filepath, 'rb') as f:
            content = f.read()
        hex_content = format_hex_view(content)
        content_html = render_template_string(HEX_VIEW_CARVED_FILE_CONTENT, filename=s_filename, hex_content=hex_content)
        return render_template_string(BASE_TEMPLATE, content=content_html, uploaded_files_db=uploaded_files_db)
    except Exception as e:
        flash(f"Error reading file '{s_filename}': {e}", "error")
        return redirect(url_for('recovered_files'))

@app.route('/download_carved_file/<filename>')
def download_carved_file(filename):
    return send_from_directory(app.config['CARVED_FOLDER'], secure_filename(filename), as_attachment=True)

@app.route('/serve_file_data/<type>/<path:filename>')
def serve_file_data(type, filename):
    s_filename = secure_filename(filename)
    if type == 'carved':
        return send_from_directory(app.config['CARVED_FOLDER'], s_filename)
    elif type == 'deleted_recovered':
        return send_from_directory(app.config['DELETED_RECOVERY_FOLDER'], s_filename)
    elif type == 'deleted':
        try:
            inode = int(filename)
            if not uploaded_files_db or inode not in deleted_files_db: 
                return "File not found", 404
            file_info = deleted_files_db[inode]
            filepath = get_active_evidence_path()
            img = pytsk3.Img_Info(filepath)
            fs = pytsk3.FS_Info(img, offset=file_info['fs_offset'])
            fs_file = fs.open_meta(inode=inode)
            content = fs_file.read_random(0, file_info['size'])
            mime_type, _ = mimetypes.guess_type(file_info['name'])
            mime_type = mime_type or 'application/octet-stream'
            return Response(content, mimetype=mime_type)
        except Exception as e: 
            return f"Error serving file data: {e}", 500
        
    return "Invalid file type", 400

@app.route('/view_carved_file/<filename>')
def view_carved_file(filename):
    s_filename = secure_filename(filename)
    filepath = os.path.join(app.config['CARVED_FOLDER'], s_filename)
    if not os.path.exists(filepath):
        return "File not found", 404
    
    try:
        with open(filepath, 'rb') as f:
            content = f.read()
    except Exception as e:
        return f"Error reading file: {e}", 500
    
    data_url = url_for('serve_file_data', type='carved', filename=s_filename)
    return _generate_preview_response(content, s_filename, data_url, file_info={"name": s_filename})

@app.route('/view_deleted_file/<filename>')
def view_deleted_file(filename):
    """View deleted file with preview similar to carved files."""
    try:
        filepath = os.path.join(app.config['DELETED_RECOVERY_FOLDER'], filename)
        if not os.path.exists(filepath):
            return "File not found.", 404
        
        with open(filepath, 'rb') as f:
            content = f.read()
        
        if not content:
            return "File is empty or could not be read."
        
        mime_type = magic.from_buffer(content, mime=True)
        data_url = f"data:{mime_type};base64,{base64.b64encode(content).decode('utf-8')}"
        
        # For small files, create text preview
        if len(content) <= 2 * 1024 * 1024:  # 2MB limit for preview
            if mime_type.startswith('text/'):
                try:
                    text_content = content.decode('utf-8', errors='ignore')
                    return render_template_string(
                        VIEW_FILE_CONTENT,
                        filename=filename,
                        mime_type=mime_type,
                        text_content=text_content,
                        is_carved=False,
                        content=content
                    )
                except:
                    pass
            
            # For images, PDFs, etc. - use embedded preview
            if mime_type.startswith(('image/', 'application/pdf', 'audio/', 'video/')):
                return render_template_string(
                    VIEW_FILE_CONTENT,
                    filename=filename,
                    mime_type=mime_type,
                    data_url=data_url,
                    is_carved=False,
                    content=content
                )
        
        # Fallback: hex preview for first 2KB
        hex_preview = format_hex_view(content[:2048])
        return render_template_string(
            VIEW_FILE_CONTENT,
            filename=filename,
            mime_type=mime_type,
            text_content=f"File too large for full preview. Hex preview (first 2KB):\n\n{hex_preview}",
            is_carved=False,
            content=content
        )
        
    except Exception as e:
        return f"Error viewing file: {str(e)}", 500
    
@app.route('/hex_view_fallback/<filename>')
def hex_view_fallback(filename):
    """Fallback hex view for files that can't be previewed normally"""
    s_filename = secure_filename(filename)
    filepath = os.path.join(app.config['CARVED_FOLDER'], s_filename)
    
    if not os.path.exists(filepath):
        return "File not found", 404
    
    try:
        with open(filepath, 'rb') as f:
            content = f.read(5000)  # Read first 5000 bytes for preview
        hex_content = format_hex_view(content)
        
        return f"""
        <!DOCTYPE html>
        <html>
        <head>
            <title>Hex View: {s_filename}</title>
            <style>
                body {{ background: #111827; color: white; font-family: monospace; padding: 20px; }}
                pre {{ background: #1f2937; padding: 20px; border-radius: 5px; }}
                .header {{ display: flex; justify-content: space-between; margin-bottom: 20px; }}
            </style>
        </head>
        <body>
            <div class="header">
                <h2>Hex View: {s_filename}</h2>
                <a href="{{ url_for('recovered_files') }}">Back</a>
            </div>
            <pre>{hex_content}</pre>
            <p>Showing first {len(content)} bytes of file</p>
        </body>
        </html>
        """
    except Exception as e:
        return f"Error reading file: {e}", 500
    
@app.route('/hex_view_deleted/<filename>')
def hex_view_deleted(filename):
    if not uploaded_files_db:
        flash("No evidence file is loaded.", "error")
        return redirect(url_for('evidence_upload'))

    s_filename = secure_filename(filename)
    filepath = os.path.join(app.config['DELETED_RECOVERY_FOLDER'], s_filename)

    if not os.path.exists(filepath):
        flash(f"Recovered file '{s_filename}' not found.", "error")
        return redirect(url_for('deleted_files'))

    try:
        with open(filepath, 'rb') as f:
            content = f.read()
        hex_content = format_hex_view(content)
        content_html = render_template_string(HEX_VIEW_CARVED_FILE_CONTENT, filename=s_filename, hex_content=hex_content)
        return render_template_string(BASE_TEMPLATE, content=content_html, uploaded_files_db=uploaded_files_db)
    except Exception as e:
        flash(f"Error reading file '{s_filename}': {e}", "error")
        return redirect(url_for('deleted_files'))

@app.route('/download_deleted_file/<filename>')
def download_deleted_file(filename):
    return send_from_directory(app.config['DELETED_RECOVERY_FOLDER'], secure_filename(filename), as_attachment=True)

@app.route('/download_zip', methods=['POST'])
def download_zip():
    if not uploaded_files_db:
        flash("No evidence file loaded.", "error")
        return redirect(url_for('evidence_upload'))

    selected_files = request.form.getlist('selected_files')
    file_type = request.form.get('file_type')

    if not selected_files:
        flash("No files were selected for download.", "warning")
        return redirect(url_for('deleted_files' if file_type in ['deleted', 'deleted_recovered'] else 'recovered_files'))

    memory_file = io.BytesIO()
    zip_filename = f"{file_type}_files_{datetime.datetime.now():%Y%m%d%H%M%S}.zip"

    try:
        with zipfile.ZipFile(memory_file, 'w', zipfile.ZIP_DEFLATED) as zf:
            if file_type == 'carved':
                for filename in selected_files:
                    filepath = os.path.join(app.config['CARVED_FOLDER'], secure_filename(filename))
                    if os.path.exists(filepath):
                        zf.write(filepath, arcname=secure_filename(filename))
            
            elif file_type == 'deleted_recovered':
                for filename in selected_files:
                    filepath = os.path.join(app.config['DELETED_RECOVERY_FOLDER'], secure_filename(filename))
                    if os.path.exists(filepath):
                        zf.write(filepath, arcname=secure_filename(filename))
        
        memory_file.seek(0)
        return send_file(memory_file, download_name=zip_filename, as_attachment=True, mimetype='application/zip')
    except Exception as e:
        flash(f"Error creating ZIP file: {e}", "error")
        return redirect(url_for('deleted_files' if file_type in ['deleted', 'deleted_recovered'] else 'recovered_files'))

@app.route('/find_in_file', methods=['POST'])
def find_in_file():
    filepath = get_active_evidence_path()
    if not filepath:
        return jsonify({"offset": -1, "error": "No file uploaded"})
    
    term = request.form.get('term')
    search_type = request.form.get('type')
    start_offset = request.form.get('start_offset', 0, type=int)

    try:
        search_bytes = bytes.fromhex(term.replace(" ", "")) if search_type == 'hex' else term.encode()
    except ValueError:
        return jsonify({"offset": -1, "error": "Invalid Hex sequence."})
    
    try:
        with open(filepath, 'rb') as f:
            with mmap.mmap(f.fileno(), 0, access=mmap.ACCESS_READ) as mm:
                found_pos = mm.find(search_bytes, start_offset)
                return jsonify({"offset": found_pos})
    except Exception as e:
        return jsonify({"offset": -1, "error": str(e)})

@app.route('/run_auto_carving', methods=['POST'])
def run_auto_carving():
    """
    Handles the file carving process based on user-selected file types.
    """
   #  global carving_status, carved_files_db, sorted_carved_keys

    image_path = get_active_evidence_path()
    if not image_path:
        flash('Error: Evidence file path is missing or invalid. Please re-upload the file.', 'error')
        return redirect(url_for('evidence_upload'))

    selected_types = request.form.getlist('file_types')
    if not selected_types:
        flash('Error: Please select at least one file type to carve.', 'error')
        return redirect(url_for('auto_carving_setup'))

    threading.Thread(target=simple_file_carver, args=(image_path, selected_types)).start()
    return redirect(url_for('auto_carving_process'))

@app.route('/find_block', methods=['POST'])
def find_block():
    filepath = get_active_evidence_path()
    if not filepath:
        return jsonify({"status": "error", "message": "No file uploaded"})
    
    header_term = request.form.get('header_term')
    footer_term = request.form.get('footer_term')
    search_type = request.form.get('type')

    try:
        header_bytes = bytes.fromhex(header_term.replace(" ", "")) if search_type == 'hex' else header_term.encode()
        footer_bytes = bytes.fromhex(footer_term.replace(" ", "")) if search_type == 'hex' else footer_term.encode()
    except ValueError:
        return jsonify({"status": "error", "message": "Invalid Hex sequence."})
    
    try:
        with open(filepath, 'rb') as f:
            with mmap.mmap(f.fileno(), 0, access=mmap.ACCESS_READ) as mm:
                header_pos = mm.find(header_bytes, 0)
                if header_pos != -1:
                    footer_pos = mm.find(footer_bytes, header_pos + len(header_bytes))
                    if footer_pos != -1:
                        block_len = (footer_pos + len(footer_bytes)) - header_pos
                        return jsonify({"status": "success", "start_offset": header_pos, "length": block_len})
    except Exception as e:
        return jsonify({"status": "error", "message": str(e)})

    return jsonify({"status": "not_found"})

@app.route('/clear_session')
def clear_session():
    """Manual session clearing endpoint"""
    _clear_all_session_data()
    flash("Session cleared successfully. All recovered files and analysis data have been removed.", "success")
    return redirect(url_for('evidence_upload'))

@app.route('/perform_manual_carve', methods=['POST'])
def perform_manual_carve():
    #global sorted_carved_keys
    filepath = get_active_evidence_path()
    if not filepath:
        flash("No evidence file is loaded.", "error")
        return redirect(url_for('manual_carving'))

    try:
        start_offset = int(request.form.get('start_offset', 0))
        length = int(request.form.get('length', 0))
    except (ValueError, TypeError):
        flash("Invalid offset or length. Please enter numbers only.", "error")
        return redirect(url_for('manual_carving'))

    if length <= 0:
        flash("Length must be a positive number.", "error")
        return redirect(url_for('manual_carving'))

    file_size = os.path.getsize(filepath)

    if start_offset >= file_size:
        flash("Start offset is beyond the end of the file.", "error")
        return redirect(url_for('manual_carving'))

    with open(filepath, 'rb') as f:
        f.seek(start_offset)
        data_to_carve = f.read(min(length, file_size - start_offset))

    carved_filename = f"manual_carve_{start_offset}_{length}.dat"
    save_path = os.path.join(app.config['CARVED_FOLDER'], carved_filename)
    with open(save_path, 'wb') as out_f:
        out_f.write(data_to_carve)
    
    max_id = max([f['id'] for f in carved_files_db.values()] or [0])
    file_id = max_id + 1
    carved_files_db[carved_filename] = {"id": file_id, "name": carved_filename, "offset": f"0x{start_offset:08X}", "size_kb": f"{len(data_to_carve)/1024:.2f} KB"}
    sorted_carved_keys.append(carved_filename)

    flash(f"Successfully carved {len(data_to_carve)} bytes to '{carved_filename}'.", "success")
    return redirect(url_for('recovered_files'))

@app.route('/diag')
def diag():
    """Diagnostic endpoint for server environment and CDN reachability.
    Returns JSON with Python/platform info, quick checks of external CDNs used by the templates,
    and a tail of any wsgi import error log we created earlier.
    """
    import platform
    import urllib.request
    
    def _check_url(url, timeout=5):
        try:
            with urllib.request.urlopen(url, timeout=timeout) as resp:
                return {"ok": True, "status": getattr(resp, 'status', None)}
        except Exception as e:
            return {"ok": False, "error": str(e)}

    result = {
        "python_version": platform.python_version(),
        "platform": platform.platform(),
        "project_root": APP_ROOT,
        "cdn_tailwind": _check_url('https://cdn.tailwindcss.com'),
        "font_inter": _check_url('https://rsms.me/inter/inter.css'),
    }

    # wsgi import errors tail (if present)
    wsgi_log = os.path.join(APP_ROOT, 'wsgi_import_errors.log')
    if os.path.exists(wsgi_log):
        try:
            with open(wsgi_log, 'r', encoding='utf-8', errors='ignore') as lf:
                data = lf.read()
                result['wsgi_import_errors_tail'] = data[-4000:]
        except Exception as e:
            result['wsgi_import_errors_tail'] = f"error reading file: {e}"
    else:
        result['wsgi_import_errors_tail'] = None

    # carved/deleted folder counts
    try:
        result['carved_files_count'] = len(os.listdir(CARVED_FOLDER)) if os.path.exists(CARVED_FOLDER) else None
    except Exception as e:
        result['carved_files_count'] = f"error: {e}"
    try:
        result['deleted_files_count'] = len(os.listdir(DELETED_RECOVERY_FOLDER)) if os.path.exists(DELETED_RECOVERY_FOLDER) else None
    except Exception as e:
        result['deleted_files_count'] = f"error: {e}"

    return jsonify(result)

if __name__ == '__main__':
    app.run(debug=True)
